#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_brandbook.py — 표준 콘텐츠 스키마(content_schema.json)를 입력받아
sky 템플릿을 편집해 브랜드북 PPTX를 생성하는 범용 빌더.

사용:
    python3 build_brandbook.py <content.json> <template.pptx> <out.pptx> [assets_dir]

STEP 7 서버(FastAPI)가 이 함수를 호출한다.
"""
import sys, json, copy, os
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ===== 템플릿 디자인 토큰 (sky.pptx 추출, 기본값) =====
TPL_INK    = RGBColor(0x25,0x25,0x25)
TPL_TEAL   = RGBColor(0x26,0x6B,0x87)
TPL_TEALL  = RGBColor(0x3A,0xA0,0xC9)
TPL_GRAY   = RGBColor(0x80,0x80,0x80)
TPL_WHITE  = RGBColor(0xFF,0xFF,0xFF)
TPL_HEADBG = RGBColor(0xED,0xF4,0xF8)
TPL_ROWALT = RGBColor(0xF4,0xF8,0xFA)
SZ_TABLE=22; SZ_TABLE_HEAD=22

# 팔레트 6종 (앱 PALETTES와 동일). 단, sky 템플릿은 배경/아이콘이 PNG라
# 텍스트·표 강조색만 팔레트로 바뀐다(배경 그래픽 색은 고정).
def _hex(s): return RGBColor(int(s[1:3],16), int(s[3:5],16), int(s[5:7],16))
PALETTES = {
  'sky':    {'teal':'#266B87','head':'#EDF4F8','rowalt':'#F4F8FA'},
  'sage':   {'teal':'#3E6B5A','head':'#EAF1ED','rowalt':'#F2F7F4'},
  'navy':   {'teal':'#2E5FA3','head':'#EAF0F8','rowalt':'#F2F6FB'},
  'crimson':{'teal':'#A8392B','head':'#F8ECEA','rowalt':'#FBF3F2'},
  'mono':   {'teal':'#333333','head':'#EFEFEF','rowalt':'#F6F6F6'},
  'forest': {'teal':'#2E7D4F','head':'#E8F2EC','rowalt':'#F2F8F4'},
  'violet': {'teal':'#6B4E97','head':'#F0EBF6','rowalt':'#F6F2FB'},
}
def apply_palette(pid):
    global TPL_TEAL, TPL_HEADBG, TPL_ROWALT
    p = PALETTES.get(pid)
    if not p: return
    TPL_TEAL   = _hex(p['teal'])
    TPL_HEADBG = _hex(p['head'])
    TPL_ROWALT = _hex(p['rowalt'])

# ===== 텍스트 헬퍼 =====
def walk(sh):
    for x in sh:
        yield x
        if x.shape_type==6: yield from walk(x.shapes)

def set_text(tf, new):
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = new
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else:
        r=p0.add_run(); r.text=new
    for p in tf.paragraphs[1:]: p._p.getparent().remove(p._p)

def set_multi(tf, lines):
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else:
        r=p0.add_run(); r.text=lines[0]
    for p in tf.paragraphs[1:]: p._p.getparent().remove(p._p)
    from pptx.text.text import _Paragraph
    for ln in lines[1:]:
        newp=copy.deepcopy(p0._p); p0._p.getparent().append(newp)
        para=_Paragraph(newp,p0._parent)
        if para.runs:
            para.runs[0].text=ln
            for r in para.runs[1:]: r._r.getparent().remove(r._r)

def by_pos(slide):
    d={}
    for x in walk(slide.shapes):
        if x.has_text_frame and x.text_frame.text.strip():
            k=(round(Emu(x.left).inches,1),round(Emu(x.top).inches,1))
            d[k]=x
    return d

def edit_by_text(slide, mapping, multimap=None):
    multimap = multimap or {}
    seen={}
    for x in walk(slide.shapes):
        if not (x.has_text_frame and x.text_frame.text.strip()): continue
        t=x.text_frame.text.strip().replace('\n',' / ').replace('\xa0',' ')
        if t in multimap:
            idx=seen.get(t,0); lines=multimap[t]
            if lines and isinstance(lines[0], list):
                set_multi(x.text_frame, lines[idx] if idx<len(lines) else lines[-1])
            else:
                set_multi(x.text_frame, lines)
            seen[t]=idx+1
        elif t in mapping:
            idx=seen.get(t,0); val=mapping[t]
            if isinstance(val, list):
                set_text(x.text_frame, val[idx] if idx<len(val) else val[-1])
            else:
                set_text(x.text_frame, val)
            seen[t]=idx+1

def clone_slide(prs, src):
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    src_part = src.part; dst_part = new_slide.part
    rid_map = {}
    for rid, rel in src_part.rels.items():
        if "image" in rel.reltype:
            rid_map[rid] = dst_part.relate_to(rel.target_part, rel.reltype)
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for blip in el.iter(qn('a:blip')):
            emb = blip.get(qn('r:embed'))
            if emb in rid_map: blip.set(qn('r:embed'), rid_map[emb])
        new_slide.shapes._spTree.append(el)
    return new_slide

KEEP_LABELS = {'성과 예측','예측한 성과의 핵심 내용을 요약해 보세요. 주요 지표를 제시해도 좋습니다.'}
def blank_bg_slide(prs, src):
    sl = clone_slide(prs, src); pics=0
    for shp in list(sl.shapes):
        if shp.shape_type==13:
            pics+=1
            if pics<=3: continue
            shp._element.getparent().remove(shp._element)
        elif shp.has_text_frame:
            if shp.text_frame.text.strip() in KEEP_LABELS: continue
            shp._element.getparent().remove(shp._element)
    return sl

def add_text(sl, x,y,w,h, runs, align=PP_ALIGN.LEFT, sp=None):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if sp: p.line_spacing=sp
        for (t,sz,b,c) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b
            r.font.color.rgb=c; r.font.name='Pretendard'
    return tb

def add_photo(sl, path, x,y,w,h):
    from PIL import Image
    iw,ih=Image.open(path).size; tr=w/h; ir=iw/ih
    p=sl.shapes.add_picture(path, Inches(x),Inches(y),Inches(w),Inches(h))
    if ir>tr:
        cw=int(ih*tr); l=(iw-cw)//2; p.crop_left=l/iw; p.crop_right=(iw-l-cw)/iw
    else:
        ch=int(iw/tr); t=(ih-ch)//2; p.crop_top=t/ih; p.crop_bottom=(ih-t-ch)/ih
    return p

def make_table_slide(prs, src, title, sub, headers, rows, col_w):
    sl = blank_bg_slide(prs, src)
    edit_by_text(sl, {'성과 예측':title,
        '예측한 성과의 핵심 내용을 요약해 보세요. 주요 지표를 제시해도 좋습니다.':sub})
    nrow=len(rows)+1; ncol=len(headers)
    gtbl = sl.shapes.add_table(nrow, ncol, Inches(1.8), Inches(4.4),
                               Inches(sum(col_w)), Inches(0.85*nrow))
    tbl = gtbl.table
    for ci,w in enumerate(col_w): tbl.columns[ci].width = Inches(w)
    for ci,h in enumerate(headers):
        c=tbl.cell(0,ci); c.fill.solid(); c.fill.fore_color.rgb=TPL_HEADBG
        tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=h; r.font.size=Pt(SZ_TABLE_HEAD); r.font.bold=True
        r.font.color.rgb=TPL_TEAL; r.font.name='Pretendard'
    for ri,row in enumerate(rows,1):
        for ci,val in enumerate(row):
            c=tbl.cell(ri,ci); c.fill.solid()
            c.fill.fore_color.rgb=(TPL_WHITE if ri%2 else TPL_ROWALT)
            tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(val); r.font.size=Pt(SZ_TABLE); r.font.bold=(ci==0)
            r.font.color.rgb=(TPL_TEAL if ci==0 else TPL_INK); r.font.name='Pretendard'
    for ri in range(nrow): tbl.rows[ri].height=Inches(0.85 if ri==0 else 0.8)
    return sl

# =========================================================
#  메인 빌더
# =========================================================
def build(content, template_path, out_path, assets_dir='.'):
    apply_palette((content.get('design') or {}).get('palette','sky'))
    prs = Presentation(template_path)
    S = list(prs.slides)
    # 템플릿 슬라이드 역할 인덱스 (sky.pptx 기준)
    T_COVER, T_TOC, T_3BOX, T_4STEP, T_TABLE3, T_2x2A, T_3BOXB, T_2x2TEAM, T_FLOW = \
        S[0], S[1], S[2], S[5], S[15], S[8], S[11], S[17], S[18]
    SRC_TABLE = S[15]   # 표/사진용 깨끗 배경
    SRC_2x2   = S[17]   # 2x2 박스
    SRC_2x2A  = S[8]    # 아이콘 2x2
    SRC_COVER = S[0]

    def P(name): 
        p=os.path.join(assets_dir, name)
        return p if os.path.exists(p) else None

    ac = content['academy']; cp = content['copy']
    logo = P(ac.get('logo','logo.png'))
    built = []   # 생성된 슬라이드를 순서대로 모음

    # ---------- 1) 표지 ----------
    cov = clone_slide(prs, SRC_COVER)
    edit_by_text(cov, {
        'MIRIFOOD': ac['name'], '작성일':'개원','작성자':'원장',
        '2080.00.00': ac.get('founded','')+'년', '사업기획팀 김미리': ac.get('principal',''),
        '이곳에 해당 자료에 관한 간단한 설명을 입력해 보세요.':
            f"{ac.get('nameSub','')} · {ac.get('region','')} {ac.get('subjects','')} 전문 학원",
    }, multimap={'상품 실적 개선을 위한  / 판매 활성화 방안':
        _two_lines(cp.get('slogan', ac['name']))})
    if logo: cov.shapes.add_picture(logo, Inches(17.5), Inches(4.5), Inches(7.0), Inches(7.0))
    built.append(cov)

    # ---------- 2) 학원 소개 (사진 + 인사말) ----------
    photos = content.get('photos',{})
    intro_photos = [P(x) for x in (photos.get('intro') or []) if P(x)]
    if cp.get('introLong') or intro_photos:
        intro = blank_bg_slide(prs, SRC_TABLE)
        edit_by_text(intro, {'성과 예측':'학원 소개',
            '예측한 성과의 핵심 내용을 요약해 보세요. 주요 지표를 제시해도 좋습니다.':
            cp.get('introShort') or cp.get('introLong','')[:40]})
        if cp.get('catch'):
            add_text(intro, 1.8, 4.6, 11.5, 0.7, [[(f'"{cp["catch"]}"', 26, True, TPL_TEAL)]])
        body = cp.get('introLong','')
        lines = _wrap(body, 28)
        add_text(intro, 1.8, 5.6, 11.5, 3.6, [[(ln, 18, False, TPL_INK)] for ln in lines], sp=1.5)
        foot = f"원장 {ac.get('principal','')}  ·  {ac.get('founded','')}년 개원  ·  {ac.get('region','')}"
        add_text(intro, 1.8, 11.6, 11.5, 0.6, [[(foot, 16, True, TPL_GRAY)]])
        if len(intro_photos)>=1: add_photo(intro, intro_photos[0], 13.6, 4.5, 11.2, 4.0)
        if len(intro_photos)>=2: add_photo(intro, intro_photos[1], 13.6, 8.7, 11.2, 3.5)
        built.append(intro)

    # ---------- 3) 핵심 약속 (S3) ----------
    pr = content.get('promise')
    if pr and pr.get('pillars'):
        box = clone_slide(prs, T_3BOX)
        m = {'내용 요약':'핵심 약속',
             '유통 확장 기반 매출 성장 전략을 요약해 보세요.': pr.get('headline','')}
        labels = ['부진 원인','활성화 방안','성과 예측']
        for i,p in enumerate(pr['pillars'][:3]):
            m[labels[i]] = p['title']
        bl = pr.get('bullets',[])
        bl_keys = ['[시설] 등 유통망 확대로 접근성 강화','SNS 인플루언서 마케팅으로 인지도 확보',
                   '구독 포인트 프로그램으로 재구매율 증대 기여','주요 방안 상세 내용을 간단히 요약해 보세요']
        for i,k in enumerate(bl_keys):
            if i < len(bl): m[k]=bl[i]
        mm = {'해당 영역의 핵심 내용을 / 작성해 보세요.':
              [p.get('lines',[p['title']]) for p in pr['pillars'][:3]]}
        edit_by_text(box, m, multimap=mm)
        built.append(box)

    # ---------- 4) 수업 특징 (S9 2x2) ----------
    feats = content.get('features',[])
    if feats:
        built.append(_build_2x2_icon(prs, SRC_2x2A, '수업 특징',
            '수업의 핵심을 정리했습니다.', feats))

    # ---------- 5) 입학 안내 (S6 4단계) ----------
    adm = content.get('admission',[])
    if adm:
        st = clone_slide(prs, T_4STEP)
        m = {'고객 행동 분석':'입학 안내',
             '병목 지점, 이탈 구간 등 고객 구매 여정에서 발견한 인사이트를 요약해 보세요.':
             ' → '.join(a['title'] for a in adm),
             '위 지점에서 발견한 핵심 문제점이나 개선 기회':
             adm[-1].get('desc','') if adm else ''}
        step_titles=['인지','관심','비교 / 검토','구매']
        for i,a in enumerate(adm[:4]): m[step_titles[i]] = a['title']
        mm = {'해당 단계에서 고객이 / 하는 행동, 사용 채널, 주요 / 지표를 입력해 보세요.':
              [_wrap(a.get('desc',''),12) for a in adm[1:4]]}
        edit_by_text(st, m, multimap=mm)
        built.append(st)

    # ---------- 6) 분반 구조 (표) ----------
    divs = content.get('divisions',[])
    if divs:
        rows=[[d['grade'],d['capacity'],d['level'],d['detail']] for d in divs]
        built.append(make_table_slide(prs, SRC_TABLE, '분반 구조',
            '입학 테스트를 통해 학년별 레벨로 배정됩니다.',
            ['학년부','정원','레벨','구성'], rows, [3.0,3.2,3.5,11.3]))

    # ---------- 7) 시간표 (여러 표) ----------
    for tt in content.get('timetables',[]):
        built.append(make_table_slide(prs, SRC_TABLE,
            tt['title'], tt.get('sub',''), tt['headers'], tt['rows'],
            tt.get('colWidths', _auto_widths(len(tt['headers'])))))

    # ---------- 8) 학습 관리 (S18 2x2) ----------
    mng = content.get('management',[])
    if mng:
        built.append(_build_2x2_team(prs, SRC_2x2, '학습 관리 시스템',
            '담당 선생님이 학습 전반을 책임지고 관리합니다.', mng))

    # ---------- 9) 특별 프로그램 (S12 3박스) ----------
    sp = content.get('specials',[])
    if sp:
        pg = clone_slide(prs, T_3BOXB)
        m = {'3~4개의 전략 키워드와 함께 핵심 전략 방향을 간단하게 요약해 보세요.':
             '학년·시기별 특별 프로그램을 운영합니다.',
             '핵심 전략 방향':'특별 프로그램',
             '전략이 지향하는 최종 목표나 비전': content['copy'].get('vision','')}
        names=['유통망 확대','고객 리텐션 강화','타겟 집중 마케팅']
        for i,s in enumerate(sp[:3]): m[names[i]] = s['title']
        mm = {'전략의 핵심 내용 요약을 / 입력해 보세요.':[[s['desc']] for s in sp[:3]]}
        edit_by_text(pg, m, multimap=mm)
        built.append(pg)

    # ---------- 10) 주요 실적 (사진 + 카드) ----------
    ach = content.get('achievements',[])
    if ach:
        sl = blank_bg_slide(prs, SRC_TABLE)
        edit_by_text(sl, {'성과 예측':'주요 실적',
            '예측한 성과의 핵심 내용을 요약해 보세요. 주요 지표를 제시해도 좋습니다.':
            '전 학년에서 성과를 만들어왔습니다.'})
        ap = P(photos.get('achievement','')) if photos.get('achievement') else None
        if ap: add_photo(sl, ap, 1.8, 4.6, 9.5, 7.8); cardx=12.2
        else: cardx=1.8
        cy=4.6
        for a in ach[:4]:
            card=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cardx),Inches(cy),Inches(12.6),Inches(1.75))
            card.fill.solid(); card.fill.fore_color.rgb=TPL_ROWALT; card.line.fill.background(); card.shadow.inherit=False
            try: card.adjustments[0]=0.08
            except: pass
            tagb=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cardx+0.4),Inches(cy+0.4),Inches(2.8),Inches(0.95))
            tagb.fill.solid(); tagb.fill.fore_color.rgb=TPL_TEAL; tagb.line.fill.background(); tagb.shadow.inherit=False
            tf=tagb.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
            r=p.add_run(); r.text=a['tag']; r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=TPL_WHITE; r.font.name='Pretendard'
            add_text(sl, cardx+3.6, cy+0.32, 8.6, 1.3,
                [[(a['headline'], 19, True, TPL_INK)],[(a.get('sub',''), 15, False, TPL_GRAY)]], sp=1.3)
            cy+=1.95
        built.append(sl)

    # ---------- 11) FAQ (S9 박스, 4개씩) ----------
    faq = content.get('faq',[])
    for i in range(0, len(faq), 4):
        chunk = faq[i:i+4]
        part = f' ({i//4+1})' if len(faq)>4 else ''
        built.append(_build_faq(prs, SRC_2x2A, '자주 묻는 질문'+part, chunk))

    # ---------- 12) 운영 규정 (S18 2x2) ----------
    rules = content.get('rules',[])
    if rules:
        built.append(_build_2x2_team(prs, SRC_2x2, '운영 규정 & 환불',
            '학원 운영 규정을 안내합니다.',
            [{'title':r['title'],'lead':'','desc':' / '.join(r['items'])} for r in rules],
            bullets=True))

    # ---------- 13) 마무리 ----------
    end = clone_slide(prs, SRC_COVER)
    edit_by_text(end, {
        'MIRIFOOD': ac['name'], '작성일':'문의','작성자':'블로그',
        '2080.00.00': ac.get('phone',''), '사업기획팀 김미리': ac.get('blog',''),
        '이곳에 해당 자료에 관한 간단한 설명을 입력해 보세요.': ac.get('address',''),
    }, multimap={'상품 실적 개선을 위한  / 판매 활성화 방안':
        _two_lines(cp.get('slogan',''))})
    if logo: end.shapes.add_picture(logo, Inches(17.5), Inches(4.5), Inches(7.0), Inches(7.0))
    built.append(end)

    # ---------- 순서 재배치: built 순서대로만 남기기 ----------
    lst = prs.slides._sldIdLst
    built_ids = set(id(s) for s in built)
    # 원본 슬라이드 중 built에 없는 것 제거, built는 순서대로
    sldid_by_slide = {}
    for sldId in list(lst):
        rId = sldId.get(qn('r:id'))
        slide = prs.slides._sldIdLst  # placeholder
    # 간단법: built의 _element 순서대로 sldIdLst 재구성
    # 각 slide의 sldId 찾기
    id_objs = list(lst)
    slide_to_sldid = {}
    for sldId, slide in zip(id_objs, prs.slides):
        slide_to_sldid[id(slide)] = sldId
    for sldId in id_objs:
        lst.remove(sldId)
    for s in built:
        sid = slide_to_sldid.get(id(s))
        if sid is not None: lst.append(sid)

    prs.save(out_path)
    return len(built)

# ===== 보조 함수 =====
def _two_lines(text):
    if len(text)<=12: return [text]
    mid=len(text)//2
    for off in range(mid, len(text)):
        if text[off]==' ': return [text[:off], text[off+1:]]
    return [text[:mid], text[mid:]]

def _wrap(text, n):
    words=text.replace('\n',' ').split(' ')
    lines=[]; cur=''
    for w in words:
        if len(cur)+len(w)+1 <= n: cur=(cur+' '+w).strip()
        else: lines.append(cur); cur=w
    if cur: lines.append(cur)
    return lines[:6] or ['']

def _auto_widths(n):
    total=24.0; return [round(total/n,1)]*n

def _build_2x2_icon(prs, src, title, sub, items):
    sl = clone_slide(prs, src)
    m = {'세부 목표들을 어떤 관점에서 구성했는지 간단하게 요약해 보세요.': sub,
         '판매 활성화 목표': title}
    pos = by_pos(sl)
    # 라벨/제목/설명 위치 (S9 기준)
    label_pos=[(3.2,4.8),(15.0,4.8),(3.2,9.6),(15.0,9.6)]
    title_pos=[(5.8,4.8),(17.6,4.8),(5.8,9.6),(17.6,9.6)]
    desc_pos =[(2.8,6.0),(14.4,6.0),(2.8,10.8),(14.4,10.8)]
    edit_by_text(sl, m)
    pos = by_pos(sl)
    for i,it in enumerate(items[:4]):
        if label_pos[i] in pos: set_text(pos[label_pos[i]].text_frame, it.get('lead','') or it['title'][:6])
        if title_pos[i] in pos: set_text(pos[title_pos[i]].text_frame, it['title'])
        lines=_wrap(it.get('desc',''),20)
        for li in range(3):
            k=(desc_pos[i][0], round(desc_pos[i][1]+0.7*li,1))
            if k in pos:
                if li<len(lines): set_text(pos[k].text_frame, lines[li])
                else: pos[k]._element.getparent().remove(pos[k]._element)  # 빈 줄 박스 제거
    return sl

def _build_2x2_team(prs, src, title, sub, items, bullets=False):
    sl = clone_slide(prs, src)
    m = {'역할 분담의 전체적인 구조를 간단하게 요약해 보세요.': sub, '역할 분담': title}
    team=['마케팅팀','영업팀','운영팀','고객관리팀']
    detail_keys=[
        ['SNS 광고 캠페인 기획 및 집행','브랜드 메시지 및 콘텐츠 제작','캠페인 성과 분석 및 최적화'],
        ['[시설] 입점 협상 및 계약','매장 진열 조건 협의 및 관리','점포별 판매 실적 모니터링'],
        ['구독 프로그램 개발 및 운영','포인트 적립 정책 관리','재고 관리 및 배송 최적화'],
        ['고객 문의 및 불만 응대','구독 고객 리텐션 관리','VOC 수집 및 개선사항 전달'],
    ]
    for i,it in enumerate(items[:4]):
        m[team[i]] = it['title']
        if bullets:
            parts = it['desc'].split(' / ')
        else:
            parts = _wrap(it.get('desc',''),16)
        for j in range(3):
            m[detail_keys[i][j]] = parts[j] if j<len(parts) else ''
    edit_by_text(sl, m)
    return sl

def _build_faq(prs, src, title, faqs):
    sl = clone_slide(prs, src)
    edit_by_text(sl, {'세부 목표들을 어떤 관점에서 구성했는지 간단하게 요약해 보세요.':
                      '입학 상담에서 가장 많이 받는 질문들을 모았습니다.',
                      '판매 활성화 목표': title})
    pos = by_pos(sl)
    label_pos=[(3.2,4.8),(15.0,4.8),(3.2,9.6),(15.0,9.6)]
    title_pos=[(5.8,4.8),(17.6,4.8),(5.8,9.6),(17.6,9.6)]
    desc_pos =[(2.8,6.0),(14.4,6.0),(2.8,10.8),(14.4,10.8)]
    for i in range(4):
        if i < len(faqs):
            f=faqs[i]
            if label_pos[i] in pos: set_text(pos[label_pos[i]].text_frame, f'Q{i+1}')
            if title_pos[i] in pos: set_text(pos[title_pos[i]].text_frame, f['q'])
            lines=_wrap(f['a'],22)
            for li in range(3):
                k=(desc_pos[i][0], round(desc_pos[i][1]+0.7*li,1))
                if k in pos:
                    if li<len(lines): set_text(pos[k].text_frame, lines[li])
                    else: pos[k]._element.getparent().remove(pos[k]._element)
        else:
            # 미사용 박스: 라벨·제목·설명 모두 제거
            for k in [label_pos[i], title_pos[i]]:
                if k in pos: pos[k]._element.getparent().remove(pos[k]._element)
            for li in range(3):
                k=(desc_pos[i][0], round(desc_pos[i][1]+0.7*li,1))
                if k in pos: pos[k]._element.getparent().remove(pos[k]._element)
    return sl


if __name__=='__main__':
    if len(sys.argv)<4:
        print('usage: build_brandbook.py <content.json> <template.pptx> <out.pptx> [assets_dir]')
        sys.exit(1)
    content=json.load(open(sys.argv[1], encoding='utf-8'))
    assets = sys.argv[4] if len(sys.argv)>4 else '.'
    n=build(content, sys.argv[2], sys.argv[3], assets)
    print(f'saved {sys.argv[3]} ({n} slides)')
