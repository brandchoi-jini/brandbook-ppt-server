# -*- coding: utf-8 -*-
"""좌표 기반 3단 리플렛 빌더 (GLT 템플릿 레이아웃을 코드로 재현).

- 세원 teal/yellow 템플릿의 배치·색·서체를 좌표로 그대로 그린다.
- 템플릿 방식과 달리 '빈 칸'이 없다: 데이터 개수만큼만 카드를 생성 → 데이터가
  얼마든 세원 잔존/빈칸 문제가 구조적으로 발생하지 않는다.
- 로고·표지사진 자리를 명시적으로 둔다(assets 없으면 안내 박스).
- 11.7 x 8.27in (A4 가로), 3등분 패널. Pretendard.

입력: to_schema_v3 가 만든 v3 스키마.
출력: 2슬라이드 PPTX (바깥면/안쪽면), 모든 도형 편집 가능.
"""
from __future__ import annotations

import io
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

import io as _io
import urllib.request as _urlreq

# URL/경로/base64 이미지를 메모리로 가져오기 (실패 시 None) — v3 빌더와 동일 방식
_IMG_CACHE = {}
def _fetch_image(url):
    if not url:
        return None
    if url in _IMG_CACHE:
        return _IMG_CACHE[url]
    try:
        if url.startswith("data:"):
            import base64 as _b64
            _, _, b64 = url.partition(",")
            data = _b64.b64decode(b64)
        elif url.startswith("http"):
            req = _urlreq.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            data = _urlreq.urlopen(req, timeout=8).read()
        else:
            with open(url, "rb") as f:
                data = f.read()
        _IMG_CACHE[url] = data
        return data
    except Exception:
        _IMG_CACHE[url] = None
        return None

# A4 가로 3단 접지. 297mm = 11.6929in. 3등분이 아니라 접지 규격으로 나눈다.
# 안으로 접혀 들어가는 패널을 3mm 좁게(97mm) 만들어야 접었을 때 맞는다.
MM = 1.0 / 25.4  # mm → inch
W_IN, H_IN = 297 * MM, 210 * MM   # 11.6929 x 8.2677 in

# 패널 폭(mm) — 바깥면과 안쪽면이 좌우 반대다.
#  안쪽면: 좌100 · 중100 · 우97 (우측이 맨 먼저 안으로 접힘)
#  바깥면: 좌97 · 중100 · 우100 (안쪽면 우측 패널의 뒷면이 좌측 97)
PANELS_INSIDE_MM  = [100, 100, 97]
PANELS_OUTSIDE_MM = [97, 100, 100]

def _panel_geom(widths_mm):
    """패널 폭 배열(mm) → [(x_in, w_in), ...] 누적 좌표."""
    out, cx = [], 0.0
    for w in widths_mm:
        w_in = w * MM
        out.append((cx, w_in))
        cx += w_in
    return out

PANEL = 100 * MM  # 기준 패널 폭(3.937in) — 내부 여백 계산 기본값
# 서체는 Pretendard SemiBold 한 가지만 쓴다(Bold/Regular 구분 폐기).
# python-pptx 는 굵기를 이름으로 지정하므로 패밀리명에 SemiBold 를 박고,
# run.font.bold 는 항상 False 로 둔다(True 면 SemiBold 위에 인조 굵게가 겹침).
BODY = 10.0       # ★본문 기본 크기 — 9pt 는 인쇄에서 작았다
MIN_BODY = 8.0    # ★리플렛 본문 가독성 하한(pt). 8pt 까지 허용.
                  #   실측: 전체 텍스트의 절반이 7.0~8.5pt 로 떨어져 인쇄 시 안 읽혔다.
                  #   이 아래로는 만들지 않는다. 안 들어가면 <내용을 줄이거나 뺀다>.
FONT = "Pretendard SemiBold"
# ★런에 bold 를 안 걸면 뷰어에 따라 굵기가 안 살아난다(전 런 bold=False 였다).
#   제목·라벨은 bold=True 로 확실히 지정한다.

# 좌우 여백 기준(모든 패널 공통). 라벨·헤더·본문·번호칩이 이 세로선에 맞춰 정렬된다.
PAD = 0.30            # 패널 안쪽 좌측 여백
LEFT = PAD            # 콘텐츠 왼쪽 기준(패널 x에 더해 사용: x + LEFT)
# 콘텐츠 폭은 가장 좁은 패널(97mm=3.819in) 기준으로 통일 → 어느 패널에서도 안 넘침
CW = 97 * MM - PAD * 2  # ≈ 3.22in
MAX_ADMISSION = 6     # 입학절차 표시 상한(워드 5단계까지 그대로 수용)
CHIP = 0.34           # 번호칩 지름
CHIP_GAP = 0.14       # 칩과 텍스트 사이 간격
TEXT_X = LEFT + CHIP + CHIP_GAP   # 칩 오른쪽 텍스트 시작
TEXT_W = CW - CHIP - CHIP_GAP     # 칩 있는 줄의 텍스트 폭

# 팔레트: 세원 teal/yellow 실측색 그대로 + 확장용 navy.
PALETTES = {
    "sewon_teal": {
        "cover_bg": "17345C",   # 표지 패널 배경(딥)
        "cover_ink": "FFFFFF",  # 표지 위 글자
        "badge": "0AA6B5",      # 과목 뱃지/번호칩(주 청록)
        "accent2": "2D74DA",    # 보조 블루(커리큘럼 라벨·관리 번호)
        "label": "0AA6B5",      # 섹션 라벨(청록)
        "head": "17345C",       # 헤더(딥네이비)
        "title": "17345C",      # 카드 제목
        "body": "41505C",       # 설명 — 5C6B76 은 인쇄에서 흐렸다(대비 5.5:1 → 8.3:1)
        "soft": "E8F7F8",       # 연한 박스(FAQ답·관리카드)
        "card_line": "E3E9EE",  # 카드 테두리
        "footer": "6B7884",     # 푸터 — 8A97A0 은 인쇄에서 거의 안 보였다
        "paper": "FFFFFF",
    },
    "sewon_yellow": {
        "cover_bg": "F5C84C",
        "cover_ink": "5C3B20",
        "badge": "F29B21",
        "accent2": "27A9B7",
        "label": "F29B21",
        "head": "5C3B20",
        "title": "5C3B20",
        "body": "716254",
        "soft": "FFF3C9",
        "card_line": "EDE4CB",
        "footer": "A08A6A",
        "paper": "FFFFFF",
    },
    "navy_amber": {
        "cover_bg": "22375F",
        "cover_ink": "FFFFFF",
        "badge": "22375F",
        "accent2": "D98A1F",
        "label": "D98A1F",
        "head": "22375F",
        "title": "22375F",
        "body": "5C6B76",
        "soft": "EEF2F8",
        "card_line": "E3E9EE",
        "footer": "8A97A0",
        "paper": "FFFFFF",
    },
}
# 약칭
PALETTES["teal"] = PALETTES["sewon_teal"]
PALETTES["yellow"] = PALETTES["sewon_yellow"]

BASE_DIR = Path(__file__).resolve().parent


# ── 유틸 ──────────────────────────────────────────────
def _rgb(h: str) -> RGBColor:
    h = h.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _set_font(run, name=None, size=None, bold=False, color=None):
    """서체를 latin·ea·cs 세 슬롯에 모두 지정한다.
    ★python-pptx 의 run.font.name 은 <a:latin> 만 채운다.
      한글은 <a:ea>(East Asian) 슬롯을 보므로, 이걸 비워두면
      테마 기본서체(Calibri)로 떨어지고 PowerPoint 가 임의의 한글 폰트
      (맑은 고딕 등)로 대체한다 → 지정한 Pretendard 가 안 나온다."""
    name = name or FONT
    f = run.font
    f.name = name
    if size is not None:
        f.size = Pt(size)
    # ★서체 자체가 SemiBold 다. 여기에 bold=True 를 얹으면 PowerPoint 가
    #   합성 볼드를 만들어 굵기가 제각각으로 보인다(서체가 오락가락하는 원인).
    #   위계는 크기와 색으로 준다 → 굵기는 항상 SemiBold 하나로 고정.
    f.bold = False
    if color is not None:
        f.color.rgb = _rgb(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        el = rPr.find(_A_NS + tag)
        if el is None:
            el = rPr.makeelement(_A_NS + tag, {})
            rPr.append(el)
        el.set("typeface", name)


def _s(v: Any, default: str = "") -> str:
    if v is None:
        return default
    if isinstance(v, str):
        return re.sub(r"[ \t]+", " ", v).strip() or default
    # ★dict/list 가 그대로 str() 되어 "{'main': '너의 가능성을…'}" 이 표지에 찍혔다.
    #   대표 키를 꺼내 문자열로 바꾼다.
    if isinstance(v, dict):
        for k in ("main", "text", "title", "head", "value", "v", "name", "desc"):
            if v.get(k):
                return _s(v.get(k), default)
        return default
    if isinstance(v, (list, tuple)):
        parts = [_s(x) for x in v]
        parts = [p for p in parts if p]
        return " · ".join(parts) or default
    return str(v).strip() or default


def _first(*vals: Any, default: str = "") -> str:
    for v in vals:
        s = _s(v)
        if s:
            return s
    return default


def _sentences(text):
    """한국어 문장 분리. 종결부호(.!?·。) 뒤에서 끊고 부호는 살린다."""
    s = _s(text)
    if not s:
        return []
    parts = re.split(r"(?<=[.!?。])\s+", s)
    return [p.strip() for p in parts if p.strip()]


def _is_complete_sentence(x):
    """완결된 서술문인지. 표지 카피 판정용.
    종결어미로 끝나야 하고, 마침표만 붙은 명사형("~하는 것.")은 미완성으로 본다."""
    x = _s(x).strip()
    if not x:
        return False
    if not re.search(r"(합니다|입니다|습니다|됩니다|드립니다|십니다|해요|이에요|예요|한다|이다)[.!?]?$", x):
        return False
    body = re.sub(r"[.!?]$", "", x)
    return not re.search(
        r"(것|점|중|뿐|까지|부터|처럼|만큼|위해|통해|대한|위한|하는|되는|있는|없는)$", body)


def _cover_copy(text, max_chars):
    """표지 한마디 전용 필터.
    ★완결된 문장만 남긴다. 하나도 없으면 빈 문자열(표지에서 생략).
      미완성 조각을 표지에 올리느니 슬로건만 보여주는 편이 낫다."""
    s = _s(text)
    if not s:
        return ""
    good = [x for x in _sentences(s) if _is_complete_sentence(x)]
    if not good:
        return ""
    out = good[0]
    for nxt in good[1:]:
        if _dwidth(out) + 1 + _dwidth(nxt) > max_chars:
            break
        out = out + " " + nxt
    return out if _dwidth(out) <= max_chars * 1.35 else ""


def _trim_sentences(text, max_chars):
    """표시폭이 max_chars 를 넘지 않도록 <문장 단위로> 줄인다.
    ★완결된 문장만 남긴다. 명사형으로 끝나는 조각("~하는 것.")은 버린다.
      표지에 미완성 문장이 올라가는 것을 막기 위함(경희궁 사례).
    - 첫 문장이 완결형이면 항상 보존한다.
    - 단어 중간을 자르거나 '…' 를 붙이지 않는다."""
    s = _s(text)
    if not s:
        return s
    sents = _sentences(s)
    if not sents:
        return s

    def _complete(x):
        x = x.strip()
        if not re.search(r"(합니다|입니다|습니다|됩니다|드립니다|십니다|해요|이에요|예요|한다|이다)[.!?]?$", x):
            return False
        # 마침표만 붙은 명사형 종결은 미완성으로 본다
        body = re.sub(r"[.!?]$", "", x)
        return not re.search(r"(것|점|중|뿐|까지|부터|처럼|만큼|위해|통해|대한|위한|하는|되는|있는|없는)$", body)

    good = [x for x in sents if _complete(x)]
    if not good:
        # 완결 문장이 없으면 전체가 max 안에 들어갈 때만 통과, 아니면 비운다
        return s if _dwidth(s) <= max_chars else ""
    out = good[0]
    for nxt in good[1:]:
        if _dwidth(out) + 1 + _dwidth(nxt) > max_chars:
            break
        out = out + " " + nxt
    return out


def _clip(text: Any, n: int, suffix: str = "…") -> str:
    s = _s(text)
    if len(s) <= n:
        return s
    return s[: max(1, n - len(suffix))].rstrip() + suffix


def _items(value: Any) -> List[Dict[str, Any]]:
    """dict({items|stages|steps}) 또는 list 를 dict 리스트로 정규화."""
    if isinstance(value, dict):
        value = value.get("items") or value.get("stages") or value.get("steps") or []
    if not isinstance(value, list):
        return []
    out = []
    for x in value:
        if isinstance(x, dict):
            out.append(x)
        elif isinstance(x, str) and x.strip():
            out.append({"title": x.strip(), "desc": ""})
    return out


def _dwidth(s: str) -> float:
    """표시 폭(한글=1.0, 영숫자·공백=0.5). 한글/영문 혼용 문장의 줄 길이를
    글자 수가 아니라 실제 차지하는 폭으로 재기 위한 근사."""
    w = 0.0
    for ch in s:
        w += 1.0 if ord(ch) > 0x2E80 else 0.5
    return w


def _wrap(text: str, per_line: int, max_lines: int = 0) -> str:
    """어절 단위 줄바꿈. ★글자를 절대 버리지 않는다.
    기존 구현은 max_lines 를 넘으면 남은 어절을 통째로 버려서
    '아이의 현재 수준을' 처럼 문장이 말없이 끊겼다(경희궁 사례).
    이제 max_lines 는 무시하고 전문을 반환한다 — 넘치는 분량은
    _fit_text 가 글자 크기를 줄여 흡수한다."""
    s = _s(text)
    if not s:
        return s
    if "\n" in s:
        return s
    limit = max(4, per_line)
    lines, cur = [], ""
    for w in s.split(" "):
        # 한 어절이 한 줄보다 길면(예: 긴 영문 단어·URL) 어쩔 수 없이 강제로 끊는다.
        while _dwidth(w) > limit:
            take = ""
            for ch in w:
                if _dwidth(take + ch) > limit:
                    break
                take += ch
            if not take:
                break
            if cur:
                lines.append(cur)
                cur = ""
            lines.append(take)
            w = w[len(take):]
        if cur and _dwidth(cur) + 0.5 + _dwidth(w) > limit:
            lines.append(cur)
            cur = w
        else:
            cur = (cur + " " + w).strip()
    if cur:
        lines.append(cur)
    return "\n".join(lines)


def _char_w(size):
    """1.0폭 글자(한글 1자)가 차지하는 <가로 길이(inch)>.
    ★단위 주의: 포인트당 비율이 아니라 절대 길이를 돌려준다.
      (이전에 비율만 돌려주다 호출부에서 size 를 곱하지 않아
       한 줄에 233자가 들어간다고 계산되는 버그가 있었다.)
    큰 글자일수록 자간이 상대적으로 넓어 계수를 키운다."""
    if size >= 22:
        return size * 0.0162
    if size >= 17:
        return size * 0.0155
    return size * 0.0148


def _trim_to_fit(text, w_in, max_h, size, *, line_spacing=1.10, pad=0.10, sep=" · "):
    """★가독성 하한(9pt)에 걸리면 글자를 더 줄일 수 없으므로 <내용을 덜어낸다>.
    여러 항목이 " · " 로 이어진 글이면 뒤 항목부터, 문장이면 뒤 문장부터 뺀다.
    문장 중간은 자르지 않는다(말이 깨지므로)."""
    t = _s(text)
    if not t:
        return t
    def fits(v):
        # ★_fit_size 가 하한(9pt)에서 멈추면 wrapped 가 max_h 를 넘긴 채 돌아온다.
        #   반환된 크기가 하한이면 줄 수로 직접 판정해야 한다.
        _sz, wrapped = _fit_size(v, w_in, max_h, size,
                                 min_size=MIN_BODY, line_spacing=line_spacing)
        n = wrapped.count("\n") + 1
        need = n * (_sz * 1.30) / 72.0 + pad
        return need <= max_h + 0.02
    if fits(t):
        return t
    parts = [x.strip() for x in t.split(sep) if x.strip()]
    while len(parts) > 1:
        parts.pop()
        cand = sep.join(parts)
        if fits(cand):
            return cand
    sents = _sentences(t) if "_sentences" in globals() else [t]
    while len(sents) > 1:
        sents.pop()
        cand = " ".join(sents)
        if fits(cand):
            return cand
    return sents[0] if sents else t


def _need_h(text, w_in, size, line_spacing=1.30, pad=0.06):
    """이 글자 크기로 text 를 담는 데 필요한 상자 높이(inch)."""
    _sz, wrapped = _fit_size(_s(text), w_in, 99.0, size,
                             min_size=size, line_spacing=line_spacing)
    n = wrapped.count("\n") + 1
    return n * (size * line_spacing) / 72.0 + pad


def _fit_block(text, w_in, max_h, size, *, min_size=MIN_BODY, line_spacing=1.10,
               pad=0.10):
    """(size, wrapped, height) 를 돌려준다.
    ★height 는 wrapped 를 담기에 <충분한> 값이다. max_h 로 잘라내지 않는다.
      (예전엔 min(max_h, …) 로 깎아서, 줄 수가 max_h 를 넘길 때
       '박스는 작은데 글자는 많은' 상태가 되어 넘침이 났다.)
      호출부는 반환된 height 를 그대로 쓰되, max_h 를 넘겼는지 보고
      글을 줄일지 결정한다."""
    # ★하한(9pt)에서 멈추면 max_h 를 넘긴 채 돌아온다 → 내용을 덜어내 맞춘다.
    text = _trim_to_fit(text, w_in, max_h, size, line_spacing=line_spacing, pad=pad)
    size, wrapped = _fit_size(text, w_in, max_h, size,
                              min_size=min_size, line_spacing=line_spacing)
    n = wrapped.count("\n") + 1
    h = n * (size * line_spacing) / 72.0 + pad
    return size, wrapped, h


def _fit_size(text, w_in, h_in, size, *, min_size=MIN_BODY, line_spacing=1.04):
    """박스(w_in × h_in)에 text 전문이 들어가는 최대 글자 크기를 찾는다.
    한 줄에 들어가는 표시 폭 ≈ (박스폭 - 여백) / (글자폭). Pretendard 기준
    글자 1.0폭 ≈ 0.0148in per pt. 잘라내는 대신 크기로 맞춘다.
    ★줄바꿈이 이미 있는 텍스트도 각 줄을 다시 랩해서 실제 줄 수를 센다."""
    s = _s(text)
    if not s:
        return size, s
    usable_w = max(0.4, w_in - 0.10)
    usable_h = max(0.15, h_in)
    hard = s.split("\n")
    cand = size
    while cand >= min_size:
        per_line = max(4, int(usable_w / _char_w(cand)))
        out, n_lines = [], 0
        for seg in hard:
            w = _wrap(seg, per_line)
            out.append(w)
            n_lines += w.count("\n") + 1
        need_h = n_lines * (cand * line_spacing) / 72.0
        if need_h <= usable_h:
            return cand, "\n".join(out)
        cand -= 0.5
    # 작성자가 넣은 줄바꿈을 지키면 안 들어가는 경우 → 줄바꿈을 풀고 다시 흘린다
    if len(hard) > 1:
        joined = " ".join(x.strip() for x in hard if x.strip())
        cand = size
        while cand >= min_size:
            per_line = max(4, int(usable_w / _char_w(cand)))
            w = _wrap(joined, per_line)
            if (w.count("\n") + 1) * (cand * line_spacing) / 72.0 <= usable_h:
                return cand, w
            cand -= 0.5
    per_line = max(4, int(usable_w / _char_w(min_size)))
    return min_size, "\n".join(_wrap(seg, per_line) for seg in hard)


# ── 도형/텍스트 ───────────────────────────────────────
def _strip_style(shape):
    """python-pptx auto_shape가 자동 삽입하는 <p:style> 제거(그림자·테마색 오염 방지)."""
    try:
        el = shape._element
        for st in el.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}style"):
            st.getparent().remove(st)
        sp = el.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}style")
        if sp is not None:
            sp.getparent().remove(sp)
    except Exception:
        pass


def _rect(slide, x, y, w, h, fill=None, line=None, radius=False, line_w=0.75):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        # ★기본 adj(0.1667)는 캡슐처럼 과하게 둥글다 → 실제 반경 0.035in 로 맞춘다.
        try:
            sh.adjustments[0] = max(0.004, min(0.05, 0.035 / max(0.01, min(w, h))))
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = _rgb(line)
        sh.line.width = Pt(line_w)
    _strip_style(sh)
    return sh


def _text(slide, text, x, y, w, h, *, size=11, color="222222", bold=False,
          align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.04,
          autofit=True, cap=None, fit=True, min_size=MIN_BODY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    txt = _s(text)
    if cap:
        txt = _clip(txt, cap)
    # ★ 크기로 맞추되, 가독성 하한(9pt)에 걸려도 안 들어가면 <내용을 덜어낸다>.
    #   예전에는 하한 아래로 계속 줄여 7pt 글자가 나왔고, 하한을 걸자 상자를 넘쳤다.
    #   문장 중간은 자르지 않고 뒤 항목·뒤 문장부터 뺀다.
    if fit and txt:
        txt = _trim_to_fit(txt, w, h, size, line_spacing=line_spacing, pad=0.02)
        size, txt = _fit_size(txt, w, h, size,
                              min_size=min_size, line_spacing=line_spacing)
        # ★하한(9pt)에서 멈춰 상자보다 커지면, 글자를 더 줄이는 대신 상자를 늘린다.
        #   (아래에 여백이 남는 자리이므로 겹치지 않는다. 실측 2건에서 0.07~0.15in 부족했다)
        _need = (txt.count("\n") + 1) * (size * 1.30) / 72.0 + 0.04
        if _need > h:
            h = _need
            box.height = Inches(h)
    lines = txt.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        # ★bold 인자를 받아놓고 무시하고 있었다(전 런 bold=False).
        #   세미볼드 단일 서체라 크기·색만으로는 제목이 본문과 구분되지 않는다.
        #   제목·라벨은 호출부가 준 bold 를 그대로 살린다.
        _set_font(r, size=size, bold=bool(bold), color=color)
    if autofit:
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
    return box


def _chip(slide, x, y, d, number, c):
    """번호 원형 칩."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(c["badge"])
    sh.line.fill.background()
    _strip_style(sh)
    tf = sh.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(number)
    _set_font(r, size=10.5, bold=False, color="FFFFFF")
    return sh


# ── 섹션 헤더 바 ──────────────────────────────────────
# ★시크릿아카데미 리플렛처럼 <어두운 바>에 제목과 한 줄 설명을 함께 담는다.
#   예전에는 작은 알약 라벨 + 큰 헤드라인에 1.2in 을 썼는데,
#   그만큼 본문 자리가 줄어 FAQ·규정이 한두 개만 들어갔다.
_SEC_SUB = {
    "상담 · 예약 안내": "문의부터 반 배정까지 진행 순서",
    "자주 묻는 질문":   "등록 전 많이 물으시는 내용",
    "우리 학원의 강점": "다른 곳과 다르게 하고 있는 것",
    "교육 과정":       "학년별로 무엇을 어떻게 하는지",
    "학습 관리":       "배우고, 확인하고, 끝까지 봅니다",
    "주요 실적":       "재원생이 만든 결과",
    "특강 및 기타 수업": "정규 수업 외에 따로 운영하는 프로그램",
    "학원 규정":       "미리 알고 오시면 좋은 기준",
    "학원 안내":       "한눈에 보는 운영 방식",
}

def _bar_head(slide, x, title, c, sub=None, y=0.30):
    """어두운 헤더 바. 그린 뒤 본문 시작 y 를 반환."""
    t = _s(title)
    if not t:
        return y
    sb = _s(sub) if sub is not None else _SEC_SUB.get(t, "")
    h = 0.72 if sb else 0.50
    _rect(slide, x + LEFT, y, CW, h, fill=c["head"], radius=True)
    _text(slide, t, x + LEFT + 0.20, y + 0.09, CW - 0.40, 0.28,
          size=12, color=c["paper"], fit=False)
    if sb:
        # ★fit=False 라서 폭을 넘으면 뒷부분이 잘렸다("학년별 수업 내용과 해내는").
        #   글자 크기를 낮춰 한 줄에 담고, 그래도 길면 부제를 생략한다.
        _ss = 8.5
        while _ss > 6.5 and _dwidth(sb) * _ss * 0.0148 > (CW - 0.44):
            _ss -= 0.25
        if _dwidth(sb) * _ss * 0.0148 <= (CW - 0.40):
            _text(slide, sb, x + LEFT + 0.20, y + 0.39, CW - 0.40, 0.24,
                  size=_ss, color="BFD3E6", fit=False)
    return y + h + 0.26


def _label(slide, x, text, c, dark=False, y=0.33):
    """섹션 라벨 — ★헤더 문구를 없앤 뒤로 이 라벨이 그 면의 유일한 제목이 됐다.
    9pt 로는 제목 구실을 못 해, 색 배경 박스에 12pt 로 키운다."""
    t = _s(text)
    if not t:
        return y
    bw = min(CW, max(1.05, _dwidth(t) * 0.155 + 0.44))
    if dark:
        _text(slide, t, x + LEFT, y, CW, 0.30, size=12, color=c["paper"], bold=True)
    else:
        _rect(slide, x + LEFT, y, bw, 0.32, fill=c["label"], radius=True)
        _text(slide, t, x + LEFT, y, bw, 0.32, size=12, color="FFFFFF", bold=True,
              align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fit=False)
    return y + 0.32


def _header(slide, x, text, c, dark=False, y=0.72, max_h=1.30):
    """대형 헤더. ★실제로 차지한 높이(in)를 반환한다.
    3줄짜리 헤더가 들어오면 0.9in 박스를 넘겨 아래 subhead 와 글자가 겹쳤다
    (경희궁 '레벨 테스트부터 반 배정까지, 네 단계로 안내합니다').
    이제 줄 수만큼 높이를 늘려 잡고, 호출부는 반환값 아래에 다음 요소를 놓는다."""
    # ★헤더 문구가 없으면 자리를 차지하지 않는다.
    #   예전에는 학원마다 같은 고정 문구("숫자보다 변화의 폭을 봅니다")를 찍었다.
    t = _s(text)
    if not t:
        return y
    size, wrapped, h = _fit_block(t, CW, max_h, 19.5,
                                  min_size=14.5, line_spacing=1.16, pad=0.10)
    h = max(0.42, h)
    _text(slide, wrapped, x + LEFT, y, CW, h,
          size=size, color=(c["paper"] if dark else c["head"]), bold=True,
          line_spacing=1.16, fit=False)
    return y + h


def _footer(slide, x, name, phone, c):
    line = f"{name}  ·  {phone}" if phone else name
    _text(slide, line, x + LEFT, 7.92, CW, 0.23,
          size=7.5, color=c["footer"], align=PP_ALIGN.LEFT)


def _place_image(slide, url, x, y, w, h, cover=True):
    """url/base64/경로 이미지를 박스에 배치. cover=True면 꽉 채우고 crop, False면 비율 유지."""
    data = _fetch_image(url)
    if not data:
        return False
    try:
        from PIL import Image as _PILImage
        im = _PILImage.open(_io.BytesIO(data))
        iw, ih = im.size
        box_ratio = w / h
        img_ratio = iw / ih
        pic = slide.shapes.add_picture(_io.BytesIO(data), Inches(x), Inches(y), Inches(w), Inches(h))
        if cover:
            if img_ratio > box_ratio:
                crop = (1 - box_ratio / img_ratio) / 2
                pic.crop_left = crop; pic.crop_right = crop
            else:
                crop = (1 - img_ratio / box_ratio) / 2
                pic.crop_top = crop; pic.crop_bottom = crop
        else:
            if img_ratio > box_ratio:
                nw = w; nh = w / img_ratio
            else:
                nh = h; nw = h * img_ratio
            pic.width = Inches(nw); pic.height = Inches(nh)
            pic.left = Inches(x); pic.top = Inches(y)
        return True
    except Exception:
        try:
            slide.shapes.add_picture(_io.BytesIO(data), Inches(x), Inches(y), Inches(w), Inches(h))
            return True
        except Exception:
            return False


def _photo_or_box(slide, path, x, y, w, h, c, label_txt, cover_mode=True):
    """이미지 있으면 배치, 없으면 점선 안내 박스."""
    if _place_image(slide, path, x, y, w, h, cover=cover_mode):
        return True
    box = _rect(slide, x, y, w, h, fill=c["soft"], line=c["card_line"], radius=True)
    try:
        box.line.dash_style = 4
    except Exception:
        pass
    _text(slide, label_txt, x, y, w, h, size=8, color=c["footer"], bold=True,
          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return False


# ── 데이터 정규화 ─────────────────────────────────────
def _feature_list(schema):
    out = []
    for it in _items(schema.get("features"))[:6]:
        t = _s(it.get("title") or it.get("name") or "")
        d = _s(it.get("desc") or it.get("description") or "")
        if t or d:
            out.append({"title": t, "desc": d})
    return out


def _curriculum_list(schema):
    cur = schema.get("curriculum") or {}
    stages = _items(cur) or _items(schema.get("targets"))
    out = []
    for it in stages[:3]:
        name = _s(it.get("name") or it.get("title") or it.get("grade"))
        # v3 스키마: 설명은 desc가 아니라 tag(목표 한마디) + items(반·과목 배열)에 있다.
        tag = _s(it.get("tag"))
        items = it.get("items")
        items_txt = ""
        if isinstance(items, list):
            # ★" · " 로 이으면 중1·중2·중3 처럼 <완결된 문장>들이 한 덩어리가 되어
            #   가독성이 크게 떨어진다. 항목마다 줄을 바꾼다.
            _its = [_s(x) for x in items if _s(x)]
            _sent = sum(1 for x in _its if re.search(r"(다|요)\.\s*$", x))
            items_txt = ("\n".join(_its) if (_sent >= 2 or len(_its) >= 3)
                         else " · ".join(_its))
        elif isinstance(items, str):
            items_txt = _s(items)
        # 기존 desc/subj 계열도 폴백으로 지원
        desc_fallback = _s(it.get("desc") or it.get("description") or it.get("subj"))
        # 표시용 설명: tag를 위(굵게), items를 아래(상세)로 합침
        desc = tag or desc_fallback
        # ★detail 을 items 배열에서만 읽어 <detail 필드가 통째로 무시>됐다.
        #   교육 과정 면에 '초등 단계 / 개념 다지기' 만 남고 본문이 사라진 원인.
        detail = items_txt or _s(it.get("detail") or it.get("body") or it.get("text"))
        # tag 와 desc 가 모두 있으면 desc 는 상세로 내린다(내용 손실 방지)
        if tag and desc_fallback and desc_fallback != tag and not detail:
            detail = desc_fallback
        if not (name or desc or detail):
            continue
        out.append({
            "name": name.replace("부", ""),
            "desc": desc,
            "detail": detail,
        })
    return out


def _mgmt_list(schema):
    mg = schema.get("management") or {}
    out = []
    # v3 스키마: management.columns = [{name, rows:[{k,v}]}]
    if isinstance(mg, dict) and isinstance(mg.get("columns"), list):
        for col in mg["columns"]:
            if not isinstance(col, dict):
                continue
            for r in (col.get("rows") or []):
                if not isinstance(r, dict):
                    continue
                k = _s(r.get("k") or r.get("key") or r.get("title"))
                v = _s(r.get("v") or r.get("value") or r.get("desc"))
                line = f"{k} — {v}" if (k and v) else (k or v)
                if line:
                    out.append(line)
        if out:
            return out[:4]
    # 폴백: items/steps/math/science 계열
    src = _items(mg)
    if not src and isinstance(mg, dict):
        for key in ("steps", "items", "math", "science"):
            src.extend(_items(mg.get(key)))
    for it in src[:4]:
        # ★'title or desc' 라서 제목이 있으면 <설명을 버렸다>.
        #   '과제'만 남고 '매 수업 후 복습 과제·오답노트 병행' 이 사라졌다.
        t = _s(it.get("title") or it.get("name") or it.get("k"))
        d = _s(it.get("desc") or it.get("description") or it.get("v") or it.get("body"))
        if not d:
            _sub = it.get("items")
            if isinstance(_sub, list):
                d = " · ".join(_s(x) for x in _sub if _s(x))
        line = f"{t} — {d}" if (t and d) else (t or d)
        if line:
            out.append(line)
    return out[:4]


# ★질문 앞에 이미 'Q1.' 'Q.' '질문 1)' 이 붙어 오는 경우가 있다.
#   렌더러가 'Q. ' 를 한 번 더 붙여 "Q. Q1. 어떤 학생에게…" 로 나왔다.
_Q_PREFIX = re.compile(r"^\s*(?:Q|q|질문)\s*\d*\s*[.)\]:·-]*\s*")
_A_PREFIX = re.compile(r"^\s*(?:A|a|답변)\s*\d*\s*[.)\]:·-]*\s*")


def _faq_list(schema):
    out = []
    for it in _items(schema.get("faq"))[:4]:
        q = _Q_PREFIX.sub("", _s(it.get("q") or it.get("question")))
        a = _A_PREFIX.sub("", _s(it.get("a") or it.get("answer") or it.get("desc")))
        if q:
            out.append({"q": q, "a": a})
    return out


def _admission_list(schema):
    """입학·상담 절차 전체를 반환한다.
    워드 컨펌본이 5단계면 리플렛도 5단계다 — 3단계 하드캡은 폐기.
    글자수도 여기서 자르지 않는다(잘림의 원인). 넘치는 분량은
    _panel_admission 이 단계 수에 맞춰 크기·간격을 줄여 흡수한다."""
    src = _items(schema.get("admission"))
    out = []
    for it in src[:MAX_ADMISSION]:
        t = _s(it.get("title") or it.get("name"))
        d = _s(it.get("desc") or it.get("body"))
        if not d:
            # v3 admission step: 내용이 items 배열에 있음
            items = it.get("items")
            if isinstance(items, list):
                d = " · ".join(_s(x) for x in items if _s(x))
        if t or d:
            out.append({"title": t, "desc": d})
    return out


_EMOJI_RE = re.compile(
    "[\U0001F300-\U0001FAFF\U00002600-\U000027BF\U0001F000-\U0001F0FF\uFE0F]+")


def _de_emoji(s):
    """이모지 제거(디자인 원칙: 이모지 금지 → 선아이콘·도형으로).
    앞머리 이모지는 통째로 떼고, 문장 중간은 공백으로 바꾼다."""
    t = _EMOJI_RE.sub("", _s(s))
    return re.sub(r"\s{2,}", " ", t).strip(" ·-—")


def _achievement_lines(schema):
    ach = schema.get("achievements") or {}
    # ★achievements 가 리스트로 바로 오면 .get 에서 죽었다.
    if not isinstance(ach, dict):
        ach = {"items": ach if isinstance(ach, list) else []}
    lines = []
    # ★실적이 groups(유형별) 로 오면 _items 가 못 읽어 <실적 면이 통째로 사라졌다>.
    #   groups 를 평면 항목으로 펼쳐 준다.
    _src = _items(ach)
    if not _src and isinstance(ach, dict) and ach.get("groups"):
        _src = []
        for g in (ach.get("groups") or []):
            for it in (g.get("items") or []):
                t = _s(it.get("title") or it.get("name"))
                ch = _s(it.get("change"))
                nt = _s(it.get("note"))
                # ★change 값 자체가 '→' 로 시작하는 경우가 있어 '→ →' 로 겹쳤다.
                ch = re.sub(r"^\s*(?:→|->|⇒)\s*", "", ch)
                _desc = " ".join(x for x in [("→ " + ch) if ch else "", nt] if x)
                _src.append({"name": t, "desc": _desc})
    for it in _src[:8]:
        name = _s(it.get("name") or it.get("title"))
        desc = _s(it.get("desc") or it.get("description"))
        # ★'중1 입학 · → 고3 …' 처럼 가운뎃점과 화살표가 겹쳐 보였다.
        if desc.startswith(("→", "->", "⇒")):
            line = f"{name} {desc}".strip() if name else desc
        else:
            line = " · ".join(x for x in [name, desc] if x)
        line = _de_emoji(line)
        if not line:
            continue
        # ★한 항목에 실적이 통째로 뭉쳐 오는 경우가 있다(경희궁: "· 주요 실적 · 2024 KCIA…").
        #   줄바꿈이나 ' · ' 로 이어붙은 덩어리는 항목별로 쪼개 읽기 쉽게 만든다.
        parts = [p.strip(" ·-—") for p in re.split(r"[\n\r]+", line)]
        if len(parts) == 1 and len(line) > 46 and line.count(" · ") >= 2:
            parts = [p.strip() for p in line.split(" · ")]
        for p in parts:
            p = re.sub(r"^주요\s*실적\s*[·:]?\s*", "", p).strip(" ·-—")
            if p:
                lines.append(p)
    head = _s(ach.get("head") if isinstance(ach, dict) else "", "함께 이룬 결과")
    return head, lines[:8]


def _subject_line(schema):
    ac = schema.get("academy") or {}
    subj = _s(ac.get("subjects"))
    subj = subj.replace(" 전문", "")
    grades = []
    for t in _items(schema.get("targets")) or _items(schema.get("curriculum")):
        g = _s(t.get("grade") or t.get("name"))
        if g:
            grades.append(g.replace("부", ""))
    gtext = " · ".join(dict.fromkeys(grades)) if grades else ""
    if gtext and subj:
        return f"{gtext}  |  {subj}"
    return gtext or subj or ""


# ── 보조 카드(면 하단 여백을 실제 정보로 채운다) ──────────
_AUX_TITLE = {"rules": "학원 규정", "faq": "자주 묻는 질문",
              "specials": "특강 및 기타 수업", "schedule": "수업 요일 · 시간"}

def _aux_rows(schema, key):
    """보조 카드에 넣을 [(제목, 설명)] 목록."""
    if key == "rules":
        src = _items(schema.get("rules")) or _items(schema.get("policy"))
        return [(_s(i.get("k") or i.get("title")), _s(i.get("v") or i.get("desc")))
                for i in src if _s(i.get("k") or i.get("title")) or _s(i.get("v") or i.get("desc"))]
    if key == "faq":
        return [(i["q"], i["a"]) for i in _faq_list(schema)]
    if key == "specials":
        return [(_s(i.get("title") or i.get("name")),
                 _s(i.get("desc") or i.get("description")))
                for i in _items(schema.get("specials")) if _s(i.get("title") or i.get("name"))]
    if key == "schedule":
        return [(_s(i.get("k")), _s(i.get("v"))) for i in _items(schema.get("schedule")) if _s(i.get("v"))]
    return []


def _aux_taken(schema):
    """이미 배정·사용된 보조 카드 키 집합.
    ★앞면(상담·예약)과 뒷면(강점)이 서로 무엇을 썼는지 몰라
      '특강 및 기타 수업' 이 양면에 똑같이 두 번 그려졌다."""
    if not isinstance(schema, dict):
        return set()
    taken = set(schema.get("_auxUsed") or set())
    taken |= set((schema.get("_auxAssign") or {}).values())
    return taken


def _aux_choose(schema, host, cands):
    """이 면(host)이 하단에 넣을 보조 카드 키를 하나 고른다.
    계획(_auxAssign)에 배정된 것이 있으면 그것을, 없으면 아직 아무도
    쓰지 않은 후보 중에서 고른다. 고른 키는 사용 처리한다."""
    if not isinstance(schema, dict):
        return None
    used = set(schema.get("_auxUsed") or set())
    planned = (schema.get("_auxAssign") or {}).get(host)
    pick = None
    if planned and planned not in used and _aux_rows(schema, planned):
        pick = planned
    else:
        own = schema.get("_ownPanel") or set()
        blocked = _aux_taken(schema)
        for k in cands:
            if k in own or k in blocked:
                continue
            if _aux_rows(schema, k):
                pick = k
                break
    if pick:
        used.add(pick)
        schema["_auxUsed"] = used
    return pick


def _draw_aux_card(slide, schema, x, c, key, y_top, y_bot):
    """연한 카드 안에 제목(굵게)+설명을 넣는다. 그린 높이를 반환."""
    rows = [(t, d) for t, d in _aux_rows(schema, key) if t or d]
    if not rows or (y_bot - y_top) < 0.95:
        return 0.0
    IW = CW - 0.44          # 카드 안쪽 폭(렌더와 측정을 같은 값으로)
    PAD_T, PAD_B, GAP = 0.20, 0.16, 0.14
    # ★설명까지 넣으면 1개만 들어가는 경우가 많았다(특강 5개 중 1개).
    #   ①글자 축소 → ②설명 접고 제목만 → 순서로 <전부> 싣는 쪽을 택한다.
    size, with_desc = BODY, True

    def _fit(size, with_desc):
        hs, keep, used = [], [], PAD_T
        for t, d in rows:
            th = _need_h(t, IW, size, pad=0.02) if t else 0.0
            dh = _need_h(d, IW, size, pad=0.02) if (d and with_desc) else 0.0
            rh = max(0.22, th) + (dh + 0.02 if dh else 0.0) + GAP
            if 0.80 + used + rh + PAD_B > (y_bot - y_top):
                break
            keep.append((t, d if with_desc else "", max(0.22, th), dh)); used += rh
        return keep, used

    def _fit_mixed(size):
        """★설명을 통째로 접으면 '운영 시간', '상담 비용' 처럼 라벨만 남아
        읽을 내용이 없는 카드가 됐다. 들어가는 항목까지는 설명을 싣고,
        남는 항목은 제목만 실어 <전부 싣기>와 <내용 싣기>를 함께 만족시킨다."""
        limit = (y_bot - y_top) - 0.80 - PAD_B
        th_all = [(_need_h(t, IW, size, pad=0.02) if t else 0.0) for t, d in rows]
        keep, used = [], PAD_T
        for i, (t, d) in enumerate(rows):
            th = max(0.22, th_all[i])
            # 뒤에 남은 항목을 제목만으로 실을 때 필요한 높이
            rest = sum(max(0.22, th_all[j]) + GAP for j in range(i + 1, len(rows)))
            dh = _need_h(d, IW, size, pad=0.02) if d else 0.0
            if dh and used + th + dh + 0.02 + GAP + rest <= limit:
                keep.append((t, d, th, dh)); used += th + dh + 0.02 + GAP
            elif used + th + GAP + rest <= limit:
                keep.append((t, "", th, 0.0)); used += th + GAP
            else:
                break
        return keep, used

    keep, used = _fit_mixed(size)
    while len(keep) < len(rows) and size > MIN_BODY:
        size -= 0.5
        keep, used = _fit_mixed(size)
    two_col = False
    if len(keep) < len(rows):
        size, with_desc = BODY, False
        keep, used = _fit(size, False)
        while len(keep) < len(rows) and size > MIN_BODY:
            size -= 0.5
            keep, used = _fit(size, False)
    # ★제목만으로도 다 안 들어가면 2열로 배치해 전부 싣는다(특강 5개 중 3개만 나왔다).
    if len(keep) < len(rows) and len(rows) >= 4:
        size = BODY
        while size > MIN_BODY:
            _rh = _need_h("가", IW / 2 - 0.10, size, pad=0.02) + 0.10
            _need_all = PAD_T + _rh * ((len(rows) + 1) // 2) + PAD_B + 0.80
            if _need_all <= (y_bot - y_top):
                break
            size -= 0.5
        _rh = _need_h("가", IW / 2 - 0.10, size, pad=0.02) + 0.10
        keep = [(t, "", _rh - 0.10, 0.0) for t, d in rows]
        used = PAD_T + _rh * ((len(rows) + 1) // 2)
        two_col = True
    if not keep:
        return 0.0
    card_h = used - GAP + PAD_B
    # ★예전에는 작은 라벨만 붙여, 위 섹션(우리 학원의 강점)에 딸린 항목처럼 보였다.
    #   서로 종속 관계가 아니므로 같은 모양의 어두운 헤더 바를 준다.
    _ttl = _AUX_TITLE.get(key, key)
    cy = _bar_head(slide, x, _ttl, c, y=y_top) - 0.04
    _rect(slide, x + LEFT, cy, CW, card_h, fill=c["soft"], radius=True)
    ry = cy + PAD_T
    if two_col:
        colw = IW / 2
        rows_n = (len(keep) + 1) // 2
        for i, (t, _d, th, _dh) in enumerate(keep):
            cxx = x + LEFT + 0.22 + (colw if i >= rows_n else 0.0)
            cyy = ry + (th + 0.10) * (i - rows_n if i >= rows_n else i)
            _text(slide, t, cxx, cyy, colw - 0.10, th,
                  size=size, color=c["title"], min_size=MIN_BODY, fit=False)
        return 0.32 + card_h
    for t, d, th, dh in keep:
        if t:
            _text(slide, t, x + LEFT + 0.22, ry, IW, th,
                  size=size, color=c["title"],
                  min_size=MIN_BODY, line_spacing=1.20)
            ry += th
        if d and dh:
            _text(slide, d, x + LEFT + 0.22, ry + 0.02, IW, dh,
                  size=size, color=c["body"], min_size=MIN_BODY, line_spacing=1.26)
            ry += dh + 0.02
        ry += GAP
    return 0.32 + card_h


# ── 바깥면 패널 ───────────────────────────────────────
def _panel_admission(slide, schema, x, c):
    ac = schema.get("academy") or {}
    admission = schema.get("admission") or {}
    hb = _bar_head(slide, x, _first(admission.get("label") if isinstance(admission, dict) else "",
                                    default="상담 · 예약 안내"), c,
                   sub=_first(admission.get("subhead") if isinstance(admission, dict) else "",
                              default="문의부터 반 배정까지 진행 순서"))
    sub_y = hb - 0.16

    steps = _admission_list(schema)
    # ★없는 절차를 지어내지 않는다. 하드코딩 3단계(전화 상담/학생 상담/등원 결정)를
    #   폴백으로 넣던 탓에 학원에 없는 절차가 인쇄물에 나갔다.
    #   데이터가 없으면 단계 부분만 비우고, 이 면의 나머지(보조 카드·연락처)는 그대로 그린다.
    if not steps:
        steps = []
    # ── 단계 수에 맞춰 자동 축소 ──
    # 세로 가용구간: subhead 아래 ~ 5.00(구분선 위). 헤더가 길면 함께 내려간다.
    # ★연락처를 실적 면으로 옮기면 이 면 하단 3인치가 통째로 비었다.
    #   연락처를 싣지 않는 경우 단계가 면 전체를 쓰도록 끝점을 내린다.
    _has_contact = "achievements" not in (schema.get("_ownPanel") or set())
    Y0, Y_END = max(1.92, sub_y + 0.44), (4.98 if _has_contact else 7.30)
    n = max(1, len(steps))
    # 내용 길이에 맞춰 단계별 높이를 배분(고정 pitch 는 04번 설명 3줄에서 넘쳤다)
    t_h0 = 0.29
    need = []
    for it in steps:
        dl = max(1, int(_dwidth(it["desc"]) / 26) + 1) if it["desc"] else 0
        need.append(t_h0 + dl * 0.20 + 0.14)
    total = sum(need)
    avail = Y_END - Y0
    scale = min(1.0, avail / total) if total > 0 else 1.0
    tight = scale < 0.92
    t_size = 11.5 if not tight else 10.5
    d_size = 9.5 if not tight else 8.5   # ★8.5pt 는 작았다
    chip_d = CHIP if not tight else 0.30

    # ★남는 높이를 간격으로 벌리면 성의 없어 보인다.
    #   간격은 읽기 좋은 최소치만 주고, 남는 자리는 <학원 규정>으로 채운다.
    gap_extra = min(0.22, max(0.0, (avail - total) / max(1, len(steps))))

    y = Y0
    for i, it in enumerate(steps):
        blk = need[i] * scale
        _chip(slide, x + LEFT, y + 0.05, chip_d, f"{i+1:02d}", c)
        t_h = t_h0 if not tight else 0.27
        _text(slide, it["title"], x + TEXT_X, y, TEXT_W, t_h,
              size=t_size, color=c["title"], bold=True)
        # 설명은 자르지 않고 남은 높이 전부를 준다
        _text(slide, it["desc"], x + TEXT_X, y + t_h + 0.02,
              TEXT_W, max(0.20, blk - t_h - 0.06), size=d_size, color=c["body"],
              line_spacing=1.06)
        y += blk + gap_extra

    # ── 남는 자리를 보조 카드로 채운다(간격으로 때우지 않는다) ──
    _aux = _aux_choose(schema, "admission", ("rules", "specials", "faq"))
    if _aux:
        _draw_aux_card(slide, schema, x, c, _aux, y + 0.30,
                       5.74 if _has_contact else 7.56)

    # ── 연락처 (짧게) ───────────────────────────────────
    # ★실적 면(뒷면)에 연락처·QR 이 이미 들어간다. 여기 또 넣으면
    #   전화·주소가 두 면에 걸쳐 두 번 나온다 → 실적 면이 없을 때만 싣는다.
    _own_c = schema.get("_ownPanel") or set()
    if "achievements" not in _own_c:
        _rect(slide, x + LEFT, 5.90, CW, 0.012, fill=c["card_line"])
        phone = _s(ac.get("phone"))
        _text(slide, f"전화  {phone}" if phone else "전화 상담", x + LEFT, 6.10,
              CW, 0.28, size=11, color=c["head"], bold=True, fit=False)
        loc = _first(ac.get("location"), ac.get("address_short"), default="")
        if loc:
            lh = min(_need_h(loc, CW, MIN_BODY, pad=0.02), 1.10)
            _text(slide, loc, x + LEFT, 6.44, CW, lh,
                  size=MIN_BODY, color=c["body"], min_size=MIN_BODY, line_spacing=1.30)


def _panel_faq(slide, schema, x, c, name):
    hb = _bar_head(slide, x, "자주 묻는 질문", c)
    faqs = _faq_list(schema)
    if not faqs:
        faqs = [
            {"q": "수업은 어떻게 배정되나요?", "a": "상담·진단 후 수준에 맞는 반을 안내합니다."},
            {"q": "결석하면 보강되나요?", "a": "운영 기준에 따라 보강 또는 자료를 제공합니다."},
        ]
    # ── 답변 전문이 들어가도록 높이를 먼저 배분한다 ──
    # 가용 세로: 1.38 ~ 7.80(푸터 위). 질문·답변 길이에 비례해 나눠 갖고,
    # 총합이 넘치면 전체를 비례 축소한다(글자를 자르지 않는다).
    Y0, Y_END = max(0.95, hb + 0.20), 7.78
    blocks = []
    for i, it in enumerate(faqs):
        q = f"Q{i+1}. {it['q']}"
        q_lines = max(1, int(_dwidth(q) / 20) + 1)
        q_h = 0.30 * q_lines
        a_lines = max(1, int(_dwidth(it["a"]) / 24) + 1)
        box_h = 0.26 + 0.21 * a_lines
        blocks.append({"q": q, "a": it["a"], "q_h": q_h, "box_h": box_h})

    GAP = 0.22
    total = sum(b["q_h"] + 0.08 + b["box_h"] for b in blocks) + GAP * max(0, len(blocks) - 1)
    avail = Y_END - Y0
    scale = min(1.0, avail / total) if total > 0 else 1.0
    gap = GAP * scale

    y = Y0
    for b in blocks:
        q_h = b["q_h"] * scale
        box_h = b["box_h"] * scale
        _text(slide, b["q"], x + LEFT, y, CW, q_h + 0.06,
              size=10, color=c["head"], bold=True, line_spacing=1.05)
        by = y + q_h + 0.06 * scale
        # ★답은 색 상자 대신 왼쪽 세로선으로 구분한다(세원 리플렛 방식)
        _rect(slide, x + LEFT, by, 0.028, box_h, fill=c["badge"])
        _text(slide, b["a"], x + LEFT + 0.16, by, CW - 0.20, box_h,
              size=MIN_BODY, color=c["body"], line_spacing=1.25, min_size=MIN_BODY)
        y = by + box_h + gap


def _panel_cover(slide, schema, x, c, assets, pw=None):
    PANEL = pw if pw else globals()["PANEL"]  # 이 패널의 실제 폭(표지=100mm)
    ac = schema.get("academy") or {}
    name = _s(ac.get("name"), "우리학원")
    phone = _s(ac.get("phone"))
    # 슬로건: 최상위 slogan / academy.slogan / brand.slogan / closing.head 모두 수용
    slogan = _first(schema.get("slogan"), ac.get("slogan"),
                    (schema.get("brand") or {}).get("slogan"),
                    (schema.get("closing") or {}).get("head"),
                    default="학생의 성장을\n함께 만들어갑니다.")
    # ※미리 _wrap 하지 않는다 — 아래 _fit_slogan 이 크기와 줄바꿈을 함께 정한다.
    # 한마디: intro.promise / intro.body / intro.head / closing.cta
    promise = _first((schema.get("intro") or {}).get("promise"),
                     (schema.get("intro") or {}).get("body"),
                     (schema.get("intro") or {}).get("head"),
                     (schema.get("closing") or {}).get("cta"),
                     default="정확한 진단 · 맞춤 수업 · 꾸준한 관리")

    # 표지 패널 배경(딥/노랑)
    _rect(slide, x, 0, PANEL, H_IN, fill=c["cover_bg"])

    # 로고 자리(좌상단) + 학원명
    logo_ok = _photo_or_box(slide, (assets or {}).get("logo"),
                            x + 0.19, 0.29, 0.72, 0.67, c, "로고", cover_mode=False)
    _text(slide, name, x + 1.02, 0.39, PANEL - 1.2, 0.5,
          size=13.5, color=c["cover_ink"], bold=True, valign=MSO_ANCHOR.MIDDLE)

    # 슬로건 — 줄 수에 맞춰 높이를 잡고, 아래 요소를 그만큼 밀어낸다.
    # (3줄 슬로건 'From Reading to Your Voice — and Beyond!' 이
    #  1.5in 고정 박스를 넘겨 소개문·사진과 겹쳤다)
    # 슬로건 — 표지의 주인공. 작게 줄이지 않고 3줄까지만 쓴다.
    # 문장부호가 없어 문장 트리밍이 안 되는 경우까지 대비해,
    # <실제로 감싼 줄 수>를 보고 어절을 하나씩 덜어낸다.
    SL_MAX_LINES = 3
    SL_MIN = 18
    SL_MAX = SL_MAX_LINES * (25 * 1.18) / 72.0 + 0.12
    sl_words = _s(slogan).replace("\n", " ").split()

    def _fit_slogan(txt):
        return _fit_block(txt, PANEL - 0.38, SL_MAX, 25,
                          min_size=SL_MIN, line_spacing=1.18, pad=0.12)

    sl_size, sl_wrapped, sl_h = _fit_slogan(slogan)
    while sl_wrapped.count("\n") + 1 > SL_MAX_LINES and len(sl_words) > 1:
        sl_words.pop()
        sl_size, sl_wrapped, sl_h = _fit_slogan(" ".join(sl_words))
    # 어절이 하나뿐인데도 넘치면(끊을 공백이 없음) 글자 단위로 줄인다
    if sl_wrapped.count("\n") + 1 > SL_MAX_LINES and len(sl_words) == 1:
        one = sl_words[0]
        while len(one) > 4 and sl_wrapped.count("\n") + 1 > SL_MAX_LINES:
            one = one[:-2]
            sl_size, sl_wrapped, sl_h = _fit_slogan(one)
    sl_h = max(0.5, sl_h)
    _text(slide, sl_wrapped, x + 0.19, 1.35, PANEL - 0.38, sl_h,
          size=sl_size, color=c["cover_ink"], bold=True,
          line_spacing=1.18, fit=False)

    # 소개 한마디 — 표지는 "읽히는 크기"가 우선이다.
    # ★슬로건과 크기 차이가 너무 벌어지지 않게, 슬로건의 <절반>을 기준으로 잡는다.
    #   (25pt 슬로건 → 12.5pt 소개문. 슬로건이 줄면 소개문도 같은 비율로 준다.)
    # ① 문장 경계에서 먼저 줄이고(첫 문장 보존) ② 그래도 넘치면 어절을 덜어낸다.
    pr_y = 1.35 + sl_h + 0.16
    PR_RATIO = 0.5
    PR_BASE = round(sl_size * PR_RATIO * 2) / 2.0        # 0.5pt 단위
    PR_MIN = max(9.5, round(SL_MIN * PR_RATIO * 2) / 2.0)
    PR_MAX_LINES = 3
    per_line = max(8, int((PANEL - 0.38 - 0.10) / _char_w(PR_BASE)))
    # ★표지에는 완결된 문장만 올린다. 길이와 무관하게 조각은 제외.
    promise = _cover_copy(promise, per_line * PR_MAX_LINES)
    pr_max_h = PR_MAX_LINES * (PR_BASE * 1.16) / 72.0 + 0.10

    def _fit_promise(txt):
        return _fit_block(txt, PANEL - 0.38, pr_max_h, PR_BASE,
                          min_size=PR_MIN, line_spacing=1.16, pad=0.10)

    pr_size, pr_wrapped, pr_h = _fit_promise(promise)
    # 넘치면 <문장 단위로> 덜어낸다(어절을 잘라 미완성 문장을 만들지 않는다)
    _pr_sents = _sentences(pr_wrapped.replace("\n", " "))
    while pr_wrapped.count("\n") + 1 > PR_MAX_LINES and len(_pr_sents) > 1:
        _pr_sents.pop()
        pr_size, pr_wrapped, pr_h = _fit_promise(" ".join(_pr_sents))
    # 한 문장인데도 3줄을 넘으면 표지에서 뺀다(잘린 문장을 올리지 않는다)
    if pr_wrapped.count("\n") + 1 > PR_MAX_LINES:
        pr_wrapped, pr_h = "", 0.0
    pr_h = max(0.0, pr_h)
    if pr_wrapped:
        _text(slide, pr_wrapped, x + 0.19, pr_y, PANEL - 0.38, pr_h,
              size=pr_size, color=c["cover_ink"], line_spacing=1.16, fit=False)

    # 표지 사진 자리 — 소개문 아래에서 시작(겹침 방지). 공간이 부족하면 생략.
    ph_y = pr_y + pr_h + 0.14
    ph_h = 6.04 - ph_y
    if ph_h >= 1.10:
        _photo_or_box(slide, (assets or {}).get("cover") or (assets or {}).get("banner"),
                      x + 0.19, ph_y, PANEL - 0.38, ph_h, c, "학원 사진", cover_mode=True)

    # 과목·학년 뱃지
    subj = _subject_line(schema)
    if subj:
        _rect(slide, x + 0.21, 6.27, PANEL - 0.42, 0.34, fill=c["badge"], radius=True)
        _text(slide, subj, x + 0.21, 6.27, PANEL - 0.42, 0.34,
              size=9.5, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER,
              valign=MSO_ANCHOR.MIDDLE)

    # 표지 하단 강점 요약
    feats = _feature_list(schema)
    if feats:
        ftext = "  ·  ".join(f["title"] for f in feats[:3])
        _text(slide, ftext, x + 0.19, 6.75, PANEL - 0.38, 0.55,
              size=10, color=c["cover_ink"])

    _text(slide, f"{name}  ·  {phone}" if phone else name, x + 0.17, 7.92,
          PANEL - 0.2, 0.23, size=7.5, color=c["cover_ink"], align=PP_ALIGN.LEFT)


# ── 안쪽면 패널 ───────────────────────────────────────
def _panel_features(slide, schema, x, c, name):
    hb0 = _bar_head(slide, x, "우리 학원의 강점", c)
    feats = _feature_list(schema)
    # ★"성장을 만드는 N가지 학습 원칙"은 학원마다 똑같이 찍히는 자동 문구다.
    #   STEP6 가 만든 헤더가 있을 때만 쓰고, 없으면 헤더 없이 간다.
    Y0, Y_END = hb0, 7.80
    # 제목이 2줄 되는 항목이 있으므로(예: '1:1 레벨 진단 및 성장 경로 설계')
    # 항목마다 필요한 높이를 재서 배분한다.
    need = []
    for it in feats:
        tl = max(1, int(_dwidth(it["title"]) / 19) + 1)
        dl = max(1, int(_dwidth(it["desc"]) / 26) + 1) if it["desc"] else 0
        need.append(tl * 0.26 + dl * 0.20 + 0.20)
    total = sum(need)
    avail = Y_END - Y0
    # ★항목이 적으면 블록을 키워 면을 채운다(간격만 벌리면 성의 없어 보인다).
    grow = 1.0
    _has_aux = bool((schema.get("_auxAssign") or {}).get("features"))
    if (not _has_aux) and total > 0 and total < avail * 0.82:
        grow = min(1.45, (avail * 0.92) / total)
        need = [n * grow for n in need]
        total = sum(need)
    scale = min(1.0, avail / total) if total > 0 else 1.0
    t_size = (12.5 if grow > 1.15 else 12) if scale >= 0.92 else 11
    # ★간격으로 벌려 채우면 성의 없어 보인다 — 읽기 좋은 최소 간격만 준다.
    gap_extra = min(0.14 if _has_aux else 0.26,
                    max(0.0, (avail - total) / max(1, len(feats))))
    y = Y0
    for i, it in enumerate(feats):
        blk = need[i] * scale
        tl = max(1, int(_dwidth(it["title"]) / 19) + 1)
        t_h = max(0.26, tl * 0.26 * scale)
        _chip(slide, x + LEFT, y, CHIP, f"{i+1:02d}", c)
        _text(slide, it["title"], x + TEXT_X, y, TEXT_W, t_h,
              size=t_size, color=c["title"], bold=True, line_spacing=1.08)
        _text(slide, it["desc"], x + TEXT_X, y + t_h + 0.03, TEXT_W,
              max(0.20, blk - t_h - 0.09), size=9.5, color=c["body"],
              line_spacing=1.06)
        # ★마지막 항목 뒤에도 간격을 더해 y 가 밀려 보조 카드가 안 들어갔다.
        y += blk + (gap_extra if i < len(feats) - 1 else 0.0)

    # ★강점 4개만 있으면 하단 2.5in 가 통째로 비었다.
    #   간격을 벌리지 않고 보조 카드(자주 묻는 질문 등)로 채운다.
    _aux = _aux_choose(schema, "features", ("faq", "rules", "specials", "schedule"))
    if _aux and y + 0.92 < 7.56:
        _draw_aux_card(slide, schema, x, c, _aux, y + 0.34, 7.56)


# ── 교육 과정 상세를 학년별로 쪼갠다 ────────────────────
# ★'중1 …다. 중2 …다. 중3 …다.' 가 한 덩어리로 찍혀
#   학년 구분 없이 빽빽하게 보였다. 학년 표기 앞에서 줄을 나눈다.
_GRADE_SPLIT = re.compile(r"(?=(?:초|중|고)\s?[1-6](?:\s|·|,))")
_GRADE_HEAD = re.compile(r"^((?:초|중|고)\s?[1-6])\s*[·,]?\s*(.*)$", re.S)

# ★생성 단계에서 다음 구간 표시로 붙는 '[중등]' 같은 꼬리표가 본문에 남아
#   "…읽기 루틴을 만든다. [중등]" 처럼 인쇄물에 그대로 찍혔다.
_GRP_TAG = re.compile(r"\s*\[\s*(?:예비)?\s*(?:초등|중등|고등|초|중|고)\s*(?:부)?\s*\]\s*")


def _grade_rows(detail):
    """상세 텍스트 → [(학년라벨, 내용)] . 학년 표기가 없으면 [("", 전체)]."""
    t = _GRP_TAG.sub(" ", _s(detail)).strip()
    if not t:
        return []
    # 줄바꿈이 이미 있으면 그것을 우선 존중한다
    base = [ln.strip() for ln in t.split("\n") if ln.strip()]
    out = []
    for ln in base:
        parts = [q.strip() for q in _GRADE_SPLIT.split(ln) if q.strip()]
        if len(parts) <= 1:
            m = _GRADE_HEAD.match(ln)
            out.append((m.group(1), m.group(2)) if m else ("", ln))
            continue
        for q in parts:
            m = _GRADE_HEAD.match(q)
            out.append((m.group(1), m.group(2)) if m else ("", q))
    return [(g, v) for g, v in out if v or g]


def _panel_curriculum(slide, schema, x, c):
    cur = _curriculum_list(schema)
    if not cur:
        # 커리큘럼 없으면 패널 라벨·헤더 자체를 안 그림(빈 잔존 없음)
        return
    hb0 = _bar_head(slide, x, "교육 과정", c)
    cin = 0.18  # 카드 내부 여백
    y = hb0
    # 카드 높이를 내용 길이로 산출하고, 넘치면 비례 축소(글자 안 자름)
    # 아래에 특강 및 기타 수업 블록이 붙으면 그 위에서 멈춘다.
    _has_sp = bool(_items(schema.get("specials")))
    # ★아래에 <수업 요일·시간>과 특강 및 기타 수업이 붙으므로 그 위에서 멈춘다
    _has_sch = bool(_items(schema.get("schedule")))
    # ★4.90 으로 묶어 학년 행이 시간표와 겹쳤다. 특강은 이 면에서 빠졌으므로
    #   시간표 자리(약 1.9in)만 남기고 나머지를 커리큘럼에 준다.
    # 시간표(제목+2행)에 약 1.3in 을 남긴다 — 중등부 행이 잘렸다.
    Y_END = (6.05 if _has_sch else 7.78)
    txt_w = CW - cin * 2
    # ★렌더는 _need_h(실측)로 그리는데 계획은 scale 로만 줄여서
    #   내용이 지면을 넘고 시간표가 밀려났다. 여기서 <크기>를 정해 둔다.
    _CSZ = BODY
    def _cur_total(csz):
        tot = 0.0
        for _it in cur:
            _nm = _s(_it["name"]); _ds = _s(_it.get("desc"))
            hh = (_need_h(_nm, CW, 10.5, pad=0.02) if _nm else 0.0) + 0.04
            if _ds:
                hh += _need_h(_ds, CW, 10.5, pad=0.02) + 0.03
            for _g, _v in _grade_rows(_it.get("detail")):
                hh += _need_h(_v, (CW - 0.52) if _g else CW, csz, pad=0.02) + 0.07
            tot += hh + 0.30
        return tot
    # ★8pt 까지 떨어져 이 면만 유독 작았다 → 9.5pt 아래로는 줄이지 않는다.
    while _CSZ > 9.5 and _cur_total(_CSZ) > (Y_END - y):
        _CSZ -= 0.5
    plan = []
    for it in cur:
        d_lines = max(1, int(_dwidth(it.get("desc") or "") / 20) + 1) if it.get("desc") else 0
        _dt = _s(it.get("detail"))
        # ★학년별로 행이 나뉘므로 행마다 높이를 잡는다(한 덩어리 계산은 어긋났다).
        _gr = _grade_rows(_dt)
        t_h = 0.0
        for _g, _v in _gr:
            _w = (CW - 0.52) if _g else CW
            t_h += _need_h(_v, _w, _CSZ, pad=0.02) + 0.07
        h = 0.56 + d_lines * 0.24 + t_h + 0.16
        # ★1.05in 최소높이를 두면 내용 2줄짜리 초등 카드가 과하게 커져
        #   중등 카드와 균형이 깨졌다(경희궁). 내용만큼만 준다.
        plan.append({"it": it, "h": max(0.82, h), "d": d_lines, "t": t_h})
    GAP = 0.35
    total = sum(p["h"] for p in plan) + GAP * max(0, len(plan) - 1)
    avail = Y_END - y
    scale = min(1.0, avail / total) if total > 0 else 1.0
    gap = GAP * scale

    # ★세원 리플렛 방식 — 카드 테두리·pill 을 걷어내고
    #   단계명(컬러 굵게) + 목표 + 상세를 얇은 구분선으로만 나눈다.
    for idx, p in enumerate(plan):
        it, ch = p["it"], p["h"] * scale
        if idx:
            _rect(slide, x + LEFT, y, CW, 0.010, fill=c["card_line"])
            y += 0.14
        nm = _s(it["name"])
        nh = _need_h(nm, CW, 10.5, pad=0.02) if nm else 0.0
        if nm:
            _text(slide, nm, x + LEFT, y, CW, nh,
                  size=10.5, color=c["label"], bold=True, fit=False)
        cy = y + nh + 0.04
        if it.get("desc"):
            dh = _need_h(_s(it["desc"]), CW, 10.5, pad=0.02)
            _text(slide, it["desc"], x + LEFT, cy, CW, dh,
                  size=10.5, color=c["title"], bold=True, min_size=MIN_BODY)
            cy += dh + 0.03
        if it.get("detail"):
            GW = 0.52          # 학년 라벨 칼럼 폭
            for _g, _v in _grade_rows(it["detail"]):
                if _g:
                    _vh = _need_h(_v, CW - GW, _CSZ, pad=0.02)
                    _text(slide, _g, x + LEFT, cy, GW - 0.06, 0.26,
                          size=_CSZ, color=c["badge"], fit=False)
                    _text(slide, _v, x + LEFT + GW, cy, CW - GW, _vh,
                          size=_CSZ, color=c["body"], min_size=MIN_BODY,
                          line_spacing=1.32)
                    cy += _vh + 0.07
                else:
                    _vh = _need_h(_v, CW, _CSZ, pad=0.02)
                    _text(slide, _v, x + LEFT, cy, CW, _vh,
                          size=_CSZ, color=c["body"], min_size=MIN_BODY,
                          line_spacing=1.32)
                    cy += _vh + 0.07
        y = max(cy, y + ch) + 0.10

    # ── 수업 요일 · 시간 (교육 과정 아래) ────────────────
    # ★강좌가 100개 넘는 학원은 표를 실을 수 없다(더케이 104개).
    #   학년별 대표 운영 패턴을 요약해 교육 과정 바로 아래에 둔다.
    sched = _items(schema.get("schedule"))
    if sched:
        # ★커리큘럼 끝과 간격이 벌어져 시간표 행이 잘렸다 → 바짝 붙인다.
        sy0 = max(y + 0.02, 5.10)
        _rect(slide, x + LEFT, sy0, CW, 0.012, fill=c["card_line"])
        _label(slide, x, "수업 요일 · 시간", c, y=sy0 + 0.10)
        sy = sy0 + 0.48
        LIM = 7.80   # ★특강 블록을 뺐으므로 하단까지 쓴다(중등부가 잘렸다)
        for it in sched[:3]:
            k = _s(it.get("k")); v = _s(it.get("v"))
            if not v: continue
            v = " · ".join([p_.strip() for p_ in v.split("·") if p_.strip()][:2])
            vh = _need_h(v, CW - 0.82, 9.5, pad=0.02)
            if sy + vh > LIM:
                break
            # ★MIN_BODY(8pt)로 찍혀 시간표만 눈에 띄게 작았다 → 본문 크기로.
            _text(slide, k, x + LEFT, sy, 0.78, max(0.20, vh),
                  size=9.5, color=c["label"], fit=False)
            _text(slide, v, x + LEFT + 0.82, sy, CW - 0.82, vh,
                  size=9.5, color=c["body"], min_size=9.0, line_spacing=1.22)
            sy += vh + 0.09

    # ★특강을 교육 과정 면 하단 자투리(약 1.1in)에 밀어 넣으니
    #   5개 중 1개만 나오고 '수업 요일·시간'까지 잘렸다.
    #   특강은 전용 면이나 '한눈에 보기' 면에서 전부 보여준다 → 여기서는 싣지 않는다.


def _panel_management(slide, schema, x, c):
    steps = _mgmt_list(schema)
    hb0 = _bar_head(slide, x, "학습 관리", c)
    sub_y = hb0 - 0.10
    if steps:
        _text(slide, _first(_s((schema.get("management") or {}).get("subhead")),
                            default=f"학습 관리 {len(steps)} STEP"),
              x + LEFT, sub_y, CW, 0.27, size=9.5, color=c["body"], bold=True)
    # 카드 높이를 내용 길이로 산출(고정 0.82in 이 길면 잘렸다).
    # 실적 블록이 아래에 붙으므로 그 위까지만 쓰고 여유를 둔다.
    Y0 = max(0.95, sub_y + 0.30)
    # ★실적이 <자기 면>을 받았으면 여기 또 넣지 않는다(같은 내용이 두 번 나왔다).
    #   자리를 못 받았을 때만 하단에 이어 싣는다.
    _own = schema.get("_ownPanel") or set()
    has_ach = bool(_achievement_lines(schema)[1]) and ("achievements" not in _own)
    Y_END = 5.72 if has_ach else 7.80
    plan = []
    for txt in steps:
        lines = max(1, int(_dwidth(txt) / 22) + 1)
        plan.append({"txt": txt, "h": max(0.62, 0.26 + lines * 0.22)})
    GAP = 0.20
    total = sum(p["h"] for p in plan) + GAP * max(0, len(plan) - 1)
    avail = Y_END - Y0
    scale = min(1.0, avail / total) if total > 0 else 1.0
    gap = GAP * scale
    # 간격은 최소치만(과도하게 벌리면 성의 없어 보인다)
    gap += min(0.24, max(0.0, (avail - total) / max(1, len(plan))))

    # ★세원 리플렛 방식 — 색 카드를 깔지 않고 번호 + 글로만 정리한다.
    y = Y0
    for i, p in enumerate(plan):
        bh = p["h"] * scale
        _chip(slide, x + LEFT, y + 0.02, CHIP, str(i + 1), c)
        _txt = _s(p["txt"])
        _k, _v = "", _txt
        _m = re.match(r"^\s*([^:：—–]{2,14})\s*[:：—–]\s*(.+)$", _txt)
        if _m:
            _k, _v = _m.group(1).strip(), _m.group(2).strip()
        if _k:
            kh = _need_h(_k, TEXT_W, 10.5, pad=0.02)
            _text(slide, _k, x + TEXT_X, y, TEXT_W, kh,
                  size=10.5, color=c["title"], bold=True, min_size=MIN_BODY, fit=False)
            vh = max(0.20, bh - kh - 0.02)
            _text(slide, _v, x + TEXT_X, y + kh + 0.02, TEXT_W, vh,
                  size=BODY, color=c["body"], min_size=MIN_BODY, line_spacing=1.34)
        else:
            _text(slide, _v, x + TEXT_X, y, TEXT_W, bh,
                  size=BODY, color=c["body"], min_size=MIN_BODY, line_spacing=1.34)
        y += bh + gap

    # 실적 — ★자기 면을 받았으면 여기 또 넣지 않는다(has_ach 로 판정).
    head, lines = _achievement_lines(schema)
    if has_ach and lines:
        _rect(slide, x + LEFT, 5.92, CW, 0.012, fill=c["card_line"])  # 구분선
        _text(slide, _de_emoji(head), x + LEFT, 6.10, CW, 0.28,
              size=10.5, color=c["head"], bold=True)
        # 줄마다 랩된 줄 수를 세어 실제 필요한 높이를 준다(푸터 7.90 위까지).
        body = "\n".join("· " + l for l in lines)
        _text(slide, body, x + LEFT, 6.42, CW, 1.44,
              size=MIN_BODY, color=c["body"], line_spacing=1.20, min_size=MIN_BODY)


# ── 예비 패널(대체 섹션) ──────────────────────────────
def _fill_leftover(slide, schema, x, c, y0, y_end):
    """면 아래에 남는 자리를 <그 학원에 있는 것>으로 채운다.
       학원마다 자료가 달라 한 가지로 정할 수 없다 →
       ①아직 못 실은 실적이 있으면 실적을 더 싣고
       ②없으면 약도, ③약도도 없으면 학원 사진을 넣는다.
       ④셋 다 없으면 그냥 비워 둔다(억지로 채우지 않는다)."""
    avail = y_end - y0
    if avail < 1.10:
        return y0
    _as = schema.get("assets") or {}

    # ① 남은 실적 줄
    rest = schema.get("_achRest") or []
    if rest:
        _text(slide, "그 밖의 성장 사례", x + LEFT, y0, CW, 0.24,
              size=9.5, color=c["label"], fit=False)
        yy = y0 + 0.28
        card_h = min(avail - 0.34, 0.24 + len(rest) * 0.30)
        _rect(slide, x + LEFT, yy, CW, card_h, fill=c["soft"], radius=True)
        iy = yy + 0.12
        for line in rest:
            if iy + 0.28 > yy + card_h:
                break
            _text(slide, _s(line), x + LEFT + 0.14, iy, CW - 0.28, 0.26,
                  size=9.5, color=c["title"], min_size=8.0)
            iy += 0.30
        schema["_achRest"] = []
        return yy + card_h + 0.16

    # ② 약도 → ③ 학원 사진
    img = _s(_as.get("map")) or ""
    cap = "오시는 길"
    if not img:
        _ph = _as.get("photos")
        if isinstance(_ph, list) and _ph:
            img, cap = _s(_ph[0]), "학원 전경"
        else:
            img = _s(_as.get("cover")) or _s(_as.get("banner"))
            cap = "학원 전경"
    if not img:
        return y0
    _text(slide, cap, x + LEFT, y0, CW, 0.24,
          size=9.5, color=c["label"], fit=False)
    ih = min(avail - 0.34, CW * 0.72)
    _photo_or_box(slide, img, x + LEFT, y0 + 0.28, CW, ih, c, cap, cover_mode=True)
    return y0 + 0.28 + ih + 0.16


def _panel_achievements(slide, schema, x, c):
    """주요 실적 — ★3단 접지에서 <뒷면>이 되는 면.
    손에 들었을 때 바로 보이므로 실적은 <한 줄씩 제목만> 싣고,
    남는 자리에 연락처와 약도 QR 을 둔다(설명까지 넣으면 줄이 길어져 자리가 없다)."""
    ach = schema.get("achievements") or {}
    if not isinstance(ach, dict):
        ach = {"items": ach if isinstance(ach, list) else []}
    ac = schema.get("academy") or {}

    # ★to_schema 가 실적을 <대입/고입/내신/성장> 유형으로 분류해 두는데
    #   _items() 로 평탄화해 유형을 버리고 긴 문장만 나열했다.
    #   시크릿 리플렛처럼 유형별 소제목 + 항목으로 묶어 훑어보게 만든다.
    groups = []
    for g in (ach.get("groups") or []):
        if not isinstance(g, dict):
            continue
        gname = _de_emoji(_s(g.get("name")))
        gitems = []
        for it in (g.get("items") or []):
            if not isinstance(it, dict):
                if isinstance(it, str) and it.strip():
                    gitems.append((it.strip(), ""))
                continue
            ttl = _de_emoji(_s(it.get("title") or it.get("name")))
            chg = _de_emoji(_s(it.get("change")))
            if ttl or chg:
                gitems.append((ttl, chg))
        if gitems:
            groups.append((gname, gitems))
    if not groups:
        flat = []
        for it in _items(ach)[:8]:
            if not isinstance(it, dict):
                continue
            ttl = _de_emoji(_s(it.get("title") or it.get("name")))
            chg = _de_emoji(_s(it.get("change")))
            if ttl or chg:
                flat.append((ttl, chg))
        if not flat:
            _h, lines = _achievement_lines(schema)
            flat = [(ln, "") for ln in lines[:8]]
        if flat:
            groups = [("", flat)]

    hb0 = _bar_head(slide, x, "주요 실적", c)
    y = hb0

    QW, FOOT = 0.56, 7.74
    qy = FOOT - QW - 0.06                 # QR 줄
    cy = qy - 0.94                        # 연락처 블록 시작
    LIMIT = cy - 0.18                     # 실적은 여기까지만
    AV = LIMIT - y

    GH, PADC, GAPI, GAPG = 0.26, 0.14, 0.09, 0.15   # 소제목·카드여백·항목·그룹 간격
    def _measure(sz, arrow=True):
        out, tot = [], 0.0
        for gname, items in groups:
            hs = []
            for ttl, chg in items:
                line = (ttl + ("  →  " + chg if (chg and arrow) else
                               ("  " + chg if chg else "")))
                hs.append((line, _need_h(line, CW - PADC * 2 - 0.12, sz, pad=0.02)))
            card = PADC + sum(h for _l, h in hs) + GAPI * max(0, len(hs) - 1) + PADC
            out.append((gname, hs, card))
            tot += (GH if gname else 0.0) + card + GAPG
        return out, (tot - GAPG if out else 0.0)

    _sz = 10.0
    meas, tot = _measure(_sz)
    while tot > AV and _sz > 8.0:
        _sz -= 0.5
        meas, tot = _measure(_sz)

    _left = []
    for gname, hs, card_h in meas:
        if y + (GH if gname else 0.0) + card_h > LIMIT + 0.12:
            _left.extend(l for l, _h in hs)      # 못 실은 실적은 뒤에서 쓴다
            continue
        if gname:
            _text(slide, gname, x + LEFT, y, CW, 0.26,
                  size=9.5, color=c["label"], fit=False)
            y += GH
        _rect(slide, x + LEFT, y, CW, card_h, fill=c["soft"], radius=True)
        iy = y + PADC
        for line, hh in hs:
            _text(slide, line, x + LEFT + PADC, iy, CW - PADC * 2, hh,
                  size=_sz, color=c["title"], min_size=8.0, line_spacing=1.22)
            iy += hh + GAPI
        y += card_h + GAPG

    # ★남는 자리 — 학원 상황에 맞게 채운다(실적 우선, 없으면 약도·사진)
    schema["_achRest"] = _left
    if LIMIT - y >= 1.10:
        y = _fill_leftover(slide, schema, x, c, y, LIMIT)

    # ── 연락처 ─────────────────────────────────────────
    _rect(slide, x + LEFT, cy, CW, 0.012, fill=c["card_line"])
    phone = _s(ac.get("phone"))
    _text(slide, f"전화  {phone}" if phone else "전화 상담", x + LEFT, cy + 0.16,
          CW, 0.28, size=11, color=c["head"], bold=True, fit=False)
    loc = _first(ac.get("location"), ac.get("address_short"), default="")
    if loc:
        _avail = qy - 0.14 - (cy + 0.50)
        if _avail >= 0.20:
            lh = min(_need_h(loc, CW, MIN_BODY, pad=0.02), _avail)
            _maxln = max(1, int((lh - 0.04) / (MIN_BODY * 1.30 / 72.0)))
            _per = max(6, int((CW - 0.10) / (MIN_BODY * 0.0148)))
            _ls = _wrap(loc, _per).split("\n")
            if len(_ls) > _maxln:
                loc = " ".join(_ls[:_maxln]).rstrip(" ,·") + "…"
            _text(slide, loc, x + LEFT, cy + 0.50, CW, lh, fit=False,
                  size=MIN_BODY, color=c["body"], line_spacing=1.30)

    # ── 약도 QR (네이버 · 카카오) ────────────────────────
    for i, lb in enumerate(("네이버", "카카오")):
        qx = x + LEFT + i * (CW / 2)
        _rect(slide, qx, qy, QW, QW, fill="FFFFFF", line=c["card_line"], radius=True)
        _text(slide, "QR", qx, qy, QW, QW, size=MIN_BODY, color=c["footer"],
              align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fit=False)
        _text(slide, lb, qx + QW + 0.06, qy + 0.16, CW / 2 - QW - 0.10, 0.26,
              size=MIN_BODY, color=c["body"], fit=False)


def _panel_rules(slide, schema, x, c):
    """학원 규정 패널(예비)."""
    rules = _items(schema.get("rules")) or _items(schema.get("policy"))
    hb = _bar_head(slide, x, "학원 규정", c)
    y = hb
    for it in rules[:6]:
        k = _s(it.get("k") or it.get("title") or it.get("name"))
        v = _s(it.get("v") or it.get("desc") or it.get("body"))
        if not (k or v):
            continue
        if k:
            _text(slide, k, x + LEFT, y, CW, 0.26,
                  size=10, color=c["title"], bold=True)
        _text(slide, v, x + LEFT, y + 0.28, CW, 0.46,
              size=8.5, color=c["body"])
        y += 0.86


def _panel_specials(slide, schema, x, c):
    """특강 및 기타 수업 단독 패널(예비)."""
    sp = _items(schema.get("specials"))
    hb0 = _bar_head(slide, x, "특강 및 기타 수업", c)
    y = hb0
    for i, it in enumerate(sp[:5]):
        t = _s(it.get("title") or it.get("name"))
        d = _s(it.get("desc") or it.get("description"))
        _chip(slide, x + LEFT, y + 0.05, CHIP, f"{i+1:02d}", c)
        _text(slide, t, x + TEXT_X, y, TEXT_W, 0.28,
              size=10.5, color=c["title"], bold=True)
        _text(slide, d, x + TEXT_X, y + 0.30, TEXT_W, 0.48,
              size=8.5, color=c["body"])
        y += 0.94


# ── 섹션 풀: 데이터가 있는 섹션만 패널에 배정 ──────────
# ★전화번호만 있어도 True 를 주던 탓에, 절차 데이터가 없는데도 면이 열리고
#   하드코딩 문구로 채워졌다. 실제 절차가 있을 때만 전용 면을 준다.
def _has_admission(s):    return bool(_admission_list(s))
def _has_faq(s):          return bool(_faq_list(s))
def _has_features(s):     return bool(_feature_list(s))
def _has_curriculum(s):   return bool(_curriculum_list(s))
def _has_management(s):   return bool(_mgmt_list(s))
def _has_achievements(s): return bool(_achievement_lines(s)[1])
def _has_specials(s):     return bool(_items(s.get("specials")))
def _has_rules(s):        return bool(_items(s.get("rules")) or _items(s.get("policy")))

# 우선순위 순. 앞 섹션이 비면 뒤 섹션이 자동으로 앞 면으로 당겨진다(A+C).
# 표지(_panel_cover)는 위치가 고정이라 이 풀에서 제외한다.
# ★자리는 5칸인데 실적이 6번이라 <실적 면이 통째로 빠졌다>(더포스둔산 12건 보유).
#   학부모가 가장 먼저 보는 정보이므로 강점 다음으로 올린다.
#   FAQ 는 지면을 많이 먹고 상담에서 답할 수 있어 뒤로 보낸다.
# ★3단 접지에서 읽는 순서
#     표지(바깥 우) → 강점(바깥 좌) → 안쪽 실적·교육과정·학습관리 → 상담·예약(바깥 중, 뒷면)
#   예전에는 상담·예약이 바깥 좌라 안으로 접혀 들어가고 강점이 뒷면이 됐다.
#   연락처·상담 안내는 뒷면에 있어야 손에 들었을 때 바로 보인다.
SECTION_POOL = [
    ("admission",    _has_admission,    _panel_admission),      # 바깥 좌 (안으로 접히는 면)
    ("achievements", _has_achievements, _panel_achievements),   # 바깥 중 = 뒷면 (실적+연락처+QR)
    ("features",     _has_features,     None),                  # 안쪽 좌 — name 인자 필요
    ("curriculum",   _has_curriculum,   _panel_curriculum),
    ("faq",          _has_faq,          None),   # name 인자 필요
    #   ★FAQ 는 답변이 길어 카드로 접으면 1~2개만 들어간다 → 전용 면을 준다.
    ("management",   _has_management,   _panel_management),
    ("specials",     _has_specials,     _panel_specials),
    ("rules",        _has_rules,        _panel_rules),
]


# ── 면 균형 배치 ────────────────────────────────────────
# ★어떤 면은 빽빽하고 어떤 면은 하단이 통째로 비었다(상담·예약 면).
#   섹션마다 필요한 높이를 먼저 재서, 넘치는 섹션은 <다음 면으로 나누고>
#   남는 면에는 짧은 섹션을 이어 붙인다.
PANEL_H_USABLE = 6.85     # 라벨·헤더·푸터를 뺀 실제 본문 높이(inch)

def _sec_lines(schema, key):
    """섹션 → [(제목, 설명)] 목록. 분량 계산과 분할에 함께 쓴다."""
    if key == "features":
        return [(i["title"], i["desc"]) for i in _feature_list(schema)]
    if key == "curriculum":
        return [(i["name"], f"{i.get('desc') or ''} {i.get('detail') or ''}".strip())
                for i in _curriculum_list(schema)]
    if key == "management":
        return [(_s(l), "") for l in _mgmt_list(schema)]
    if key == "faq":
        return [(i["q"], i["a"]) for i in _faq_list(schema)]
    if key == "admission":
        return [(_s(i.get("title") or i.get("k")), _s(i.get("desc") or i.get("v")))
                for i in _admission_list(schema)]
    if key == "achievements":
        return [(_s(l), "") for l in _achievement_lines(schema)[1]]
    if key == "specials":
        return [(_s(i.get("title") or i.get("name")),
                 _s(i.get("desc") or i.get("description") or i.get("text")))
                for i in _items(schema.get("specials"))]
    if key == "rules":
        src = _items(schema.get("rules")) or _items(schema.get("policy"))
        return [(_s(i.get("k") or i.get("title")), _s(i.get("v") or i.get("desc"))) for i in src]
    return []

def _sec_height(rows):
    """[(제목, 설명)] → 예상 높이(inch)."""
    h = 0.0
    for t, d in rows:
        h += _need_h(t, CW - 0.10, BODY, pad=0.02)
        if d:
            h += _need_h(d, CW - 0.60, MIN_BODY, pad=0.02) + 0.04
        h += 0.14
    return h

def _plan_panels(schema, cand, slots):
    """cand(섹션 후보) → slots 개 면에 배치한 계획.
    ★면은 5칸인데 섹션이 8개면 6번째부터 <통째로 버려졌다>(특강이 사라진 원인).
      앞 (slots-1) 칸은 전용 면으로, 남는 섹션은 전부 마지막 면에 모아 싣는다.
    반환: [(kind, payload)]  kind='sec' → (key, draw_fn) / kind='combo' → parts
    """
    if len(cand) <= slots:
        if isinstance(schema, dict):
            schema["_ownPanel"] = {k for k, _f in cand}
            schema["_auxAssign"] = {}
        return [("sec", c) for c in cand]
    own = cand[:slots - 1]
    rest = [k for k, _fn in cand[slots - 1:]]
    own_keys = {k for k, _f in own}
    # ★남는 섹션을 전부 '한눈에 보기' 한 면에 몰면 그 면만 빽빽해지고
    #   다른 면 하단은 빈다. 짧은 섹션은 여유 있는 면 아래에 붙인다.
    #     학원 규정 → 상담·예약 면 / 자주 묻는 질문 → 강점 면
    #   시크릿 리플렛처럼 주요 섹션은 전용 면을 갖고,
    #   짧은 정보(규정·FAQ)는 여유 있는 면 <하단 카드>로 접어 넣는다.
    aux = {}
    for host, guest in (("admission", "rules"), ("features", "specials"),
                        ("curriculum", "rules"), ("features", "rules")):
        if host in aux or guest not in rest or host not in own_keys:
            continue
        aux[host] = guest
        rest.remove(guest)
    plan = [("sec", c) for c in own]
    parts = [(t, ls) for t, ls in _combo_parts(schema, rest) if ls]
    if parts:
        plan.append(("combo", parts))
    if isinstance(schema, dict):
        schema["_ownPanel"] = own_keys
        schema["_auxAssign"] = aux
    return plan


def _pick_sections(schema, slots, name):
    """채울 수 있는 섹션을 우선순위대로 slots개 고른다.
    - 데이터 없는 섹션은 건너뛰고 뒤 섹션을 당겨온다 → 빈 면이 생기지 않는다.
    - 후보가 slots보다 적으면 있는 만큼만(억지 확대·빈칸 없음).
    반환: [(key, draw_fn(slide, x, c)), ...]"""
    picked = []
    for key, has, fn in SECTION_POOL:
        if len(picked) >= slots:
            break
        try:
            if not has(schema):
                continue
        except Exception:
            continue
        if key == "features":
            draw = lambda sl, x, c: _panel_features(sl, schema, x, c, name)
        elif key == "faq":
            draw = lambda sl, x, c: _panel_faq(sl, schema, x, c, name)
        else:
            draw = (lambda f: lambda sl, x, c: f(sl, schema, x, c))(fn)
        picked.append((key, draw))
    # ★자기 면을 받은 섹션을 여기서 확정한다.
    #   build() 에서 기록하면 패널이 이미 그려진 뒤라 늦다(실적이 두 번 나왔다).
    if isinstance(schema, dict):
        schema["_ownPanel"] = {k for k, _f in picked}
    return picked


def _panel_combo(slide, schema, x, c, parts):
    """여러 섹션을 한 면에 모아 그린다.
    ★예전에는 전부 '· 한 줄' 불릿으로 흘려 성의 없어 보였다.
      섹션마다 [라벨 + 연한 카드]를 주고, 카드 안에서 제목(굵게)과
      설명(본문)을 구분해 다른 면과 위계를 맞춘다."""
    hb = _bar_head(slide, x, "학원 안내", c)
    Y0, Y_END = max(0.95, hb + 0.24), 7.62

    # (제목, 설명) 으로 분해 — _combo_parts 는 '제목 — 설명' 으로 넘긴다
    blocks = []
    for title, lines in parts:
        rows = []
        for l in lines:
            t = _s(l)
            if not t:
                continue
            m = re.split(r"\s+[—–]\s+", t, maxsplit=1)
            if len(m) == 2:
                rows.append((m[0].strip(), m[1].strip()))
            else:
                rows.append((t, ""))
        if rows:
            blocks.append((title, rows))
    if not blocks:
        return

    # ★블록 제목 아래에 부제 한 줄을 넣는다. 다른 면은 헤더 바에 부제가 붙는데
    #   이 면만 제목뿐이라 위계가 어긋나 보였다.
    LAB_H, PAD_T, PAD_B, GAP_ROW, GAP_SEC = 0.46, 0.18, 0.14, 0.12, 0.26
    IW = CW - 0.44
    AVAIL = Y_END - Y0

    def _plan(bl, size, nodesc=()):
        """블록 목록 → (측정치, 전체높이).
        nodesc 에는 <블록 인덱스> 또는 <(블록, 행) 쌍> 이 들어간다.
        ★예전에는 블록 인덱스만 받아 그 블록의 설명을 통째로 접었다.
          FAQ와 학습 관리 둘만 있으면 한쪽은 반드시 라벨만 남아
          '수업 방식 / 근거 확인 / 내신 대비' 처럼 읽을 내용이 사라졌다.
          이제 행 단위로 접어 앞 항목의 설명을 지킨다."""
        out, total = [], 0.0
        for bi, (ti, rows) in enumerate(bl):
            hs = []
            for ri, (t, d) in enumerate(rows):
                _qa = (bl[bi][0] == "자주 묻는 질문")
                _skip = (bi in nodesc) or ((bi, ri) in nodesc)
                th = _need_h(("Q. " + t) if _qa else t, IW, size, pad=0.02) if t else 0.0
                dh = (_need_h(d, IW - (0.20 if _qa else 0.0), size, pad=0.02)
                      if (d and not _skip) else 0.0)
                hs.append((max(0.21, th), dh))
            _qgap = 0.06 if (bl[bi][0] == "자주 묻는 질문") else 0.0
            card = PAD_T + sum(a + (b + 0.02 if b else 0.0) for a, b in hs) \
                   + (GAP_ROW + _qgap) * max(0, len(rows) - 1) + PAD_B
            out.append((ti, rows, hs, card))
            total += LAB_H + card + GAP_SEC
        return out, (total - GAP_SEC if bl else 0.0)

    # ① 글자 크기를 줄여 맞춘다
    size, nodesc = BODY, set()
    while True:
        meas, total = _plan(blocks, size, nodesc)
        if total <= AVAIL or size <= MIN_BODY:
            break
        size -= 0.5
    # ② 그래도 넘치면 <가장 긴 블록부터 하나씩> 설명을 접는다.
    #    ★전부 한꺼번에 접으니 제목만 남고 하단 2.4in 가 비었다.
    #   ★가장 긴 블록부터 접으니 <FAQ 답변>이 먼저 날아갔다.
    #     답변·설명이 정보의 핵심인 블록은 마지막에 접는다.
    # ★설명이 곧 내용인 블록은 마지막까지 접지 않는다.
    #   예전에는 FAQ만 보호해서, FAQ와 같은 면에 놓인 <학습 관리>가 먼저 접혔다.
    #   그 결과 "과제 / 오답 관리 / 성취도 점검 / 학부모 소통" 라벨만 남아
    #   무엇을 어떻게 관리하는지 알 수 없는 면이 되었다.
    _FOLD_LAST = ("자주 묻는 질문", "학습 관리", "특강 및 기타 수업")

    def _fold_targets():
        """접을 수 있는 (블록, 행) 목록. 뒤쪽 행부터, 보호 블록은 나중에."""
        plain, keep = [], []
        for bi, (ti, rows) in enumerate(blocks):
            bucket = keep if ti in _FOLD_LAST else plain
            for ri in range(len(rows) - 1, -1, -1):
                t, d = rows[ri]
                if d and (bi, ri) not in nodesc and bi not in nodesc:
                    bucket.append((bi, ri))
        return plain + keep

    while total > AVAIL:
        _tg = _fold_targets()
        if not _tg:
            break
        nodesc.add(_tg[0])
        size = BODY
        while True:
            meas, total = _plan(blocks, size, nodesc)
            if total <= AVAIL or size <= MIN_BODY:
                break
            size -= 0.5
    # ③ 최후 — 여전히 넘치면 뒤 항목부터 덜어낸다(잘린 카드를 만들지 않는다)
    if total > AVAIL:
        trim = [(t, list(r)) for t, r in blocks]
        while total > AVAIL and sum(len(r) for _t, r in trim) > len(trim):
            longest = max(range(len(trim)), key=lambda i: len(trim[i][1]))
            if len(trim[longest][1]) <= 1:
                break
            trim[longest][1].pop()
            meas, total = _plan(trim, size, nodesc)

    y = Y0
    for ti, rows, hs, card_h in meas:
        _text(slide, ti, x + LEFT, y, CW, 0.26,
              size=10.5, color=c["label"], bold=True, fit=False)
        _sub_t = _SEC_SUB.get(ti, "")
        if _sub_t:
            _text(slide, _sub_t, x + LEFT, y + 0.24, CW, 0.20,
                  size=8.5, color=c["body"], fit=False)
        cy = y + LAB_H
        _rect(slide, x + LEFT, cy, CW, card_h, fill=c["soft"], radius=True)
        ry = cy + PAD_T
        # ★질문과 답이 같은 자리·같은 색이라 구분이 안 됐다.
        #   질문은 진한 색, 답은 컬러 세로바 + 들여쓰기로 확실히 나눈다.
        _is_qa = (ti == "자주 묻는 질문")
        _ind = 0.20 if _is_qa else 0.0
        for (t, d), (th, dh) in zip(rows, hs):
            if t:
                _text(slide, (("Q. " + t) if _is_qa else t),
                      x + LEFT + 0.22, ry, IW, th,
                      size=size, color=c["head"],
                      min_size=MIN_BODY, line_spacing=1.20)
                ry += th
            # ★설명을 접은 경우(dh=0) 높이 0 상자를 만들지 않는다.
            if d and dh:
                if _is_qa:
                    _rect(slide, x + LEFT + 0.24, ry + 0.04, 0.026, dh - 0.02,
                          fill=c["badge"])
                _text(slide, d, x + LEFT + 0.22 + _ind, ry + 0.02, IW - _ind, dh,
                      size=size, color=c["body"], min_size=MIN_BODY,
                      line_spacing=1.26)
                ry += dh + 0.02
            ry += GAP_ROW + (0.06 if _is_qa else 0.0)
        y = cy + card_h + GAP_SEC


def _combo_parts(schema, keys):
    """섹션 키 목록 → _panel_combo 가 쓸 [(소제목, [줄...])] 로 변환."""
    out = []
    for k in keys:
        if k == "achievements":
            head, lines = _achievement_lines(schema)
            out.append(("주요 실적", lines))
        elif k == "specials":
            # ★제목만 싣던 것을 '제목 — 설명' 으로 바꿔 무슨 수업인지 보이게 한다.
            _sl = []
            for i in _items(schema.get("specials")):
                _t = _s(i.get("title") or i.get("name"))
                _d = _s(i.get("desc") or i.get("description") or i.get("text"))
                if _t:
                    _sl.append(f"{_t} — {_d}" if _d else _t)
            out.append(("특강 및 기타 수업", _sl))   # 부제는 _COMBO_SUB 참조
        elif k == "management":
            out.append(("학습 관리", _mgmt_list(schema)))
        elif k == "faq":
            out.append(("자주 묻는 질문",
                        [f"{i['q']} — {i['a']}" for i in _faq_list(schema)]))
        elif k == "rules":
            src = _items(schema.get("rules")) or _items(schema.get("policy"))
            out.append(("학원 규정",
                        [f"{_s(i.get('k') or i.get('title'))} {_s(i.get('v') or i.get('desc'))}".strip()
                         for i in src]))
        elif k == "curriculum":
            out.append(("교육 과정",
                        [f"{i['name']} {i.get('desc') or ''} {i.get('detail') or ''}".strip()
                         for i in _curriculum_list(schema)]))
        elif k == "features":
            out.append(("우리 학원의 강점",
                        [f"{i['title']} — {i['desc']}" for i in _feature_list(schema)]))
        elif k == "admission":
            out.append(("상담 · 등록 안내",
                        [f"{i['title']} {i['desc']}".strip() for i in _admission_list(schema)]))
    return [(t, l) for t, l in out if l]


# ── 빌드 ──────────────────────────────────────────────
def _force_theme_font(prs, name=FONT):
    """테마의 majorFont/minorFont 를 지정 서체로 바꾼다.
    기본 테마는 Calibri(한글 글리프 없음)라서, 혹시 서체 지정이 누락된
    텍스트가 있으면 PowerPoint 가 임의의 한글 폰트로 대체해 버린다.
    테마까지 맞춰두면 그런 경우에도 Pretendard 로 떨어진다.
    ※theme1.xml 은 python-pptx 가 파싱하지 않는 일반 Part 라서
      blob(바이트)을 직접 치환한다."""
    try:
        for part in prs.part.package.iter_parts():
            if "theme" not in str(part.partname):
                continue
            xml = part.blob.decode("utf-8", errors="ignore")

            def _fix(block):
                block = re.sub(r'(<a:latin[^/>]*typeface=")[^"]*(")',
                               r"\1" + name + r"\2", block)
                block = re.sub(r'(<a:ea[^/>]*typeface=")[^"]*(")',
                               r"\1" + name + r"\2", block)
                block = re.sub(r'(<a:cs[^/>]*typeface=")[^"]*(")',
                               r"\1" + name + r"\2", block)
                return block

            for tag in ("majorFont", "minorFont"):
                m = re.search(r"<a:%s>.*?</a:%s>" % (tag, tag), xml, re.S)
                if m:
                    xml = xml[:m.start()] + _fix(m.group(0)) + xml[m.end():]
            part._blob = xml.encode("utf-8")
    except Exception:
        pass  # 테마 수정 실패는 치명적이지 않다(런 단위 지정이 이미 있음)


def build(schema: Dict[str, Any], palette: str = "sewon_teal",
          out: Union[str, os.PathLike, io.BytesIO, None] = None,
          assets: Optional[Dict[str, str]] = None):
    c = PALETTES.get(palette, PALETTES["sewon_teal"])

    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank = prs.slide_layouts[6]

    ac = schema.get("academy") or {}
    name = _s(ac.get("name"), "우리학원")
    phone = _s(ac.get("phone"))

    # assets: 스키마 내장 + 파라미터 병합(파라미터 우선)
    _assets = dict(schema.get("assets") or {})
    if isinstance(assets, dict):
        _assets.update(assets)

    # 접지 좌표 계산
    out_geom = _panel_geom(PANELS_OUTSIDE_MM)   # [(x,w)] 좌97·중100·우100
    in_geom  = _panel_geom(PANELS_INSIDE_MM)    # [(x,w)] 좌100·중100·우97

    # ── 섹션 배정: 표지 1면 + 콘텐츠 면 ──
    # 데이터 있는 섹션만 우선순위대로 채우고, 비면 뒤 섹션을 앞으로 당긴다.
    # → 빈 면이 생기지 않는다(A: 풀+동적배치, C: 예비섹션 승격).
    #
    # 3단 접지는 면 단위로만 쓸 수 있다. 바깥면 [섹션·섹션·표지] 3면,
    # 안쪽면 [섹션·섹션·섹션] 3면. 안쪽면은 3면이 다 찰 때만 만든다.
    # 따라서 실제로 쓸 수 있는 섹션 수는 2개(1장) 또는 5개(2장) 뿐이다.
    # 후보를 5개까지 뽑아보고, 3~4개면 상위 2개만 써서 1장으로 낸다.
    cand = _pick_sections(schema, 8, name)
    if len(cand) >= 5:
        # ★예전에는 cand[:5] 로 잘라 6번째 이후 섹션을 통째로 버렸고,
        #   면별 분량을 재지 않아 한 면은 빽빽, 한 면은 하단이 비었다.
        #   이제 분량을 재서 넘치는 섹션은 다음 면으로 이어 싣는다.
        _plan = _plan_panels(schema, cand, 5)
        picked = []
        for kind, payload in _plan:
            if kind == "sec":
                picked.append(payload)
            else:
                picked.append(("combo",
                               lambda sl, x, c, _p=payload: _panel_combo(sl, schema, x, c, _p)))
    elif len(cand) >= 3:
        # 3~4개: 1장으로 낸다. 콘텐츠 면이 2개뿐이므로
        # 첫 섹션은 제 면으로, 나머지 전부를 한 면에 합쳐 버리지 않는다.
        rest = [k for k, _ in cand[1:]]
        parts = _combo_parts(schema, rest)
        picked = cand[:1]
        if parts:
            picked = picked + [("combo",
                                lambda sl, x, c, _p=parts: _panel_combo(sl, schema, x, c, _p))]
    else:
        picked = cand              # 1~2개: 바깥면만

    # ★어떤 섹션이 <자기 면>을 받았는지 기록한다.
    #   패널 하단에 같은 내용을 또 붙이지 않기 위해 필요하다(실적이 두 번 나왔다).
    if isinstance(schema, dict):
        schema["_ownPanel"] = {k for k, _fn in picked if isinstance(k, str)}

    def _draw(slide, idx, x):
        """picked[idx] 를 그린다. 후보가 모자라면 아무것도 안 그림."""
        if 0 <= idx < len(picked):
            picked[idx][1](slide, x, c)
            _footer(slide, x, name, phone, c)
            return True
        return False

    n = len(picked)
    if n >= 5:
        outer_secs, inner_secs = 2, 3    # 2장: 바깥[섹·섹·표지] + 안쪽[섹·섹·섹]
    else:
        # 1장: 바깥[섹·섹·표지]. 표지는 항상 있어야 하므로 콘텐츠는 최대 2면.
        # (3면짜리 combo 케이스는 위에서 이미 2섹션+combo=3 이므로 여기서 2로 잘림 →
        #  combo 가 밀리지 않도록 n==3 이면 combo 를 두 번째 면에 둔다)
        outer_secs, inner_secs = min(2, n), 0

    # ── 슬라이드 1: 바깥면 (섹션 / 섹션 / 표지) ──
    s1 = prs.slides.add_slide(blank)
    (ax, aw), (fx, fw), (cx, cw) = out_geom
    if outer_secs > 0:
        _draw(s1, 0, ax)
    if outer_secs > 1:
        _rect(s1, fx, 0, 0.008, H_IN, fill=c["card_line"])   # 접는 선
        _draw(s1, 1, fx)
    _rect(s1, cx, 0, 0.008, H_IN, fill=c["card_line"])       # 접는 선
    _panel_cover(s1, schema, cx, c, _assets, pw=cw)          # 표지는 위치 고정

    # ── 슬라이드 2: 안쪽면 ──
    # 안쪽면은 3면이 다 차거나(3섹션) 아예 안 만든다 → 백지 면이 남지 않는다.
    if inner_secs > 0:
        s2 = prs.slides.add_slide(blank)
        (ix, iw), (mx, mw), (gx, gw) = in_geom
        _draw(s2, outer_secs, ix)
        if inner_secs > 1:
            _rect(s2, mx, 0, 0.008, H_IN, fill=c["card_line"])
            _draw(s2, outer_secs + 1, mx)
        if inner_secs > 2:
            _rect(s2, gx, 0, 0.008, H_IN, fill=c["card_line"])
            _draw(s2, outer_secs + 2, gx)

    _force_theme_font(prs, FONT)
    prs.core_properties.title = f"{name} 3단 리플렛"
    prs.core_properties.subject = "YouMeanI coordinate leaflet"

    if out is None:
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf
    prs.save(out)
    return out
