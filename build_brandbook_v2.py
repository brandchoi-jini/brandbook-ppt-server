#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_brandbook_v2.py — 유미니 상담용 PPT(신규 12장) 팔레트별 빌더.
텍스트 자리표시자만 교체(디자인·이미지·그라데이션 100% 보존).
규모별 스마트 채움: 데이터 없는 학년/섹션 슬라이드는 통째 삭제, 카드는 있는 만큼만.

표준 콘텐츠 스키마(content):
  academy{name,slogan,region,subjects,phone,kakao,address}
  copy{catch,identity}
  strengths[{title,desc}]           핵심 강점
  achievements[str]  growth[str]    주요 실적 / 성장 사례
  divisions{elem[],mid[],high[]}    각 [{name,desc}]
  management[{title,desc}]          학습관리 요소(출결·과제·성취도·상담)
  timetable{headers[],rows[[]]}
  specials[{title,desc}]            특별 프로그램
  faq[{q,a}]
  design{palette: green|blue|orange}
"""
import os, copy, json, argparse, io
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.oxml.ns import qn

TEMPLATES = {
    'green':  '유미니_상담용_ppt_그린.pptx',
    'blue':   '유미니_상담용_ppt_블루.pptx',
    'orange': '유미니_상담용_ppt_주황.pptx',
}
# 팔레트별 카드 수용량 / 슬라이드 인덱스
CAP = {
    'green':  {'strength':3, 'div':3, 'mng':4, 'faq':3, 'special':3},
    'blue':   {'strength':4, 'div':4, 'mng':4, 'faq':4, 'special':4},
    'orange': {'strength':3, 'div':3, 'mng':4, 'faq':5, 'special':3},
}
# 슬라이드 역할 인덱스 — 팔레트마다 순서가 다르다 (템플릿 실측)
#   blue  : 0표지 1정체성 2강점 3실적 4초 5중 6고 7시간표 8관리 9특강 10FAQ 11등록
#   green : 0표지 1정체성 2강점 3실적 4초 5중 6고 7관리  8시간표 9특강 10FAQ 11등록  ← 7·8 반대
#   orange: blue와 동일
IDX_BY_PAL = {
    'blue':   dict(cover=0, ident=1, strength=2, ach=3, elem=4, mid=5, high=6, tt=7, mng=8, special=9, faq=10, info=11),
    'orange': dict(cover=0, ident=1, strength=2, ach=3, elem=4, mid=5, high=6, tt=7, mng=8, special=9, faq=10, info=11),
    'green':  dict(cover=0, ident=1, strength=2, ach=3, elem=4, mid=5, high=6, mng=7, tt=8, special=9, faq=10, info=11),
}
IDX = IDX_BY_PAL['blue']   # 기본값 (build()에서 팔레트에 맞게 교체)

# 슬라이드 제목 플레이스홀더 → 학원명이 들어간 제목으로 교체
TITLE_MAP = {
    '학원의 정체성과 교육철학을 보여주세요': '{name}의 교육철학',
    '학원의 정체성을 보여주는 \n문구를 한줄로 작성해주세요.': '{name}의 교육철학',
    '학원의 정체성을 보여주는\n문구를 한줄로 작성해주세요.': '{name}의 교육철학',
    '학원의 핵심 강점을 보여주세요': '{name}의 핵심 강점',
    '학원의 핵심을\n적어주세요': '{name}의 핵심 강점',
    '학원의 주요 실적과 성장 사례를 보여주세요': '{name}의 주요 실적',
    '학원의 주요 실적과 성장 사례를 적어주세요': '{name}의 주요 실적',
    '학원의 주요 실적과 성장 사례를\n보여주세요': '{name}의 주요 실적',
    '학원의 레벨별 반구성과 학습 체계': '{name} 커리큘럼',
    '학원의 반구성과 학습 체계': '{name} 커리큘럼',
    '반구성과 학습 체계': '{name} 커리큘럼',
    '학원 반구성': '커리큘럼',
    '학원의 반 구성': '커리큘럼',
    '학원의 시간표를 알려주세요': '{name}의 시간표',
    '학원의\n시간표를 알려주세요': '{name}의 시간표',
    '학습의 관리 시스템을 적어주세요': '{name}의 학습 관리 시스템',
    '학원의 학습 관리 시스템': '{name}의 학습 관리 시스템',
    '학원의 특별 프로그램': '{name}만의 특별 프로그램',
    '학원의 등록 안내를 적어주세요': '{name} 등록 안내',
    '00학원의 체계적인 관리 시스템': '{name}의 체계적인 관리 시스템',
    '00학원만의 프로그램': '{name}만의 프로그램',
    '학원의 학습 관리 시스템을 적어주세요': '{name}의 학습 관리 시스템',
}

# 리드문(설명 문구) 플레이스홀더 → 데이터 기반 문장
LEAD_KEYS = [
    '레벨별 반구성과 학습체게에 대해 자세하게 적어주세요',
    '레벨별 반구성과 학습체계에 대해 자세하게 적어주세요',
    '세부 내용을 적어주세요',
    '세부 설명을 적어주세요',
    '학원의 입학 절차를 적어주세요',
    "지금까지 달성한 구체적인 성과를 보여주며",
]

def clear_leads(prs, leads=None):
    """안 채워진 리드문 플레이스홀더를 비우거나 대체"""
    leads = leads or {}
    for s0 in prs.slides:
        for sh in s0.shapes:
            if not sh.has_text_frame: continue
            t = sh.text_frame.text
            if not t.strip(): continue
            for k in LEAD_KEYS:
                if k in t:
                    _set(sh.text_frame, leads.get(k, ''))
                    break

def fill_titles(prs, name):
    """모든 슬라이드의 제목 플레이스홀더를 학원명으로 교체"""
    if not name: return
    for s in prs.slides:
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            t = sh.text_frame.text.strip()
            if not t: continue
            # 개행 유무 모두 대응
            for k, v in TITLE_MAP.items():
                kk = k.replace('\\n', '\n')
                if t == kk or t.replace('\n', '') == kk.replace('\n', ''):
                    _set(sh.text_frame, v.format(name=name))
                    break

def walk(shapes):
    for x in shapes:
        yield x
        if x.shape_type == 6:
            yield from walk(x.shapes)

def _set(tf, txt):
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = txt
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else:
        r = p0.add_run(); r.text = txt
    for p in tf.paragraphs[1:]: p._p.getparent().remove(p._p)

def _set_multi(tf, lines):
    from pptx.text.text import _Paragraph
    lines = [l for l in lines if l is not None] or ['']
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else:
        r = p0.add_run(); r.text = lines[0]
    for p in tf.paragraphs[1:]: p._p.getparent().remove(p._p)
    for ln in lines[1:]:
        newp = copy.deepcopy(p0._p); p0._p.getparent().append(newp)
        para = _Paragraph(newp, p0._parent)
        if para.runs:
            para.runs[0].text = ln
            for r in para.runs[1:]: r._r.getparent().remove(r._r)

def norm(t):
    return t.strip().replace('\n', ' / ').replace('\xa0', ' ')

def edit(slide, mapping=None, multimap=None):
    mapping = mapping or {}; multimap = multimap or {}
    # 키도 정규화(\xa0 제거) 해서 비교
    mp = {norm(k): v for k, v in mapping.items()}
    mm = {norm(k): v for k, v in multimap.items()}
    seen = {}
    for x in walk(slide.shapes):
        if not (x.has_text_frame and x.text_frame.text.strip()): continue
        t = norm(x.text_frame.text)
        if t in mm:
            i = seen.get(t, 0); blocks = mm[t]
            blk = blocks[i] if (blocks and isinstance(blocks[0], list)) else blocks
            if blocks and isinstance(blocks[0], list) and i >= len(blocks): blk = ['']
            _set_multi(x.text_frame, blk); seen[t] = i + 1
        elif t in mp:
            i = seen.get(t, 0); v = mp[t]
            if isinstance(v, list): _set(x.text_frame, v[i] if i < len(v) else '')
            else: _set(x.text_frame, v)
            seen[t] = i + 1

def remove_slide(prs, slide):
    lst = prs.slides._sldIdLst
    slides = list(prs.slides)
    xmls = list(lst)
    for idx, sl in enumerate(slides):
        if sl is slide:
            lst.remove(xmls[idx]); return True
    return False

# ── 학원 로고·사진 삽입 ──
def _fetch(url, timeout=8):
    """URL에서 이미지 바이트를 받아온다. 실패하면 None."""
    if not url or not str(url).startswith('http'):
        return None
    try:
        import urllib.request
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=timeout) as r:
            data = r.read()
        if len(data) < 100:
            return None
        return io.BytesIO(data)
    except Exception as e:
        print(f'[img] fetch fail: {e}')
        return None

def put_logo(slide, url, left_in, top_in, max_w_in=1.5, max_h_in=0.7):
    """슬라이드에 로고를 얹는다. 실패하면 '로고 넣을 자리' 표시를 남긴다."""
    bio = _fetch(url)
    if bio:
        try:
            pic = slide.shapes.add_picture(bio, Inches(left_in), Inches(top_in))
            rw = max_w_in / (pic.width / 914400)
            rh = max_h_in / (pic.height / 914400)
            r = min(rw, rh, 1.0)
            pic.width = int(pic.width * r)
            pic.height = int(pic.height * r)
            return True
        except Exception as e:
            print(f'[img] logo fail: {e}')
    _placeholder(slide, left_in, top_in, max_w_in, max_h_in, '로고')
    return False

def swap_photo(slide, url, target_shape, label='학원 사진'):
    """템플릿 일러스트 자리를 학원 사진으로 교체. 실패하면 자리표시만 남긴다."""
    if target_shape is None:
        return False
    L, T = target_shape.left, target_shape.top
    Wd, Ht = target_shape.width, target_shape.height
    bio = _fetch(url)
    if bio:
        try:
            target_shape._element.getparent().remove(target_shape._element)
            slide.shapes.add_picture(bio, L, T, width=Wd, height=Ht)
            return True
        except Exception as e:
            print(f'[img] swap fail: {e}')
    # 실패: 원본 일러스트는 그대로 두고, 자리 안내를 그 위 중앙에 작게 얹는다
    _w = min(Wd / 914400 * 0.75, 4.0)
    _h = 0.55
    _l = (L + Wd / 2) / 914400 - _w / 2
    _t = (T + Ht / 2) / 914400 - _h / 2
    _placeholder(slide, _l, _t, _w, _h, label + ' 넣을 자리')
    return False

def _placeholder(slide, left_in, top_in, w_in, h_in, label):
    """사진·로고를 못 받아왔을 때 '여기에 넣으세요' 안내 박스."""
    try:
        from pptx.util import Inches as _I
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     _I(left_in), _I(top_in), _I(w_in), _I(h_in))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        box.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        box.line.width = Pt(1)
        box.line.dash_style = 2   # dashed
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f'📷 {label}'
        r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        return True
    except Exception as e:
        print(f'[img] placeholder fail: {e}')
        return False

def fill_table(slide, headers, rows):
    for sh in slide.shapes:
        if not sh.has_table: continue
        tbl = sh.table
        need = len(headers)

        # ① 남는 열 삭제 (헤더보다 열이 많으면 빈 칸이 생긴다)
        while len(tbl.columns) > need and len(tbl.columns) > 1:
            _del_col(tbl, len(tbl.columns) - 1)

        ncol = len(tbl.columns)
        for c in range(min(ncol, need)):
            _set(tbl.cell(0, c).text_frame, str(headers[c]))
        for ri, row in enumerate(rows, start=1):
            if ri >= len(tbl.rows): _add_row(tbl)
            for c in range(ncol):
                _set(tbl.cell(ri, c).text_frame, str(row[c]) if c < len(row) else '')

        # ② 남는 행 삭제 (비워두면 빈 줄이 그대로 보인다)
        while len(tbl.rows) > len(rows) + 1 and len(tbl.rows) > 2:
            _del_row(tbl, len(tbl.rows) - 1)
        return True
    return False

def _del_col(tbl, idx):
    try:
        gridCol = tbl._tbl.find(qn('a:tblGrid')).findall(qn('a:gridCol'))[idx]
        gridCol.getparent().remove(gridCol)
        for tr in tbl._tbl.findall(qn('a:tr')):
            tcs = tr.findall(qn('a:tc'))
            if idx < len(tcs):
                tr.remove(tcs[idx])
    except Exception:
        pass

def _del_row(tbl, idx):
    try:
        trs = tbl._tbl.findall(qn('a:tr'))
        if 0 <= idx < len(trs):
            tbl._tbl.remove(trs[idx])
    except Exception:
        pass

def _add_row(tbl):
    last = tbl._tbl.findall(qn('a:tr'))[-1]
    new = copy.deepcopy(last)
    for tc in new.findall(qn('a:tc')):
        for t in tc.iter(qn('a:t')): t.text = ''
    tbl._tbl.append(new)

def _wrap(text, n=18):
    """카드에 넣을 텍스트를 문장·어절 경계로 나눈다. (글자 수로 무작정 자르면
    '확인하며 / , 당일' 처럼 쉼표가 줄 앞에 오고 단어가 깨진다)"""
    text = (text or '').strip()
    if not text:
        return ['']
    # 이미 줄바꿈이 있으면 그대로 존중
    if '\n' in text:
        return [l.strip() for l in text.split('\n') if l.strip()] or ['']

    # 1) 문장 단위로 먼저 쪼갠다
    import re as _re
    sents = [s.strip() for s in _re.split(r'(?<=[.!?다요])\s+', text) if s.strip()]
    if not sents:
        sents = [text]

    out = []
    for s in sents:
        if len(s) <= n * 1.6:
            out.append(s)
            continue
        # 2) 긴 문장은 어절(공백) 경계로 접는다
        words = s.split(' ')
        cur = ''
        for w in words:
            if not cur:
                cur = w
            elif len(cur) + 1 + len(w) <= n:
                cur += ' ' + w
            else:
                out.append(cur)
                cur = w
        if cur:
            out.append(cur)
    return out or ['']

def _titles(items, n):   return [(items[i]['title'] if i < len(items) else '') for i in range(n)]
def _names(items, n):    return [(items[i]['name']  if i < len(items) else '') for i in range(n)]
def _descs(items, n, max_lines=4, max_chars=110):
    """카드 설명문. 카드가 감당할 수 있는 만큼만 넣는다.
    (긴 서술문을 통째로 넣으면 칸을 뚫고 나가거나 잘려서 문장이 깨진다)"""
    out = []
    for i in range(n):
        if i >= len(items):
            out.append([''])
            continue
        d = (items[i].get('desc') or '').strip()
        if not d:
            out.append([''])
            continue
        # 1) 글자 수 상한 — 문장 경계에서 컷
        d = _cap(d, max_chars)
        # 2) 줄로 접고, 줄 수 상한
        lines = _wrap(d)
        if len(lines) > max_lines:
            lines = lines[:max_lines]
            # 마지막 줄이 문장 중간이면 말줄임
            if not lines[-1].rstrip().endswith(('.', '다', '요', '!', '?')):
                lines[-1] = lines[-1].rstrip() + '…'
        out.append(lines or [''])
    return out

def _normalize_ppt_schema(content):
    """v3 스키마(curriculum/special/dict-achievements)를 PPT build()가 쓰는 키로 정규화.
    설명문 길이도 제한해 넘침을 방지한다. 원본을 훼손하지 않도록 복사본에 채운다."""
    c = dict(content)
    # 강점 desc 길이 제한 (표지 하단 캡션은 짧아야 함)
    st = _dicts(c.get('strengths', []))
    for s in st:
        s['desc'] = _cap(s.get('desc',''), 46)
    c['strengths'] = st

    # divisions 없으면 curriculum에서 학년별로 구성
    if not c.get('divisions'):
        curric = _dicts(c.get('curriculum', []))
        gm = _grade_map(curric)
        div = {}
        for gk, key in [('초등','elem'), ('중등','mid'), ('고등','high')]:
            cc = gm.get(gk)
            if cc:
                desc = cc.get('content') or cc.get('goal') or ''
                div[key] = [{
                    'name': cc.get('title', f'{gk} 과정'),
                    'level': cc.get('subject',''),
                    'desc': _cap(desc, 90),
                }]
        c['divisions'] = div

    # specials <- special (v3 키명 차이)
    if not c.get('specials'):
        sp = _dicts(c.get('special', []))
        norm_sp = []
        for s in sp:
            desc = s.get('desc','')
            if not desc and isinstance(s.get('items'), list):
                desc = ' · '.join(s['items'][:3])
            norm_sp.append({'title': s.get('title',''), 'desc': _cap(desc, 70)})
        c['specials'] = norm_sp

    # management desc 길이 제한
    mng = _dicts(c.get('management', []))
    for m in mng:
        m['desc'] = _cap(m.get('desc',''), 60)
    c['management'] = mng

    # achievements/growth: dict 리스트 -> 문자열 리스트
    def _to_strs(lst):
        out = []
        for x in (lst or []):
            if isinstance(x, dict):
                v = x.get('text') or x.get('tag') or x.get('title') or ''
            else:
                v = str(x)
            v = v.strip()
            if v: out.append(_cap(v, 40))
        return out
    c['achievements'] = _to_strs(c.get('achievements'))
    c['growth'] = _to_strs(c.get('growth'))

    return c


def build(content, out_path, template_dir='.'):
    global IDX
    content = _normalize_ppt_schema(content)
    design = content.get('design') or {}
    pal = design.get('palette', 'green')
    if pal not in TEMPLATES: pal = 'green'
    IDX = IDX_BY_PAL.get(pal, IDX_BY_PAL['blue'])   # ★ 팔레트마다 슬라이드 순서가 다르다
    prs = Presentation(os.path.join(template_dir, TEMPLATES[pal]))
    S = list(prs.slides)
    cap = CAP[pal]
    ac = content.get('academy', {}); cp = content.get('copy', {})
    name = ac.get('name', '학원명')
    st = content.get('strengths', [])

    # ★ 슬라이드 제목 플레이스홀더를 학원명으로 (본문 채우기 전에)
    fill_titles(prs, name)

    # ★ 학원 로고·사진 — 받아지면 넣고, 못 받으면 "넣을 자리" 표시를 남긴다
    imgs = content.get('images') or {}
    _logo = imgs.get('logo') or ''
    _photos = [p for p in (imgs.get('photos') or []) if p]

    # 표지 좌하단 로고 — 슬라이드 실제 크기 기준으로 배치 (템플릿마다 크기가 다름)
    _SW = prs.slide_width / 914400.0
    _SH = prs.slide_height / 914400.0
    _lw = _SW * 0.12          # 로고 폭 = 슬라이드 폭의 12%
    _lh = _lw * 0.45
    put_logo(S[IDX['cover']], _logo,
             left_in=_SW * 0.045,
             top_in=_SH - _lh - (_SH * 0.05),
             max_w_in=_lw, max_h_in=_lh)

    # 정체성 슬라이드 중앙 일러스트 → 학원 사진
    _sl = S[IDX['ident']]
    _big = None
    for _sh in _sl.shapes:
        if hasattr(_sh, 'image'):
            _a = (_sh.width / 914400) * (_sh.height / 914400)
            if _a > 5 and (_big is None or _a > (_big.width/914400)*(_big.height/914400)):
                _big = _sh
    if _big is not None:
        swap_photo(_sl, _photos[0] if _photos else '', _big, label='학원 사진')

    # [0] 표지 — 큰 제목 자리에는 학원명만. (catch/정체성 같은 긴 문장을 넣으면 화면을 뚫는다)
    _slogan = ac.get('slogan') or cp.get('catch') or ''
    edit(S[IDX['cover']], {
        '학원명': name, '제목을': name, '제목을\n입력해 주세요': name,
        '입력해 주세요': name, '제목을 / 입력해 주세요': name,
        '슬로건을 입력해주세요': _slogan,
        '슬로건을 입력해 주세요': _slogan,
        '슬로건을 입력해주세요 슬로건을 입력해주세요 슬로건을 입력해주세요': _slogan,
    })

    # [1] 정체성 — 팔레트별 자리표시자 상이
    ident_common = {
        '학원의 정체성을 보여주는  / 문구를 한줄로 작성해주세요.': cp.get('identity',''),
        '학원의 정체성을 보여주는 / 문구를 한줄로 작성해주세요.': cp.get('identity',''),
        '학원의 교육 방향과 / 정체성을 / 한 줄로 정의해 주세요': cp.get('catch',''),
    }
    if pal == 'green':
        edit(S[1], {**ident_common,
            '교육 역량 1': st[0]['title'] if len(st)>0 else '',
            '교육 역량 2': st[1]['title'] if len(st)>1 else '',
            '교육 역량 3': st[2]['title'] if len(st)>2 else '',
            '[학원만의 교육 방향을 / 적어주세요': [ (st[i]['desc'] if i<len(st) else '') for i in range(3) ],
        })
    elif pal == 'orange':
        edit(S[1], {**ident_common,
            '교육 방향 1': st[0]['title'] if len(st)>0 else '',
            '교육 방향 2': st[1]['title'] if len(st)>1 else '',
            '교육 방향 3': st[2]['title'] if len(st)>2 else '',
            '[학원만의 / 교육 방향을 / 적어주세요': [ (st[i]['desc'] if i<len(st) else '') for i in range(3) ],
            '학원만의 / 교육 방향을 / 적어주세요': [ (st[i]['desc'] if i<len(st) else '') for i in range(3) ],
        })
    else:  # blue: 정체성 슬라이드는 "한 줄 정의 / 교육철학 2" 2열
        edit(S[1], {
            '학원에 대한 한 줄 정의': cp.get('catch','') or '한 줄 정의',
            '교육 철학 2': (st[0]['title'] if st else '교육 철학'),
        }, multimap={
            '세부 내용을 입력해 주세요 / 세부 내용을 입력해 주세요 / 세부 내용을 입력해 주세요 / 세부 내용을 입력해 주세요':
                [ _wrap(cp.get('identity','')), _wrap(st[0]['desc'] if st else '') ],
        })

    # [2] 핵심 강점
    if pal == 'green':
        edit(S[2], {'강점을 적어주세요': _titles(st,3)},
             multimap={'학원의 핵심 강정을 설명해 주세요 / 학원의 핵심 강정을 설명해 주세요': _descs(st,3)})
    elif pal == 'orange':
        edit(S[2], {'핵심 강점 1': st[0]['title'] if len(st)>0 else '',
                    '핵심 강점 2': st[1]['title'] if len(st)>1 else '',
                    '핵심 강점 3': st[2]['title'] if len(st)>2 else ''},
             multimap={'학원의 핵심 강정을 설명해 주세요 / 학원의 핵심 강정을 설명해 주세요 / 학원의 핵심 강정을 설명해 주세요': _descs(st,3)})
    else:  # blue 4개
        edit(S[2], {'핵심 강점 1': st[0]['title'] if len(st)>0 else '',
                    '핵심 강점 2': st[1]['title'] if len(st)>1 else '',
                    '핵심 강점 3': st[2]['title'] if len(st)>2 else '',
                    '핵심 강점 4': st[3]['title'] if len(st)>3 else ''},
             multimap={'강점에 대한 세부 설명을 적어주세요 / 강점에 대한 세부 설명을 적어주세요': _descs(st,4)})

    # [3] 주요 실적 + 성장 사례
    ach = content.get('achievements', []); growth = content.get('growth', [])
    if ach or growth:
        _ach_head = cp.get('positioning','') or '한 명도 놓치지 않는 관리로 성장 사례를 만들어갑니다.'
        edit(S[3], {'학원의 주요 실적과 성장 사례를 적어주세요': _ach_head,
                    '학원의 주요 실적과 성장 사례를 / 적어주세요': _ach_head})
        if pal == 'green':
            details = [([ach[i]] if i<len(ach) else ['']) for i in range(3)] + \
                      [([growth[i]] if i<len(growth) else ['']) for i in range(3)]
            edit(S[3], multimap={'세부 내용을 입력해 주세요': details})
        else:  # blue/orange: "입시 성과를 구체적으로 적어주세요" 6줄
            lines = (ach + growth)[:6]
            edit(S[3], multimap={
                '입시 성과를 구체적으로 적어주세요 / 입시 성과를 구체적으로 적어주세요 / 입시 성과를 구체적으로 적어주세요 / 입시 성과를 구체적으로 적어주세요 / 입시 성과를 구체적으로 적어주세요 / 입시 성과를 구체적으로 적어주세요':
                    lines if lines else ['']})
    else:
        remove_slide(prs, S[3])

    # [4][5][6] 학년별 반구성 — 스마트: 데이터 있는 학년만
    div = content.get('divisions', {})
    for key, si in [('elem', IDX['elem']), ('mid', IDX['mid']), ('high', IDX['high'])]:
        items = div.get(key, [])
        if not items:
            remove_slide(prs, S[si]); continue
        ncap = cap['div']
        if pal == 'green':
            edit(S[si], {'반 이름을 적어주세요': _names(items, 3)},
                 multimap={'학습체계를 구체적으로 적어주세요 / 학습체계를 구체적으로 적어주세요 / 학습체계를 구체적으로 적어주세요': _descs(items, 3)})
        elif pal == 'blue':
            edit(S[si], {'반 이름을 적어주세요': _names(items, 4)},
                 multimap={'세부적인 내용을 적어주세요': _descs(items, 4, max_lines=2, max_chars=60)})
        else:  # orange
            edit(S[si], {'반 이름을 적어주세요': _names(items, 3),
                         '반 레벨을 설명해 주세요': [(items[i].get('level', items[i].get('name','')) if i<len(items) else '') for i in range(3)]},
                 multimap={'학습체계를 구체적으로 적어주세요 / 학습체계를 구체적으로 적어주세요 / 학습체계를 구체적으로 적어주세요 / 학습체계를 구체적으로 적어주세요 / 학습체계를 구체적으로 적어주세요 / 학습체계를 구체적으로 적어주세요': _descs(items,3)})

    # [7] 학습 관리
    mng = content.get('management', [])
    if mng:
        labels = ['출결','과제','성취도','상담']
        m = {}
        for i, lb in enumerate(labels):
            m[lb] = mng[i]['title'] if i < len(mng) else lb
        if pal == 'green':
            m['요소 1'] = ''; m['요소 2'] = ''; m['요소 3'] = ''; m['요소 4'] = ''
            edit(S[IDX['mng']], m, multimap={'구체적인 설명을 적어주세요 / 구체적인 설명을 적어주세요': _descs(mng,4)})
        elif pal == 'blue':
            edit(S[IDX['mng']], m, multimap={'구체적인 설명을 적어주세요 / 구체적인 설명을 적어주세요 / 구체적인 설명을 적어주세요': _descs(mng,4)})
        else:
            m2 = {'한 줄로 정의해주세요': [ (mng[i]['desc'].split('.')[0] if i<len(mng) else '') for i in range(4) ]}
            edit(S[IDX['mng']], {**m, **{}}, multimap={'구체적인 설명을 적어주세요 / 구체적인 설명을 적어주세요 / 구체적인 설명을 적어주세요 / 구체적인 설명을 적어주세요': _descs(mng,4)})
    else:
        remove_slide(prs, S[IDX['mng']])

    # [8] 시간표
    tt = content.get('timetable', {})
    if tt and tt.get('rows'):
        fill_table(S[IDX['tt']], tt.get('headers', ['','월','화','수','목','금','토','일']), tt['rows'])
    else:
        remove_slide(prs, S[IDX['tt']])

    # [9] 특별 프로그램
    sp = content.get('specials', [])
    if sp:
        lead = {'프로그램의 한 줄 정의를 설명해주세요': cp.get('specialsLead',''),
                '학원의 특별 프로그램을 / 적어주세요': cp.get('specialsLead','') or '학원의 특별 프로그램'}
        if pal == 'green':
            edit(S[9], {**lead, '프로그램명': _titles(sp,3)},
                 multimap={'세부적인 설명을 적어주세요 / 세부적인 설명을 적어주세요 / 세부적인 설명을 적어주세요 / 세부적인 설명을 적어주세요': _descs(sp,3)})
        elif pal == 'blue':
            edit(S[9], {'프로그램명': _titles(sp,4)},
                 multimap={'세부 설명을 적어주세요 / 세부 설명을 적어주세요 / 세부 설명을 적어주세요': _descs(sp,4)})
        else:
            edit(S[9], {**lead, '핵심 프로그램명': _titles(sp,3)},
                 multimap={'세부적인 설명을 적어주세요 / 세부적인 설명을 적어주세요 / 세부적인 설명을 적어주세요 / 세부적인 설명을 적어주세요': _descs(sp,3)})
    else:
        remove_slide(prs, S[IDX['special']])

    # [10] FAQ — 답변이 길면 칸을 뚫으므로 길이 제한
    faq = content.get('faq', [])
    if faq:
        n = cap['faq']
        # 팔레트별 답변 최대 길이 (템플릿 칸 크기 실측 기준)
        A_MAX = {'green': 150, 'blue': 130, 'orange': 110}.get(pal, 130)
        Q_MAX = 40
        _fq = [{'q': _cap(x.get('q',''), Q_MAX),
                'a': _cap(x.get('a',''), A_MAX),
                'detail': _cap(x.get('detail',''), 90)} for x in faq]
        if pal == 'green':
            edit(S[10], {'질문을 입력해 주세요': [(_fq[i]['q'] if i<len(_fq) else '') for i in range(3)],
                         'A. 답변을 적어주세요': [('A. '+_fq[i]['a'] if i<len(_fq) else '') for i in range(3)]},
                 multimap={'세부 설명을 추가해 주세요 / 세부 설명을 추가해 주세요 / 세부 설명을 추가해 주세요 / 세부 설명을 추가해 주세요':
                           [(_wrap(_fq[i].get('detail','')) if i<len(_fq) else ['']) for i in range(3)]})
        elif pal == 'blue':
            edit(S[10], {'질문을 적어주세요': [(_fq[i]['q'] if i<len(_fq) else '') for i in range(4)]},
                 multimap={'답변을 적어주세요 / 답변을 적어주세요': [(_wrap(_fq[i]['a']) if i<len(_fq) else ['']) for i in range(4)]})
        else:  # orange 5개
            edit(S[10], {'질문의 내용을 적어주세요': [(_fq[i]['q'] if i<len(_fq) else '') for i in range(5)],
                         '질문에 대한 상세한 답변을 적어주세요': [(_fq[i]['a'] if i<len(_fq) else '') for i in range(5)]})
    else:
        remove_slide(prs, S[IDX['faq']])

    # [11] 등록 안내
    edit(S[IDX['info']], {
        '상담 전화번호를 적어주세요': ac.get('phone',''),
        '상담 전화번호를 / 적어주세요': ac.get('phone',''),
        '카카오톡 채널을 적어주세요': ac.get('kakao',''),
        '카카오톡 채널을 / 적어주세요': ac.get('kakao',''),
        '주소를 적어주세요': ac.get('address',''),
        '주소를 / 적어주세요': ac.get('address',''),
    })

    # 남은 placeholder 보강: 핵심(Key point)·등록안내 등 매핑 안 된 안내문을 데이터로 교체
    _posit = cp.get('positioning','') or cp.get('identity','') or ac.get('slogan','')
    _reg_lines = []
    _rules = content.get('operatingRules')
    if isinstance(_rules, dict):
        if _rules.get('refund'): _reg_lines.append('환불: '+_rules['refund'])
        if _rules.get('refundBefore'): _reg_lines.append('수강 전 환불: '+_rules['refundBefore'])
        if _rules.get('refundAfter'): _reg_lines.append('수강 후 환불: '+_rules['refundAfter'])
    _reg_txt = '\n'.join(_reg_lines) if _reg_lines else (
        '상담 신청 후 학습 상담과 레벨 테스트를 거쳐 반이 배정됩니다.' +
        (('\n문의: '+ac.get('phone','')) if ac.get('phone') else ''))
    _PH_FILL = {
        '학원의 핵심을 적어주세요': _posit,
        '학원의 핵심을\n적어주세요': _posit,
        '학원의 등록 안내를 적어주세요': _reg_txt,
        '학원의 정체성과 교육철학을 보여주세요': cp.get('identity','') or _posit,
    }
    def _norm(s): return s.replace('\xa0',' ').replace('\n',' ').strip()
    _PH_NORM = { _norm(k):v for k,v in _PH_FILL.items() }
    for _s in S:
        for _sh in _s.shapes:
            if not _sh.has_text_frame: continue
            _val = _PH_NORM.get(_norm(_sh.text_frame.text))
            if _val:
                _p = _sh.text_frame.paragraphs[0]
                if _p.runs:
                    _p.runs[0].text = _val
                    for _r in _p.runs[1:]: _r.text = ''
                else:
                    _p.add_run().text = _val
                for _ex in _sh.text_frame.paragraphs[1:]:
                    _ex._p.getparent().remove(_ex._p)

    # 넘침 방지: 각 텍스트박스를 폭 기반으로 자동 축소
    for _s in list(prs.slides):
        for _sh in _s.shapes:
            if not _sh.has_text_frame: continue
            txt = (_sh.text_frame.text or '').strip()
            if not txt: continue
            try:
                w_in = _in(_sh.width); h_in = _in(_sh.height)
            except Exception:
                continue
            if w_in <= 0.2 or h_in <= 0.1: continue

            # 현재 폰트 크기
            cur = None
            for _p in _sh.text_frame.paragraphs:
                for _r in _p.runs:
                    if _r.font.size: cur = _r.font.size.pt; break
                if cur: break
            if not cur: continue

            # 실제 수용 글자수 = (가로 글자수) × (줄 수)
            # 한글 1글자 ≈ 폰트pt, 1inch = 72pt. 줄높이 ≈ 폰트pt × 1.35
            chars_per_line = max(1, int((w_in * 72) / max(cur, 1) * 0.95))
            lines = max(1, int((h_in * 72) / (max(cur, 1) * 1.35)))
            cap_chars = max(4, chars_per_line * lines)

            if len(txt) > cap_chars:
                _autosize(_sh, base_pt=cur, max_chars_at_base=cap_chars, min_pt=8)

    # ★ 남은 리드문 플레이스홀더 처리 — 데이터가 있으면 채우고, 없으면 비운다
    div = content.get('divisions') or {}
    _lead_div = ''
    _all = (div.get('elem') or []) + (div.get('mid') or []) + (div.get('high') or [])
    _descs_all = [d.get('desc','') for d in _all if isinstance(d, dict) and (d.get('desc') or '').strip()]
    if _descs_all:
        _lead_div = _descs_all[0][:90]
    _mg = content.get('management') or []
    _lead_mng = (_mg[0].get('desc','')[:80] if _mg and isinstance(_mg[0], dict) else '')
    _sp = content.get('specials') or []
    _lead_sp = (_sp[0].get('desc','')[:80] if _sp and isinstance(_sp[0], dict) else '')
    _adm = content.get('admission') or []
    _lead_adm = ' → '.join((a.get('step') or '') for a in _adm if isinstance(a, dict) and a.get('step'))

    clear_leads(prs, {
        '레벨별 반구성과 학습체게에 대해 자세하게 적어주세요': _lead_div,
        '레벨별 반구성과 학습체계에 대해 자세하게 적어주세요': _lead_div,
        '세부 내용을 적어주세요': _lead_mng,
        '세부 설명을 적어주세요': _lead_sp,
        '학원의 입학 절차를 적어주세요': _lead_adm,
        '지금까지 달성한 구체적인 성과를 보여주며': '',
    })

    prs.save(out_path)
    return out_path

if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('content_json'); ap.add_argument('out_pptx')
    ap.add_argument('--templates', default='.')
    a = ap.parse_args()
    content = json.load(open(a.content_json, encoding='utf-8'))
    print('저장:', build(content, a.out_pptx, a.templates))


# ============================================================
# 카탈로그(6장) / 리플렛(2장) 빌더
# ============================================================
CATALOG_TEMPLATES = {
    'gray_gray':  '유미니_카탈로그_디자인_-_그레이__그레이.pptx',
    'gray_green': '유미니_카탈로그_디자인_-_그레이__그린.pptx',
    'blue':       '유미니_카탈로그_디자인_-_블루.pptx',
}
LEAFLET_TEMPLATES = {
    'gray_green': '그레이_그린.pptx',
    'gray_blue':  '그레이_블루.pptx',
    'blue_blue':  '블루_블루.pptx',
}
# 앱 팔레트 → 카탈로그/리플렛 색
PAL_TO_CATALOG = {'forest':'gray_green','sage':'gray_green','navy':'blue','violet':'blue','crimson':'gray_gray','mono':'gray_gray'}
PAL_TO_LEAFLET = {'forest':'gray_green','sage':'gray_green','navy':'blue_blue','violet':'blue_blue','crimson':'gray_blue','mono':'gray_blue'}

# 교육 특징 6개 자리표시자(제목→설명) — 카탈로그/리플렛 공통
FEATURE_SLOTS = [
    ('학년 담임 통합', '전 과목 진도와 생활을 함께 관리합니다.'),
    ('과목별 전담 강사', '5년 이상 경력의 전문 강사진입니다.'),
    ('주간 통합 리포트', '전 과목 학습 현황을 한 페이지로 공유합니다.'),
    ('과목 간 연계 학습', '독해 · 지문을 교차로 활용합니다.'),
    ('시험 4주 대비', '전 과목 동시 부스터를 가동합니다.'),
    ('초등 · 중등 · 고등 진학 연결', '기초부터 시작해 고등까지 이어갑니다.'),
]


# ============================================================
# v3 좌표 기반 빌더 (카탈로그·리플렛) — 예시 완전 제거 + 학년 편차 스마트 채움
# ============================================================

CATALOG_TEMPLATES = {
    'gray':  '유미니_카탈로그_디자인_-_그레이__그레이.pptx',
    'green': '유미니_카탈로그_디자인_-_그레이__그린.pptx',
    'blue':  '유미니_카탈로그_디자인_-_블루.pptx',
}
LEAFLET_TEMPLATES = {
    'green': '그레이_그린.pptx',
    'blue':  '그레이_블루.pptx',
    'navy':  '블루_블루.pptx',
}

def _in(v): return round(Emu(v).inches, 1)

from pptx.util import Pt, Inches

def _autosize(shape, base_pt=None, min_pt=9, max_chars_at_base=None):
    """텍스트 길이에 따라 폰트 크기를 줄여 넘침 방지. base_pt 미지정 시 현재 폰트 사용."""
    if shape is None or not shape.has_text_frame: return
    tf = shape.text_frame
    txt = tf.text or ''
    n = len(txt)
    # 기준 폰트
    cur = None
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size: cur = r.font.size.pt; break
        if cur: break
    base = base_pt or cur or 14
    if max_chars_at_base is None:
        max_chars_at_base = 40
    if n > max_chars_at_base:
        scale = (max_chars_at_base / n) ** 0.5
        new = max(min_pt, round(base * scale))
    else:
        new = base
    for p in tf.paragraphs:
        for r in p.runs:
            try: r.font.size = Pt(new)
            except Exception: pass
    try:
        tf.word_wrap = True
    except Exception: pass

def _cap(s, n):
    """설명문 길이 제한 — 문장 경계 우선."""
    s = (s or '').strip()
    if len(s) <= n: return s
    cut = s[:n]
    for sep in ['. ', '다 ', '다. ', '요 ', '요. ', ' · ', ' — ', ', ']:
        i = cut.rfind(sep)
        if i > n*0.5: return cut[:i+1].strip()
    return cut.strip() + '…'

def _find_pos(slide, left, top, tol=0.35):
    best, bestd = None, 999
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        l, t = _in(sh.left), _in(sh.top)
        d = abs(l-left) + abs(t-top)
        if d <= tol*2 and d < bestd:
            best, bestd = sh, d
    return best

def _put(shape, txt):
    if shape is None or txt is None: return
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = txt
        for r in p.runs[1:]: r.text = ''
    else:
        p.add_run().text = txt
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)

def _put_multi(shape, lines):
    if shape is None: return
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    ref = p0.runs[0] if p0.runs else None
    if p0.runs:
        p0.runs[0].text = lines[0] if lines else ''
        for r in p0.runs[1:]: r.text = ''
    else:
        p0.add_run().text = lines[0] if lines else ''
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    for ln in lines[1:]:
        np = tf.add_paragraph(); r = np.add_run(); r.text = ln
        if ref is not None:
            try:
                r.font.size = ref.font.size; r.font.bold = ref.font.bold
                if ref.font.color and ref.font.color.type is not None:
                    r.font.color.rgb = ref.font.color.rgb
            except Exception: pass

def _clear_region(slide, xmin, xmax, tmin, tmax):
    for sh in list(slide.shapes):
        if not sh.has_text_frame: continue
        l, t = _in(sh.left), _in(sh.top)
        if xmin <= l <= xmax and tmin <= t <= tmax:
            _put(sh, '')

def _dicts(lst):
    """리스트에서 dict 원소만 남긴다. 문자열은 {'title':s,'desc':s}로, 그 외는 버림."""
    if not isinstance(lst, list): return []
    out = []
    for x in lst:
        if isinstance(x, dict): out.append(x)
        elif isinstance(x, str) and x.strip(): out.append({'title': x.strip()[:20], 'desc': x.strip(), 'q': x.strip(), 'a': '', 'text': x.strip()})
    return out

def _grade_map(curric):
    gm = {}
    for c in curric:
        g = c.get('grade','')
        key = '초등' if '초' in g else '중등' if '중' in g else '고등' if '고' in g else g
        gm[key] = c
    return gm


def build_catalog(content, out_path, template_dir='.'):
    design = content.get('design') or {}
    color = design.get('palette', 'green')
    if color not in CATALOG_TEMPLATES: color = 'green'
    prs = Presentation(os.path.join(template_dir, CATALOG_TEMPLATES[color]))
    S = list(prs.slides)
    ac = content.get('academy', {}); cp = content.get('copy', {})
    feats = _dicts(content.get('features', []) or content.get('strengths', []))
    strengths = _dicts(content.get('strengths', []))
    faq = _dicts(content.get('faq', []))
    curric = _dicts(content.get('curriculum', []))
    mgmt = _dicts(content.get('management', []))
    special = _dicts(content.get('special', []))
    ach = _dicts(content.get('achievements', []))
    rules = content.get('operatingRules')
    if not isinstance(rules, dict): rules = {}
    name = ac.get('name', '')

    # S0 표지
    _put(_find_pos(S[0], 1.0, 0.8), cp.get('catch','') or ac.get('slogan',''))
    for i,(l,t) in enumerate([(4.3,8.5),(6.8,8.5),(9.3,8.5)]):
        _put(_find_pos(S[0], l, t), strengths[i]['title'] if i < len(strengths) else '')
    _put(_find_pos(S[0], 6.0,10.1), ac.get('phone',''))

    # S1 PHILOSOPHY + 교육특징 6
    _put(_find_pos(S[1], 1.2,1.4), f"왜 {name}인가 ?")
    _put(_find_pos(S[1], 1.2,2.4), cp.get('identity','') or cp.get('positioning',''))
    for i,(l,t) in enumerate([(2.3,4.6),(2.3,6.3),(2.3,8.0),(9.0,4.6),(9.0,6.3),(9.0,8.0)]):
        ts=_find_pos(S[1], l, t); ds=_find_pos(S[1], l, t+0.4)
        if i < len(feats):
            _put(ts, feats[i].get('title','')); _put(ds, feats[i].get('desc',''))
        else:
            _put(ts, ''); _put(ds, '')

    # S2 CURRICULUM 3컬럼 (초/중/고) — 학년 편차 스마트
    col_x  = {'초등':1.0, '중등':5.8, '고등':10.8}
    col_x2 = {'초등':2.2, '중등':7.0, '고등':12.0}
    gm = _grade_map(curric)
    for gk in ['초등','중등','고등']:
        x1, x2 = col_x[gk], col_x2[gk]
        c = gm.get(gk)
        title_sh = _find_pos(S[2], x1, 3.05, tol=0.4)
        if c:
            _put(title_sh, c.get('title', f'{gk} 과정'))
            _clear_region(S[2], x2-0.6, x2+0.6, 3.6, 4.95)
            cont = _find_pos(S[2], x2, 3.7, tol=0.3)
            lines=[]
            if c.get('content'):
                first=c['content'].split('.')[0]
                if len(first)>60: first=c['content'][:60]
                cur=''
                for w in first.split(' '):
                    if len(cur)+len(w)+1<=20: cur=(cur+' '+w).strip()
                    else: lines.append(cur); cur=w
                if cur: lines.append(cur)
                lines=lines[:4]
            _put_multi(cont, lines or [''])
            _put(_find_pos(S[2], x2, 5.2, tol=0.3), c.get('subject',''))
            _put(_find_pos(S[2], x2, 5.65, tol=0.3), c.get('goal',''))
        else:
            _clear_region(S[2], x1-0.3, x2+0.5, 2.9, 5.9)

    # 특별프로그램 2블록 → special 있으면, 없으면 관리로 대체
    sp_titles=[_find_pos(S[2],5.0,7.1),_find_pos(S[2],10.7,7.1)]
    sp_bodies=[_find_pos(S[2],4.7,7.7,tol=0.5),_find_pos(S[2],11.0,7.7,tol=0.5)]
    if special:
        for i in range(2):
            if i < len(special):
                sp=special[i]
                _put(sp_titles[i], sp.get('title',''))
                _put_multi(sp_bodies[i], sp.get('items',[]) if isinstance(sp.get('items'),list) else [sp.get('desc','')])
            else:
                _put(sp_titles[i],''); _put_multi(sp_bodies[i],[''])
    elif mgmt:
        _put(sp_titles[0], '학습 관리')
        _put_multi(sp_bodies[0], [m.get('title','') for m in mgmt[:4]])
        _put(sp_titles[1], ''); _put_multi(sp_bodies[1], [''])
    else:
        for t in sp_titles: _put(t,'')
        for b in sp_bodies: _put_multi(b,[''])

    # S3 관리시스템 6칸 (management + features 보충)
    _put(_find_pos(S[3], 1.7,3.1,tol=0.5), cp.get('positioning','') or '한 명도 놓치지 않는 관리 시스템으로 함께합니다.')
    mgmt_fill=list(mgmt)
    for f in feats:
        if len(mgmt_fill)>=6: break
        if not any(m.get('title')==f.get('title') for m in mgmt_fill):
            mgmt_fill.append({'title':f.get('title',''),'desc':f.get('desc','')})
    mslots=[(2.2,4.5),(5.3,4.6),(2.2,6.3),(5.3,6.3),(2.2,8.1),(5.3,8.1)]
    mdesc =[(1.6,5.1),(4.8,5.1),(1.6,6.9),(4.8,6.9),(1.6,8.7),(4.8,8.7)]
    for i,((tl,tt),(dl,dt)) in enumerate(zip(mslots,mdesc)):
        ts=_find_pos(S[3],tl,tt,tol=0.4); ds=_find_pos(S[3],dl,dt,tol=0.5)
        if i < len(mgmt_fill):
            _put(ts,mgmt_fill[i].get('title','')); _put(ds,mgmt_fill[i].get('desc',''))
        else:
            _put(ts,''); _put(ds,'')
    # 실적
    for i,(l,t) in enumerate([(8.7,3.7),(11.3,3.7),(8.7,4.3),(11.3,4.3)]):
        sh=_find_pos(S[3],l,t,tol=0.4)
        if sh and sh.text_frame.text.strip() not in ('초등','중등','고등'):
            _put(sh, ach[i].get('text','') if i<len(ach) else '')
    if not ach:
        for l,t in [(8.1,3.7),(10.7,3.7),(8.1,4.3)]:
            sh=_find_pos(S[3],l,t,tol=0.25)
            if sh and sh.text_frame.text.strip() in ('초등','중등','고등'): _put(sh,'')
    # 규정
    for (l,t),v in {(9.1,6.0):rules.get('homework'),(9.1,6.4):rules.get('absence'),
                    (9.1,6.7):rules.get('withdrawal'),(9.1,8.2):rules.get('refundBefore'),
                    (9.1,8.6):rules.get('refundAfter'),(9.1,8.9):rules.get('attendance')}.items():
        if v: _put(_find_pos(S[3], l, t, tol=0.3), v)

    # S4 입학안내
    _put(_find_pos(S[4], 8.4,2.6,tol=0.6), '주소 : ' + ac.get('address',''))
    kk=_find_pos(S[4], 9.4,3.6,tol=0.6)
    if kk and ac.get('kakao'): _put(kk, ac.get('kakao'))

    # S5 FAQ 6칸
    for i,(l,t) in enumerate([(1.2,3.3),(1.2,5.4),(1.2,7.4),(8.2,3.3),(8.2,5.4),(8.2,7.4)]):
        q=_find_pos(S[5], l, t, tol=0.3); a=_find_pos(S[5], l+0.5, t+0.5, tol=0.5)
        if i < len(faq):
            _put(q, f"Q.  {faq[i]['q']}"); _put(a, faq[i]['a'].replace('\n',' '))
        else:
            _put(q, ''); _put(a, '')

    prs.save(out_path)
    return out_path


def build_leaflet(content, out_path, template_dir='.'):
    design = content.get('design') or {}
    color = design.get('palette', 'green')
    if color not in LEAFLET_TEMPLATES: color = 'green'
    prs = Presentation(os.path.join(template_dir, LEAFLET_TEMPLATES[color]))
    S = list(prs.slides)
    ac = content.get('academy', {}); cp = content.get('copy', {})
    feats = _dicts(content.get('features', []) or content.get('strengths', []))
    strengths = _dicts(content.get('strengths', []))
    faq = _dicts(content.get('faq', []))
    curric = _dicts(content.get('curriculum', []))
    special = _dicts(content.get('special', []))
    adm = _dicts(content.get('admission', []))
    mgmt = _dicts(content.get('management', []))
    ach = _dicts(content.get('achievements', []))
    rules = content.get('operatingRules')
    if not isinstance(rules, dict): rules = {}
    name = ac.get('name', '')

    s0, s1 = S[0], S[1]
    # S0 FAQ 6
    for i,qt in enumerate([2.0,3.3,4.6,5.9,7.2,8.5]):
        q=_find_pos(s0,0.7,qt,tol=0.3); a=_find_pos(s0,1.0,qt+0.3,tol=0.4)
        if i < len(faq):
            _put(q, f"Q.  {faq[i]['q']}"); _put(a, faq[i]['a'].replace('\n',' '))
        else:
            _put(q, ''); _put(a, '')
    # S0 입학안내 3
    for i,at in enumerate([2.1,3.2,4.4]):
        if i < len(adm):
            _put(_find_pos(s0,6.5,at,tol=0.3), adm[i].get('step',''))
            _put(_find_pos(s0,8.2,at,tol=0.4), adm[i].get('desc','')[:20])
    _put(_find_pos(s0,5.7,6.7,tol=0.5), '주소 : '+ac.get('address',''))
    kk=_find_pos(s0,6.0,7.4,tol=0.5)
    if kk and ac.get('kakao'): _put(kk, '카카오톡 : '+ac.get('kakao'))
    # S0 표지
    _put(_find_pos(s0,11.0,0.9,tol=0.4), (cp.get('catch') or ac.get('slogan','')).replace(' ','\n',1))
    for i,(l,t) in enumerate([(11.0,8.2),(12.5,8.2),(13.9,8.2)]):
        sh=_find_pos(s0,l,t,tol=0.4)
        if sh and i < len(strengths): _put(sh, strengths[i].get('title','').replace(' ','\n'))
    _put(_find_pos(s0,11.7,9.4,tol=0.4), ac.get('phone',''))

    # S1 PHILOSOPHY + 교육특징 6
    _put(_find_pos(s1,0.5,0.8,tol=0.3), f"왜 {name}인가 ?")
    _put(_find_pos(s1,0.5,1.7,tol=0.4), cp.get('identity','') or cp.get('positioning','')[:40])
    for i,ft in enumerate([3.4,4.6,5.7,6.8,8.0,9.1]):
        ts=_find_pos(s1,1.3,ft,tol=0.25); ds=_find_pos(s1,1.3,ft+0.3,tol=0.25)
        if i < len(feats):
            _put(ts,feats[i].get('title','')); _put(ds,feats[i].get('desc','')[:30])
        else:
            _put(ts,''); _put(ds,'')

    # S1 CURRICULUM 세로 3단
    col_tops={'초등':(1.9,2.4,3.4,3.7),'중등':(4.4,4.8,5.8,6.1),'고등':(6.8,7.2,8.2,8.5)}
    gm=_grade_map(curric)
    for gk,(tt,ct,st,gt) in col_tops.items():
        c=gm.get(gk)
        _clear_region(s1,6.5,7.5,ct-0.1,ct+0.7)
        title_sh=_find_pos(s1,6.0,tt,tol=0.3)
        if c:
            _put(title_sh, c.get('title',f'{gk} 과정'))
            if c.get('content'):
                _put(_find_pos(s1,6.8,ct,tol=0.2), c['content'].split('.')[0][:40])
            _put(_find_pos(s1,6.8,st,tol=0.2), c.get('subject',''))
            _put(_find_pos(s1,6.8,gt,tol=0.2), c.get('goal',''))
        else:
            _clear_region(s1,5.9,7.5,tt-0.1,gt+0.3)

    # S1 특별프로그램
    sp=_find_pos(s1,5.8,9.6,tol=0.4)
    if special:
        items=special[0].get('items',[]) if isinstance(special[0],dict) else []
        _put(sp, '\n'.join(items[:4]) if items else '')
    elif strengths:
        _put(sp, '\n'.join(s.get('title','') for s in strengths[:4]))
    else:
        _put(sp, '')

    # S1 세번째 패널: 학습관리
    mgmt_fill=list(mgmt)
    for f in feats:
        if len(mgmt_fill)>=6: break
        if not any(m.get('title')==f.get('title') for m in mgmt_fill):
            mgmt_fill.append({'title':f.get('title',''),'desc':f.get('desc','')})
    mg=[(11.6,2.9,11.2,3.3),(13.7,2.9,13.3,3.3),(11.6,4.1,11.2,4.5),
        (13.7,4.1,13.3,4.5),(11.6,5.3,11.2,5.7),(13.7,5.3,13.3,5.7)]
    for i,(tl,tt,dl,dt) in enumerate(mg):
        ts=_find_pos(s1,tl,tt,tol=0.3); ds=_find_pos(s1,dl,dt,tol=0.4)
        if i < len(mgmt_fill):
            _put(ts,mgmt_fill[i].get('title','')); _put(ds,mgmt_fill[i].get('desc','')[:34])
        else:
            _put(ts,''); _put(ds,'')
    _put(_find_pos(s1,11.2,1.9,tol=0.4), cp.get('positioning','')[:50] or '한 명도 놓치지 않는 관리 시스템으로 함께합니다.')
    for i,(l,t) in enumerate([(11.6,6.9),(13.2,6.9),(11.6,7.4)]):
        sh=_find_pos(s1,l,t,tol=0.3)
        if sh and sh.text_frame.text.strip() not in ('초등','중등','고등'):
            _put(sh, ach[i].get('text','') if i<len(ach) else '')
    if not ach:
        for l,t in [(11.2,6.9),(12.9,6.9),(11.2,7.4)]:
            sh=_find_pos(s1,l,t,tol=0.2)
            if sh and sh.text_frame.text.strip() in ('초등','중등','고등'): _put(sh,'')
    for (l,t),v in {(11.9,8.4):rules.get('homework'),(11.9,8.6):rules.get('absence'),
                    (11.9,8.8):rules.get('withdrawal'),(11.9,9.8):rules.get('refundBefore'),
                    (11.9,10.0):rules.get('refundAfter'),(11.9,10.2):rules.get('attendance')}.items():
        if v: _put(_find_pos(s1,l,t,tol=0.2), v)

    prs.save(out_path)
    return out_path
