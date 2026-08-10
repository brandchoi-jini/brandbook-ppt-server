"""리플렛 스킨 'card_navy' — 수학귀신 시안 기준.

A4 가로 3단 접지 11.6929 x 8.2677in.
기존 build_leaflet_coord 의 그리기 기본기(_text/_rect/_photo_or_box)와
데이터 읽기(_admission_list 등)를 그대로 재사용하고, 배치만 새로 잡는다.

시안에서 뽑은 설계
  패널 x = 0.40 / 4.30 / 8.20, 콘텐츠 폭 3.10, 접지선 3.90 / 7.80
  색   네이비 #13233A · 틸 #3E8A80 · 연회색 #F6F8F9
  글자 라벨 10 · 헤더 15 · 본문 9.5~10 · 보조 8.5~9 · 표지명 27
  행   ①키-값 2열  ②번호칩 카드  ③틸 알약+값  ④연회색 행

키-값 2열 행이 핵심이다. 제목 위·설명 아래로 쌓는 기존 방식은 한 항목에
0.80in 이 들어가 학원 규정 4항목이 면에 들어가지 못했다. 좌우로 놓으면
0.30in 이면 되므로 같은 자리에 내용까지 담긴다.
"""

import io
import os
import re
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from build_leaflet_coord import (
    _s, _items, _first, _text, _rect, _need_h, _photo_or_box,
    _admission_list, _achievement_lines, _feature_list, _mgmt_list,
    _faq_list, W_IN, H_IN, MIN_BODY, _force_theme_font, FONT, _set_font,
    _strip_style,
)

# ── 색 ────────────────────────────────────────────────
NAVY = "13233A"
TEAL = "3E8A80"
SOFT = "F6F8F9"
WHITE = "FFFFFF"
INK = "1F2933"
BODY = "41505C"
MUTED = "8A97A2"
LINE = "DFE5E9"
ON_NAVY = "E8EDF2"

# ── 좌표 ──────────────────────────────────────────────
PX = (0.40, 4.30, 8.20)      # 패널 콘텐츠 시작 x
CWID = 3.10                  # 콘텐츠 폭
FOLDS = (3.90, 7.80)         # 접지선
Y_TOP = 0.62                 # 라벨 시작
Y_BOT = 7.78                 # 콘텐츠 하한 (푸터 7.98 위로 여유)

LABEL_PT = 10.0
HEAD_PT = 15.0
BODY_PT = 9.5
SMALL_PT = 9.0
TINY_PT = 8.5


# ══════════ 기본 블록 ══════════

def _hair(slide, x, y, w=CWID, color=LINE):
    _rect(slide, x, y, w, 0.008, fill=color)


def _label(slide, x, y, text, color=TEAL):
    _text(slide, text, x, y, CWID, 0.22, size=LABEL_PT, color=color,
          bold=True, fit=False)


def _head(slide, x, y, text, color=INK, w=CWID):
    """패널 헤더(15pt, 최대 2줄). 실제로 쓴 높이를 반환."""
    t = _s(text)
    if not t:
        return 0.0
    h = max(0.30, _need_h(t, w, HEAD_PT, line_spacing=1.22, pad=0.02))
    h = min(h, 0.86)
    _text(slide, t, x, y, w, h, size=HEAD_PT, color=color, bold=True,
          line_spacing=1.22, min_size=12.0)
    return h


def _kv(slide, x, y, k, v, kw=0.82, size=BODY_PT, kcolor=INK, vcolor=BODY,
        rule=True, w=CWID):
    """키-값 좌우 2열 행. 다음 y 를 반환.
    ★쌓기 방식(0.80in)의 40% 높이로 같은 내용을 담는다."""
    k, v = _s(k), _s(v)
    if not (k or v):
        return y
    if rule:
        _hair(slide, x, y, w)
        y += 0.08
    vh = max(0.22, _need_h(v, w - kw - 0.06, size, line_spacing=1.22, pad=0.02))
    if k:
        _text(slide, k, x, y, kw, 0.24, size=size, color=kcolor, bold=True,
              fit=False)
    if v:
        _text(slide, v, x + kw, y, w - kw - 0.04, vh, size=size, color=vcolor,
              line_spacing=1.22, min_size=MIN_BODY)
    return y + max(0.24, vh) + 0.08


def _pill(slide, x, y, text, w=1.14, h=0.24, fill=TEAL, color=WHITE,
          size=TINY_PT):
    """양끝이 반원인 캡슐. ★coord 의 radius=True 는 반경 0.035in 여서
    시안의 캡슐과 달리 각져 보였다. adjustment 0.5 = 높이의 절반 반경."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    shp.line.fill.background()
    # ★python-pptx 가 자동으로 넣는 <p:style>(그림자·테마색)을 제거한다.
    #   shadow.inherit=False 만으로는 테마 effectRef 가 남아 그림자가 찍혔다.
    _strip_style(shp)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    par = tf.paragraphs[0]
    par.alignment = PP_ALIGN.CENTER
    run = par.add_run()
    run.text = _s(text)
    # ★폭은 블록 안에서 통일한다. 긴 라벨은 폰트를 줄여 캡슐 안에 넣는다.
    _sz = size
    while _sz > 8.0 and _txt_w(_s(text), _sz) > (w - 0.14):
        _sz -= 0.5
    _set_font(run, size=_sz, bold=False, color=color)


def _txt_w(text, size):
    """표시 폭(inch) 근사. 한글 1자 ≈ size pt, 영숫자·공백 ≈ size*0.55pt."""
    w = 0.0
    for ch in _s(text):
        w += (size if ord(ch) > 0x2000 else size * 0.55)
    return w / 72.0


def _pill_w_uniform(labels, size=TINY_PT, lo=1.02, hi=1.60):
    """★한 블록 안의 캡슐 폭은 하나로 통일한다.
    길이마다 다르면 오른쪽 값의 시작선이 어긋나 지저분해 보인다.
    가장 긴 라벨을 기준으로 정하고, 상한을 넘으면 폰트로 흡수한다."""
    if not labels:
        return lo
    need = max(_txt_w(l, size) for l in labels if _s(l)) + 0.24
    return max(lo, min(hi, need))


def _pill_row(slide, x, y, pill, value, w=CWID, pw=None):
    """틸 알약 + 오른쪽 값. 다음 y 를 반환."""
    pill, value = _s(pill), _s(value)
    if not (pill or value):
        return y
    if pw is None:
        pw = 1.14
    _hair(slide, x, y, w)
    y += 0.06
    vh = max(0.24, _need_h(value, w - pw - 0.16, SMALL_PT + 0.5,
                           line_spacing=1.24, pad=0.02))
    if pill:
        _pill(slide, x, y, pill, w=pw)
    if value:
        _text(slide, value, x + pw + 0.12, y - 0.01, w - pw - 0.14, vh,
              size=SMALL_PT + 0.5, color=BODY, line_spacing=1.24,
              min_size=MIN_BODY)
    return y + max(0.26, vh) + 0.08


def _step_card(slide, x, y, no, title, desc, w=CWID, h=None):
    """번호칩 + 흰 카드. 다음 y 를 반환."""
    title, desc = _s(title), _s(desc)
    chip_w = 0.62
    tw = w - chip_w - 0.20
    dh = _need_h(desc, tw, TINY_PT, line_spacing=1.26, pad=0.02) if desc else 0.0
    ch = h or max(0.74, 0.28 + dh + 0.12)
    _rect(slide, x, y, w, ch, fill=WHITE, line=LINE, radius=True)
    _rect(slide, x, y, chip_w, ch, fill=NAVY, radius=True)
    _text(slide, str(no), x, y + (ch - 0.34) / 2.0, chip_w, 0.34,
          size=HEAD_PT, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
          fit=False)
    _text(slide, title, x + chip_w + 0.14, y + 0.13, tw, 0.24,
          size=11.0, color=INK, bold=True, min_size=9.5)
    if desc:
        _text(slide, desc, x + chip_w + 0.14, y + 0.39, tw, max(0.24, dh),
              size=TINY_PT, color=BODY, line_spacing=1.26, min_size=MIN_BODY)
    return y + ch + 0.10


def _feat_card(slide, x, y, pill, desc, w=CWID, pw=None):
    """틸 알약이 상단에 걸린 흰 카드. 다음 y 를 반환."""
    pill, desc = _s(pill), _s(desc)
    dh = _need_h(desc, w - 0.36, 10.0, line_spacing=1.30, pad=0.02) if desc else 0.0
    ch = max(0.68, 0.40 + dh + 0.10)
    _rect(slide, x, y + 0.12, w, ch, fill=WHITE, line=LINE, radius=True)
    if pill:
        pw = pw or min(w - 0.44, 1.60)
        _pill(slide, x + (w - pw) / 2.0, y, pill, w=pw, h=0.26, size=10.0)
    if desc:
        _text(slide, desc, x + 0.18, y + 0.42, w - 0.36, max(0.24, dh),
              size=10.0, color=BODY, line_spacing=1.30,
              align=PP_ALIGN.CENTER, min_size=MIN_BODY)
    return y + ch + 0.16


def _soft_row_h(label, desc, right="", w=CWID):
    """_soft_row 가 차지할 높이. ★자리 계산을 고정값(0.62)으로 했다가
    설명이 2줄인 행에서 푸터와 겹쳤다."""
    rw = 1.16 if _s(right) else 0.0
    dw = w - 0.32 - rw
    dh = _need_h(_s(desc), dw, BODY_PT, line_spacing=1.24, pad=0.02) if _s(desc) else 0.0
    return max(0.56, 0.30 + dh + 0.12) + 0.10


def _soft_row(slide, x, y, label, desc, right="", w=CWID):
    """연회색 행: 라벨 + 설명(+오른쪽 과목). 다음 y 를 반환."""
    label, desc, right = _s(label), _s(desc), _s(right)
    rw = 1.16 if right else 0.0
    dw = w - 0.32 - rw
    dh = _need_h(desc, dw, BODY_PT, line_spacing=1.24, pad=0.02) if desc else 0.0
    ch = max(0.56, 0.30 + dh + 0.12)
    _rect(slide, x, y, w, ch, fill=SOFT, radius=True)
    _text(slide, label, x + 0.16, y + 0.07, w - 0.32 - rw, 0.24,
          size=11.5, color=TEAL, bold=True, min_size=10.0)
    if right:
        _text(slide, right, x + w - rw - 0.14, y + 0.08, rw, 0.22,
              size=BODY_PT, color=TEAL, bold=True, align=PP_ALIGN.RIGHT,
              fit=False)
    if desc:
        _text(slide, desc, x + 0.16, y + 0.33, dw, max(0.22, dh),
              size=BODY_PT, color=BODY, line_spacing=1.24, min_size=MIN_BODY)
    return y + ch + 0.08


def _quote(slide, x, y, text, w=CWID):
    """네이비 인용 박스. 다음 y 를 반환."""
    t = _s(text)
    if not t:
        return y
    th = _need_h(t, w - 0.36, BODY_PT, line_spacing=1.34, pad=0.02)
    ch = max(0.72, th + 0.30)
    _rect(slide, x, y, w, ch, fill=NAVY, radius=True)
    _text(slide, t, x + 0.18, y + 0.15, w - 0.36, max(0.30, th),
          size=BODY_PT, color=WHITE, bold=True, line_spacing=1.34,
          min_size=MIN_BODY)
    return y + ch + 0.14


def _stat_box(slide, x, y, big, cap, w=CWID):
    """연회색 강조 박스(큰 숫자 + 캡션). 다음 y 를 반환."""
    big, cap = _s(big), _s(cap)
    if not (big or cap):
        return y
    ch = 0.92 if cap else 0.62
    _rect(slide, x, y, w, ch, fill=SOFT, radius=True)
    _text(slide, big, x + 0.18, y + 0.12, w - 0.36, 0.30, size=14.0,
          color=INK, bold=True, min_size=11.0)
    if cap:
        _text(slide, cap, x + 0.18, y + 0.46, w - 0.36, 0.40, size=SMALL_PT,
              color=BODY, line_spacing=1.26, min_size=MIN_BODY)
    return y + ch + 0.14


def _qr_png(url: str) -> Optional[io.BytesIO]:
    """QR 이미지 생성. 라이브러리가 없거나 실패하면 None."""
    u = _s(url)
    if not u:
        return None
    try:
        import qrcode
        img = qrcode.make(u)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf
    except Exception:
        return None


def _qr_grid(slide, x, y, items, w=CWID):
    """QR 2열 배치. items = [(캡션, url)]. 다음 y 를 반환."""
    items = [(c, u) for c, u in items if _s(u)]
    if not items:
        return y
    side, gap = 1.10, 0.16
    for i, (cap, url) in enumerate(items[:4]):
        cx = x + (side + gap) * (i % 2)
        cy = y + (side + 0.36) * (i // 2)
        png = _qr_png(url)
        if png is not None:
            _rect(slide, cx, cy, side, side, fill=WHITE, line=LINE)
            slide.shapes.add_picture(png, Inches(cx + 0.04), Inches(cy + 0.04),
                                     Inches(side - 0.08), Inches(side - 0.08))
        else:
            _rect(slide, cx, cy, side, side, fill=SOFT, line=LINE)
            _text(slide, "QR", cx, cy + side / 2 - 0.14, side, 0.26,
                  size=SMALL_PT, color=MUTED, align=PP_ALIGN.CENTER, fit=False)
        _text(slide, cap, cx, cy + side + 0.04, side, 0.20, size=TINY_PT,
              color=MUTED, align=PP_ALIGN.CENTER, fit=False)
    rows = (len(items[:4]) + 1) // 2
    return y + (side + 0.36) * rows + 0.06


# ══════════ 데이터 읽기 ══════════

def _basic(schema) -> Dict[str, Any]:
    raw = schema.get("_raw")
    b = raw.get("basic") if isinstance(raw, dict) else None
    return b if isinstance(b, dict) else {}


def _intro_ch(schema) -> Dict[str, Any]:
    raw = schema.get("_raw")
    c = raw.get("introChannel") if isinstance(raw, dict) else None
    return c if isinstance(c, dict) else {}


def _addr(schema) -> str:
    """★academy.location 은 '운양역 4번 출구 도보 1분…' 같은 오시는 길 설명이라
    주소 칸에 넣으면 안 된다. 사업장 주소를 먼저 쓴다."""
    ac = schema.get("academy") or {}
    return _first(_basic(schema).get("businessAddress"),
                  ac.get("address"), ac.get("address_short"),
                  ac.get("location"), default="")


def _ops(schema) -> Dict[str, Any]:
    """유미니 원본의 operations. 서버가 schema['_raw'] 로 붙여 준다."""
    raw = schema.get("_raw")
    if isinstance(raw, dict):
        ops = raw.get("operations")
        if isinstance(ops, dict):
            return ops
    ops = schema.get("operations")
    return ops if isinstance(ops, dict) else {}


def _hour_rows(schema) -> List[Tuple[str, str]]:
    """평일·토요일·일요일 운영시간 → 키-값 행."""
    ot = (_ops(schema).get("operatingTime") or {})
    if not isinstance(ot, dict):
        return []
    out = []
    for key, ko in (("weekdays", "평일"), ("saturday", "토요일"),
                    ("sunday", "일요일")):
        v = ot.get(key) or {}
        if not isinstance(v, dict):
            continue
        st, en = _s(v.get("start")), _s(v.get("end"))
        if st and en:
            out.append((ko, f"{st} – {en}"))
        elif not st and not en:
            out.append((ko, "휴원"))
    return out


def _fee_rows(schema) -> List[Tuple[str, str]]:
    """상담·테스트 소요시간과 비용 → 키-값 행. 납부 계좌는 싣지 않는다."""
    op = _ops(schema)
    out = []

    def _one(t_key, f_key, ko):
        t, f = _s(op.get(t_key)), op.get(f_key)
        if not t and f in (None, ""):
            return
        parts = []
        if t:
            parts.append(f"{t}분")
        # ★금액은 싣지 않는다(요청). 무료 여부만 알린다.
        if f in (0, "0"):
            parts.append("무료")
        if parts:
            out.append((ko, " · ".join(parts)))

    _one("consultingTime", "consultingFee", "상담")
    _one("testTime", "testFee", "레벨테스트")
    return out


_URL_RE = re.compile(r"https?://\S+")
_TAIL_RE = re.compile(r"(?:[·,/]\s*)?\[?\s*확인\s*필요\s*\]?[^·,/]*$")


def _clean_tail(text):
    """★앱이 60자에서 잘라 '고3T · [확인 필요] 테스트' 처럼 끝나는 경우가 있다.
    끊긴 꼬리와 매달린 구분자를 떼어낸다."""
    t = _s(text)
    if not t:
        return ""
    t = _TAIL_RE.sub("", t)
    return t.rstrip(" ·,/-—").strip()


def _no_url(text):
    """설명에 그대로 실려 오는 주소를 지운다(리플렛에는 QR 이 따로 있다)."""
    t = _URL_RE.sub("", _s(text))
    return re.sub(r"\s{2,}", " ", t).rstrip(" ·,/-—.").strip()


def _lead_sentence(text, n=46):
    """표지 슬로건에서 첫 문장만. ★"'너의 가능성을 믿는, 날 믿어봐!' 따뜻한
    관심과 열정으로…" 처럼 두 문장이 붙어 와 표지가 빽빽했다."""
    t = _s(text).strip()
    if not t:
        return ""
    # 여는 인용부호로 시작하면 닫는 부호까지를 한 문장으로 본다
    if t[0] in "\"'“‘":
        close = t.find(t[0] if t[0] in "\"'" else {"“": "”", "‘": "’"}[t[0]], 1)
        if 4 < close < n + 8:
            return t[:close + 1].strip()
    for mark in ("!", "?", ". ", "다.", "요."):
        i = t.find(mark)
        if 4 < i < n + 4:
            return t[:i + len(mark)].strip()
    return t if len(t) <= n else _clip_sent(t, n)


def _clip_sent(text, n):
    """문장 단위로 잘라 n자 이내로. ★소개문이 길어 패널이 빽빽했다."""
    t = _s(text)
    if len(t) <= n:
        return t
    cut = t[:n]
    for mark in ("다. ", "다.", ". "):
        i = cut.rfind(mark)
        if i > n * 0.5:
            return cut[:i + len(mark)].strip()
    return cut.rstrip() + "…"


def _ach_rows(schema) -> List[Tuple[str, str]]:
    """실적 → (라벨, 값). ★items 는 {title, change, note} 구조다.
    기존 헬퍼가 문자열로 합친 것을 다시 쪼개다 보니
    '2021~2025년 매년 수상' 같은 설명이 알약 라벨로 들어갔다."""
    ach = schema.get("achievements") or {}
    src = _items(ach)
    if not src and isinstance(ach, dict):
        for g in (ach.get("groups") or []):
            src.extend(_items(g))
    out = []
    for it in src[:4]:
        label = _s(it.get("title") or it.get("name"))
        val = " · ".join(v for v in [_s(it.get("change")), _s(it.get("note"))] if v)
        if not val:
            val = _s(it.get("desc") or it.get("body"))
        if label or val:
            out.append((label, val))
    return out


def _rule_rows(schema) -> List[Tuple[str, str]]:
    src = _items(schema.get("rules")) or _items(schema.get("policy"))
    out = []
    for it in src[:6]:
        k = _s(it.get("k") or it.get("title") or it.get("name") or it.get("key"))
        v = _s(it.get("v") or it.get("desc") or it.get("value") or it.get("body"))
        if k or v:
            out.append((k, v))
    return out


def _mgmt_rows(schema) -> List[Tuple[str, str]]:
    """_mgmt_list 는 '제목 — 설명' 문자열을 준다. 알약 행으로 쓰려면 쪼갠다."""
    out = []
    for line in _mgmt_list(schema)[:6]:
        s = _s(line)
        if " — " in s:
            k, v = s.split(" — ", 1)
        else:
            k, v = s, ""
        out.append((k.strip(), v.strip()))
    return out


def _target_rows(schema) -> List[Tuple[str, str, str]]:
    """수업 대상 → (학년, 설명, 과목)."""
    tg = schema.get("targets") or {}
    out = []
    for it in _items(tg)[:4]:
        g = _s(it.get("grade") or it.get("title") or it.get("name"))
        # ★desc 는 항상 빈 문자열로 오고 실제 내용은 subj(개설 반 나열)에 있다.
        d = _clean_tail(_first(it.get("desc"), it.get("body"),
                                it.get("subj"), default=""))
        sj = _clean_tail(_s(it.get("subject") or it.get("subjects")))
        if sj == d:
            sj = ""
        if g or d:
            out.append((g, d, sj))
    return out


def _level_rows(schema) -> List[Tuple[str, str]]:
    """반 구성 → (구간 라벨, 반 이름 나열).
    ★실제 스키마는 classes = {head, groups:[{grade, rows:[{name,desc}]}]} 다.
      items 를 찾고 있어서 반 구성 블록이 늘 비어 있었다."""
    cl = schema.get("classes") or {}
    out = []
    # ★수업 대상 블록과 같은 내용이면 싣지 않는다(잘품에서 학년 목록이
    #   두 면에 똑같이 나왔다). 이미 쓴 반 이름을 모아 비교한다.
    seen = set()
    for _g, _d, _sj in _target_rows(schema):
        for tok in re.split(r"[·,/]", _s(_d) + " " + _s(_sj)):
            tok = tok.strip()
            if tok:
                seen.add(tok)
    groups = cl.get("groups") if isinstance(cl, dict) else None
    if isinstance(groups, list) and groups:
        for g in groups[:3]:
            if not isinstance(g, dict):
                continue
            head = _s(g.get("grade") or g.get("name") or g.get("title"))
            # '확인 필요' 같은 판정 꼬리가 그룹 라벨로 올라오는 경우 제외
            if not head or "확인" in head:
                continue
            names = []
            for r in (g.get("rows") or g.get("items") or []):
                nm = _s(r.get("name") or r.get("title")) if isinstance(r, dict) else _s(r)
                if nm:
                    names.append(nm)
            body = _clean_tail(" · ".join(names))
            # 수업 대상에 이미 나온 반 이름이 대부분이면 중복이므로 건너뛴다
            if names and seen:
                dup = sum(1 for n in names if n.strip() in seen)
                if dup >= max(1, int(len(names) * 0.6)):
                    continue
            if head or body:
                out.append((head, body))
        return out
    # 폴백: items 형태로 오는 경우
    for it in _items(cl)[:3]:
        head = _s(it.get("name") or it.get("subject") or it.get("title"))
        body = _clean_tail(_s(it.get("desc") or it.get("body")))
        if head or body:
            out.append((head, body))
    return out


def _specials_line(schema) -> str:
    names = []
    for it in _items(schema.get("specials")):
        t = _s(it.get("title") or it.get("name"))
        if t:
            names.append(t)
    return " · ".join(names)


def _channels(schema) -> List[Tuple[str, str]]:
    """★유미니 실제 키는 naverMapUrl · kakaoMapUrl · blogUrl · snsUrl ·
    kakaoChannel 이다. kakaoChatUrl·instagramUrl 을 찾고 있어 못 읽었다."""
    ac = schema.get("academy") or {}
    ch = _intro_ch(schema)

    def _u(v):
        # ★kakaoChannel 은 문자열이 아니라 {profileId, chatUrl} 객체로 온다.
        if isinstance(v, dict):
            v = v.get("chatUrl") or v.get("url") or v.get("profileId") or ""
        v = _s(v)
        if not v:
            return ""
        if v.startswith("http"):
            return v
        if v.startswith("@"):            # 카카오 채널 ID
            return "https://pf.kakao.com/" + v.lstrip("@")
        return ""

    out = [
        ("네이버 지도", _u(_first(ac.get("naverMapUrl"), ch.get("naverMapUrl"),
                             default=""))),
        ("카카오톡 상담", _u(ch.get("kakaoChannel") or ch.get("kakaoMapUrl"))),
        ("네이버 블로그", _u(ch.get("blogUrl"))),
        ("SNS", _u(ch.get("snsUrl"))),
    ]
    return [(c, u) for c, u in out if u]


# ══════════ 앞면 ══════════

def _front_contact(slide, schema, x):
    """① 상담 문의 — 전화·주소·운영시간·예약 채널."""
    ac = schema.get("academy") or {}
    _label(slide, x, Y_TOP, "상담 문의")
    y = Y_TOP + 0.24
    region = _s(ac.get("region"))
    name = _s(ac.get("name"), "우리학원")
    y += _head(slide, x, y, (region + "\n" + name) if region else name) + 0.10

    y = _kv(slide, x, y, "전화", _s(ac.get("phone")))
    y = _kv(slide, x, y, "주소", _addr(schema))

    hours = _hour_rows(schema)
    if hours:
        y += 0.12
        _label(slide, x, y, "운영시간")
        y += 0.24
        for k, v in hours:
            y = _kv(slide, x, y, k, v, kw=0.94)

    # ── 약도 ──
    # ★학원마다 약도 유무가 다르다. 있으면 남는 자리를 계산해 넣고,
    #   없으면 이 블록 자체를 건너뛴다(빈 칸을 만들지 않는다).
    _as = schema.get("assets") or {}
    _map = _s(_as.get("map"))
    _ph = _as.get("photos")
    _photo = _s(_ph[0]) if isinstance(_ph, list) and _ph else _s(_as.get("cover"))
    chs = [(c, u) for c, u in _channels(schema) if _s(u)]
    _qr_h = ((1.10 + 0.36) * ((len(chs[:4]) + 1) // 2) + 0.30) if chs else 0.0
    # ★약도 자리는 항상 확보한다(요청). 약도가 있으면 약도, 없으면 학원 사진,
    #   둘 다 없으면 자리만 남겨 나중에 이미지를 채울 수 있게 한다.
    _img = _map or _photo
    _cap = "오시는 길" if _map else ("학원 전경" if _photo else "오시는 길")
    _avail = Y_BOT - 0.56 - _qr_h - y - 0.30
    _mh = max(1.05, min(1.90, _avail))
    if _mh >= 1.00:
        y += 0.14
        _label(slide, x, y, _cap)
        y += 0.25
        if _img:
            _photo_or_box(slide, _img, x, y, CWID, _mh,
                          {"soft": SOFT, "body": MUTED, "card_line": LINE},
                          _cap, cover_mode=True)
        else:
            _rect(slide, x, y, CWID, _mh, fill=SOFT, line=LINE, radius=True)
            _text(slide, "약도 이미지", x, y + _mh / 2 - 0.14, CWID, 0.26,
                  size=SMALL_PT, color=MUTED, align=PP_ALIGN.CENTER, fit=False)
        y += _mh + 0.14

    if chs and y < Y_BOT - _qr_h - 0.20:
        y += 0.10
        _label(slide, x, y, "예약 · 채널")
        y += 0.24
        y = _qr_grid(slide, x, y, chs)

    tail = _first((schema.get("copy") or {}).get("contactLine"),
                  default="학생의 현재 위치를 정확히 확인하는 자리입니다. 부담 없이 문의해 주세요.")
    if y < Y_BOT - 0.50:
        _text(slide, tail, x, min(y + 0.10, Y_BOT - 0.46), CWID, 0.44,
              size=BODY_PT, color=BODY, line_spacing=1.30, min_size=MIN_BODY)


def _front_admission(slide, schema, x):
    """② 입학 안내 — 단계 카드 + 등록 안내 + 인용."""
    ad = schema.get("admission") or {}
    _label(slide, x, Y_TOP, "입학 안내")
    y = Y_TOP + 0.24
    y += _head(slide, x, y, _first(ad.get("head") if isinstance(ad, dict) else "",
                                   default="테스트와 상담을 거쳐 맞는 반을 배정합니다")) + 0.14

    steps = _admission_list(schema)
    fees = _fee_rows(schema)
    # ★표지에 쓰는 coverLine 을 여기 또 쓰면 같은 문장이 두 면에 나온다.
    _cp = schema.get("copy") or {}
    quote = _first(_cp.get("admissionQuote"), _cp.get("promise"), default="")

    # 아래 블록에 필요한 자리를 먼저 떼어 둔다
    reserve = 0.0
    if fees:
        reserve += 0.30 + 0.42 * len(fees)      # ★2줄로 늘어나는 행이 있었다
    if quote:
        reserve += 0.92
    # ★잘품에서 '레벨테스트 30분·무료' 줄이 푸터와 겹쳐 인쇄 시 잘렸다.
    y_steps_end = max(y + 1.0, min(Y_BOT - reserve - 0.24, 7.30))

    if steps:
        each = (y_steps_end - y) / max(1, len(steps)) - 0.12
        each = max(0.62, min(1.02, each))
        for i, st in enumerate(steps):
            y = _step_card(slide, x, y, i + 1, st.get("title"),
                           _no_url(st.get("desc")), h=each)

    y = max(y, y_steps_end - 0.10)
    if fees:
        _label(slide, x, y, "등록 안내")
        y += 0.24
        for k, v in fees:
            y = _kv(slide, x, y, k, v, kw=1.00, size=SMALL_PT)
    if quote and y < Y_BOT - 0.70:
        _quote(slide, x, min(y + 0.10, Y_BOT - 0.94), quote)


def _front_cover(slide, schema, x, assets):
    """③ 표지 — 네이비 전면."""
    ac = schema.get("academy") or {}
    cp = schema.get("copy") or {}
    _rect(slide, FOLDS[1], 0.0, W_IN - FOLDS[1], H_IN, fill=NAVY)

    y = 0.66
    logo = _s(assets.get("logo"))
    if logo:
        _rect(slide, x, y, 1.22, 1.22, fill=WHITE, radius=True)
        _photo_or_box(slide, logo, x + 0.07, y + 0.07, 1.08, 1.08,
                      {"soft": WHITE, "body": MUTED, "card_line": LINE},
                      "", cover_mode=False)
        y += 1.34
    else:
        y += 0.20

    tags = " · ".join(t for t in [_s(ac.get("grades")), _s(ac.get("subjects"))] if t)
    if tags:
        _text(slide, tags, x, y, CWID, 0.22, size=BODY_PT, color=TEAL,
              bold=True, fit=False)
        y += 0.28

    _text(slide, _s(ac.get("name"), "우리학원"), x, y, CWID, 0.62, size=27.0,
          color=WHITE, bold=True, min_size=18.0)
    y += 0.78

    slogan = _first(cp.get("slogan"), ac.get("slogan"), schema.get("slogan"),
                    default="")
    # ★"'너의 가능성을 믿는, 날 믿어봐!' 따뜻한 관심과 열정으로…" 처럼
    #   두 문장이 붙어 오면 표지가 빽빽해진다 → 첫 문장만 쓴다.
    slogan = _lead_sentence(slogan, 46)
    if slogan:
        sh = min(1.20, max(0.44, _need_h(slogan, CWID, HEAD_PT,
                                         line_spacing=1.34, pad=0.02)))
        _text(slide, slogan, x, y, CWID, sh, size=HEAD_PT, color=WHITE,
              bold=True, line_spacing=1.34, min_size=12.0)
        y += sh + 0.24

    _rect(slide, x, y, 1.20, 0.014, fill=TEAL)
    y += 0.18

    cover = _s(cp.get("coverLine"))
    if cover:
        ch = max(0.30, _need_h(cover, CWID, 10.5, line_spacing=1.32, pad=0.02))
        _text(slide, cover, x, y, CWID, min(ch, 0.70), size=10.5, color=TEAL,
              bold=True, line_spacing=1.32, min_size=MIN_BODY)
        y += min(ch, 0.70) + 0.12

    voice = _s(cp.get("directorVoice") or cp.get("quote"))
    if voice:
        _text(slide, voice, x, y, CWID, 0.26, size=BODY_PT, color=ON_NAVY,
              min_size=MIN_BODY)
        y += 0.36

    # ★copy.principles 가 없는 학원이 많다 → 강점 이름으로 대신 채운다.
    _bul = [_s(v) for v in (cp.get("principles") or []) if _s(v)]
    if not _bul:
        _bul = [_s(f.get("title")) for f in _feature_list(schema) if _s(f.get("title"))]
    if not _bul:
        _bul = [k for k, _v in _mgmt_rows(schema) if k]
    for b in _bul[:3]:
        _text(slide, "· " + b, x, y, CWID, 0.24, size=BODY_PT, color=TEAL,
              min_size=MIN_BODY)
        y += 0.28

    _text(slide, _s(ac.get("phone")), x, 7.14, CWID, 0.26, size=12.0,
          color=WHITE, bold=True, fit=False)
    _text(slide, _addr(schema), x, 7.42, CWID, 0.36, size=TINY_PT,
          color=ON_NAVY, line_spacing=1.24, min_size=MIN_BODY)


# ══════════ 뒷면 ══════════

def _back_why(slide, schema, x):
    """① 왜 ○○인가 — 인용 + 이야기 + 실적 + 강조."""
    ac = schema.get("academy") or {}
    cp = schema.get("copy") or {}
    _label(slide, x, Y_TOP, "왜 " + _s(ac.get("name"), "우리학원") + "인가")
    y = Y_TOP + 0.24
    y += _head(slide, x, y, _first(cp.get("identity"), cp.get("slogan"),
                                   default="")) + 0.14

    y = _quote(slide, x, y, _s(cp.get("promise") or cp.get("coverLine")))

    intro = _s((schema.get("intro") or {}).get("body") or cp.get("intro"))
    if intro:
        intro = _clip_sent(intro, 240)
        ih = min(1.60, max(0.40, _need_h(intro, CWID, 10.0,
                                         line_spacing=1.52, pad=0.02)))
        _text(slide, intro, x, y, CWID, ih, size=10.0, color=BODY,
              line_spacing=1.52, min_size=MIN_BODY)
        y += ih + 0.20

    rows = _ach_rows(schema)
    if rows:
        _label(slide, x, y, "주요 실적")
        y += 0.24
        _pw = _pill_w_uniform([k for k, _v in rows])
        for k, v in rows:
            y = _pill_row(slide, x, y, k, v, pw=_pw)

    st = (schema.get("achievements") or {})
    hl = st.get("highlight") if isinstance(st, dict) else None
    if isinstance(hl, dict) and y < Y_BOT - 1.00:
        y = _stat_box(slide, x, min(y + 0.06, Y_BOT - 0.98),
                      _s(hl.get("title")), _s(hl.get("desc")))
    # ★남는 자리를 사진으로 채우던 기능은 뺐다. assets.photos 에 규정 문서
    #   스캔이 섞여 들어와 '학원 전경'으로 인쇄되는 사고가 났다.
    #   빈 자리가 잘못된 이미지보다 낫다.


def _back_features(slide, schema, x):
    """② 교육 특징 — 알약 카드 + 수업 대상."""
    cp = schema.get("copy") or {}
    _label(slide, x, Y_TOP, "교육 특징")
    y = Y_TOP + 0.24
    y += _head(slide, x, y, _s(cp.get("featureHead") or schema.get("feat_head"))
               or "세분하고, 첨삭하고, 확인하고 넘어갑니다") + 0.14

    feats = _feature_list(schema)[:4]
    targets = _target_rows(schema)
    # ★고정값(0.62/행)으로 자리를 잡다가 설명이 2줄인 행에서 푸터와 겹쳤다.
    #   실제 높이를 재서 확보하고, 감당 못 하면 뒤 항목부터 덜어낸다.
    t_h = [_soft_row_h(g, d, sj) for g, d, sj in targets]
    reserve = (0.30 + sum(t_h)) if targets else 0.0
    while targets and reserve > 3.40:
        targets.pop(); t_h.pop()
        reserve = 0.30 + sum(t_h)
    y_end = Y_BOT - reserve - 0.10

    _fpw = _pill_w_uniform([f.get("title") for f in feats], size=10.0,
                           lo=1.30, hi=2.50)
    for f in feats:
        if y > y_end - 0.60:
            break
        y = _feat_card(slide, x, y, f.get("title"), f.get("desc"), pw=_fpw)

    if targets:
        y = max(y, y_end - 0.06)
        _label(slide, x, y, "수업 대상 · 과목")
        y += 0.24
        for i, (g, d, sj) in enumerate(targets):
            if y + t_h[i] > Y_BOT:            # 하한 강제 — 겹치기 전에 멈춘다
                break
            y = _soft_row(slide, x, y, g, d, right=sj)


def _pill_row_h(pill, value, w=CWID, pw=1.14):
    """_pill_row 가 차지할 높이(간격 제외)."""
    vh = max(0.24, _need_h(_s(value), w - pw - 0.16, SMALL_PT + 0.5,
                           line_spacing=1.24, pad=0.02))
    return 0.06 + max(0.26, vh)


def _back_manage(slide, schema, x):
    """③ 반 구성 · 수업 관리 — 레벨 + 관리 알약 + 특강 한 줄.

    ★특강을 면 하단에 고정해 두어, 반 구성이 생략된 학원(잘품)에서
      수업 관리와 특강 사이가 2인치 넘게 비었다.
      이제 위에서부터 흐르게 놓고, 남는 여백은 행 간격으로 고르게 나눈다.
    """
    cp = schema.get("copy") or {}
    levels = _level_rows(schema)
    mrows = _mgmt_rows(schema)
    sp = _specials_line(schema)

    # 반 구성이 없으면 제목에서도 뺀다(빈 항목을 제목이 약속하지 않게)
    _title = "반 구성 · 수업 관리" if levels else "수업 관리"
    _label(slide, x, Y_TOP, _title)
    y = Y_TOP + 0.24
    y += _head(slide, x, y, _s(cp.get("classHead"))
               or ("학년별 단계에 맞춰 반을 배정합니다" if levels
                   else "배우고, 확인하고, 끝까지 봅니다")) + 0.14

    # ── 필요한 높이를 먼저 다 재고, 남는 만큼만 간격으로 나눈다 ──
    lv_h = []
    for _hd, bd in levels:
        lv_h.append(0.42 + max(0.26, _need_h(bd, CWID - 0.32, 12.0,
                                             line_spacing=1.24, pad=0.02)))
    _pw = _pill_w_uniform([k for k, _v in mrows]) if mrows else 1.14
    mg_h = [_pill_row_h(k, v, pw=_pw) for k, v in mrows]
    sp_h = _need_h(sp, CWID, BODY_PT, line_spacing=1.30, pad=0.02) if sp else 0.0

    need = sum(lv_h) + sum(mg_h) + sp_h
    need += 0.24 * (1 if levels and mrows else 0)      # 수업 관리 라벨
    need += 0.35 if sp else 0.0                        # 특강 라벨
    slots = max(1, len(lv_h) + len(mg_h))
    slack = max(0.0, (Y_BOT - y) - need)
    gap = min(0.26, 0.08 + slack / (slots + 2))        # 과하게 벌리지 않는다

    for i, (hd, bd) in enumerate(levels):
        if y + lv_h[i] > Y_BOT:
            break
        _rect(slide, x, y, CWID, lv_h[i], fill=SOFT, radius=True)
        _text(slide, hd, x + 0.16, y + 0.09, CWID - 0.32, 0.22,
              size=BODY_PT, color=TEAL, bold=True, min_size=MIN_BODY)
        _text(slide, bd, x + 0.16, y + 0.34, CWID - 0.32, lv_h[i] - 0.42,
              size=12.0, color=INK, bold=True, min_size=9.5)
        y += lv_h[i] + gap

    if mrows:
        if levels:
            _label(slide, x, y, "수업 관리")
            y += 0.26
        for i, (k, v) in enumerate(mrows):
            if y + mg_h[i] > Y_BOT - (sp_h + 0.35 if sp else 0.0):
                break
            y = _pill_row(slide, x, y, k, v, pw=_pw)
            y += max(0.0, gap - 0.08)

    if sp:
        y = min(y + 0.10, Y_BOT - sp_h - 0.28)
        _label(slide, x, y, "특강 및 기타 수업")
        y += 0.26
        _text(slide, sp, x, y, CWID, max(0.26, sp_h), size=BODY_PT,
              color=BODY, line_spacing=1.30, min_size=MIN_BODY)


# ══════════ 조립 ══════════

def _new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _folds(slide):
    for fx in FOLDS:
        _rect(slide, fx, 0.0, 0.006, H_IN, fill="E3E8EC")


def _footer(slide, schema, panels=(0, 1)):
    ac = schema.get("academy") or {}
    tail = " · ".join(t for t in [_s(ac.get("name")), _s(ac.get("phone"))] if t)
    if not tail:
        return
    for i in panels:
        _text(slide, tail, PX[i], 7.98, CWID, 0.20, size=8.0,
              color=MUTED, fit=False)


def build(schema: Dict[str, Any], palette: str = "card_navy",
          out: Union[str, os.PathLike, io.BytesIO, None] = None,
          assets: Optional[Dict[str, str]] = None):
    _assets = dict(schema.get("assets") or {})
    if isinstance(assets, dict):
        _assets.update(assets)

    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    _force_theme_font(prs, FONT)

    # 앞면
    s1 = _new_slide(prs)
    _folds(s1)
    _front_contact(s1, schema, PX[0])
    _front_admission(s1, schema, PX[1])
    _front_cover(s1, schema, PX[2], _assets)
    _footer(s1, schema, panels=(0, 1))

    # 뒷면
    s2 = _new_slide(prs)
    _folds(s2)
    _back_why(s2, schema, PX[0])
    _back_features(s2, schema, PX[1])
    _back_manage(s2, schema, PX[2])
    _footer(s2, schema, panels=(0, 1, 2))

    if out is None:
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf
    prs.save(out)
    return out
