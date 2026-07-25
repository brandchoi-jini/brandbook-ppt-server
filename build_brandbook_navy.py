# -*- coding: utf-8 -*-
"""
build_brandbook_navy.py — 브랜드북 navy 스킨 (PPT 12장)
기준점: 더바른수학전문학원 초등부 소개서(네이비강조)

입력: v3 표준스키마 (book 스킨과 동일 계약)
출력: .pptx (경로 또는 BytesIO)

스마트 채움:
  - 데이터 없는 섹션은 슬라이드 자체를 만들지 않는다
  - 학년(초/중/고)은 실제 있는 것만 렌더
  - 카드 개수는 데이터 수에 맞춰 3열/2열/1열 자동
  - 페이지 번호는 실제 생성된 순서로 다시 매긴다
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from navy_core import (
    _rgb, _latin_font, FONT,
    G, FS, get_palette, blank_slide, page_head, add_box, add_text, add_pill,
    place_image, clean, fit_size, gv, as_dicts, nonempty, wrap_words,
    force_theme_font,
)

SLIDE_W = 13.333
SLIDE_H = 7.5

# 모든 박스 공통 라운드 반경 (살짝 둥근 정도로 통일)
R = 0.055


# ══════════════════════════════════════════════════════════════════
# 여백 채움 헬퍼
# ══════════════════════════════════════════════════════════════════
def _photo_pool(d):
    """본문에 쓸 수 있는 사진 목록. 표지 사진은 중복을 피해 뒤로 돌린다."""
    a = d.get("assets") or {}
    pool = []
    for p in (a.get("photos") or []):
        if p and p not in pool:
            pool.append(p)
    for k in ("cover", "photo"):
        v = a.get(k)
        if v and v not in pool:
            pool.append(v)
    return pool


def take_photo(d, used):
    """아직 안 쓴 사진 하나. 없으면 None (있는 것만 재사용은 하지 않는다)."""
    for p in _photo_pool(d):
        if p not in used:
            used.add(p)
            return p
    return None


def fill_gap(s, d, P, used, y_from, x=None, w=None, y_to=6.55,
             min_h=1.15, caption=None):
    if used is None:
        used = set()
    """
    페이지 아래쪽 빈 공간에 사진을 넣는다.
    남은 높이가 min_h 미만이면 아무것도 하지 않는다(억지로 늘리지 않음).
    사진이 없으면 얇은 구분선만 그어 마무리한다.
    """
    x = G.ML if x is None else x
    w = G.CW if w is None else w
    h = y_to - y_from
    if h < min_h:
        return False

    src = take_photo(d, used)
    if not src:
        # 사진이 없으면 억지로 늘리지 않는다.
        # 다만 여백이 크게 남으면 마감선을 그어 '미완성처럼' 보이지 않게 한다.
        if h >= 1.6:
            add_box(s, x, y_from + 0.1, w, 0.012, fill=P["line"])
        return False

    h = min(h, 2.9)
    # cover 는 비율 유지를 위해 영역 밖으로 넘칠 수 있다(슬라이드 이탈).
    # 여백 채움에서는 영역 안에 들어오는 contain 방식을 쓴다.
    add_box(s, x, y_from, w, h, fill=P["card"], radius=R)
    ok = place_image(s, src, x + 0.06, y_from + 0.06,
                     w - 0.12, h - 0.12, cover=False)
    if ok is None:
        return False
    if caption:
        add_text(s, x + 0.24, y_from + h - 0.44, w - 0.48, 0.3, caption,
                 size=FS.small, bold=True, color=P["onDark"],
                 line_spacing=1.0, shrink=False, clip=False)
    return True


# ══════════════════════════════════════════════════════════════════
# 01. 표지
# ══════════════════════════════════════════════════════════════════
def s_cover(prs, d, P, used=None):
    s = blank_slide(prs)
    b = d.get("basic", {})
    name = clean(b.get("name") or "학원")
    slogan = clean(b.get("slogan") or gv(d, "identity", "slogan") or "")
    target = clean(b.get("target") or "")
    subline = clean(b.get("subline") or gv(d, "identity", "oneline") or "")
    phone = clean(b.get("phone") or "")
    addr = clean(b.get("address") or "")
    logo = gv(d, "assets", "logo")
    photo = gv(d, "assets", "cover") or gv(d, "assets", "photo")

    # 우측 딥네이비 패널 + 사진
    panel_x = 7.66
    panel_w = SLIDE_W - panel_x
    add_box(s, panel_x, 0, panel_w, SLIDE_H, fill=P["deep"])
    if photo:
        place_image(s, photo, panel_x, 0, panel_w, SLIDE_H, cover=True)
        if used is not None:
            used.add(photo)

    # 로고
    y = 0.47
    if logo:
        place_image(s, logo, 0.85, y, 1.7, 1.0, cover=False)
        y = 1.88
    else:
        y = 0.9

    # 대상 배지
    if target:
        w = min(3.0, 0.42 + len(target) * 0.135)
        add_pill(s, 0.6, y, w, 0.29, target, fill=P["primary"],
                 color=P["onDark"], size=FS.small)
        y += 0.67
    else:
        y += 0.2

    # 슬로건 (2줄까지)
    if slogan:
        lines = wrap_words(slogan, 17)[:2]
        add_text(s, 0.6, y, 6.6, 1.45, lines, size=FS.cover_ttl, bold=True,
                 color=P["text"], line_spacing=1.28)
        y += 1.62

    add_text(s, 0.6, y, 5.6, 0.35, name, size=FS.cover_sub, bold=True,
             color=P["text"], line_spacing=1.0)
    y += 0.57

    if subline:
        add_text(s, 0.6, y, 6.5, 0.42, subline, size=FS.body,
                 color=P["muted"], line_spacing=1.3)

    foot = "  |  ".join([x for x in (phone, addr) if x])
    if foot:
        add_text(s, 0.6, 6.65, 6.6, 0.25, foot, size=FS.small,
                 color=P["muted"], line_spacing=1.0)
    return s


# ══════════════════════════════════════════════════════════════════
# 02. 강점 카드 (2~6개, 3열 그리드)
# ══════════════════════════════════════════════════════════════════
def s_strengths(prs, d, P, page, used=None):
    items = as_dicts(d.get("strengths") or gv(d, "identity", "strengths"))
    if not items:
        return None
    items = items[:6]
    s = blank_slide(prs)
    hb = page_head(s, P,
              "핵심 강점",
              gv(d, "identity", "headline") or "우리 학원이 만드는 변화",
              gv(d, "identity", "lead"),
              page)

    n = len(items)
    cols = 3 if n >= 3 else n
    rows = -(-n // cols)
    gap = 0.34
    cw = (G.CW - gap * (cols - 1)) / cols
    ch = 1.73 if rows >= 2 else 2.1
    top = max(G.y_body, hb)

    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        x = G.ML + c * (cw + gap)
        y = top + r * (ch + 0.4)
        add_box(s, x, y, cw, ch, fill=P["card"], line=P["line"], radius=R)
        add_text(s, x + 0.23, y + 0.19, 0.6, 0.25, f"{i+1:02d}",
                 size=FS.card_num, bold=True, color=P["primary"], line_spacing=1.0)
        ttl = clean(it.get("title") or it.get("name") or "")
        add_text(s, x + 0.23, y + 0.56, cw - 0.46, 0.33, ttl,
                 size=fit_size(ttl, FS.card_ttl, cw - 0.46, 0.36),
                 bold=True, color=P["text"], line_spacing=1.15)
        body = clean(it.get("desc") or it.get("body") or it.get("text") or "")
        if body:
            add_text(s, x + 0.25, y + 1.02, cw - 0.5, ch - 1.22, body,
                     size=fit_size(body, FS.body, cw - 0.5, ch - 1.22),
                     color=P["muted"], line_spacing=1.4)

    # 카드 아래가 비면 학원 사진으로 채운다
    bottom = top + rows * ch + (rows - 1) * 0.4
    fill_gap(s, d, P, used, bottom + 0.42)
    return s


# ══════════════════════════════════════════════════════════════════
# 03. 철학 / 왜 우리인가 (좌 체크리스트 + 우 강조박스)
# ══════════════════════════════════════════════════════════════════
def s_philosophy(prs, d, P, page, used=None):
    ph = d.get("philosophy") or {}
    points = as_dicts(ph.get("points"), "title")
    intro = clean(ph.get("intro") or gv(d, "identity", "intro") or "")
    note = ph.get("note") or {}
    if not points and not intro:
        return None
    s = blank_slide(prs)
    hb = page_head(s, P, "학원 소개",
              clean(ph.get("headline") or "공부하는 방법이 결과를 만듭니다"),
              None, page)

    if intro:
        # 포인트가 없으면 오른쪽 박스도 없으므로 전폭으로 크게 배치
        if points:
            add_text(s, G.ML, max(1.88, hb), 5.47, 1.04, intro, size=FS.body,
                     color=P["muted"], line_spacing=1.5)
        else:
            iy = max(1.88, hb)
            # 실제 글 분량만큼만 높이를 잡는다 (남는 공간은 사진으로)
            ilines = wrap_words(intro, max(24, int(G.CW * 72 / (FS.lead * 1.55))))
            ih = min(2.6, len(ilines) * FS.lead * 1.7 / 72.0 + 0.1)
            add_text(s, G.ML, iy, G.CW, ih, intro,
                     size=FS.lead, color=P["text"], line_spacing=1.7)
            # 아래 여백에 학원 사진
            if not points:
                fill_gap(s, d, P, used, iy + ih + 0.42)

    y = 3.39
    for i, p in enumerate(points[:4]):
        add_pill(s, G.ML, y, 0.5, 0.29, f"{i+1:02d}", fill=P["primary"],
                 color=P["onDark"], size=FS.small, radius=0.15)
        add_text(s, 1.25, y, 4.9, 0.29, clean(p.get("title") or ""),
                 size=FS.body + 0.75, bold=True, color=P["text"],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        y += 0.725

    # 우측 강조 박스 (포인트나 노트가 있을 때만)
    has_note = any(clean(note.get(k) or "") for k in ("title", "body", "footer"))
    if not (points or has_note):
        return s
    bx, by, bw, bh = 7.08, 1.82, 5.42, 4.17
    add_box(s, bx, by, bw, bh, fill=P["card"], line=P["line"], radius=R)
    ny = by + 0.47
    nt = clean(note.get("title") or "")
    if nt:
        add_text(s, bx + 0.47, ny, bw - 0.94, 0.36, nt, size=FS.card_ttl,
                 bold=True, color=P["primary"], line_spacing=1.2)
        ny += 0.57
    nb = clean(note.get("body") or "")
    if nb:
        add_text(s, bx + 0.47, ny, bw - 0.94, 1.56, nb, size=FS.lead,
                 color=P["text"], line_spacing=1.6)
        ny += 2.04
    nf = clean(note.get("footer") or "")
    if nf:
        add_text(s, bx + 0.47, ny, bw - 0.94, 0.65, nf, size=FS.body,
                 bold=True, color=P["accent"], line_spacing=1.45)
    return s


# ══════════════════════════════════════════════════════════════════
# 03-2. 수업 대상 (학년별 과목·반 구성)
# ══════════════════════════════════════════════════════════════════
GRADE_ORDER = ["초등", "중등", "고등", "특강", "기타"]


def s_targets(prs, d, P, page, used=None):
    items = as_dicts(d.get("targets"), "grade")
    items = [it for it in items
             if clean(it.get("grade") or it.get("subj") or it.get("desc"))]
    if not items:
        return None
    items = items[:4]
    s = blank_slide(prs)
    hb = page_head(s, P, "수업 대상",
                   clean(d.get("targetsHead") or "학년별 수업 대상"),
                   "학년과 과목에 맞춰 반을 나누어 운영합니다.", page)

    top = max(G.y_body, hb)
    avail = 6.5 - top
    n = len(items)
    gap = 0.24
    ch = min(1.5, (avail - gap * (n - 1)) / n)

    for i, it in enumerate(items):
        y = top + i * (ch + gap)
        add_box(s, G.ML, y, G.CW, ch, fill=P["card"], line=P["line"], radius=R)

        grade = clean(it.get("grade") or "")
        gw = min(1.6, max(0.9, 0.44 + len(grade) * 0.15))
        add_pill(s, G.ML + 0.26, y + (ch - 0.33) / 2, gw, 0.33, grade,
                 fill=P["primary"], color=P["onDark"], size=FS.small)

        tx = G.ML + 0.26 + gw + 0.3
        tw = G.CW - (tx - G.ML) - 0.3
        subj = clean(it.get("subj") or "")
        desc = clean(it.get("desc") or "")

        if subj and desc:
            add_text(s, tx, y + 0.24, tw, 0.3, subj,
                     size=FS.card_ttl, bold=True, color=P["text"],
                     line_spacing=1.1)
            add_text(s, tx, y + 0.64, tw, ch - 0.86, desc,
                     size=fit_size(desc, FS.body, tw, ch - 0.86),
                     color=P["muted"], line_spacing=1.45)
        else:
            one = subj or desc
            add_text(s, tx, y, tw, ch, one,
                     size=fit_size(one, FS.body + 0.75, tw, ch - 0.3),
                     bold=bool(subj), color=P["text"],
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
    return s


# ══════════════════════════════════════════════════════════════════
# 04. 커리큘럼 (학년별 카드 + 태그 필)
# ══════════════════════════════════════════════════════════════════
def s_curriculum(prs, d, P, page, used=None):
    stages = as_dicts(d.get("curriculum"), "name")
    if not stages:
        return None
    stages = stages[:3]
    s = blank_slide(prs)
    hb = page_head(s, P, "교육 과정",
              "학년별 진도와 수업 방식은 한 흐름 안에서 설계합니다",
              gv(d, "curriculumLead") or
              "현재 수준에 맞춰 다음 단계까지 자연스럽게 연결합니다.",
              page)

    n = len(stages)
    gap = 0.31
    cw = (G.CW - gap * (n - 1)) / n
    top, ch = max(G.y_body, hb), 4.06

    for i, st in enumerate(stages):
        x = G.ML + i * (cw + gap)
        add_box(s, x, top, cw, ch, fill=P["card"], line=P["line"], radius=R)

        gname = clean(st.get("name") or st.get("grade") or "")
        gw = min(cw - 0.5, 0.5 + len(gname) * 0.135)
        add_pill(s, x + 0.25, top + 0.3, gw, 0.31, gname,
                 fill=P["primary"], color=P["onDark"], size=FS.small)
        ttl = clean(st.get("title") or "")
        add_text(s, x + 0.25, top + 0.88, cw - 0.5, 0.6, ttl,
                 size=fit_size(ttl, FS.stage_ttl, cw - 0.5, 0.62),
                 bold=True, color=P["text"], line_spacing=1.25)
        body = clean(st.get("desc") or st.get("body") or "")
        if body:
            add_text(s, x + 0.25, top + 1.67, cw - 0.5, 1.1, body,
                     size=fit_size(body, FS.body, cw - 0.5, 1.1),
                     color=P["muted"], line_spacing=1.5)

        tags = [clean(t) for t in (st.get("tags") or []) if clean(t)][:3]
        if tags:
            tw = (cw - 0.5 - 0.2) / max(len(tags), 1)
            for j, t in enumerate(tags):
                add_pill(s, x + 0.25 + j * (tw + 0.1), top + 3.02,
                         tw, 0.29, t, fill=P["card2"], color=P["primary"],
                         size=FS.small, radius=0.2)

    foot = clean(d.get("curriculumNote") or
                 "반 이름은 학생의 우열이 아니라, 현재 필요한 학습 내용과 속도를 구분하는 기준입니다.")
    add_text(s, G.ML, 6.46, G.CW, 0.38, foot, size=FS.body,
             color=P["muted"], line_spacing=1.3)
    return s


# ══════════════════════════════════════════════════════════════════
# 05. 특별 프로그램 (2x2)
# ══════════════════════════════════════════════════════════════════
def s_specials(prs, d, P, page, used=None):
    items = as_dicts(d.get("specials"))
    if not items:
        return None
    items = items[:4]
    s = blank_slide(prs)
    hb = page_head(s, P, "특별 프로그램",
              "수업 밖의 시간도 공부하는 힘으로 바꿉니다", None, page)

    colw, rowh = 6.15, 2.13
    for i, it in enumerate(items):
        r, c = divmod(i, 2)
        x = G.ML + c * colw
        y = max(1.98, hb) + r * rowh
        add_text(s, x, y, 0.45, 0.42, f"{i+1:02d}", size=FS.card_ttl,
                 bold=True, color=P["accent"], line_spacing=1.0)
        add_text(s, x + 0.58, y, 5.05, 0.34,
                 clean(it.get("title") or it.get("name") or ""),
                 size=FS.card_ttl, bold=True, color=P["text"],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        body = clean(it.get("desc") or it.get("body") or "")
        if body:
            add_text(s, x + 0.58, y + 0.44, 5.05, 0.86, body, size=FS.body,
                     color=P["muted"], line_spacing=1.5)
    return s


# ══════════════════════════════════════════════════════════════════
# 06. 학습관리 (2x2 카드, 좌 키워드 / 우 설명)
# ══════════════════════════════════════════════════════════════════
def s_management(prs, d, P, page, used=None):
    items = as_dicts(d.get("management"))
    if not items:
        return None
    items = items[:4]
    s = blank_slide(prs)
    hb = page_head(s, P, "학습 관리",
              "작은 빈틈을 제때 채워야 어려워지지 않습니다",
              "수업과 가정이 같은 정보를 보고, 필요한 시점에 학습 방향을 조정합니다.",
              page)

    cw, ch, gx, gy = 5.62, 1.51, 6.15, 1.92
    for i, it in enumerate(items):
        r, c = divmod(i, 2)
        x = G.ML + c * gx
        y = max(2.14, hb) + r * gy
        add_box(s, x, y, cw, ch, fill=P["card"], line=P["line"], radius=R)
        key = clean(it.get("key") or it.get("title") or "")
        add_pill(s, x + 0.25, y + 0.22, 0.94, 0.33, key, fill=P["primary"],
                 color=P["onDark"], size=FS.small, radius=0.18)
        body = clean(it.get("desc") or it.get("body") or "")
        add_text(s, x + 1.36, y + 0.2, cw - 1.61, ch - 0.4, body,
                 size=fit_size(body, FS.body, cw - 1.61, ch - 0.4),
                 color=P["text"], line_spacing=1.5, anchor=MSO_ANCHOR.MIDDLE)
    return s


# ══════════════════════════════════════════════════════════════════
# 07. 성장 경로 (단계 카드 + 화살표)
# ══════════════════════════════════════════════════════════════════
def s_growth(prs, d, P, page, used=None):
    steps = as_dicts(d.get("growth"), "name")
    if not steps:
        return None
    steps = steps[:4]
    s = blank_slide(prs)
    hb = page_head(s, P, "성장 로드맵",
              "지금의 공부 습관이 다음 단계의 경쟁력으로 이어집니다",
              gv(d, "growthLead"), page)

    n = len(steps)
    aw = 0.45
    cw = (G.CW - aw * (n - 1)) / n
    top, ch = max(2.45, hb), 2.45
    for i, st in enumerate(steps):
        x = G.ML + i * (cw + aw)
        add_box(s, x, top, cw, ch, fill=P["card"], line=P["line"], radius=R)
        add_text(s, x + 0.21, top + 0.36, cw - 0.42, 0.29,
                 clean(st.get("name") or st.get("stage") or ""),
                 size=FS.body + 0.75, bold=True, color=P["primary"],
                 line_spacing=1.0)
        body = st.get("desc") or st.get("body") or ""
        if isinstance(body, list):
            body = [clean(b) for b in body if clean(b)]
        add_text(s, x + 0.21, top + 0.94, cw - 0.42, 1.2, body,
                 size=FS.body, color=P["text"], line_spacing=1.5)
        if i < n - 1:
            add_text(s, x + cw, top + 0.88, aw, 0.42, "→", size=FS.stage_ttl,
                     bold=True, color=P["accent"], align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    note = clean(d.get("growthNote") or "")
    if note:
        add_text(s, 0.94, 5.62, 11.46, 0.67, note, size=FS.lead,
                 color=P["muted"], align=PP_ALIGN.CENTER, line_spacing=1.5)
    return s


# ══════════════════════════════════════════════════════════════════
# 08. 실적 (큰 숫자 카드 3단)
# ══════════════════════════════════════════════════════════════════
def s_achievements(prs, d, P, page, used=None):
    items = as_dicts(d.get("achievements"))
    if not items:
        return None
    items = items[:3]
    s = blank_slide(prs)
    hb = page_head(s, P, "주요 실적",
              clean(gv(d, "achievementsHead") or "숫자로 확인하는 학습의 결과"),
              gv(d, "achievementsLead"), page)

    n = len(items)
    gap = 0.34
    cw = (G.CW - gap * (n - 1)) / n
    top, ch = max(2.08, hb), 2.29
    for i, it in enumerate(items):
        x = G.ML + i * (cw + gap)
        add_box(s, x, top, cw, ch, fill=P["card"], line=P["line"], radius=R)
        add_text(s, x + 0.21, top + 0.26, cw - 0.42, 0.68,
                 clean(it.get("value") or it.get("num") or ""),
                 size=FS.big_num, bold=True, color=P["primary"], line_spacing=1.0)
        add_text(s, x + 0.21, top + 0.99, cw - 0.42, 0.31,
                 clean(it.get("title") or it.get("label") or ""),
                 size=FS.card_ttl, bold=True, color=P["text"], line_spacing=1.0)
        body = clean(it.get("desc") or it.get("body") or "")
        if body:
            add_text(s, x + 0.37, top + 1.46, cw - 0.73, 0.7, body,
                     size=fit_size(body, FS.body, cw - 0.73, 0.7),
                     color=P["muted"], line_spacing=1.45)

    notes = [clean(x) for x in (d.get("achievementsNote") or []) if clean(x)]
    if notes:
        add_box(s, G.ML, 4.9, G.CW, 1.04, fill=P["card2"], radius=R)
        add_text(s, 0.94, 5.05, 11.25, 0.75, notes[:2], size=FS.body,
                 color=P["text"], line_spacing=1.6, space_after=4)
    return s


# ══════════════════════════════════════════════════════════════════
# 09. 시간표 (그룹 수만큼 동적, 넘치면 다음 슬라이드로 분할)
# ══════════════════════════════════════════════════════════════════
MAX_ROWS_PER_COL = 6


def _style_table(tbl, P, col_widths, header_h=0.46, row_h=0.4):
    """표 스타일. 밴딩/테마 제거하고 팔레트 색으로 직접 칠한다."""
    from navy_core import _rgb
    # 기본 밴딩 끄기 (테마색 오염 방지)
    tblPr = tbl._tbl.find(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}tblPr')
    if tblPr is not None:
        tblPr.set('bandRow', '0')
        tblPr.set('firstRow', '1')

    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    for ri, row in enumerate(tbl.rows):
        row.height = Inches(header_h if ri == 0 else row_h)
        for ci, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = _rgb(P["primary"])
            else:
                cell.fill.fore_color.rgb = _rgb(
                    "FFFFFF" if ri % 2 else P["card"])


def _set_cell(cell, text, P, size, bold=False, color=None,
              align=PP_ALIGN.LEFT):
    from navy_core import _rgb, _latin_font, FONT
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.15
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    r = p.add_run()
    r.text = clean(text)
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = _rgb(color or P["text"])
    _latin_font(r)


def s_timetables(prs, d, P, page, used=None):
    """시간표 = 진짜 table 객체 (PowerPoint에서 행 추가·수정 가능).
    행이 많으면 잘라내지 않고 '(이어서)' 표로 분할한다."""
    groups = as_dicts(d.get("timetables"), "group")
    groups = [g for g in groups if g.get("rows")]
    if not groups:
        return []

    ROWS_PER_TABLE = 10

    # 그룹을 표 단위로 쪼갠다 (행 손실 없음)
    units = []
    for grp in groups:
        rows = as_dicts(grp.get("rows"), "name")
        gname = clean(grp.get("group") or grp.get("title") or "")
        for i in range(0, len(rows), ROWS_PER_TABLE):
            chunk = rows[i:i + ROWS_PER_TABLE]
            label = gname if i == 0 else f"{gname} (이어서)"
            units.append({"label": label, "rows": chunk})

    made = []
    pairs = [units[i:i + 2] for i in range(0, len(units), 2)]
    for ci, pair in enumerate(pairs):
        s = blank_slide(prs)
        hb = page_head(s, P, "수업 시간표",
                       "수업 시간표" if len(pairs) == 1
                       else f"수업 시간표 ({ci+1})",
                       "반 편성과 시간표는 학기별로 조정될 수 있으니 "
                       "상담 시 확인해 주세요.",
                       page + ci)

        n = len(pair)
        gap = 0.4
        tw = (G.CW - gap * (n - 1)) / n
        top = max(2.05, hb + 0.4)
        tbl_bottom = top

        for gi, u in enumerate(pair):
            x = G.ML + gi * (tw + gap)
            rows = u["rows"]
            if not rows:
                continue

            add_text(s, x, top - 0.42, tw, 0.3, u["label"],
                     size=FS.card_ttl, bold=True, color=P["text"],
                     line_spacing=1.0)

            avail = 6.55 - top
            # 행이 적으면 넉넉히, 많으면 촘촘히 (지면을 고르게 채운다)
            row_h = (avail - 0.46) / max(len(rows), 1)
            row_h = min(0.62, max(0.3, row_h))
            th = 0.46 + row_h * len(rows)

            gf = s.shapes.add_table(len(rows) + 1, 2, Inches(x), Inches(top),
                                    Inches(tw), Inches(th))
            tbl = gf.table
            c1 = tw * 0.58
            _style_table(tbl, P, [c1, tw - c1], 0.46, row_h)

            _set_cell(tbl.cell(0, 0), "반", P, FS.small, True, P["onDark"])
            _set_cell(tbl.cell(0, 1), "수업 요일 · 시간", P, FS.small, True,
                      P["onDark"], PP_ALIGN.RIGHT)

            for ri, row in enumerate(rows, start=1):
                nm = clean(row.get("name") or row.get("class") or "")
                tm = clean(row.get("time") or row.get("when") or "")
                _set_cell(tbl.cell(ri, 0), nm, P,
                          fit_size(nm, FS.body, c1 - 0.24, row_h, min_size=8.0),
                          True, P["text"])
                _set_cell(tbl.cell(ri, 1), tm, P, FS.body_sm, False,
                          P["muted"], PP_ALIGN.RIGHT)
            tbl_bottom = max(tbl_bottom, top + th)
        # 표 아래가 많이 비면 사진으로 채운다
        fill_gap(s, d, P, used, tbl_bottom + 0.45, min_h=1.35)
        made.append(s)
    return made


# ══════════════════════════════════════════════════════════════════
# 10. 입학 절차 (번호 카드)
# ══════════════════════════════════════════════════════════════════
def s_admission(prs, d, P, page, used=None):
    steps = as_dicts(d.get("admission"))
    if not steps:
        return None
    steps = steps[:4]
    s = blank_slide(prs)
    hb = page_head(s, P, "입학 안내",
              "처음 한 걸음부터 학원이 함께 안내합니다", None, page)

    n = len(steps)
    gap = 0.45
    cw = (G.CW - gap * (n - 1)) / n
    top, ch = max(2.29, hb), 2.71
    for i, st in enumerate(steps):
        x = G.ML + i * (cw + gap)
        add_box(s, x, top, cw, ch, fill=P["card"], line=P["line"], radius=R)
        add_text(s, x + 0.21, top + 0.26, 0.52, 0.42, str(i + 1),
                 size=FS.big_num * 0.72, bold=True, color=P["accent"],
                 line_spacing=1.0)
        add_text(s, x + 0.21, top + 0.89, cw - 0.42, 0.36,
                 clean(st.get("title") or st.get("name") or ""),
                 size=FS.card_ttl, bold=True, color=P["text"], line_spacing=1.1)
        body = clean(st.get("desc") or st.get("body") or "")
        if body:
            add_text(s, x + 0.21, top + 1.46, cw - 0.42, 1.0, body,
                     size=fit_size(body, FS.body, cw - 0.42, 1.0),
                     color=P["muted"], line_spacing=1.5)

    note = clean(d.get("admissionNote") or "")
    if note:
        add_text(s, G.ML, 5.68, G.CW, 0.47, note, size=FS.lead,
                 color=P["muted"], line_spacing=1.4)
    return s


# ══════════════════════════════════════════════════════════════════
# 11. FAQ (2x2 카드)
# ══════════════════════════════════════════════════════════════════
def s_faq(prs, d, P, page, used=None):
    items = as_dicts(d.get("faq"), "q")
    items = [x for x in items if clean(x.get("q") or x.get("question"))]
    if not items:
        return None
    items = items[:4]
    s = blank_slide(prs)
    hb = page_head(s, P, "자주 묻는 질문",
              "등록 전에 가장 많이 물어보시는 내용을 정리했습니다", None, page)

    cw, ch, gx, gy = 5.62, 1.79, 6.15, 2.13
    for i, it in enumerate(items):
        r, c = divmod(i, 2)
        x = G.ML + c * gx
        y = max(1.93, hb) + r * gy
        add_box(s, x, y, cw, ch, fill=P["card"], line=P["line"], radius=R)
        q = clean(it.get("q") or it.get("question") or "")
        a = clean(it.get("a") or it.get("answer") or "")
        add_text(s, x + 0.21, y + 0.18, 0.35, 0.29, "Q.", size=FS.body + 0.75,
                 bold=True, color=P["primary"], line_spacing=1.0)
        add_text(s, x + 0.61, y + 0.18, cw - 0.88, 0.31, q,
                 size=fit_size(q, FS.body + 1.5, cw - 0.88, 0.33),
                 bold=True, color=P["text"], line_spacing=1.1)
        add_text(s, x + 0.21, y + 0.73, 0.35, 0.29, "A.", size=FS.body + 0.75,
                 bold=True, color=P["accent"], line_spacing=1.0)
        add_text(s, x + 0.61, y + 0.69, cw - 0.88, ch - 0.9, a,
                 size=fit_size(a, FS.body, cw - 0.88, ch - 0.9),
                 color=P["muted"], line_spacing=1.5)

    rows_n = -(-len(items) // 2)
    bottom = 1.93 + rows_n * gy - (gy - ch)
    fill_gap(s, d, P, used, bottom + 0.42)
    return s


# ══════════════════════════════════════════════════════════════════
# 12. 연락처 / 마무리
# ══════════════════════════════════════════════════════════════════
def s_contact(prs, d, P, page, used=None):
    s = blank_slide(prs)
    b = d.get("basic", {})
    ct = d.get("contact") or {}
    hb = page_head(s, P, "등록 문의",
              clean(ct.get("headline") or b.get("slogan") or "상담을 기다리고 있습니다"),
              None, page)

    qr = gv(d, "assets", "qr")
    mapimg = gv(d, "assets", "map")

    phone = clean(b.get("phone") or "")
    addr = clean(b.get("address") or "")

    # 좌측: 연락처 → 약도 순. 약도는 '오시는 길'이므로 크게 보여준다.
    ty = max(1.93, hb)
    if phone:
        add_text(s, G.ML, ty, 1.86, 0.25, "전화 문의", size=FS.small,
                 bold=True, color=P["primary"], line_spacing=1.0)
        add_text(s, G.ML, ty + 0.32, 4.17, 0.4, phone, size=FS.header * 0.78,
                 bold=True, color=P["text"], line_spacing=1.0)
        ty += 0.85
    if addr:
        add_text(s, G.ML, ty, 4.9, 0.68, wrap_words(addr, 22)[:2],
                 size=FS.body, color=P["muted"], line_spacing=1.5)
        ty += 0.78

    if mapimg:
        # 약도: 좌측 하단을 넉넉히 사용
        mh = min(2.55, 6.35 - ty - 0.3)
        if mh >= 1.0:
            add_text(s, G.ML, ty, 2.2, 0.25, "오시는 길", size=FS.small,
                     bold=True, color=P["primary"], line_spacing=1.0)
            ty += 0.34
            mh = min(2.55, 6.35 - ty)
            add_box(s, G.ML, ty, 5.42, mh, fill=P["card"],
                    line=P["line"], radius=R)
            place_image(s, mapimg, G.ML + 0.06, ty + 0.06,
                        5.42 - 0.12, mh - 0.12, cover=False)
            if used is not None:
                used.add(mapimg)

    if qr:
        # QR: 우측 안내 박스 위쪽 여백에 작게
        place_image(s, qr, G.W - G.MR - 1.25, max(1.93, hb) - 0.02,
                    1.25, 1.25, cover=False)
        if used is not None:
            used.add(qr)

    # 우측 안내 박스 (QR이 있으면 그만큼 아래에서 시작)
    bx = 6.67
    by = (max(1.93, hb) + 1.42) if qr else 1.82
    bw = 5.83
    bh = 5.88 - by
    add_box(s, bx, by, bw, bh, fill=P["card"], line=P["line"], radius=R)
    add_text(s, bx + 0.52, by + 0.52, bw - 1.04, 0.38,
             clean(ct.get("title") or "상담 전에 알려주시면 좋아요"),
             size=FS.card_ttl, bold=True, color=P["primary"], line_spacing=1.1)

    asks = [clean(x) for x in (ct.get("asks") or []) if clean(x)][:4]
    ay = by + 1.25
    for i, a in enumerate(asks):
        add_pill(s, bx + 0.52, ay, 0.33, 0.29, str(i + 1), fill=P["primary"],
                 color=P["onDark"], size=FS.small, radius=0.5)
        add_text(s, bx + 1.02, ay, bw - 1.6, 0.31, a, size=FS.body,
                 color=P["text"], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        ay += 0.585

    closing = clean(ct.get("closing") or "")
    if closing:
        add_text(s, bx + 0.52, max(ay + 0.1, by + bh - 0.75), bw - 1.04, 0.33,
                 closing, size=FS.body, bold=True, color=P["accent"],
                 line_spacing=1.3)
    return s


# ══════════════════════════════════════════════════════════════════
# 조립
# ══════════════════════════════════════════════════════════════════
def build(data, out=None, palette=None):
    """
    data    : v3 표준스키마 dict
    out     : 파일경로 또는 None(BytesIO 반환)
    palette : 'navy_gold' | 'charcoal_gold' | 'forest_gold'
    """
    P = get_palette(palette)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    d = data or {}
    used = set()          # 이미 쓴 사진 (중복 배치 방지)
    s_cover(prs, d, P, used)

    page = 2
    plan = [s_philosophy, s_achievements, s_strengths, s_targets,
            s_curriculum, s_specials, s_management, s_growth]
    for fn in plan:
        if fn(prs, d, P, page, used) is not None:
            page += 1

    made = s_timetables(prs, d, P, page, used)
    page += len(made)

    for fn in (s_admission, s_faq):
        if fn(prs, d, P, page, used) is not None:
            page += 1

    s_contact(prs, d, P, page, used)

    force_theme_font(prs)
    if out is None:
        import io
        bio = io.BytesIO()
        prs.save(bio)
        bio.seek(0)
        return bio
    prs.save(out)
    return out


if __name__ == "__main__":
    import json
    import sys
    src = sys.argv[1] if len(sys.argv) > 1 else "sample.json"
    pal = sys.argv[2] if len(sys.argv) > 2 else "navy_gold"
    with open(src, encoding="utf-8") as f:
        data = json.load(f)
    print(build(data, f"navy_{pal}.pptx", pal))
