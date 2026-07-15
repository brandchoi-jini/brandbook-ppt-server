#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
브랜드북 PPT — 컨설팅 보고서 스타일 (SG 보고서 톤)
흰 배경 + 네이비. 러닝헤드/페이지번호, 큰 제목+부제, 3열 비교표, 인용구 강조.
사용: python build_report.py <content.json> <raw.json> <out.pptx> [palette]
"""
import sys, json, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PALETTES = {
    "navy":  {"main":"1E3A5F","accent":"2E5FA3","soft":"F1F5FA","soft2":"E4EBF3","line":"D8E0EA","muted":"7A899B"},
    "navy2": {"main":"152B47","accent":"3E6FA8","soft":"EEF3F9","soft2":"E0E8F1","line":"D3DEEA","muted":"73838F"},
}
BLACK="222A35"; WHITE="FFFFFF"

def C(h): return RGBColor.from_string(h)

class Report:
    def __init__(self, content, raw, palette="navy"):
        self.ct=content; self.raw=raw
        self.P=PALETTES.get(palette,PALETTES["navy"])
        self.prs=Presentation()
        self.prs.slide_width=Inches(10); self.prs.slide_height=Inches(5.625)
        self.blank=self.prs.slide_layouts[6]
        self.W=10; self.H=5.625; self.MX=0.4
        self.CW=self.W-self.MX*2
        self.page=0
        b=raw.get("basic",{})
        self.name=(b.get("name") or "").strip()
        smap={"math":"수학","science":"과학","korean":"국어","english":"영어","social":"사회"}
        self.subjects=[smap.get(k,k) for k,v in (b.get("subjects") or {}).items() if v]

    # ── 헬퍼 ──
    def _strip(self, shp):
        el=shp._element
        for st in el.findall(qn('p:style')): el.remove(st)

    def _slide(self):
        s=self.prs.slides.add_slide(self.blank)
        f=s.background.fill; f.solid(); f.fore_color.rgb=C(WHITE)
        return s

    def _rect(self,s,x,y,w,h,color,line=None,rounded=False,lw=1.0):
        shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                               Inches(x),Inches(y),Inches(w),Inches(h))
        self._strip(shp)
        if color: shp.fill.solid(); shp.fill.fore_color.rgb=C(color)
        else: shp.fill.background()
        if line: shp.line.color.rgb=C(line); shp.line.width=Pt(lw)
        else: shp.line.fill.background()
        shp.shadow.inherit=False
        return shp

    def _line(self,s,x,y,w,color=None,wpt=1.0):
        ln=s.shapes.add_connector(2,Inches(x),Inches(y),Inches(x+w),Inches(y))
        self._strip(ln)
        ln.line.color.rgb=C(color or self.P["line"]); ln.line.width=Pt(wpt)
        ln.shadow.inherit=False
        return ln

    def _t(self,s,x,y,w,h,runs,size=11,bold=False,color=BLACK,align=PP_ALIGN.LEFT,
           anchor=MSO_ANCHOR.TOP,ls=1.15,spacing=0,font="맑은 고딕"):
        tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
        tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
        for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
        p=tf.paragraphs[0]; p.alignment=align
        if ls: p.line_spacing=ls
        if isinstance(runs,str): runs=[(runs,{})]
        for txt,opt in runs:
            r=p.add_run(); r.text=txt
            r.font.size=Pt(opt.get("size",size)); r.font.bold=opt.get("bold",bold)
            r.font.color.rgb=C(opt.get("color",color)); r.font.name=font
            if spacing or opt.get("spacing"):
                self._set_spacing(r, opt.get("spacing",spacing))
        return tb

    def _set_spacing(self, run, pts):
        rPr=run._r.get_or_add_rPr(); rPr.set('spc', str(int(pts*100)))

    # 러닝헤드 + 페이지번호 + [핵심메시지 헤드라인] + 정보라벨
    def _head(self, s, section, title, subtitle=None, headline=None):
        self.page+=1
        # 상단 얇은 러닝헤드(정보) + 페이지
        self._t(s,self.MX,0.2,6.0,0.3,section,size=9.5,bold=True,color=self.P["accent"],
                spacing=1.2,anchor=MSO_ANCHOR.MIDDLE)
        self._t(s,self.W-0.7,0.2,0.3,0.3,str(self.page),size=9.5,bold=True,
                color=self.P["muted"],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
        self._line(s,self.MX,0.5,self.CW,self.P["line"],0.75)
        if headline:
            # 핵심메시지형: 큰 헤드라인 + 작은 정보라벨(제목) 병기
            self._t(s,self.MX,0.62,self.CW,0.45,'"'+headline+'"',size=20,bold=True,
                    color=self.P["main"],anchor=MSO_ANCHOR.MIDDLE)
            self._t(s,self.MX,1.12,self.CW,0.28,
                    [(title,{"bold":True,"color":self.P["accent"]}),
                     (("   ·   "+subtitle) if subtitle else "",{"color":self.P["muted"]})],
                    size=10.5)
            return 1.5  # 본문 시작 y
        else:
            self._t(s,self.MX,0.62,self.CW,0.45,title,size=19,bold=True,color=self.P["main"],
                    anchor=MSO_ANCHOR.MIDDLE)
            if subtitle:
                self._t(s,self.MX,1.12,self.CW,0.3,subtitle,size=11,color=self.P["muted"])
            return 1.6

    @staticmethod
    def _fit(text,base,maxc):
        n=len(text or "")
        return base if n<=maxc else max(base*0.72, round(base*maxc/n,1))

    # ══════ 슬라이드 ══════
    def cover(self):
        s=self._slide(); P=self.P
        # 좌측 얇은 네이비 바
        self._rect(s,0,0,0.18,self.H,P["main"])
        # 우상단 태그
        self._rect(s,8.5,0.35,1.1,0.32,P["soft"])
        self._t(s,8.5,0.35,1.1,0.32,"BRANDBOOK",size=8,bold=True,color=P["accent"],
                align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=0.5)
        # 학원명
        self._t(s,0.6,1.6,8.5,1.0,self.name,size=40,bold=True,color=P["main"],
                anchor=MSO_ANCHOR.MIDDLE)
        # 슬로건
        slg=(self.ct.get("brand") or {}).get("slogan","")
        if slg:
            self._t(s,0.6,2.75,8.5,0.4,slg,size=15,bold=True,color=P["accent"])
        # 과목
        if self.subjects:
            self._t(s,0.6,3.2,8.5,0.3,"  ·  ".join(self.subjects)+"  전문",size=11,
                    color=P["muted"],spacing=0.5)
        # 하단 구분선 + 위치/연락처
        self._line(s,0.6,4.55,8.8,P["line"],0.75)
        ct=self.ct.get("contact") or {}
        parts=[]
        if ct.get("address"): parts.append([("위치  ",{"bold":True,"color":P["main"]}),(ct["address"],{"color":P["muted"]})])
        if parts:
            self._t(s,0.6,4.7,8.8,0.3,parts[0],size=10)
        if ct.get("phone"):
            self._t(s,0.6,5.0,8.8,0.3,[("연락처  ",{"bold":True,"color":P["main"]}),(ct["phone"],{"color":P["muted"]})],size=10)

    def intro(self):
        s=self._slide(); P=self.P; b=self.ct.get("brand") or {}
        # 러닝헤드 + 제목 + 부제
        self.page+=1
        self._t(s,self.MX,0.2,6.0,0.3,"학원 소개",size=9.5,bold=True,color=P["accent"],
                spacing=1.2,anchor=MSO_ANCHOR.MIDDLE)
        self._t(s,self.W-0.7,0.2,0.3,0.3,str(self.page),size=9.5,bold=True,
                color=P["muted"],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
        self._line(s,self.MX,0.5,self.CW,P["line"],0.75)
        self._t(s,self.MX,0.6,self.CW,0.45,self.name+" 소개",size=19,bold=True,
                color=P["main"],anchor=MSO_ANCHOR.MIDDLE)
        # 소제목 태그
        tag=self.ct.get("intro_tag","")
        if tag:
            self._t(s,self.MX,1.2,4.0,0.3,"[ "+tag+" ]",size=11,bold=True,color=P["accent"])
        # 핵심 문장 (크게)
        key = b.get("identity","")
        self._t(s,self.MX,1.5,9.0,0.5,key,size=16,bold=True,color=P["main"])
        # 설명 (2줄)
        self._t(s,self.MX,2.05,9.0,0.9,b.get("intro",""),size=10,color="4A5563",ls=1.4)
        # ── 하단 3축 매트릭스 ──
        mtx=self.ct.get("intro_matrix")
        if mtx:
            axes=mtx["axes"]; n=len(axes)
            labW=1.6; gap=0.15
            colW=(self.CW-labW-gap*n)/n
            headY=3.35; headH=0.42
            # 컬러 헤더 (3축)
            for ci,ax in enumerate(axes):
                cx=self.MX+labW+ci*(colW+gap)
                self._rect(s,cx,headY,colW,headH,P["accent"])
                self._t(s,cx,headY,colW,headH,f"{ci+1}. {ax}",size=11,bold=True,color=WHITE,
                        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
            # 행들
            rowY=headY+headH+0.1
            rows=mtx["rows"]
            rowH=(self.H-rowY-0.3)/len(rows)
            for ri,row in enumerate(rows):
                y=rowY+ri*rowH
                if ri>0: self._line(s,self.MX,y,self.CW,P["line"],0.5)
                self._t(s,self.MX,y,labW,rowH,row["label"],size=9.5,bold=True,
                        color=P["muted"],anchor=MSO_ANCHOR.MIDDLE)
                for ci,v in enumerate(row["values"]):
                    cx=self.MX+labW+ci*(colW+gap)
                    self._t(s,cx+0.1,y,colW-0.2,rowH,v,size=10,color=BLACK,
                            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,ls=1.2)

    def strengths(self):
        items=(self.ct.get("strengths") or [])[:6]
        if not items: return
        s=self._slide(); P=self.P
        self._head(s,"학원 소개","핵심 강점","학생 한 명 한 명에게 닿는 여섯 가지 방식")
        cols=3; gap=0.25; cw=(self.CW-gap*(cols-1))/cols; ch=1.5; y0=1.7
        for i,it in enumerate(items):
            r=i//cols; c=i%cols; x=self.MX+c*(cw+gap); y=y0+r*(ch+0.2)
            # 번호 + 상단 구분선 (카드 대신 정돈된 구획)
            self._t(s,x,y,0.5,0.35,f"{i+1:02d}",size=15,bold=True,color=P["accent"])
            self._line(s,x,y+0.42,cw,P["line"],0.75)
            self._t(s,x,y+0.5,cw,0.35,it.get("t",""),size=12.5,bold=True,color=P["main"])
            self._t(s,x,y+0.88,cw,ch-0.95,it.get("d",""),size=10,color=BLACK,ls=1.25)

    def achievements(self):
        A=self.ct.get("achievements") or []
        if not A: return
        s=self._slide(); P=self.P
        hl=(self.ct.get("headlines") or {}).get("achievements")
        self._head(s,"학원 성과","주요 실적","진학·합격으로 증명하는 결과",headline=hl)
        short={"중등 진학":"중등","고등 진학":"고등","재원생 전원":"대입"}
        # 큰 번호 + 태그 + 내용, 좌측 굵은 세로바로 구획 (박스 없음)
        n=len(A); rowH=(self.H-1.75-0.3)/n; y0=1.75
        for i,a in enumerate(A):
            y=y0+i*rowH
            # 좌측 큰 번호 (연한 회색)
            self._t(s,self.MX,y,1.2,rowH,f"{i+1:02d}",size=40,bold=True,color=P["soft2"] if "soft2" in P else P["line"],
                    anchor=MSO_ANCHOR.MIDDLE)
            # 굵은 세로바
            self._rect(s,self.MX+1.35,y+rowH*0.2,0.06,rowH*0.6,P["accent"])
            # 태그
            self._t(s,self.MX+1.6,y,1.5,rowH,short.get(a.get("tag",""),a.get("tag","")),
                    size=16,bold=True,color=P["main"],anchor=MSO_ANCHOR.MIDDLE)
            # 내용
            self._t(s,self.MX+3.2,y,self.CW-3.2,rowH,a.get("text",""),
                    size=13,color=BLACK,anchor=MSO_ANCHOR.MIDDLE,ls=1.3)
            if i<n-1:
                self._line(s,self.MX,y+rowH,self.CW,P["line"],0.5)

    def subjects_targets(self):
        cls=self.ct.get("classes") or {}; gk=list(cls.keys())
        if not gk: return
        s=self._slide(); P=self.P
        hl=(self.ct.get("headlines") or {}).get("subjects")
        self._head(s,"교육 과정","수업 대상 · 과목",
                   "과목 "+" · ".join(self.subjects)+"  /  대상 초등~고등",headline=hl)
        ranges={"초등":"초4 ~ 초6","중등":"중1 ~ 중3","고등":"고1 ~ 고3"}
        n=len(gk); gap=0.5; cw=(self.CW-gap*(n-1))/n; y0=1.75
        colTop=y0; colBot=self.H-0.35
        for i,g in enumerate(gk):
            x=self.MX+i*(cw+gap)
            # 학년 큰 라벨 + 범위(작게 옆)
            self._t(s,x,colTop,cw,0.5,
                    [(g+"  ",{"bold":True,"color":P["main"],"size":22}),
                     (ranges.get(g,""),{"color":P["accent"],"size":11})],
                    anchor=MSO_ANCHOR.BOTTOM)
            # 굵은 밑줄 (액센트)
            self._rect(s,x,colTop+0.6,0.5,0.05,P["accent"])
            self._line(s,x+0.5,colTop+0.625,cw-0.5,P["line"],0.75)
            # 반 리스트 (세로 넉넉히)
            items=cls[g]; ly=colTop+0.9
            step=min(0.55,(colBot-ly)/max(len(items),1))
            for c in items:
                nm=c.get("n","")
                self._t(s,x,ly,0.25,0.35,"·",size=14,bold=True,color=P["accent"])
                self._t(s,x+0.28,ly,cw-0.3,0.4,nm,size=self._fit(nm,12,18),color=BLACK,
                        anchor=MSO_ANCHOR.MIDDLE)
                ly+=step

    def curriculum(self):
        cur=self.ct.get("curriculum_math")
        if not cur: return
        s=self._slide(); P=self.P
        hl=(self.ct.get("headlines") or {}).get("curriculum")
        self._head(s,"교육 과정","연계 커리큘럼","초등부터 대입까지, 흐름을 끊지 않는 설계",headline=hl)
        tracks=list(cur.keys()); n=len(tracks)
        gap=0.3; cw=(self.CW-gap*(n-1))/n; y0=1.7; ch=3.55
        for i,t in enumerate(tracks):
            tr=cur[t]; x=self.MX+i*(cw+gap)
            # 헤더 (학년+범위)
            self._rect(s,x,y0,cw,0.5,P["main"])
            self._t(s,x,y0,cw,0.5,[(t+"  ",{"bold":True,"color":WHITE,"size":13}),
                                   ("("+tr.get("range","")+")",{"color":"C4D2E4","size":9.5})],
                    align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
            self._rect(s,x,y0+0.5,cw,ch-0.5,WHITE,line=P["line"])
            # 목표
            self._t(s,x+0.2,y0+0.62,cw-0.4,0.4,tr.get("goal",""),size=10.5,bold=True,
                    color=P["accent"],ls=1.1)
            # 단계
            ly=y0+1.2
            for st in tr.get("steps",[]):
                self._t(s,x+0.2,ly,cw-0.4,0.26,st.get("label",""),size=10,bold=True,color=P["main"])
                self._t(s,x+0.2,ly+0.22,cw-0.4,0.38,st.get("desc",""),size=9,color="555F6B",ls=1.05)
                ly+=0.64
            # 전환점 강조: 액센트 배경 박스 + 화살표
            if tr.get("turn"):
                turnH=0.72
                turnY=y0+ch-turnH-0.15
                self._rect(s,x+0.12,turnY,cw-0.24,turnH,P["soft"])
                self._rect(s,x+0.12,turnY,0.06,turnH,P["accent"])
                self._t(s,x+0.28,turnY+0.08,cw-0.44,0.24,
                        [("↓ 전환점",{"bold":True,"color":P["accent"],"size":9.5})])
                self._t(s,x+0.28,turnY+0.3,cw-0.44,turnH-0.35,tr["turn"].get("desc",""),
                        size=9,bold=True,color=P["main"],ls=1.1)

    def timetable(self):
        DAY={"MONDAY":"월","TUESDAY":"화","WEDNESDAY":"수","THURSDAY":"목",
             "FRIDAY":"금","SATURDAY":"토","SUNDAY":"일"}
        courses=((self.raw.get("timetable") or {}).get("courses") or [])
        tt=[]
        for co in courses:
            slots=((co.get("slots") or {}).get("slots") or [])
            times=", ".join(DAY.get(x.get("day"),"")+" "+str(x.get("start","")) for x in slots)
            tt.append({"name":(co.get("courseName") or "").strip(),
                       "teacher":(co.get("staffName") or "").strip(),
                       "room":(co.get("room") or "").strip(),"times":times,
                       "raw":co.get("courseName") or ""})
        if not tt: return
        def grade(nm):
            if re.search(r"초등|초\d",nm): return "초등"
            if re.search(r"방학특강",nm): return "특강·기타"
            if re.search(r"약술|물리|화학|생명|지구|논술",nm): return "특강·기타"
            if re.search(r"중등|중학|중\d",nm): return "중등"
            if re.search(r"고등|고\d|통합과학",nm): return "고등"
            return "특강·기타"
        def draw(rows,suffix):
            s=self._slide(); P=self.P
            self._head(s,"교육 과정","시간표 · "+suffix,None)
            heads=["수업 과목","요일 · 시간","강사","강의실"]
            nrows=len(rows)+1; colw=[1.9,3.4,0.95,0.95]
            # 세로 공간(1.55~5.3)에 맞춰 행 높이 자동 (최소 0.4, 최대 0.7)
            availH=self.H-1.55-0.35
            rowH=min(0.7,max(0.4,availH/nrows))
            tblH=rowH*nrows
            tblY=1.55+(availH-tblH)/2 if tblH<availH else 1.55
            tbl=s.shapes.add_table(nrows,4,Inches(self.MX),Inches(tblY),
                                   Inches(sum(colw)),Inches(tblH)).table
            tbl.first_row=False; tbl.horz_banding=False
            for ci,w in enumerate(colw): tbl.columns[ci].width=Inches(w)
            for ri in range(nrows): tbl.rows[ri].height=Inches(rowH)
            for ci,h in enumerate(heads):
                cell=tbl.cell(0,ci); cell.fill.solid(); cell.fill.fore_color.rgb=C(P["main"])
                self._cell(cell,h,10.5,True,WHITE)
            for ri,row in enumerate(rows,start=1):
                nm=re.sub(r"[월화수목금토일()0-9:.\s]+$","",row["name"]) or row["name"]
                vals=[nm,row["times"],row["teacher"],row["room"] or "-"]
                for ci,v in enumerate(vals):
                    cell=tbl.cell(ri,ci); cell.fill.solid()
                    cell.fill.fore_color.rgb=C(P["soft"] if ri%2 else WHITE)
                    self._cell(cell,v,10,False,BLACK)
        if len(tt)<=11: draw(tt,"전체")
        else:
            order=["초등","중등","고등","특강·기타"]; byG={}
            for c in tt: byG.setdefault(grade(c["raw"]),[]).append(c)
            for g in order:
                if byG.get(g): draw(byG[g],g)

    def _cell(self,cell,text,size,bold,color):
        for m,v in (("margin_left",0.06),("margin_right",0.06),("margin_top",0.02),("margin_bottom",0.02)):
            setattr(cell,m,Inches(v))
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=text
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=C(color); r.font.name="맑은 고딕"

    def specials(self):
        sp=self.ct.get("specials") or []
        if not sp: return
        s=self._slide(); P=self.P
        self._head(s,"특별 과정","특강 · 특별 프로그램","정규 수업 외 개별 성장 지원")
        cols=2; gap=0.25; cw=(self.CW-gap)/cols; rows=(len(sp)+1)//2
        ch=min(1.7,(3.5)/rows-0.15); y0=1.7
        for i,p in enumerate(sp):
            r=i//cols; c=i%cols; x=self.MX+c*(cw+gap); y=y0+r*(ch+0.2)
            self._t(s,x,y,0.5,0.35,f"{i+1:02d}",size=14,bold=True,color=P["accent"])
            self._line(s,x,y+0.4,cw,P["line"],0.75)
            self._t(s,x,y+0.48,cw,0.35,p.get("t",""),size=12,bold=True,color=P["main"])
            if p.get("d"):
                self._t(s,x,y+0.84,cw,ch-0.9,p["d"],size=10,color=BLACK,ls=1.25)

    def management(self):
        mg=self.ct.get("management") or {}; subs=list(mg.keys())
        if not subs: return
        s=self._slide(); P=self.P
        hl=(self.ct.get("headlines") or {}).get("management")
        by=self._head(s,"학습 관리","과목별 학생 관리","편성·과제·테스트·상담",headline=hl)
        # SG식 매트릭스: 행=과목, 열=관리 항목
        # 항목 키 수집 (순서 유지)
        cols=[]
        for it in mg[subs[0]]:
            cols.append(it.get("k",""))
        ncol=len(cols)
        # 레이아웃: 좌측 과목 라벨열 + 항목 열들
        labW=1.1
        colW=(self.CW-labW)/ncol
        headY=by+0.05; headH=0.42
        # 헤더 (컬러바)
        self._rect(s,self.MX,headY,labW,headH,P["soft"])
        self._t(s,self.MX,headY,labW,headH,"과목",size=10,bold=True,color=P["main"],
                align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        for ci,ck in enumerate(cols):
            cx=self.MX+labW+ci*colW
            self._rect(s,cx,headY,colW,headH,P["main"])
            self._t(s,cx,headY,colW,headH,ck,size=10,bold=True,color=WHITE,
                    align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        # 행 (과목별)
        rowY=headY+headH
        rowH=(self.H-rowY-0.35)/len(subs)
        for ri,sub in enumerate(subs):
            y=rowY+ri*rowH
            # 과목 라벨
            self._rect(s,self.MX,y,labW,rowH,P["accent"])
            self._t(s,self.MX,y,labW,rowH,sub,size=13,bold=True,color=WHITE,
                    align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
            for ci,it in enumerate(mg[sub]):
                cx=self.MX+labW+ci*colW
                self._rect(s,cx,y,colW,rowH,P["soft"] if ri%2 else WHITE,line=P["line"],lw=0.5)
                self._t(s,cx+0.12,y+0.08,colW-0.24,rowH-0.16,it.get("v",""),
                        size=8.5,color=BLACK,ls=1.15)
            rowY_=y

    def admission(self):
        ad=self.ct.get("admission") or []
        if not ad: return
        s=self._slide(); P=self.P
        self._head(s,"등록 안내","입학 절차","상담부터 등원까지")
        n=len(ad); gap=0.4; cw=min(3.6,(self.CW-gap*(n-1))/n)
        totW=cw*n+gap*(n-1); x0=self.MX+(self.CW-totW)/2; y=2.0; ch=2.3
        for i,a in enumerate(ad):
            x=x0+i*(cw+gap); emph=(i==n-1)
            self._rect(s,x,y,cw,ch,P["soft"] if emph else WHITE,
                       line=P["accent"] if emph else P["line"],lw=1.5 if emph else 1)
            self._t(s,x,y+0.35,cw,0.7,f"STEP {i+1}",size=26,bold=True,color=P["accent"],
                    align=PP_ALIGN.CENTER)
            self._t(s,x,y+1.15,cw,0.4,a.get("step",""),size=13,bold=True,color=P["main"],
                    align=PP_ALIGN.CENTER)
            self._t(s,x+0.25,y+1.55,cw-0.5,ch-1.6,a.get("desc",""),size=9.5,color=self.P["muted"],
                    align=PP_ALIGN.CENTER,ls=1.2)
            if i<n-1:
                self._t(s,x+cw,y,gap,ch,"›",size=22,bold=True,color=P["line"],
                        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

    def policy(self):
        rules=self.ct.get("rules") or []
        if not rules: return
        s=self._slide(); P=self.P
        self._head(s,"운영 지침","학원 관리 지침","출결·과제·결석·퇴원·환불 기준")
        y0=1.7; rowH=(3.6)/len(rules)
        for i,r in enumerate(rules):
            y=y0+i*rowH
            if i>0: self._line(s,self.MX,y,self.CW,P["line"],0.5)
            self._rect(s,self.MX,y+rowH*0.15,1.1,rowH*0.55,P["main"])
            self._t(s,self.MX,y+rowH*0.15,1.1,rowH*0.55,r.get("k",""),size=11,bold=True,
                    color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
            self._t(s,self.MX+1.35,y,self.CW-1.4,rowH,r.get("v",""),size=11,color=BLACK,
                    anchor=MSO_ANCHOR.MIDDLE,ls=1.2)

    def faq(self):
        fa=(self.ct.get("faq") or [])[:4]
        if not fa: return
        s=self._slide(); P=self.P
        self._head(s,"안내","자주 묻는 질문")
        cols=2; gap=0.3; cw=(self.CW-gap)/cols; ch=1.75; y0=1.65
        for i,f in enumerate(fa):
            r=i//cols; c=i%cols; x=self.MX+c*(cw+gap); y=y0+r*(ch+0.25)
            self._rect(s,x,y,cw,ch,P["soft"])
            self._t(s,x+0.25,y+0.2,cw-0.5,0.6,
                    [("Q  ",{"bold":True,"color":P["accent"]}),(f.get("q",""),{"bold":True,"color":P["main"]})],
                    size=11,ls=1.15)
            self._t(s,x+0.25,y+0.85,cw-0.5,ch-0.95,
                    [("A  ",{"bold":True,"color":P["muted"]}),(f.get("a",""),{"color":BLACK})],
                    size=9.5,ls=1.2)

    def contact(self):
        ct=self.ct.get("contact") or {}
        s=self._slide(); P=self.P
        self._rect(s,0,0,self.W,self.H,P["main"])
        self._t(s,0.6,0.9,8.8,0.35,"CONTACT",size=11,bold=True,color="7FA8D8",spacing=1.5)
        self._t(s,0.6,1.3,8.8,0.7,self.name,size=30,bold=True,color=WHITE)
        y=2.6
        def row(label,val):
            nonlocal y
            self._t(s,0.6,y,8.8,0.4,[(label+"    ",{"bold":True,"color":WHITE}),(val,{"color":"C4D2E4"})],size=12)
            y+=0.55
        if ct.get("phone"): row("전화",ct["phone"])
        if ct.get("address"):
            row("주소",ct["address"]+("  ("+ct["parking"]+")" if ct.get("parking") else ""))
        if ct.get("hours"): row("운영",", ".join(h["k"]+" "+h["v"] for h in ct["hours"]))
        self._line(s,0.6,y+0.05,8.8,"33507A",0.75); y+=0.25
        chans=" ".join(l["k"]+" "+l["v"] for l in (ct.get("links") or []))
        if chans:
            self._t(s,0.6,y,8.8,0.6,chans,size=10,color="9FB4CE",ls=1.4)

    def build(self,out):
        self.cover(); self.intro(); self.strengths(); self.achievements()
        self.subjects_targets(); self.curriculum(); self.timetable()
        self.specials(); self.management(); self.admission(); self.policy()
        self.faq(); self.contact()
        self.prs.save(out); return out

def main():
    content=json.load(open(sys.argv[1],encoding="utf-8"))
    raw=json.load(open(sys.argv[2],encoding="utf-8"))
    out=sys.argv[3] if len(sys.argv)>3 else "out.pptx"
    pal=sys.argv[4] if len(sys.argv)>4 else "navy"
    Report(content,raw,pal).build(out)
    print("생성:",out,"| 팔레트:",pal)

if __name__=="__main__":
    main()
