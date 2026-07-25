# -*- coding: utf-8 -*-
"""
navy_core.py — 브랜드북 navy 스킨 공통 코어
기준점: 더바른수학전문학원 초등부 소개서(네이비강조) / 3단리플렛(네이비골드)

기존 build_brandbook_v3.py / build_brandbook_book.py 는 건드리지 않는다.
navy 스킨은 v3 스키마를 그대로 입력받는 별도 렌더러다.
"""
import base64
import io
import re
import urllib.request

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FONT = "Pretendard SemiBold"

# ── 팔레트 (기준점 실측값) ────────────────────────────────────────────
PALETTES = {
    "navy_gold": {
        "key": "navy_gold",
        "label": "네이비 & 골드",
        "primary":   "18324A",  # 네이비 (라벨/번호/강조)
        "deep":      "172D43",  # 딥네이비 (표지 패널/CTA 배경)
        "accent":    "B79A58",  # 골드 (포인트/CTA)
        "card":      "F3F7F8",  # 카드 배경
        "card2":     "E8EEF2",  # 보조 카드 배경
        "line":      "D8E1E4",  # 카드 테두리/구분선
        "text":      "252A2E",  # 본문 먹색
        "muted":     "69747C",  # 보조 텍스트
        "onDark":    "FFFFFF",
        "onDarkSub": "C9D5DE",
    },
    "charcoal_gold": {
        "key": "charcoal_gold",
        "label": "차콜 & 골드",
        "primary": "2B2F33", "deep": "1C1F22", "accent": "B79A58",
        "card": "F4F5F6", "card2": "EAECEE", "line": "DCDFE2",
        "text": "252A2E", "muted": "6B7075",
        "onDark": "FFFFFF", "onDarkSub": "D2D6DA",
    },
    "forest_gold": {
        "key": "forest_gold",
        "label": "포레스트 & 골드",
        "primary": "1F4739", "deep": "16342A", "accent": "B79A58",
        "card": "F2F7F4", "card2": "E6EFEA", "line": "D5E2DB",
        "text": "252A2E", "muted": "68746E",
        "onDark": "FFFFFF", "onDarkSub": "C7D8CF",
    },
}
DEFAULT_PALETTE = "navy_gold"


def get_palette(name=None):
    if isinstance(name, dict):
        return name
    return PALETTES.get(name or DEFAULT_PALETTE, PALETTES[DEFAULT_PALETTE])


# ── 폰트 크기 규격 (기준점 실측) ──────────────────────────────────────
class FS:
    eyebrow   = 9.5     # 영문 아이브로 라벨
    header    = 24.0    # 페이지 헤더
    header_sm = 21.0    # 긴 헤더 축소
    lead      = 13.0    # 리드 문장
    card_num  = 11.5    # 카드 번호(01/02)
    card_ttl  = 15.0    # 카드 제목
    body      = 11.5    # 본문 (가독성 위해 상향)
    body_sm   = 10.5    # 본문 축소
    small     = 9.5     # 캡션/태그
    tiny      = 8.0     # 페이지 번호
    big_num   = 30.0    # 큰 숫자(실적)
    stage_ttl = 17.0    # 커리큘럼 단계 제목
    pill      = 9.5     # 필/태그
    cover_ttl = 30.0    # 표지 타이틀
    cover_sub = 13.5    # 표지 서브


# ── 그리드 (13.33 x 7.5in) ───────────────────────────────────────────
class G:
    W = 13.333
    H = 7.5
    ML = 0.6                  # 좌 마진
    MR = 0.6                  # 우 마진
    CW = 13.333 - 1.2         # 콘텐츠 폭 12.13
    y_eyebrow = 0.35
    y_header  = 0.65
    y_lead    = 1.23
    y_body    = 1.98          # 리드 있을 때
    y_body_nolead = 1.88      # 리드 없을 때
    y_page    = 7.08


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def strip_style(shape):
    """python-pptx auto_shape 가 자동 삽입하는 <p:style> 제거.
    이걸 안 지우면 테마색·그림자가 오염된다."""
    sp = shape._element
    for st in sp.findall(qn('p:style')):
        sp.remove(st)


def no_line(shape):
    shape.line.fill.background()


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
            radius=None, shape_type=None):
    """도형 하나. fill/line 은 hex 문자열 또는 None(투명)."""
    st = shape_type or (MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE)
    sh = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    strip_style(sh)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(fill)
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = _rgb(line)
        sh.line.width = Pt(line_w)
    else:
        no_line(sh)
    if radius is not None and st == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    sh.shadow.inherit = False
    return sh


def add_text(slide, x, y, w, h, text, size=FS.body, bold=False,
             color="252A2E", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.3, space_after=0, wrap=True, shrink=True,
             min_size=None, clip=True):
    """
    텍스트 박스. text 는 문자열 또는 문자열 리스트(여러 문단).

    shrink=True (기본): 박스에 안 들어가면 폰트를 자동으로 줄인다.
    clip=True  (기본): 최소 크기로도 안 들어가면 문장 경계에서 잘라내고 '…' 처리.
                       핵심 카피는 clip=False 로 보호할 것.
    """
    lines = text if isinstance(text, (list, tuple)) else [text]
    lines = [clean(t) for t in lines if t is not None and clean(t) != ""]
    if not lines:
        lines = [""]

    if shrink and h and w:
        joined = "\n".join(lines)
        size = fit_size(joined, size, w, h, min_size=min_size)
        if clip:
            lines = _clip_to_box(lines, size, w, h, line_spacing)

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    for i, t in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = t
        f = r.font
        f.name = FONT
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = _rgb(color)
        _latin_font(r)
    return tb


def _est_lines(text, size_pt, width_in, cpl_factor=1.85):
    if not text:
        return 0
    cpl = max(1, int(width_in * 72 / (size_pt * cpl_factor) * 1.85))
    return max(1, -(-len(text) // cpl))


def _clip_to_box(lines, size, width_in, height_in, line_spacing=1.3):
    """
    최소 폰트로도 넘치면 들어가는 만큼만 남긴다.
    문장 경계(。.!?) 우선, 없으면 어절 경계에서 자르고 '…' 을 붙인다.
    (콘텐츠 원칙: 첫 문장은 보존, 설명·부연만 트리밍)
    """
    max_lines = max(1, int(height_in * 72 / (size * line_spacing)))
    out, used = [], 0
    for ln in lines:
        if used >= max_lines:
            break
        need = _est_lines(ln, size, width_in)
        if used + need <= max_lines:
            out.append(ln)
            used += need
            continue

        room = max_lines - used
        cpl = max(1, int(width_in * 72 / (size * 1.85) * 1.85))
        budget = max(8, room * cpl - 1)

        head = ln[:budget]
        cut = max(head.rfind("다. "), head.rfind(". "),
                  head.rfind("! "), head.rfind("? "))
        if cut > budget * 0.45:
            out.append(ln[:cut + 1].strip())
        else:
            sp = head.rfind(" ")
            head = head[:sp] if sp > budget * 0.5 else head
            out.append(head.rstrip(" ,·-") + "…")
        used = max_lines
        break
    return out or [""]


def _latin_font(run):
    """한글 폰트가 latin/ea 양쪽에 걸리도록."""
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', FONT)


def add_pill(slide, x, y, w, h, text, fill, color="FFFFFF",
             size=FS.pill, bold=True, radius=0.5):
    """둥근 필 배지 (배경 + 중앙 텍스트)."""
    add_box(slide, x, y, w, h, fill=fill, radius=radius)
    add_text(slide, x, y, w, h, text, size=size, bold=bold, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0, shrink=False, clip=False)


def clean(s):
    """non-breaking space 등 정리. 한글 텍스트에 \xa0 빈번."""
    if s is None:
        return ""
    s = str(s)
    s = s.replace("\xa0", " ").replace("\u200b", "")
    s = re.sub(r"[ \t]+", " ", s)
    return s.strip()


# ── 텍스트 길이 기반 자동 축소 ────────────────────────────────────────
def fit_size(text, base, width_in, height_in, min_size=None, cpl_factor=1.85):
    """
    글자수 기준으로 폰트 크기를 낮춘다.
    한글은 폭이 넓어 1글자 ≈ size pt. cpl_factor 로 보정.
    가독성 하한은 9.0pt (인쇄물 기준). 그 아래로는 내리지 않는다.
    """
    READABLE_MIN = 9.0
    text = clean(text)
    if not text:
        return base
    min_size = max(READABLE_MIN, min_size or base * 0.75)
    if min_size > base:
        min_size = base
    size = base
    while size > min_size:
        cpl = max(1, int(width_in * 72 / (size * cpl_factor) * 1.85))
        lines = 0
        for seg in text.split("\n"):
            lines += max(1, -(-len(seg) // cpl))
        need = lines * size * 1.35 / 72.0
        if need <= height_in:
            return round(size, 2)
        size -= 0.5
    return round(min_size, 2)


def wrap_words(text, max_chars):
    """어절 단위 줄바꿈. 단어 중간에서 안 끊기게."""
    text = clean(text)
    words = text.split(" ")
    lines, cur = [], ""
    for w in words:
        if not cur:
            cur = w
        elif len(cur) + 1 + len(w) <= max_chars:
            cur += " " + w
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


# ── 페이지 공통 요소 ─────────────────────────────────────────────────
def page_head(slide, P, eyebrow, header, lead=None, page_no=None):
    """섹션 라벨(한글, 작게·회색) → 헤더 → 리드 3단 + 우하단 페이지번호.
    ★영문 아이브로 금지 원칙 — .upper() 폐기, 한글 라벨을 그대로 쓴다.
    반환값: 콘텐츠 시작 가능 y (헤더가 2줄이면 그만큼 내려간다)."""
    if eyebrow:
        add_text(slide, G.ML, G.y_eyebrow, 5.0, 0.25, clean(eyebrow),
                 size=FS.eyebrow, bold=True, color=P["primary"],
                 line_spacing=1.0, shrink=False, clip=False)

    y = G.y_header
    if header:
        htxt = clean(header)
        hs = FS.header if len(htxt) <= 30 else FS.header_sm
        hlines = wrap_words(htxt, max(14, int(G.CW * 72 / (hs * 1.62))))[:2]
        hh = max(0.56, len(hlines) * hs * 1.18 / 72.0 + 0.04)
        add_text(slide, G.ML, y, G.CW, hh, hlines,
                 size=hs, bold=True, color=P["text"], line_spacing=1.18,
                 clip=False)
        y += hh + 0.06

    if lead:
        ltxt = clean(lead)
        llines = wrap_words(ltxt, max(20, int(G.CW * 72 / (FS.lead * 1.55))))[:2]
        lh = len(llines) * FS.lead * 1.35 / 72.0 + 0.04
        add_text(slide, G.ML, y, G.CW - 0.2, lh, llines,
                 size=FS.lead, color=P["muted"], line_spacing=1.35)
        y += lh

    if page_no is not None:
        add_text(slide, G.W - G.MR - 0.57, G.y_page, 0.57, 0.19,
                 f"{page_no:02d}", size=FS.tiny, bold=True,
                 color=P["muted"], align=PP_ALIGN.RIGHT, line_spacing=1.0,
                 shrink=False, clip=False)

    return max(y + 0.34, G.y_body)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def force_theme_font(prs, name=None):
    """테마 majorFont/minorFont 를 지정 서체로 바꾼다.
    서체 지정이 누락된 텍스트가 있어도 Calibri(한글 글리프 없음)로
    떨어지지 않게 하는 안전장치. theme1.xml 은 python-pptx 가 파싱하지
    않는 일반 Part 라 blob 을 직접 치환한다."""
    name = name or FONT
    try:
        for part in prs.part.package.iter_parts():
            if "theme" not in str(part.partname):
                continue
            xml = part.blob.decode("utf-8", errors="ignore")

            def _fix(b):
                for slot in ("latin", "ea", "cs"):
                    b = re.sub(r'(<a:%s[^/>]*typeface=")[^"]*(")' % slot,
                               r"\1" + name + r"\2", b)
                return b

            for tag in ("majorFont", "minorFont"):
                m = re.search(r"<a:%s>.*?</a:%s>" % (tag, tag), xml, re.S)
                if m:
                    xml = xml[:m.start()] + _fix(m.group(0)) + xml[m.end():]
            part._blob = xml.encode("utf-8")
    except Exception:
        pass


# ── 이미지 로딩 (URL / base64 data URL / 경로 / bytes) ───────────────
def fetch_image(src, timeout=12):
    """v3와 동일 계약: http(s) URL, data:base64, 로컬경로, bytes 모두 지원."""
    if not src:
        return None
    try:
        if isinstance(src, (bytes, bytearray)):
            return io.BytesIO(src)
        s = str(src).strip()
        if s.startswith("data:"):
            b64 = s.split(",", 1)[1] if "," in s else ""
            return io.BytesIO(base64.b64decode(b64))
        if s.startswith("http://") or s.startswith("https://"):
            req = urllib.request.Request(s, headers={"User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(req, timeout=timeout) as r:
                return io.BytesIO(r.read())
        if len(s) > 200 and not s.startswith("/"):
            # 순수 base64 문자열
            return io.BytesIO(base64.b64decode(s))
        with open(s, "rb") as f:
            return io.BytesIO(f.read())
    except Exception:
        return None


def place_image(slide, src, x, y, w, h, cover=True):
    """
    cover=True  : 영역을 꽉 채우되 비율 유지(넘치는 부분은 자름 대신 중앙 배치)
    cover=False : 비율 유지하며 영역 안에 들어가게(로고 찌그러짐 방지)
    """
    stream = fetch_image(src)
    if stream is None:
        return None
    try:
        from PIL import Image
        stream.seek(0)
        im = Image.open(stream)
        iw, ih = im.size
        stream.seek(0)
    except Exception:
        stream.seek(0)
        return slide.shapes.add_picture(stream, Inches(x), Inches(y),
                                        Inches(w), Inches(h))

    if not iw or not ih:
        return None
    ar_img = iw / ih
    ar_box = w / h
    if cover:
        if ar_img > ar_box:
            nh, nw = h, h * ar_img
        else:
            nw, nh = w, w / ar_img
    else:
        if ar_img > ar_box:
            nw, nh = w, w / ar_img
        else:
            nh, nw = h, h * ar_img
    nx = x + (w - nw) / 2
    ny = y + (h - nh) / 2
    try:
        return slide.shapes.add_picture(stream, Inches(nx), Inches(ny),
                                        Inches(nw), Inches(nh))
    except Exception:
        return None


# ── 스키마 안전 접근 ─────────────────────────────────────────────────
def gv(d, *keys, default=None):
    """중첩 dict 안전 접근."""
    cur = d
    for k in keys:
        if not isinstance(cur, dict):
            return default
        cur = cur.get(k)
        if cur is None:
            return default
    return cur


def as_dicts(items, name_key="title"):
    """리스트 원소가 문자열이어도 dict로 정규화."""
    out = []
    for it in (items or []):
        if isinstance(it, dict):
            out.append(it)
        elif isinstance(it, str) and it.strip():
            out.append({name_key: it.strip()})
    return out


def nonempty(v):
    if v is None:
        return False
    if isinstance(v, str):
        return bool(v.strip())
    if isinstance(v, (list, tuple, dict)):
        return len(v) > 0
    return True
