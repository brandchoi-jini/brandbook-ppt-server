# -*- coding: utf-8 -*-
"""
브랜드북 빌더 v3 (좌표 기반 python-pptx)
- 세원 16장 시안(4번 톤)을 코드로 재현
- 팔레트 3종: teal_blue / navy_amber / green_orange
- 학원별 데이터 편차 스마트 채움:
  * 데이터 없는 섹션 슬라이드 자동 생략
  * 시간표는 그룹 수만큼 동적 생성, 행 많으면 행간 축소 + 자동 페이지 분할
  * 본문 넘침 시 폰트 auto-fit(길이 기반 pt 단계 축소)
슬라이드 크기 13.3 x 7.5 in. 서체 Pretendard(미설치 환경은 대체폰트).
"""
# ★배포 확인용 — Railway에 새 코드가 올라갔는지 /  응답에서 바로 볼 수 있게 한다
BUILD_VERSION = "2026-07-29k"

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import math
import io as _io, urllib.request as _urlreq

EMU_IN = 914400

# URL/경로 이미지를 메모리로 가져오기 (실패 시 None)
_IMG_CACHE = {}
def _fetch_image(url):
    if not url: return None
    if url in _IMG_CACHE: return _IMG_CACHE[url]
    try:
        if url.startswith("data:"):
            # 담당자가 PC에서 올린 이미지(base64 data URL)
            import base64 as _b64
            _, _, b64 = url.partition(",")
            data = _b64.b64decode(b64)
        elif url.startswith("http"):
            req = _urlreq.Request(url, headers={"User-Agent":"Mozilla/5.0"})
            data = _urlreq.urlopen(req, timeout=8).read()
        else:
            with open(url,"rb") as f: data = f.read()
        _IMG_CACHE[url] = data
        return data
    except Exception:
        _IMG_CACHE[url] = None
        return None

def place_image(slide, url, x, y, w, h, cover=True):
    """url 이미지를 (x,y,w,h) 박스에 넣음. cover=True면 박스를 꽉 채우고 넘치는 부분 crop.
    성공 True / 실패 False (실패 시 호출부가 회색 박스 유지)."""
    data = _fetch_image(url)
    if not data: return False
    try:
        from PIL import Image as _PILImage
        im = _PILImage.open(_io.BytesIO(data))
        iw, ih = im.size
        box_ratio = w/h; img_ratio = iw/ih
        pic = slide.shapes.add_picture(_io.BytesIO(data), Inches(x), Inches(y), Inches(w), Inches(h))
        if cover:
            # 박스를 꽉 채우도록 crop (가운데 기준)
            if img_ratio > box_ratio:
                # 이미지가 더 넓음 → 좌우 crop
                crop = (1 - box_ratio/img_ratio)/2
                pic.crop_left = crop; pic.crop_right = crop
            else:
                crop = (1 - img_ratio/box_ratio)/2
                pic.crop_top = crop; pic.crop_bottom = crop
        else:
            # 비율 유지해 박스 안에 전체가 들어가게(로고 찌그러짐 방지). 좌·상단 기준 정렬.
            if img_ratio > box_ratio:
                nw = w; nh = w/img_ratio
            else:
                nh = h; nw = h*img_ratio
            pic.width = Inches(nw); pic.height = Inches(nh)
            pic.left = Inches(x); pic.top = Inches(y)
        return True
    except Exception:
        try:
            slide.shapes.add_picture(_io.BytesIO(data), Inches(x), Inches(y), Inches(w), Inches(h))
            return True
        except Exception:
            return False

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Pretendard SemiBold"

# ★한글 서체 적용: python-pptx 의 font.name 은 <a:latin> 만 채운다.
#   한글은 <a:ea>(East Asian) 슬롯을 보므로 이걸 비워두면 테마 기본서체
#   (Calibri, 한글 글리프 없음)로 떨어져 PowerPoint 가 맑은 고딕 등으로
#   임의 대체한다. 세 슬롯을 모두 지정해야 지정한 서체가 실제로 나온다.
def _apply_font(run, name=None):
    name = name or FONT
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _force_theme_font(prs, name=None):
    """테마 majorFont/minorFont 를 지정 서체로 바꾼다.
    서체 지정이 누락된 텍스트가 있어도 Calibri(한글 없음)로 떨어지지 않게
    하는 안전장치. theme1.xml 은 python-pptx 가 파싱하지 않는 일반 Part 라
    blob 을 직접 치환한다."""
    name = name or FONT
    try:
        for part in prs.part.package.iter_parts():
            if "theme" not in str(part.partname):
                continue
            xml = part.blob.decode("utf-8", errors="ignore")
            def _fix(b):
                for slot in ("latin", "ea", "cs"):
                    b = re.sub(r'(<a:%s[^/>]*typeface=")[^"]*(")' % slot,
                               r"\1" + name + r"\2", b)
                return b
            for tag in ("majorFont", "minorFont"):
                m = re.search(r"<a:%s>.*?</a:%s>" % (tag, tag), xml, re.S)
                if m:
                    xml = xml[:m.start()] + _fix(m.group(0)) + xml[m.end():]
            part._blob = xml.encode("utf-8")
    except Exception:
        pass

# ---------------- 팔레트 ----------------
# pill_bg/pill_ink = 박스 소제목(바탕/글씨). 인쇄에서 흐려지지 않게 진한 바탕 + 흰 글씨.
# card = 박스 바탕(연한 색), soft/soft2 = 태그·구분용 옅은 톤.
PALETTES = {
    "navy_blue":  {"accent": "1F3350", "accent2": "3A6EA5", "navy": "1F3350",
                   "soft": "E4EBF5", "soft2": "EAF0F8", "card": "EDF2F9",
                   "pill_bg": "1F3350", "pill_ink": "FFFFFF"},
    "teal_blue":  {"accent": "12727E", "accent2": "3A6EA5", "navy": "1F3350",
                   "soft": "E0EDEF", "soft2": "E9EFF7", "card": "E9F2F3",
                   "pill_bg": "12727E", "pill_ink": "FFFFFF"},
    "navy_amber": {"accent": "22375F", "accent2": "C07A15", "navy": "22375F",
                   "soft": "E4EAF3", "soft2": "F8EDD8", "card": "EDF1F8",
                   "pill_bg": "22375F", "pill_ink": "FFFFFF"},
    "green_orange":{"accent": "1B6B4A", "accent2": "C96A2C", "navy": "1D3F30",
                   "soft": "DFEEE6", "soft2": "F8E7D8", "card": "E8F2EC",
                   "pill_bg": "1B6B4A", "pill_ink": "FFFFFF"},
}

INK = "2D3748"        # 본문 진한먹색
INK_STRONG = "1A2233" # 헤더
LABEL = "5B6B85"      # 라벨 회색
LINE = "E3E8EF"
CARD = "F6F8FB"
CARD_LINE = "E7ECF3"
MUTED = "9AA7BD"
PH_BG = "DBE4EF"      # 사진 자리 배경

# 시간표 학년 태그 색 (팔레트 무관 고정 구분색은 accent 계열로 매핑)
def group_color(pal, group):
    g = (group or "").strip()
    if g.startswith("초"): return pal["accent"]
    if g.startswith("중"): return pal["accent2"]
    if g.startswith("고"): return pal["navy"]
    if "특" in g:          return "7C5CFF"
    return "E08A2B"


# ---------------- 저수준 헬퍼 ----------------
def _rgb(h): return RGBColor.from_string(h)

def _set_bg_white(slide):
    # 슬라이드 배경 흰색
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb("FFFFFF")

def _no_line(shape):
    shape.line.fill.background()

def _fill(shape, hexc):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hexc)
    _no_line(shape)

def _strip_style(shape):
    """python-pptx auto_shape가 자동 삽입하는 <p:style>(테마색/그림자 오염) 제거"""
    sp = shape._element
    for st in sp.findall(qn('p:style')):
        sp.remove(st)

def _qr_png(url):
    """링크를 QR 이미지(BytesIO)로. qrcode 패키지가 없으면 None을 돌려주고 자리 박스를 쓴다."""
    try:
        import qrcode
        from io import BytesIO
        qr = qrcode.QRCode(version=None, box_size=8, border=1,
                           error_correction=qrcode.constants.ERROR_CORRECT_M)
        qr.add_data(str(url).strip()); qr.make(fit=True)
        img = qr.make_image(fill_color="#22375F", back_color="white").convert("RGB")
        buf = BytesIO(); img.save(buf, format="PNG"); buf.seek(0)
        return buf
    except Exception:
        return None


CORNER_R = 0.035    # 모서리 반경(인치) — 살짝만. 0.07은 여전히 둥글어 보였다.

def card_bg(pal):
    """박스 바탕 — 팔레트별 연한 색(없으면 기본 회백색)"""
    return pal.get("card", CARD)


def rect(slide, x, y, w, h, hexc=None, radius=False, line_hex=None, line_w=None):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    _strip_style(s)
    if radius:
        # PowerPoint 기본 adj(0.1667)는 캡슐처럼 과하게 둥글다 →
        # 실제 반경이 CORNER_R 인치가 되도록 짧은 변 기준으로 환산한다.
        try:
            adj = CORNER_R / max(0.01, min(w, h))
            s.adjustments[0] = max(0.004, min(0.05, adj))
        except Exception:
            pass
    if hexc:
        _fill(s, hexc)
    else:
        s.fill.background()
        _no_line(s)
    if line_hex:
        s.line.color.rgb = _rgb(line_hex)
        s.line.width = Pt(line_w or 1)
    else:
        if not hexc:
            _no_line(s)
    # 그림자 제거
    s.shadow.inherit = False
    return s

def hline(slide, x, y, w, hexc=LINE, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = _rgb(hexc)
    ln.line.width = Pt(weight)
    return ln

def vline(slide, x, y, h, hexc=LINE, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x), Inches(y+h))
    ln.line.color.rgb = _rgb(hexc)
    ln.line.width = Pt(weight)
    return ln

def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        line_spacing=1.0, wrap=True, space_after=None):
    """
    runs: list of (text, size_pt, bold, color_hex) 또는 문자열
    여러 문단은 runs 안에 dict {'para':[(...)]}로 넣거나, 리스트의 리스트로 전달
    여기서는 간단히: runs=[(text,size,bold,color), ...] = 한 문단 내 여러 run
    문단 구분이 필요하면 paras() 사용
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    if space_after is not None: p.space_after = Pt(space_after)
    if isinstance(runs, str):
        runs = [(runs, 18, False, INK)]
    for (t, sz, b, c) in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b
        r.font.color.rgb = _rgb(c); _apply_font(r)
    return tb

def paras(slide, x, y, w, h, plist, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          line_spacing=1.0, para_space=4, wrap=True, hang=0.0):
    """plist: list of paragraphs; each = list of (text,size,bold,color) runs.
    hang: >0이면 hanging indent(인치) — 둘째 줄부터 그만큼 들여써서 마커 뒤에 맞춤."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for para in plist:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(para_space)
        if hang > 0:
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(int(hang*EMU_IN)))
            pPr.set('indent', str(int(-hang*EMU_IN)))
        for (t, sz, b, c) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = _rgb(c); _apply_font(r)
    return tb

def bullets(slide, x, y, w, h, items, size=13, color=INK, dot_hex="0EA5B7",
            line_spacing=1.15, para_space=3):
    """불릿 목록 - 점은 색 있는 마커로. hanging indent로 둘째 줄을 텍스트 시작선에 맞춤."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    # 들여쓰기 폭: 글자 크기에 비례 (점+공백 만큼)
    indent_in = 0.02 + size*0.0135   # 대략 점"• " 폭
    marL = Emu(int(indent_in*EMU_IN))
    ind = Emu(int(-indent_in*EMU_IN))
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing
        p.space_after = Pt(para_space)
        # hanging indent: 문단 왼쪽 marL, 첫 줄만 -indent(점이 왼쪽으로)
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(int(indent_in*EMU_IN)))
        pPr.set('indent', str(int(-indent_in*EMU_IN)))
        r0 = p.add_run(); r0.text = "• "
        r0.font.size = Pt(size); r0.font.bold = True; r0.font.color.rgb = _rgb(dot_hex); _apply_font(r0)
        r1 = p.add_run(); r1.text = it
        r1.font.size = Pt(size); r1.font.color.rgb = _rgb(color); _apply_font(r1)
    return tb


# ---------------- auto-fit (길이 기반 pt 축소) ----------------
def fit_size(text, base, min_size, cap_per_size):
    """
    text 길이에 따라 base에서 단계적으로 줄임.
    cap_per_size: 해당 pt에서 '한 줄에 편안한 글자수 * 허용 줄수' 상한(대략).
    넘으면 1pt씩 내려 min_size까지.
    """
    n = len(text)
    size = base
    while size > min_size and n > cap_per_size(size):
        size -= 0.5
    return size


def _disp_width(text):
    """한글=1.0, 영문/숫자/기호=0.55 로 계산한 표시 폭(글자수 환산)."""
    w = 0.0
    for ch in text:
        o = ord(ch)
        # 한글 음절/자모 + CJK
        if 0xAC00 <= o <= 0xD7A3 or 0x3130 <= o <= 0x318F or 0x4E00 <= o <= 0x9FFF:
            w += 1.0
        elif ch in " ·":
            w += 0.4
        else:
            w += 0.55
    return w

def wrap_lines(text, width_in, size_pt):
    """텍스트박스 폭(in)과 폰트 크기(pt)로 실제 줄바꿈 수를 근사한다.
    Pretendard 기준 한글 한 글자 폭 ≈ size_pt*0.0139 in."""
    if not text:
        return 1
    char_in = size_pt * 0.0139
    chars_per_line = max(6, width_in / char_in)
    return max(1, math.ceil(_disp_width(text) / chars_per_line))


def fit_box(text, w, h, base, min_size, line_spacing=1.32, step=0.5):
    """텍스트박스 폭·높이 안에 들어오도록 글자 크기를 낮춘다(문장은 자르지 않는다)."""
    t = str(text or "")
    if not t: return base
    sz = base
    while sz > min_size:
        lines = sum(max(1, wrap_lines(ln, w, sz)) for ln in t.split("\n"))
        if lines * (sz*0.0139*1.62) * line_spacing <= h:
            return sz
        sz -= step
    return min_size


def _chunk(lst, n):
    """n개씩 나눈다. 개수가 많아도 잘라 버리지 않고 슬라이드를 나누기 위한 헬퍼."""
    n = max(1, int(n))
    return [lst[i:i+n] for i in range(0, len(lst), n)]

def _part(base, pi, parts):
    """여러 장으로 나뉠 때 라벨에 (1/2) 표기"""
    return base if parts <= 1 else f"{base} ({pi}/{parts})"


# ---------------- 공통 슬라이드 프레임 ----------------
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg_white(s)
    return s

PADX = 0.62      # 좌우 여백(인치)
PADY = 0.58      # 상단 여백

def label(slide, pal, text, x=PADX, y=PADY):
    """작은 라벨 + 세로바 (한글만)"""
    vline(slide, x+0.02, y+0.02, 0.20, pal["accent"], weight=3.2)
    # ★상자 폭 6.0 고정이라 오른쪽 열 텍스트와 <상자가> 겹쳤다 → 글자 폭만큼
    _lw = min(6.0, max(0.9, _disp_width(str(text or ""))*0.19 + 0.25))
    txt(slide, x+0.14, y-0.02, _lw, 0.3, [(text, 12, True, LABEL)], anchor=MSO_ANCHOR.TOP)

def header(slide, text, x=PADX, y=None, size=27, w=11.9):
    if y is None: y = PADY + 0.30
    t = str(text or "").strip()
    # ★헤더는 <한 문장>만 쓴다. 마침표로 끝난 <완결된 문장>이 둘 이상이면 첫 문장만 남긴다.
    #   (문장 도중을 자르는 것이 아니므로 말이 끊기지 않는다)
    if t:
        _sents = [x_ for x_ in re.split(r"(?<=[.!?])\s+", t) if x_.strip()]
        if len(_sents) > 1:
            t = _sents[0].strip()
        # ★예전에는 길면 쉼표에서 끊었다. 그래서 "수업이 끝난 그날 밤, 배운 내용을
        #   학부모님께 전달합니다" 가 "수업이 끝난 그날 밤" 으로 잘려 <말이 끊긴 채>
        #   슬라이드에 실렸다. 한마디는 문장이므로 도중에 자르지 않는다.
        #   길면 글자 크기만 낮춘다. (길이 조절은 앱이 <뜻을 살려 다시 쓰는> 방식으로 한다)
        if wrap_lines(t, w-0.14, 19) > 2:
            _ok = False
            for _s in (18, 17, 16, 15, 14):
                if wrap_lines(t, w-0.14, _s) <= 2:
                    size = min(size, _s)
                    _ok = True
                    break
            if not _ok:
                # ★여기서 문장을 자르지 않는다. 잘린 말은 어떤 경우에도 슬라이드에
                #   올리지 않는다("수업이 끝난 그날 밤" 처럼 끊겨 나가던 문제).
                #   길이 조절은 앱이 <뜻을 살려 짧게 다시 쓰는> 방식으로 처리한다.
                #   여기서는 최소 크기까지만 낮추고 그대로 담는다.
                size = min(size, 14)
        t = t.rstrip(" ,·—–")
    # 두 줄을 넘어가면 아래 본문과 겹친다 → 크기를 낮춰 두 줄 안에 담는다
    while size > 19 and wrap_lines(t, w-0.14, size) > 2:
        size -= 1
    # 제목 왼쪽을 라벨 글자 시작선(세로바 오른쪽, +0.14)에 맞춤
    # ★상자 높이 0.95 고정이면 한 줄일 때 아래 요소와 상자가 겹친다 → 실제 줄 수만큼
    _hh = max(1, wrap_lines(t, w-0.14, size))*(size*0.0139*1.35) + 0.10
    txt(slide, x+0.14, y, w-0.14, _hh, [(t, size, True, INK_STRONG)], line_spacing=1.12)


# ======================================================================
#  각 슬라이드 빌더
# ======================================================================

def slide_cover(prs, pal, d):
    s = new_slide(prs)
    a = d["academy"]
    assets = d.get("assets", {}) or {}
    photos = assets.get("photos", []) or []
    logo = assets.get("logo", "")
    # 우측 사진 영역 — 위아래 여백을 줘 표지가 답답하지 않게
    PH_X, PH_Y = 7.15, 0.62
    PH_W, PH_H = SLIDE_W - PH_X - 0.72, SLIDE_H - 0.62*2
    rect(s, PH_X, PH_Y, PH_W, PH_H, PH_BG)
    cover_img = assets.get("cover") or ""          # 담당자가 '표지 사진'으로 지정한 것만
    if cover_img and place_image(s, cover_img, PH_X, PH_Y, PH_W, PH_H, cover=True):
        pass
    else:
        txt(s, PH_X, PH_Y+PH_H/2-0.2, PH_W, 0.4,
            [("학생 학습 사진", 12, False, "9FB0C6")], align=PP_ALIGN.CENTER)
    # 좌측
    lx = 0.85
    name_y = 2.25
    # 로고 — 사진과 윗변을 맞추고 크기를 키움
    if logo:
        LOGO_W, LOGO_H = 2.1, 1.0
        if place_image(s, logo, lx, PH_Y, LOGO_W, LOGO_H, cover=False):
            pass  # 비율 유지
    # 슬로건
    slo = a.get("slogan", "")
    ssize = fit_size(slo, 18, 12, lambda sz: int(6.2/(sz*0.017)))
    txt(s, lx, 1.85, 6.3, 0.5, [(slo, ssize, False, pal["accent"])],
        align=PP_ALIGN.LEFT, wrap=True)
    # 학원명 — ★두 줄이 되는 긴 이름에서 과목 뱃지와 겹쳤다(고정 3.75) → 아래로 흐르게
    nm = str(a.get("name","") or "")
    BADGE_H, CONTACT_Y = 0.62, 5.75
    def _nm_h(sz):
        return max(1, wrap_lines(nm, 6.4, sz))*(sz*0.0139*1.30) + 0.10
    nsz = 52.0
    while nsz > 26:
        if wrap_lines(nm, 6.4, nsz) <= 2 and name_y + _nm_h(nsz) + 0.16 + BADGE_H <= CONTACT_Y - 0.20:
            break
        nsz -= 2
    nh = _nm_h(nsz)
    txt(s, lx, name_y, 6.4, nh, [(nm, nsz, True, pal["navy"])], line_spacing=1.0)
    # 과목 뱃지 — 학원명 아래로
    _bw = _badge_w(a.get("subjects",""))
    by = min(name_y + nh + 0.16, CONTACT_Y - BADGE_H - 0.20)
    badge = rect(s, lx, by, _bw, BADGE_H, pal["pill_bg"], radius=True)
    txt(s, lx, by, _bw, BADGE_H,
        [(a.get("subjects",""), 17, True, "FFFFFF")], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    # 연락처 (각 한 줄) — 전화가 두 개면 둘 다, 카카오 채널도 있으면 함께
    _tel = a.get("phone","")
    if a.get("phone2"):
        _tel = f"{_tel} · {a['phone2']}" if _tel else a["phone2"]
    _c_paras = [
        [("위치   ", 12, True, pal["navy"]), (a.get("location",""), 12, False, "667788")],
        [("연락처   ", 12, True, pal["navy"]), (_tel, 12, False, "667788")],
    ]
    if a.get("kakao"):
        _c_paras.append([("카카오   ", 12, True, pal["navy"]), (a.get("kakao",""), 12, False, "667788")])
    paras(s, lx, 5.75, 6.4, 1.5, _c_paras, line_spacing=1.2, para_space=6)
    return s

def _badge_w(text):
    # 한글 표시폭 반영 — len() 은 한글을 과소평가해 뱃지를 넘쳤다
    return max(1.9, 0.52 + _disp_width(str(text or ""))*0.23)

def slide_intro(prs, pal, d):
    """학원 소개 — 한마디 / 슬로건 / 소개문 / 강점 박스를 한 장에."""
    intro = d.get("intro") or {}
    body  = str(intro.get("body","")).strip()
    _fsrc = d.get("features") or []
    # ★features 가 {items:[...]} 로 오면 .get 에서 죽었다.
    if isinstance(_fsrc, dict):
        _fsrc = _fsrc.get("items") or []
    feats = [f for f in _fsrc if isinstance(f, dict) and str(f.get("desc","")).strip()][:6]
    if not body and not feats: return None
    a = d.get("academy", {})
    s = new_slide(prs)
    label(s, pal, "학원 소개")

    head = str(intro.get("head","")).strip()
    if head:
        header(s, head, size=26)
        hl = max(1, wrap_lines(head, 11.76, 26))     # 한마디가 두 줄이면 그만큼 아래에서 시작
        y = 0.88 + hl*0.50 + 0.16
    else:
        y = 1.10

    slogan = str(a.get("slogan","")).strip()
    if slogan:
        ssz = 15.5
        while ssz > 12.5 and wrap_lines(slogan, 11.76, ssz) > 1:
            ssz -= 0.5
        txt(s, PADX+0.14, y, 11.76, 0.38, [(slogan, ssz, True, pal["accent"])], line_spacing=1.1)
        y += 0.44

    if slogan and body:
        _key = re.sub(r"[^가-힣A-Za-z0-9]", "", slogan)
        _ss = [x.strip() for x in re.split(r"(?<=[.!?])\s+", body.replace("\n", " ")) if x.strip()]
        if len(_ss) > 1 and len(_key) >= 6:
            _f = re.sub(r"[^가-힣A-Za-z0-9]", "", _ss[0])
            if _key in _f or _f.startswith(_key[:8]):
                body = " ".join(_ss[1:])

    grid_top = 2.45
    if body:
        bsz = 14.0
        while bsz > 11.5 and wrap_lines(body, 11.76, bsz) > 3:
            bsz -= 0.5
        bl = max(1, wrap_lines(body, 11.76, bsz))
        bh = bl*(bsz*0.0139*1.40)*1.5 + 0.08
        txt(s, PADX+0.14, y, 11.76, bh+0.2, [(body, bsz, False, INK)], line_spacing=1.5)
        grid_top = y + bh + 0.26

    if feats:
        n = len(feats)
        ncol = 3 if n >= 5 else (n if n <= 3 else 2)
        nrow = (n+ncol-1)//ncol
        gap, gapy = 0.26, 0.22
        gw = (11.9-(ncol-1)*gap)/ncol
        gh = min(2.2, (7.05-grid_top-gapy*(nrow-1))/max(nrow,1))
        # ★배지 폭이 글자 수마다 달라 6칸이 들쭉날쭉했다.
        #   가장 긴 제목 기준으로 폭을 맞춰 통일감을 준다.
        _PW = min(gw-0.3, max(0.5 + _disp_width(str(f.get("title","")))*0.20
                              for f in feats))
        for i, f in enumerate(feats):
            r = i//ncol; c = i%ncol
            in_row = min(ncol, n - r*ncol)
            off = ((ncol-in_row)*(gw+gap))/2
            x = PADX + off + c*(gw+gap); yy = grid_top + r*(gh+gapy)
            rect(s, x, yy, gw, gh, card_bg(pal), radius=True, line_hex=CARD_LINE, line_w=1)
            title = str(f.get("title",""))
            pw = _PW
            rect(s, x+(gw-pw)/2, yy+0.18, pw, 0.40, pal["pill_bg"], radius=True)
            txt(s, x+(gw-pw)/2, yy+0.19, pw, 0.38, [(title, 12.5, True, pal["pill_ink"])],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            desc = str(f.get("desc",""))
            dh = gh-0.72
            dsz = fit_box(desc, gw-0.44, dh, 13.5, 9.5, line_spacing=1.32)
            txt(s, x+0.22, yy+0.66, gw-0.44, dh, [(desc, dsz, False, INK)],
                align=PP_ALIGN.CENTER, line_spacing=1.32)
    return s


def slide_achievements(prs, pal, d):
    """주요 실적 — 항목이 많으면 <슬라이드를 나눈다>.
    ★예전에는 한 장에 밀어 넣어 텍스트가 카드 밖(하단 8.49in)으로 나갔다.
      슬라이드 높이는 7.5in 이므로 인쇄하면 잘린다."""
    ach = d.get("achievements")
    if not ach: return None
    if not isinstance(ach, dict):
        ach = {"items": ach if isinstance(ach, list) else [ach]}

    def _norm(lst):
        """★항목이 문자열이나 None 으로 오면 .get 에서 죽었다 → dict 로 맞춘다."""
        out = []
        for it in (lst or []):
            if isinstance(it, dict):
                out.append(it)
            elif isinstance(it, str) and it.strip():
                out.append({"title": it.strip(), "change": "", "note": ""})
        return out

    groups = []
    for g in (ach.get("groups") or []):
        if not isinstance(g, dict):
            continue
        its = _norm(g.get("items"))
        if its:
            groups.append({"name": g.get("name", ""), "items": its})
    if not groups:
        items = _norm(ach.get("items"))
        if not items: return None
        groups = [{"name": "", "items": items}]
    MAXI, MAXG = 5, 3          # 카드당 항목 수 / 한 장당 카드 수
    pages, cur = [], []
    for g in groups:
        its = g["items"]
        for i in range(0, len(its), MAXI):
            if len(cur) >= MAXG:
                pages.append(cur); cur = []
            cur.append({"name": g.get("name",""), "items": its[i:i+MAXI]})
    if cur: pages.append(cur)
    # ★마지막 장에 카드가 1개만 남아 지면이 통째로 비었다(성장 사례 1건).
    #   앞 장에 자리가 있으면 끌어올려 합친다.
    if len(pages) >= 2 and len(pages[-1]) == 1 and len(pages[-2]) < MAXG + 1:
        pages[-2].extend(pages.pop())
    last = None
    for pi, gs in enumerate(pages, 1):
        last = _achievements_page(prs, pal, d, gs, pi, len(pages))
    return last


def _achievements_page(prs, pal, d, groups, pi=1, parts=1):
    ach = d.get("achievements") or {}
    # ★achievements 가 리스트로 바로 오면 .get 에서 죽었다.
    if not isinstance(ach, dict):
        ach = {}
    s = new_slide(prs)
    label(s, pal, _part("주요 실적", pi, parts))
    header(s, ach.get("head","") if pi == 1 else "", size=26)

    # ★하단 사진 띠는 쓰지 않는다. 가로로 납작하게 잘려(1.55인치) 내용이 안 보이고
    #   카드 영역만 좁아졌다. 실적은 카드로만 보여준다.
    top = 2.25 if (pi == 1 and str(ach.get("head","")).strip()) else 1.60
    bottom = 7.05

    n = len(groups)
    cols = min(n, 3 if n <= 3 else 2)
    rows = (n+cols-1)//cols
    gapx, gapy = 0.30, 0.24
    gw = (11.9-(cols-1)*gapx)/cols
    gh_max = (bottom-top-(rows-1)*gapy)/rows
    # 내용 분량으로 카드 높이를 정한다(빈 공간이 크게 남지 않게)
    _need = 1.0
    for g in groups:
        _rows = 0.0
        for x in g["items"]:
            _nt = str(x.get("note") or "")
            _nl = max(1, wrap_lines(_nt, gw-0.82, 11.5)) if _nt else 0
            _rows += 0.28 + _nl*0.26 + 0.14
        _need = max(_need, 0.92 + _rows + 0.26)
    gh = min(gh_max, max(2.1, _need))
    top = top + max(0.0, (bottom-top - (gh*rows + gapy*(rows-1)))/2)
    colors = [pal["accent"], pal["navy"], pal["accent2"], pal["accent"]]

    for gi, g in enumerate(groups):
        c = gi % cols; r = gi // cols
        x = PADX + c*(gw+gapx); y = top + r*(gh+gapy)
        kc = colors[gi % len(colors)]
        rect(s, x, y, gw, gh, card_bg(pal), radius=True, line_hex=CARD_LINE, line_w=1)
        # 유형 이름 — 진한 바탕 배지 + 아이콘
        gname = str(g.get("name","")).strip() or "주요 실적"
        _ICON = {"고입 실적":"", "대입 실적":"", "내신·성적 향상":"", "성장 사례":""}
        ico = _ICON.get(gname, "★")
        _lab = f"{ico}  {gname}"
        pw = min(gw-0.4, 0.6 + _disp_width(_lab)*0.20)
        rect(s, x+(gw-pw)/2, y+0.26, pw, 0.48, pal["pill_bg"], radius=True)
        tsz = 14 if _disp_width(_lab) < 12 else 12.5
        txt(s, x+(gw-pw)/2, y+0.27, pw, 0.46, [(_lab, tsz, True, pal["pill_ink"])],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        # 항목들
        area_y0 = y + 0.92
        area_h  = (y + gh - 0.22) - area_y0
        # 각 항목 높이를 내용으로 계산해 블록 전체를 세로 가운데에 둔다
        # ★높이 계산(11.5pt 고정)과 실제 렌더 크기가 어긋나 카드를 넘쳤다 → 같은 값으로
        _NSZ = 11.5
        def _calc(nsz):
            rs = []
            for it in g["items"]:
                _n = str(it.get("note") or "")
                _nl = max(1, wrap_lines(_n, gw-0.76, nsz)) if _n else 0
                rs.append(0.30 + _nl*(nsz*0.0139*1.42) + 0.14)
            return rs
        _rows = _calc(_NSZ)
        while sum(_rows) > area_h and _NSZ > 9.0:
            _NSZ -= 0.5
            _rows = _calc(_NSZ)
        _blk = sum(_rows)
        rowsy = area_y0 + max(0.0, (area_h - _blk)/2)
        for ii, it in enumerate(g["items"]):
            title  = str(it.get("title") or it.get("name") or "").strip()
            change = str(it.get("change") or "").strip()
            note   = str(it.get("note") or "").strip()
            if not change and not note:
                ls = [t.strip() for t in str(it.get("desc","")).split("\n") if t.strip()]
                change = ls[0] if ls else ""
                note = " / ".join(ls[1:])
            iy2 = rowsy + sum(_rows[:ii])
            # 제목 줄 — 왼줄 맞춤. 길면 크기를 줄여 한 줄에 담는다.
            _full = ("• " + title + ("  →  " if (title and change) else "") + change)
            tsz2 = 12.5
            while tsz2 > 10.0 and _disp_width(_full)*tsz2*0.0139 > (gw-0.52):
                tsz2 -= 0.5
            runs = [("•  ", tsz2, True, kc)]
            # ★화살표가 없어 '중1 입학   고3 수능까지…' 처럼 두 문장이 붙어 읽혔다.
            if title:  runs.append((title, tsz2, True, INK_STRONG))
            if change:
                if title: runs.append(("  →  ", tsz2, True, kc))
                runs.append((change, tsz2, True, kc))
            paras(s, x+0.26, iy2, gw-0.52, 0.30, [runs],
                  align=PP_ALIGN.LEFT, line_spacing=1.15, wrap=True)
            if note:
                # ★상자 높이(nl*0.26)와 블록 계산식(nsz*0.0139*1.42)이 달라 항목이 겹쳤다 → 통일
                nl = max(1, wrap_lines(note, gw-0.82, _NSZ))
                txt(s, x+0.52, iy2+0.30, gw-0.82, nl*(_NSZ*0.0139*1.42)+0.04,
                    [(note, _NSZ, False, "5A6A7A")], align=PP_ALIGN.LEFT, line_spacing=1.26)

    return s


def _stage_colors(pal, n):
    seq = [pal["accent"], pal["accent2"], pal["navy"]]
    if n <= 3:
        return seq[:n] if n > 1 else [pal["accent"]]
    out = []
    for i in range(n):
        out.append(seq[min(i, len(seq)-1)])
    return out


def _tint_for(pal, c):
    if c == pal["accent"]: return pal["soft"]
    if c == pal["accent2"]: return pal["soft2"]
    return "E7ECF4"


def slide_targets(prs, pal, d):
    """수업 대상 · 커리큘럼 — 위는 학년별 반 편성, 아래는 단계별 커리큘럼. 한 장에 묶는다."""
    tg = d.get("targets") or {}
    cl = d.get("classes") or {}
    cu = d.get("curriculum") or {}
    rows_src = tg.get("items") or []
    stages   = cu.get("stages") or []
    if not rows_src and not stages: return None

    s = new_slide(prs)
    label(s, pal, "수업 대상 · 과목")
    head = str(tg.get("head","") or cu.get("head","")).strip()
    if head:
        header(s, head, size=26)
        # ★대상 블록이 아래로 처져 커리큘럼 배지와 붙어 보였다 → 위로 올린다.
        top = 1.92
    else:
        top = 1.42

    colors = _stage_colors(pal, max(len(rows_src), len(stages), 1))
    photos = tg.get("_photos") or []

    # ── 위: 학년별 반 편성 ──
    y = top
    if rows_src:
        n = len(rows_src)
        # ★커리큘럼 카드 위에 배지가 떠 있어(카드 상단 -0.52in) 대상 블록과 겹쳤다.
        #   단계가 있으면 대상 영역을 그만큼 위에서 끊는다.
        area = (3.72 if stages else 7.05) - top
        rh = min(0.92, area/max(n,1))
        pw_col, gap = 1.15, 0.28
        has_ph = bool(photos)
        ph_w = 1.55 if has_ph else 0.0
        tw = 11.9 - pw_col - gap - (ph_w + 0.25 if has_ph else 0)
        for i, it in enumerate(rows_src):
            yy = y + i*rh
            kc = colors[i % len(colors)]
            gname = str(it.get("grade","")).strip()
            if gname:
                rect(s, PADX, yy+(rh-0.40)/2, pw_col, 0.40, _tint_for(pal, kc), radius=True)
                txt(s, PADX, yy+(rh-0.40)/2, pw_col, 0.40, [(gname, 13, True, kc)],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            subj = str(it.get("subj","")).strip()
            if subj:
                ssz = fit_box(subj, tw-0.1, rh-0.16, 15.5, 11.5, line_spacing=1.3)
                txt(s, PADX+pw_col+gap, yy+0.08, tw, rh-0.16,
                    [(subj, ssz, True, kc)], line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)
            if has_ph:
                px = 11.9+PADX-ph_w
                rect(s, px, yy+0.08, ph_w, rh-0.20, PH_BG, radius=True)
                u = photos[i % len(photos)]
                if u: place_image(s, u, px, yy+0.08, ph_w, rh-0.20, cover=True)
            if i < n-1:
                hline(s, PADX, yy+rh, 11.9, LINE, 1)
        y = top + n*rh + 0.16

    # ── 아래: 단계별 커리큘럼 ──
    if stages:
        if rows_src:
            hline(s, PADX, y, 11.9, LINE, 1)
            y += 0.72          # ★배지가 카드 위로 떠 있으므로 구분선과 넉넉히 띄운다
        n = len(stages)
        gap = 0.26
        # ★단계가 2개뿐인 학원(중등·고등만 운영)은 카드가 가로로 늘어나 밋밋했다.
        #   카드 폭에 상한을 두고 가운데로 모아, 3단계 학원과 같은 밀도로 보이게 한다.
        CARD_MAX = 4.15
        gw = min(CARD_MAX, (11.9-(n-1)*gap)/max(n,1))
        x0 = PADX + max(0.0, (11.9 - (gw*n + gap*(n-1)))/2)
        # ★계단식 — 뒤 단계일수록 카드를 위로 올려 <올라가는 흐름>이 보이게 한다.
        #   (평평하게 늘어놓으면 단계가 아니라 목록처럼 읽힌다)
        RISE = 0.0 if n <= 1 else min(0.52, (7.05 - y - 2.60) / max(1, n-1))
        step_colors = _stage_colors(pal, n)
        for i, st_ in enumerate(stages):
            x = x0 + i*(gw+gap)
            c = step_colors[i]
            cy = y + (n-1-i)*RISE          # 첫 단계가 가장 아래
            gh = 7.05 - cy
            rect(s, x, cy, gw, gh, card_bg(pal), radius=True, line_hex=CARD_LINE, line_w=1)
            if i > 0:      # 카드 사이 진행 화살표 — 올라가는 방향
                txt(s, x-gap-0.10, cy+0.34, gap+0.20, 0.40,
                    [("↗" if RISE > 0.05 else "→", 15, True, c)],
                    align=PP_ALIGN.CENTER, wrap=False)
            nm = str(st_.get("name","")).strip()
            tag = str(st_.get("tag","")).strip()
            iy = cy + 0.24
            if tag:
                # ★배지는 카드 위에 띄워 단계 이름과 겹치지 않게 한다
                pw = min(gw-0.4, 0.5 + _disp_width(tag)*0.19)
                by = cy - 0.52
                if by < PADY + 0.10:
                    by = cy + 0.10
                    iy = by + 0.46
                rect(s, x+0.22, by, pw, 0.36, pal["pill_bg"], radius=True)
                txt(s, x+0.30, by, pw, 0.36, [(tag, 12.5, True, pal["pill_ink"])],
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            if nm:
                txt(s, x+0.22, iy, gw-0.44, 0.44, [(nm, 17, True, INK_STRONG)], line_spacing=1.15)
                iy += 0.46
            items = [str(v).strip() for v in (st_.get("items") or []) if str(v).strip()]
            if items:
                ih = cy + gh - iy - 0.18
                isz = 13.5
                while isz > 9.5 and sum(max(1, wrap_lines(v, gw-0.62, isz)) for v in items)*(isz*0.0139*1.62)*1.30 > ih:
                    isz -= 0.5
                bullets(s, x+0.22, iy, gw-0.44, ih, items, size=isz,
                        color=INK, dot_hex=c, line_spacing=1.22, para_space=4)
    return s


def slide_timetable(prs, pal, d, group, rows, part=None, parts=None):
    s = new_slide(prs)
    label(s, pal, "시간표")
    gc = group_color(pal, group)
    # 학년 태그 + 헤더
    tag_w = 0.5 + len(group)*0.28
    rect(s, PADX, PADY+0.30, tag_w, 0.5, gc, radius=True)
    txt(s, PADX, PADY+0.35, tag_w, 0.42, [(group, 15, True, "FFFFFF")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ttl = "시간표" + (f"  ({part}/{parts})" if parts and parts>1 else "")
    txt(s, PADX+tag_w+0.25, PADY+0.28, 8.0, 0.6, [(ttl, 26, True, INK_STRONG)])
    # 표
    _draw_table(s, pal, gc, rows, top=1.95)
    return s

def _draw_table(s, pal, header_hex, rows, top):
    # ★ 진짜 표(python-pptx table 객체)로 생성 — 원장이 셀 단위로 편집 가능.
    #    (기존 박스+텍스트 레이어 방식 폐기)
    cols_ratio = [0.30, 0.34, 0.18, 0.18]  # 반 / 요일시간 / 강사 / 강의실
    tbl_w = 11.9
    headers = ["반", "요일 · 시간", "강사", "강의실"]
    n = len(rows)
    ncol = len(headers)

    # 행 높이 자동조절: 세로영역 ÷ (행수+헤더), 0.30~0.62 제한
    avail = 5.0
    rh = max(0.30, min(0.62, avail/(n+1)))
    hh = min(0.5, rh+0.04)
    body_font = 15 if n <= 7 else (14 if n <= 10 else 13)
    tbl_h = hh + rh*n

    gf = s.shapes.add_table(n+1, ncol, Inches(PADX), Inches(top), Inches(tbl_w), Inches(tbl_h))
    tbl = gf.table
    # 기본 표 스타일(밴딩/강조 행) 끄기 → 색을 직접 제어
    tbl.first_row = False
    tbl.horz_banding = False
    tbl.first_col = False
    # 기본 테마 표스타일(테두리/밴딩 색 오염) 제거 → "스타일 없음"으로
    try:
        tblPr = tbl._tbl.find(qn('a:tblPr'))
        if tblPr is not None:
            styleId = tblPr.find(qn('a:tableStyleId'))
            if styleId is not None:
                tblPr.remove(styleId)
    except Exception:
        pass

    # 열 너비
    for ci, ratio in enumerate(cols_ratio):
        tbl.columns[ci].width = Inches(ratio*tbl_w)
    # 헤더 행 높이
    tbl.rows[0].height = Inches(hh)
    for ri in range(1, n+1):
        tbl.rows[ri].height = Inches(rh)

    def _style_cell(cell, text, *, fs, bold, color, fill_hex, align_left=True):
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(fill_hex)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if align_left else PP_ALIGN.CENTER
        run = p.add_run(); run.text = text
        run.font.size = Pt(fs); run.font.bold = bold
        run.font.color.rgb = _rgb(color); _apply_font(run)

    # 헤더
    for ci, htext in enumerate(headers):
        _style_cell(tbl.cell(0, ci), htext, fs=14, bold=True,
                    color="FFFFFF", fill_hex=header_hex)
    # 바디 (지브라: 짝수 흰색 / 홀수 연회색)
    for ri, row in enumerate(rows):
        zebra = CARD if (ri % 2 == 1) else "FFFFFF"
        for ci in range(ncol):
            cell_text = row[ci] if ci < len(row) else ""
            color = MUTED if cell_text.strip() == "-" else INK
            _style_cell(tbl.cell(ri+1, ci), cell_text, fs=body_font, bold=False,
                        color=color, fill_hex=zebra)

    # 표 테두리(가로 구분선만 은은하게) — 셀 하단 라인
    _table_hlines(tbl, LINE)
    return gf


def _table_hlines(tbl, hex_line):
    """각 셀 하단에 가는 가로선 추가(표 편집성 유지하면서 구분감)."""
    from pptx.oxml.ns import qn as _qn
    for row in tbl.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            # 아래쪽 테두리 lnB
            for tag in ("a:lnB",):
                existing = tcPr.find(_qn(tag))
                if existing is not None:
                    tcPr.remove(existing)
                ln = tcPr.makeelement(_qn(tag), {"w": "6350", "cap": "flat"})  # 6350 EMU ≈ 0.5pt
                fill = ln.makeelement(_qn("a:solidFill"), {})
                clr = fill.makeelement(_qn("a:srgbClr"), {"val": hex_line})
                fill.append(clr); ln.append(fill)
                tcPr.append(ln)

def slide_specials(prs, pal, d):
    sp = d.get("specials")
    if not sp: return None
    if not isinstance(sp, dict):
        sp = {"items": sp if isinstance(sp, list) else [sp]}
    # ★항목이 문자열이나 None 으로 오면 .get 에서 죽었다 → dict 로 맞춘다.
    _its = []
    for it in (sp.get("items") or []):
        if isinstance(it, dict):
            if str(it.get("title") or it.get("name") or "").strip():
                _its.append(it)
        elif isinstance(it, str) and it.strip():
            _its.append({"title": it.strip(), "desc": ""})
    if not _its: return None
    sp = dict(sp); sp["items"] = _its
    pages = _chunk(sp["items"], 6)
    last = None
    for pi, items in enumerate(pages, 1):
        last = _specials_page(prs, pal, sp, items, pi, len(pages))
    return last

def _specials_page(prs, pal, sp, items, pi, parts):
    s = new_slide(prs)
    label(s, pal, _part("특강 및 기타 수업", pi, parts))
    header(s, sp.get("head","") if pi == 1 else "", size=26)
    n = len(items)
    gx, gy = PADX, 2.25
    # 5~6개는 3열 2행
    if n >= 5:
        cols, rows = 3, 2
        gapx, gapy = 0.28, 0.24
        cw = (11.9-(cols-1)*gapx)/cols
        ch = (7.05-gy-(rows-1)*gapy)/rows
        for i, it in enumerate(items):
            r = i//cols; c = i%cols
            in_row = min(cols, n - r*cols)
            off = ((cols-in_row)*(cw+gapx))/2
            x = gx + off + c*(cw+gapx); y = gy + r*(ch+gapy)
            rect(s, x, y, cw, ch, card_bg(pal), radius=True, line_hex=CARD_LINE, line_w=1)
            paras(s, x+0.26, y+0.28, cw-0.52, 0.52, [[
                (it.get("no","")+"  ", 16, True, pal["accent"]),
                (it["title"], 16, True, INK_STRONG)]], line_spacing=1.2)
            dsz = fit_box(it.get("desc",""), cw-0.52, ch-1.05, 13, 10, line_spacing=1.32)
            txt(s, x+0.26, y+0.88, cw-0.52, ch-1.05,
                [(it.get("desc",""), dsz, False, INK)], line_spacing=1.32)
        return s
    # 개수에 맞춰 배치: 1개=전폭 / 2개=좌우 / 3개=가로 3단 / 4개=사분면
    if n <= 3:
        gap = 0.34
        cw = (11.9 - (n-1)*gap)/n
        ch = 4.5
        for i, it in enumerate(items):
            x = gx + i*(cw+gap)
            rect(s, x, gy, cw, ch, card_bg(pal), radius=True, line_hex=CARD_LINE, line_w=1)
            paras(s, x+0.3, gy+0.35, cw-0.6, 0.6, [[
                (it.get("no","")+"  ", 20, True, pal["accent"]),
                (it["title"], 19, True, INK_STRONG)]], line_spacing=1.2)
            dsize = 14.5
            while dsize > 11.0 and wrap_lines(it.get("desc",""), cw-0.6, dsize)*(dsize*0.0139*1.62)*1.35 > (ch-1.5):
                dsize -= 0.5
            txt(s, x+0.3, gy+1.15, cw-0.6, ch-1.5, [(it.get("desc",""), dsize, False, INK)], line_spacing=1.35)
        return s
    cw, ch = 11.9/2, 2.15
    vline(s, gx+cw, gy, ch*2, LINE, 1)
    hline(s, gx, gy+ch, 11.9, LINE, 1)
    for i, it in enumerate(items):
        r=i//2; c=i%2
        x = gx + c*cw + 0.3; y = gy + r*ch + 0.3
        paras(s, x, y, cw-0.6, 0.5, [[
            (it.get("no","")+"  ", 20, True, pal["accent"]),
            (it["title"], 19, True, INK_STRONG)]])
        txt(s, x, y+0.65, cw-0.6, ch-1.0, [(it.get("desc",""), 14.5, False, INK)], line_spacing=1.35)
    return s

def slide_management(prs, pal, d):
    mg = d.get("management")
    if not mg or not mg.get("columns"): return None
    # ★과목이 3개·4개여도 버리지 않는다. 한 장에 2개씩 싣고 나머지는 다음 장으로.
    _cols = [c for c in mg["columns"] if c.get("rows")]
    # 과목 2개까지 한 장. 3개 이상이면 2개씩 나눠 설명 폭을 지킨다.
    pages = _chunk(_cols, 2)
    if not pages: return None
    last = None
    for pi, cols in enumerate(pages, 1):
        last = _management_page(prs, pal, mg, cols, pi, len(pages))
    return last

def _management_page(prs, pal, mg, cols, pi, parts):
    s = new_slide(prs)
    label(s, pal, _part("과목별 학생 관리", pi, parts))
    header(s, mg.get("head","") if pi == 1 else "", size=26)
    single = (len(cols) == 1)
    key_colors = [pal["accent"], pal["accent2"], pal["navy"], pal["accent"]]

    # 타입 선택: mg["style"] = "row" | "box" | "auto"
    # auto → 과목 1개면 박스(오른쪽 여백 방지), 2개면 줄 타입, 3개 이상이면 열이 좁아지므로 박스
    style = (mg.get("style") or "auto").lower()
    if style in ("auto", "box"):
        style = "row"      # 세원 양식으로 통일: 항목 pill 왼쪽 + 설명 오른쪽

    if style == "box":
        return _management_box(s, pal, cols, single, key_colors)
    return _management_row(s, pal, cols, single, key_colors)


def _management_row(s, pal, cols, single, key_colors):
    # 줄 타입: 좌 키박스 + 우 설명. 과목 1개면 전체폭 대신 왼쪽 정리(오른쪽 여백 축소)
    total = 11.9; gap = 0.7
    # ★과목이 하나면 전체폭 한 줄로 늘어나 오른쪽이 허전하고 줄 간격만 벌어졌다.
    #   항목이 4개 이상이면 두 칸으로 나눠, 2과목 학원과 같은 밀도로 보이게 한다.
    if single and len(cols) == 1 and len(cols[0].get("rows") or []) >= 4:
        _rows = list(cols[0]["rows"])
        _half = (len(_rows) + 1)//2
        cols = [{"name": cols[0].get("name",""), "rows": _rows[:_half]},
                {"name": "", "rows": _rows[_half:]}]
        single = False
    if single:
        colw = total        # 과목 1개면 전체폭을 써서 설명이 넉넉하게
        col_x = [PADX]
    else:
        colw = (total-gap)/2
        col_x = [PADX, PADX+colw+gap]
        vline(s, PADX+colw+gap/2, 2.15, 4.4, LINE, 1)
    max_rows = max(len(c["rows"]) for c in cols)
    # ★항목이 위에 몰리고 아래가 텅 비었다(2줄짜리 학원). 실제 줄 수로 간격을 잡고,
    #   남는 공간은 위아래로 나눠 슬라이드 상하 균형을 맞춘다.
    rows_top = 2.98; rows_bottom = 7.0
    _avail = rows_bottom - rows_top
    pitch = min(1.15, _avail/max(max_rows,1))
    _used = pitch * max_rows
    rows_top += max(0.0, (_avail - _used) / 2.0)
    for idx, col in enumerate(cols):
        x = col_x[idx]; kc = key_colors[idx]
        _cn = str(col.get("name") or "").strip()
        # ★슬라이드 라벨이 이미 "과목별 학생 관리"다. 과목이 하나뿐인 학원에서
        #   그 아래 또 "학습관리"라고 쓰면 같은 말이 두 번 나온다 → 소제목을 생략한다.
        #   (수학·과학처럼 과목명이 들어간 소제목은 구분에 필요하므로 남긴다)
        if re.fullmatch(r"학습\s*관리(\s*시스템)?", _cn or ""):
            _cn = ""
        if _cn:
            txt(s, x, 2.15, colw, 0.5, [(_cn, 19, True, kc)])
        hline(s, x, 2.72, colw, LINE, 1)
        kw = 1.15
        for j, row in enumerate(col["rows"]):
            yy = rows_top + j*pitch
            vv = str(row.get("v") or "")
            vsize = fit_size(vv, 14, 11, lambda sz: int(((colw-kw-0.2)/(sz*0.017))*2.4))
            # 본문은 위에서부터 흐르게 하고, 키 단추를 본문 첫 줄 높이에 맞춘다
            LINE_H = vsize * 0.0175 * 1.28          # 한 줄 높이(인치 근사)
            txt(s, x+kw+0.2, yy, colw-kw-0.2, pitch-0.05, [(vv, vsize, False, INK)],
                line_spacing=1.28, anchor=MSO_ANCHOR.TOP)
            pill_h = 0.44
            pill_y = yy + max(0.0, (LINE_H - pill_h)/2)   # 첫 줄 중앙과 단추 중앙 일치
            rect(s, x, pill_y, kw, pill_h, kc, radius=True)
            ktext = str(row.get("k") or "")
            ksize = 12.5 if len(ktext) <= 5 else (11 if len(ktext) <= 7 else 10)
            txt(s, x+0.05, pill_y+0.02, kw-0.1, pill_h-0.04, [(ktext, ksize, True, "FFFFFF")],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    return s


def _management_box(s, pal, cols, single, key_colors):
    # 박스 타입: 항목마다 카드. 과목 1개면 2열 그리드로 채워 오른쪽 여백 방지.
    area_top = 2.20 if len(cols) <= 2 else 2.02
    area_bottom = 7.05 if len(cols) <= 2 else 7.18
    if single:
        col = cols[0]; kc = key_colors[0]
        rows = col["rows"]
        # 제목
        txt(s, PADX, area_top, 11.9, 0.5, [(col["name"], 19, True, kc)])
        hline(s, PADX, area_top+0.55, 11.9, LINE, 1)
        grid_top = area_top + 0.80
        # 2열 그리드
        gcols = 2 if len(rows) >= 3 else 1
        gap = 0.5
        cw = (11.9 - gap*(gcols-1)) / gcols
        nrows = (len(rows)+gcols-1)//gcols
        ch = min(1.9, (area_bottom - grid_top - gap*(nrows-1)) / max(nrows,1))
        for i, row in enumerate(rows):
            r = i // gcols; c = i % gcols
            x = PADX + c*(cw+gap)
            y = grid_top + r*(ch+gap)
            _mgmt_card(s, pal, kc, pal.get("soft","F2F5F4"), x, y, cw, ch, row)
        return s
    # 2~4과목: 과목마다 한 열, 열 안에서 항목 카드 세로 배치
    n = len(cols)
    total = 11.9
    gap = 0.7 if n <= 2 else (0.45 if n == 3 else 0.34)
    colw = (total - gap*(n-1))/n
    col_x = [PADX + i*(colw+gap) for i in range(n)]
    soft_pool = [pal.get("soft","F2F5F4"), pal.get("soft2", pal.get("soft","F2F5F4"))]
    name_size = 19 if n <= 2 else (17 if n == 3 else 15)
    for idx, col in enumerate(cols):
        x = col_x[idx]; kc = key_colors[idx % len(key_colors)]
        txt(s, x, area_top, colw, 0.5, [(col["name"], name_size, True, kc)], line_spacing=1.15)
        hline(s, x, area_top+(0.55 if n <= 2 else 0.48), colw, LINE, 1)
        if idx < n-1 and n >= 3:
            vline(s, x+colw+gap/2, area_top, area_bottom-area_top, LINE, 1)
        rows = col["rows"]
        grid_top = area_top + (0.80 if n <= 2 else 0.62)
        vgap = 0.3 if n <= 2 else 0.14
        ch = min(1.15, (area_bottom - grid_top - vgap*(len(rows)-1)) / max(len(rows),1))
        for j, row in enumerate(rows):
            y = grid_top + j*(ch+vgap)
            _mgmt_card(s, pal, kc, soft_pool[idx % 2], x, y, colw, ch, row)
    return s


def _mgmt_card(s, pal, kc, soft_bg, x, y, w, h, row):
    # 카드 배경 없음: 소제목 pill(컬러)만 + 설명. (연한 박스 배경 제거)
    ktext = str(row.get("k") or "")
    ksize = 13 if len(ktext) <= 5 else (11.5 if len(ktext) <= 7 else 10.5)
    # 키 pill (열이 좁으면 pill과 글자를 함께 줄인다)
    kw = min(w-0.15, 0.42 + _disp_width(ktext)*0.20)
    while kw > w-0.15 and ksize > 9:
        ksize -= 0.5; kw = min(w-0.15, 0.42 + _disp_width(ktext)*0.18)
    ph = 0.34 if h < 0.95 else 0.40
    rect(s, x, y+0.06, kw, ph, kc, radius=True)
    txt(s, x, y+0.06, kw, ph, [(ktext, ksize, True, "FFFFFF")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    # 설명 — 칸 높이 안에 들어오도록 자동 축소(넘쳐서 아래 칸을 침범하지 않게)
    vv = str(row.get("v") or "")
    top_off = 0.44 if h < 0.95 else 0.54
    vsize = fit_box(vv, w-0.12, max(0.22, h-top_off), 13.5, 9.0, line_spacing=1.26)
    txt(s, x+0.02, y+top_off, w-0.10, max(0.22, h-top_off), [(vv, vsize, False, INK)],
        line_spacing=1.26, anchor=MSO_ANCHOR.TOP)

def slide_admission(prs, pal, d):
    ad = d.get("admission")
    if not ad or not ad.get("steps"): return None
    s = new_slide(prs)
    label(s, pal, "입학 절차")
    header(s, ad.get("head",""), size=26)
    steps = ad["steps"]
    n = len(steps)
    top = 2.20
    bottom = 7.10
    avail = bottom - top
    slot = avail / n
    has_bridge = bool(ad.get("bridge"))

    num_size = 58 if n <= 2 else (48 if n == 3 else 40)
    title_size = 19 if n <= 3 else 17

    text_x = PADX + 1.7
    text_w = 11.9 + PADX - text_x
    title_h = 0.40
    sep_gap = 0.14

    for i, st in enumerate(steps):
        y = top + i*slot
        hline(s, PADX, y, 11.9, LINE, 1)
        # 브릿지 — 구분선 우측 끝에 작은 캡션(왼쪽 설명과 좌우로 분리 → 겹칠 수 없음)
        if i > 0 and has_bridge:
            txt(s, 11.9+PADX-4.4, y-0.30, 4.4, 0.28,
                [("↓ ", 12, False, MUTED), (ad["bridge"], 11.5, True, pal["accent"])],
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)

        items = st.get("items", [])
        body = "  ".join(items) if items else ""
        # 설명 세로 영역: 슬롯에서 제목·여백 제외(브릿지는 다음 슬롯 상단 캡션이라 이 슬롯을 안 먹음)
        desc_top = y + sep_gap + title_h + 0.06
        desc_bottom = y + slot - 0.14
        desc_h = max(0.30, desc_bottom - desc_top)

        # 폰트를 desc_h 안에 들어오도록 자동 축소
        base_size = 14 if n <= 3 else 12.5
        min_size = 10.5
        line_sp = 1.30
        size = base_size
        while size > min_size:
            lines = wrap_lines(body, text_w, size)
            need = lines * (size*0.0139*1.55) * line_sp
            if need <= desc_h:
                break
            size -= 0.5

        # 제목 + 본문이 차지하는 실제 높이를 계산해, 숫자를 그 묶음의 세로 중앙에 맞춘다
        body_h = 0.0
        if body:
            lines = wrap_lines(body, text_w, size)
            body_h = lines * (size*0.0139*1.55) * line_sp
        block_top = y + sep_gap
        block_h = title_h + (0.06 + body_h if body else 0)
        block_mid = block_top + block_h/2

        txt(s, PADX, block_mid-0.62, 1.5, 1.24, [(st["no"], num_size, True, pal["accent"])],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 왼쪽에 큰 숫자가 있으므로 제목의 "1단계" 같은 중복 표기는 뗀다
        _t = re.sub(r"^\s*\d+\s*단계\s*", "", str(st.get("title","")))
        txt(s, text_x, y+sep_gap, text_w, 0.5, [(_t, title_size, True, INK_STRONG)])
        if body:
            txt(s, text_x, desc_top, text_w, desc_h,
                [("– ", size, False, INK), (body, size, False, INK)],
                line_spacing=line_sp, anchor=MSO_ANCHOR.TOP)
    return s

def slide_rules(prs, pal, d):
    ru = d.get("rules")
    if not ru: return None
    if not isinstance(ru, dict):
        ru = {"items": ru if isinstance(ru, list) else []}
    # ★항목이 None/문자열로 오면 .get 에서 죽었다 → dict 로 정규화한다.
    items = []
    for _it in (ru.get("items") or []):
        if isinstance(_it, dict):
            if str(_it.get("k") or _it.get("title") or "").strip() or str(_it.get("v") or "").strip():
                items.append(_it)
        elif isinstance(_it, str) and _it.strip():
            items.append({"k": _it.strip(), "v": ""})
    if not items: return None
    s = new_slide(prs)
    n = len(items)
    rx = 4.3; rw = 11.9+PADX-rx
    kw = 1.15
    vw = rw - kw - 0.42   # 설명 폭 (오른쪽 여백 확보)

    # 본문 영역(헤더 없는 슬라이드): 슬라이드 세로 정중앙(7.5in의 절반=3.75)을 기준으로 중앙 정렬
    AREA_TOP = 0.75
    AREA_BOTTOM = 7.4
    AREA_MID = 3.75   # 슬라이드(13.3×7.5) 세로 정중앙 — 항목이 몇 개든 이 지점을 블록 중앙으로

    # 각 항목 실제 높이(설명 줄 수 반영, 타이트하게)
    vsizes = []; text_h = []
    for it in items:
        vv = str(it.get("v") or "")
        vsize = fit_size(vv, 15, 12, lambda sz: int(((vw)/(sz*0.017))*2.0))
        lines = wrap_lines(vv, vw, vsize)
        th = lines * (vsize*0.0139*1.55) * 1.3      # 텍스트 자체 높이
        vsizes.append(vsize)
        text_h.append(max(0.34, th))

    # 항목 간 간격: 항목이 많으면 좁게, 적으면 넉넉하게(상하 중앙은 유지)
    content_h = sum(text_h)
    # 사용 가능한 세로에서 텍스트를 뺀 나머지를 (n+? ) 간격으로 배분하되, 상한/하한
    free = (AREA_BOTTOM - AREA_TOP) - content_h
    gap = free / (n)                     # 항목 사이+양끝을 고르게
    gap = max(0.14, min(0.9, gap))       # 너무 붙거나 너무 벌어지지 않게
    block = content_h + gap*(n-1)

    # 블록이 영역보다 크면(항목 매우 많음) 간격을 더 줄여 강제로 맞춤
    if block > (AREA_BOTTOM - AREA_TOP):
        gap = max(0.08, ((AREA_BOTTOM - AREA_TOP) - content_h) / max(n-1,1))
        block = content_h + gap*(n-1)

    top = AREA_MID - block/2
    if top < AREA_TOP: top = AREA_TOP
    mid = top + block/2

    # 좌 제목 — 블록 세로 중앙에 맞춤
    label(s, pal, "학원 규정", x=PADX, y=mid-1.02)
    # ★앱이 head를 빈 문자열로 보내면 <한마디를 넣지 않겠다>는 뜻이다.
    #   예전에는 or 를 써서 빈 값을 "없음"으로 보고 "학원 규정 안내"를 대신 찍었다.
    # ★"학원 규정"은 라벨만 쓴다. 한마디를 붙이면 규정처럼 읽혀 오해를 부른다.
    _rhead = ""
    if _rhead:
        _rsz = fit_box(_rhead, 3.35, 1.7, 31, 19, line_spacing=1.16)
        txt(s, PADX, mid-0.70, 3.35, 1.9,
            [(_rhead, _rsz, True, INK_STRONG)],
            line_spacing=1.16, anchor=MSO_ANCHOR.TOP)

    yy = top
    for i, it in enumerate(items):
        th = text_h[i]
        cyi = yy + th/2
        rect(s, rx, cyi-0.23, kw, 0.46, pal["pill_bg"], radius=True)
        txt(s, rx, cyi-0.23, kw, 0.46, [(str(it.get("k") or ""), 12.5, True, pal["pill_ink"])],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        txt(s, rx+kw+0.28, cyi-th/2-0.05, vw, th+0.1, [(str(it.get("v") or ""), vsizes[i], False, INK)],
            line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)
        yy += th + gap
    return s

def slide_faq(prs, pal, d):
    fq = d.get("faq")
    if not fq: return None
    if not isinstance(fq, dict):
        fq = {"items": fq if isinstance(fq, list) else []}
    # ★항목이 None/문자열로 오면 렌더에서 죽었다 → dict 로 정규화한다.
    _fi = []
    for _it in (fq.get("items") or []):
        if isinstance(_it, dict):
            if str(_it.get("q") or "").strip() or str(_it.get("a") or "").strip():
                _fi.append(_it)
        elif isinstance(_it, str) and _it.strip():
            _fi.append({"q": _it.strip(), "a": ""})
    if not _fi: return None
    fq = dict(fq); fq["items"] = _fi
    pages = _chunk(_fi, 6)
    last = None
    for pi, items in enumerate(pages, 1):
        last = _faq_page(prs, pal, fq, items, pi, len(pages))
    return last

def _faq_page(prs, pal, fq, _items, pi, parts):
    fq = dict(fq); fq["items"] = _items
    s = new_slide(prs)
    label(s, pal, _part("자주 묻는 질문", pi, parts))
    # ★"자주 묻는 질문"도 라벨만 쓴다.
    header(s, "", size=26)
    items = fq["items"][:6]
    gx, gy = PADX, 2.15
    gap = 0.3
    ncol = 2 if len(items) <= 4 else 3
    cw = (11.9-gap*(ncol-1))/ncol
    nrow = (len(items)+ncol-1)//ncol
    ch = min(2.30, (7.12-gy-0.22*(nrow-1))/max(nrow,1))
    for i, it in enumerate(items):
        r=i//ncol; c=i%ncol
        x = gx + c*(cw+gap); y = gy + r*(ch+0.22)
        rect(s, x, y, cw, ch, card_bg(pal), radius=True, line_hex=CARD_LINE, line_w=1)
        # Q
        q = it["q"]; a = it["a"]
        qsize = fit_box(q, cw-0.62, 0.74, 15, 11, line_spacing=1.25)
        paras(s, x+0.28, y+0.24, cw-0.56, 0.8,
              [[("Q. ", qsize, True, pal["accent2"]), (q, qsize, True, INK_STRONG)]],
              line_spacing=1.25, hang=0.32)
        asize = fit_box("A. " + a, cw-0.56, ch-1.10, 13.5, 9.5, line_spacing=1.32)
        paras(s, x+0.28, y+0.98, cw-0.56, ch-1.10,
              [[("A. ", asize, True, pal["accent"]), (a, asize, False, "4A5568")]],
              line_spacing=1.3, hang=0.30)
    return s

def slide_closing(prs, pal, d):
    if d.get("_noClosing"): return None      # 담당자가 목차에서 '상담 안내'를 끈 경우
    s = new_slide(prs)
    cl = d.get("closing", {})
    a = d["academy"]
    label(s, pal, "지도 · 오시는 길")
    # 우 지도
    mx = 6.9
    mw, mh = 11.9+PADX-mx, 4.35
    rect(s, mx, 0.9, mw, mh, PH_BG, radius=True)
    _map = (d.get("assets") or {}).get("map") or ""
    if not (_map and place_image(s, _map, mx, 0.9, mw, mh, cover=False)):
        txt(s, mx, 2.85, mw, 0.4, [("지도 · 오시는 길", 12, False, "8CA0BB")], align=PP_ALIGN.CENTER)
    # 좌 카피 (쉼표까지 한 줄 / 다음 줄)
    head = cl.get("head","")
    lines = head.split("\n")
    hl = cl.get("highlight","")
    para_runs = []
    # 카피 글자 크기를 좌측 폭(5.9in) 안에 들어오도록 맞춘다 — 그대로 두면 페이지를 넘친다
    HW = 5.9
    hsz = 30.0
    while hsz > 18.0 and any(_disp_width(ln)*hsz*0.0139 > HW for ln in lines if ln.strip()):
        hsz -= 1.0
    for ln in lines:
        runs=[]
        if hl and hl in ln:
            before, after = ln.split(hl,1)
            if before: runs.append((before, hsz, True, INK_STRONG))
            runs.append((hl, hsz, True, pal["accent"]))
            runs.append((after, hsz, True, INK_STRONG))
        else:
            runs.append((ln, hsz, True, INK_STRONG))
        para_runs.append(runs)
    _hh = max(1.6, len([l for l in lines if l.strip()]) * hsz*0.0139*1.62*1.32 + 0.2)
    paras(s, PADX, 1.0, HW, _hh, para_runs, line_spacing=1.32, para_space=2, wrap=True)
    # 연락처(각 한 줄) — 값 있는 것만
    contact_paras = []
    _rows = []
    if a.get("phone"):  _rows.append(("상담 문의", a.get("phone","")))
    if a.get("phone2"): _rows.append(("학원 전화", a.get("phone2","")))
    if a.get("kakao"):  _rows.append(("카카오 채널", a.get("kakao","")))
    if a.get("address_short") or a.get("location"):
        _rows.append(("위치", a.get("address_short","") or a.get("location","")))
    if a.get("hours"):  _rows.append(("운영시간", a.get("hours","")))
    csize = 15 if len(_rows) <= 3 else 13
    for k, v in _rows:
        contact_paras.append([(k + "    ", csize, True, pal["navy"]), (v, csize, False, INK)])
    if contact_paras:
        paras(s, PADX, 2.55, 6.0, 2.05, contact_paras, line_spacing=1.28,
              para_space=(8 if len(_rows) <= 3 else 4))
    # ── QR: 자료에 주소가 있는 채널만 만든다(없으면 아예 그리지 않는다) ──
    links = a.get("links") or []
    if not links:
        _auto = []
        for k, lab in (("blog","네이버 블로그"), ("homepage","홈페이지"),
                       ("instagram","인스타그램"), ("kakao","카카오 채널")):
            if a.get(k): _auto.append({"label": lab, "url": a[k]})
        links = _auto
    links = [l for l in links if str(l.get("url","")).strip()][:3]
    for i, l in enumerate(links):
        qx = PADX + i*1.5
        rect(s, qx, 4.7, 0.95, 0.95, "FFFFFF", radius=True, line_hex=LINE, line_w=1)
        buf = _qr_png(l["url"])
        if buf is not None:
            try:
                s.shapes.add_picture(buf, Inches(qx+0.06), Inches(4.76), Inches(0.83), Inches(0.83))
            except Exception:
                buf = None
        if buf is None:
            txt(s, qx, 5.05, 0.95, 0.3, [("QR", 11, False, "9FB0C6")], align=PP_ALIGN.CENTER)
        txt(s, qx-0.15, 5.72, 1.25, 0.3, [(str(l.get("label","")), 11, False, "667788")],
            align=PP_ALIGN.CENTER)
    # CTA
    if cl.get("cta"):
        rect(s, PADX, 6.52, 11.9, 0.66, pal["navy"], radius=True)
        txt(s, PADX, 6.57, 11.9, 0.56, [(cl["cta"], 20, True, "FFFFFF")],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


# ======================================================================
#  메인 빌드
# ======================================================================
MAX_ROWS_PER_TT = 13  # 시간표 한 장 최대 행수

def build(data, palette="teal_blue", out="brandbook.pptx"):
    pal = PALETTES.get(palette, PALETTES["teal_blue"])
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # ── 브랜드북 목차 ──
    # 표지 / ①학원 소개 / ②주요 실적 / ③수업 강점 / ④수업 대상·반 구성·커리큘럼 /
    # ⑤수업 방식·수업 관리 / ⑥시간표 / ⑦특강 및 기타 수업 / ⑧입학 절차 / ⑨학원 규정 /
    # ⑩자주 묻는 질문 / ⑪상담 안내
    slide_cover(prs, pal, data)
    slide_intro(prs, pal, data)          # ① 학원 소개 (슬로건·소개문·강점 한 장)
    slide_achievements(prs, pal, data)   # ② 주요 실적 (유형별 카드 + 실적 사진)
    slide_targets(prs, pal, data)        # ③ 수업 대상 · 커리큘럼 (한 장)
    slide_management(prs, pal, data)     # ④ 수업 관리 (과목별)

    # ⑥ 시간표: 그룹 수만큼, 행 많으면 자동 분할
    for tt in data.get("timetables", []):
        group = tt.get("group","")
        rows = tt.get("rows", [])
        if not rows: continue
        if len(rows) <= MAX_ROWS_PER_TT:
            slide_timetable(prs, pal, data, group, rows)
        else:
            chunks = [rows[i:i+MAX_ROWS_PER_TT] for i in range(0,len(rows),MAX_ROWS_PER_TT)]
            for pi, ch in enumerate(chunks, 1):
                slide_timetable(prs, pal, data, group, ch, part=pi, parts=len(chunks))

    slide_specials(prs, pal, data)       # ⑦ 특강 및 기타 수업
    slide_admission(prs, pal, data)      # ⑧ 입학 절차
    slide_rules(prs, pal, data)          # ⑨ 학원 규정 (출결·보강·환불)
    slide_faq(prs, pal, data)            # ⑩ 자주 묻는 질문
    slide_closing(prs, pal, data)        # ⑪ 상담 안내 (위치·연락처)

    _force_theme_font(prs)
    prs.save(out)   # out은 경로(str) 또는 file-like(BytesIO) 모두 가능
    return out


if __name__ == "__main__":
    import json, sys
    src = sys.argv[1] if len(sys.argv)>1 else "sewon.json"
    pal = sys.argv[2] if len(sys.argv)>2 else "teal_blue"
    out = sys.argv[3] if len(sys.argv)>3 else "brandbook.pptx"
    data = json.load(open(src, encoding="utf-8"))
    build(data, pal, out)
    print("saved", out)
