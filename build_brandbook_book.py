# -*- coding: utf-8 -*-
"""
브랜드북 빌더 — book 스킨 (책 펼침 디자인)
- 참고: 더바른 수학학원 초등부 16장
- v3 빌더(build_brandbook_v3)의 저수준 헬퍼/팔레트/시간표 로직을 재사용하고,
  프레임과 16장 표현방식만 '책 펼침' 스타일로 새로 그린다.
- v3 표준 스키마를 그대로 입력받음(표현만 다른 렌더러). build(schema, palette, out).
- 슬라이드 13.3 x 7.5in, 서체 Pretendard.

공통 프레임(모든 슬라이드):
  상단 책곡선 2개(하늘색 얇은선) + 중앙 세로 접힘선 + 하단 이중선
  좌측: 라벨(16px 700) + 대형헤더(58px 800) + 하단 요약(17px 700)
  우측: 콘텐츠 영역
  푸터: 좌(학원명) / 우(슬로건 or 날짜)
배경 흰색 #f4f5f7.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import math

# v3 헬퍼 재사용 (rect/txt/paras/bullets/place_image/_strip_style/시간표 등)
import build_brandbook_v3 as V3
from build_brandbook_v3 import (
    rect, txt, paras, bullets, hline, vline, place_image,
    _rgb, _strip_style, _fill, _no_line, wrap_lines, fit_size,
)

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Pretendard"
EMU_IN = 914400

# ---------------- book 팔레트 (v3 3종에 매핑) ----------------
# teal_blue: 청록 주 + 블루 보조 (원형 호는 이 두 색만)
BOOK_PALETTES = {
    "teal_blue":   {"accent": "0EA5B7", "accent2": "2B7FFF", "arc": "8ECAE6"},
    "navy_amber":  {"accent": "22375F", "accent2": "D98A1F", "arc": "9DB4D8"},
    "green_orange":{"accent": "1F7A55", "accent2": "E07B39", "arc": "8FC9A9"},
}

PAGE_BG   = "F4F5F7"
INK       = "2D3748"    # 본문/요약
HEAD_INK  = "333A44"    # 대형 헤더
LABEL_INK = "5B6B85"    # 라벨
CARD      = "FFFFFF"
CARD_LINE = "E4E8EE"
MUTED     = "9AA7BD"
FOLD      = "DDE3EA"    # 중앙 접힘선
DBL       = "CBD3DE"    # 하단 이중선
PH_BG     = "DBE4EF"

# 레이아웃 좌표(인치)
MX = 0.72              # 좌우 여백
TOP_ARC_Y = 0.30
FOLD_X = SLIDE_W/2
LEFT_W = FOLD_X - MX - 0.35   # 좌측 컬럼 폭 (중앙 접힘선을 넘지 않게)
RIGHT_X = FOLD_X + 0.35
RIGHT_W = SLIDE_W - RIGHT_X - MX
BODY_TOP = 1.15
BODY_BOT = 6.55        # 푸터 위 한계


# ======================================================================
#  공통 프레임
# ======================================================================
def _bg(slide, hexc=PAGE_BG):
    bg = slide.background; bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(hexc)

def _new(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    return s

def _curve(slide, cx, cy, w, h, hexc, weight=1.4):
    """상단 책 곡선: ARC 도형(얇은 선). w,h는 바운딩 박스."""
    s = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(cx), Inches(cy), Inches(w), Inches(h))
    _strip_style(s)
    s.fill.background()
    s.line.color.rgb = _rgb(hexc); s.line.width = Pt(weight)
    s.shadow.inherit = False
    return s

def _frame(slide, pal, d, footer_right=None):
    ac = d.get("academy", {})
    arc = pal["arc"]
    # 상단 책곡선 2개 (완만한 아치 두 겹)
    _curve(slide, 1.1, -0.55, SLIDE_W-2.2, 1.5, arc, 1.4)
    _curve(slide, 1.7, -0.35, SLIDE_W-3.4, 1.2, arc, 1.1)
    # 중앙 세로 접힘선
    vline(slide, FOLD_X, TOP_ARC_Y+0.55, BODY_BOT-0.55, FOLD, weight=1.2)
    # 하단 이중선
    hline(slide, MX, 6.78, SLIDE_W-2*MX, DBL, weight=1.1)
    hline(slide, MX, 6.86, SLIDE_W-2*MX, DBL, weight=1.1)
    # 푸터
    txt(slide, MX, 6.95, 6.0, 0.4, [(ac.get("name",""), 11, True, MUTED)],
        anchor=MSO_ANCHOR.MIDDLE)
    if footer_right is not None:
        fr = footer_right
    else:
        # 슬로건이 짧은 한 줄일 때만 푸터에. 길면(포지셔닝 문장) 비움.
        slg = (ac.get("slogan","") or "").strip()
        fr = slg if 0 < len(slg) <= 22 else ""
    if fr:
        txt(slide, SLIDE_W-6.72, 6.95, 6.0, 0.4, [(fr, 11, True, MUTED)],
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def _left(slide, pal, label_text, head, summary):
    """좌측 컬럼: 라벨 + 대형헤더(58px≈29pt) + 하단 요약."""
    ac = pal["accent"]
    # 라벨 + 짧은 세로바
    vline(slide, MX+0.02, BODY_TOP+0.04, 0.26, ac, weight=3.4)
    txt(slide, MX+0.16, BODY_TOP, LEFT_W-0.2, 0.34,
        [(label_text, 13, True, LABEL_INK)])
    # 대형 헤더 (자동 축소: 글자수 기준)
    hlen = max(len(ln) for ln in head.split("\n")) if head else 0
    hsize = 30 if hlen <= 8 else (26 if hlen <= 11 else (22 if hlen <= 15 else 19))
    txt(slide, MX+0.16, BODY_TOP+0.55, LEFT_W-0.16, 2.6,
        [(head, hsize, True, HEAD_INK)], line_spacing=1.12)
    # 하단 요약
    if summary:
        txt(slide, MX+0.16, 5.45, LEFT_W-0.16, 1.15,
            [(summary, 12, True, INK)], line_spacing=1.28)


def _circle_icon(slide, cx, cy, r, fill_hex, glyph="", glyph_color="FFFFFF", glyph_size=18):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
    _strip_style(s); _fill(s, fill_hex); s.shadow.inherit = False
    if glyph:
        txt(slide, cx-r, cy-r-0.02, 2*r, 2*r, [(glyph, glyph_size, True, glyph_color)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s

def _ring(slide, cx, cy, r, line_hex, weight=2.2):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
    _strip_style(s); s.fill.background()
    s.line.color.rgb = _rgb(line_hex); s.line.width = Pt(weight)
    s.shadow.inherit = False
    return s

def _card(slide, x, y, w, h, fill=CARD, line=CARD_LINE, radius=True):
    s = rect(slide, x, y, w, h, fill, radius=radius, line_hex=line, line_w=1.0)
    s.shadow.inherit = False
    return s


def _short_slogan(text, limit=30):
    """표지용 슬로건: 한 줄 카피가 아니면(문장이 길면) 첫 구절만 취함.
    빈 문자열이면 호출부가 학원명으로 폴백."""
    t = (text or "").strip().replace("\n", " ")
    if not t:
        return ""
    if len(t) <= limit:
        return t
    # 구두점 경계 우선(문장 단위). 너무 짧게 잘리는 공백 분리는 후순위.
    for sep in [". ", ".", "!", "?", " — ", " - ", ",", ";"]:
        if sep in t:
            head = t.split(sep)[0].strip()
            if 8 <= len(head) <= limit:
                return head
    # 공백 분리: limit 근처까지 단어를 모아 자연스럽게 자름
    words = t.split(" ")
    acc = ""
    for w in words:
        if len(acc) + len(w) + 1 > limit:
            break
        acc = (acc + " " + w).strip()
    if len(acc) >= 8:
        return acc
    return t[:limit].rstrip() + "…"


def _fmt_phone(p):
    """숫자만 있는 전화번호에 하이픈. 형식 불명이면 원본 유지."""
    d = "".join(ch for ch in (p or "") if ch.isdigit())
    if len(d) == 11 and d.startswith("010"):
        return f"{d[:3]}-{d[3:7]}-{d[7:]}"
    if len(d) == 10:
        return f"{d[:2]}-{d[2:6]}-{d[6:]}" if d.startswith("02") else f"{d[:3]}-{d[3:6]}-{d[6:]}"
    if len(d) == 9 and d.startswith("02"):
        return f"{d[:2]}-{d[2:5]}-{d[5:]}"
    return p or ""


def _clean_desc(title, desc):
    """강점/항목 설명이 제목과 같으면(원본에 desc가 비어 title 복사된 경우) 빈 값 처리."""
    t = (title or "").strip()
    dv = (desc or "").strip()
    if not dv or dv == t:
        return ""
    return dv


def _real_items(items):
    """내용 있는 항목만 남김 (빈 칸 방지 공통 필터)."""
    out = []
    for it in (items or []):
        if not isinstance(it, dict):
            continue
        if any((str(it.get(k,"")).strip()) for k in ("title","desc","grade","subj","no","q","a","k","v","name")):
            out.append(it)
    return out


def _trim_sentence(text, limit):
    """limit 근처에서 문장 경계로 자연스럽게 자름. 조각 단어가 남지 않게."""
    t = (text or "").strip()
    if len(t) <= limit:
        return t
    cut = t[:limit]
    # 1순위: 문장 끝(마침표 등) — limit의 45% 이후면 채택
    for sep in [". ", ".", "! ", "!", "? ", "?"]:
        p = cut.rfind(sep)
        if p >= limit * 0.45:
            return t[:p+1].strip()
    # 2순위: 어미(습니다/입니다/합니다 등) 뒤
    for end in ["습니다", "입니다", "합니다", "됩니다", "니다"]:
        p = cut.rfind(end)
        if p >= limit * 0.45:
            return t[:p+len(end)].strip() + "."
    # 3순위: 쉼표
    p = cut.rfind(",")
    if p >= limit * 0.5:
        return t[:p].strip() + "…"
    # 4순위: 공백(조각 단어 방지 위해 마지막 어절 통째 제거)
    p = cut.rfind(" ")
    if p >= limit * 0.5:
        return t[:p].strip() + "…"
    return cut.rstrip() + "…"


# ======================================================================
#  16장 표현방식 (book 스킨)
# ======================================================================

# 1) 표지 = 4모서리 아이콘 + 중앙 서브/타이틀/배지/학원명
def s01_cover(prs, pal, d):
    s = _new(prs)
    ac, ac2, arc = pal["accent"], pal["accent2"], pal["arc"]
    _bg(s, PAGE_BG)
    # 상단 곡선
    _curve(s, 1.1, -0.6, SLIDE_W-2.2, 1.6, arc, 1.4)
    _curve(s, 1.7, -0.4, SLIDE_W-3.4, 1.3, arc, 1.1)
    academy = d.get("academy", {})
    # 4모서리 아이콘 (작은 원)
    corners = [(1.5,1.5,ac),(SLIDE_W-1.5,1.5,ac2),(1.5,6.0,ac2),(SLIDE_W-1.5,6.0,ac)]
    for cx,cy,c in corners:
        _ring(s, cx, cy, 0.34, c, 2.0)
        _circle_icon(s, cx, cy, 0.10, c)
    # 중앙 배치
    sub = academy.get("subjects","") or "브랜드북"
    txt(s, 0, 2.55, SLIDE_W, 0.5, [(sub, 16, True, ac)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 표지 제목: 슬로건이 길면(한 줄 카피가 아니면) 첫 구절만, 그래도 길면 학원명으로 폴백
    slg = (academy.get("slogan","") or "").strip()
    title = _short_slogan(slg) or academy.get("name","")
    tlen = len(title)
    tsize = 40 if tlen<=12 else (34 if tlen<=18 else (27 if tlen<=24 else 22))
    txt(s, 0.9, 2.95, SLIDE_W-1.8, 1.5, [(title, tsize, True, HEAD_INK)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)
    # 배지
    bw = 2.2
    b = rect(s, (SLIDE_W-bw)/2, 4.45, bw, 0.5, ac, radius=True); b.shadow.inherit=False
    txt(s, (SLIDE_W-bw)/2, 4.48, bw, 0.44, [("BRAND BOOK", 12, True, "FFFFFF")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 학원명 (로고 있으면 로고 위에 학원명)
    _cvr_assets = d.get("assets", {}) or {}
    _logo = _cvr_assets.get("logo", "")
    if _logo and _place_img(s, _logo, (SLIDE_W-1.7)/2, 5.05, 1.7, 0.72, cover=False):
        txt(s, 0, 5.80, SLIDE_W, 0.5, [(academy.get("name",""), 17, True, INK)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        txt(s, 0, 5.15, SLIDE_W, 0.6, [(academy.get("name",""), 22, True, INK)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 하단 이중선 + 연락처
    hline(s, MX, 6.78, SLIDE_W-2*MX, DBL, 1.1)
    hline(s, MX, 6.86, SLIDE_W-2*MX, DBL, 1.1)
    loc = (academy.get("location","") or academy.get("address_short","") or "").strip()
    if len(loc) > 40:
        loc = loc[:40].rstrip() + "…"
    info = " · ".join([x for x in [loc, _fmt_phone(academy.get("phone",""))] if x])
    isz = 11 if len(info) <= 52 else (10 if len(info) <= 62 else 9)
    txt(s, MX, 6.95, SLIDE_W-2*MX, 0.4, [(info, isz, True, MUTED)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


# 2) SWOT 사분면 (원 + 2x2) — 학원 소개/강점 요약
def s02_intro(prs, pal, d):
    s = _new(prs)
    intro = d.get("intro", {}); feats = d.get("features", [])[:4]
    _frame(s, pal, d)
    # 소개 body가 비면(원본만 있고 인터뷰 전) 강점 제목을 한 줄 요약으로 대신
    body = (intro.get("body","") or "").strip()
    if not body and feats:
        titles = [f.get("title","") for f in feats if f.get("title","")]
        if titles:
            body = " · ".join(titles[:3]) + " 을 지향합니다"
    _left(s, pal, "학원 소개", intro.get("head","우리 학원을\n소개합니다"),
          _trim_sentence(body, 88))
    # 우측: 중앙 원 + 2x2 개념
    ac, ac2 = pal["accent"], pal["accent2"]
    cx = RIGHT_X + RIGHT_W/2; cy = 3.6
    _ring(s, cx, cy, 0.62, ac, 2.4)
    txt(s, cx-0.62, cy-0.62, 1.24, 1.24, [("핵심\n가치", 14, True, ac)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    # 2x2 카드: 강점을 쓰되, 4개 미만이면 기본 가치로 보충(빈 카드 금지)
    DEFAULTS = [{"title":"개념","desc":""},{"title":"이해","desc":""},
                {"title":"습관","desc":""},{"title":"성취","desc":""}]
    labels = _real_items(feats)[:4]
    for dft in DEFAULTS:
        if len(labels) >= 4: break
        if dft["title"] not in [l.get("title","") for l in labels]:
            labels.append(dft)
    cw, ch = RIGHT_W/2-0.25, 1.75
    pos = [(RIGHT_X, BODY_TOP+0.15),(RIGHT_X+cw+0.5, BODY_TOP+0.15),
           (RIGHT_X, BODY_TOP+0.15+ch+0.4),(RIGHT_X+cw+0.5, BODY_TOP+0.15+ch+0.4)]
    cols = [ac, ac2, ac2, ac]
    for i,(x,y) in enumerate(pos):
        _card(s, x, y, cw, ch)
        c = cols[i]
        rect(s, x, y, 0.10, ch, c, radius=False)
        ttl = labels[i].get("title","")
        txt(s, x+0.28, y+0.18, cw-0.4, 0.5, [(ttl, 16, True, HEAD_INK)],
            anchor=(MSO_ANCHOR.MIDDLE if not _clean_desc(ttl, labels[i].get("desc","")) else MSO_ANCHOR.TOP))
        dv = _clean_desc(ttl, labels[i].get("desc",""))
        if dv:
            txt(s, x+0.28, y+0.72, cw-0.5, ch-0.85, [(dv[:60], 11.5, False, INK)],
                line_spacing=1.22)
    return s


# 3) 원형 연결 4개 (사람 아이콘 + 중앙 플러스 + 하단 캡션) — 네 가지 힘 / 강점
def s03_four(prs, pal, d):
    feats = _real_items(d.get("features", []))[:4]
    n = len(feats)
    if n == 0:
        return None
    s = _new(prs)
    KO = {1:"한 가지 힘",2:"두 가지 힘",3:"세 가지 힘",4:"네 가지 힘"}
    _frame(s, pal, d)
    _left(s, pal, "우리의 강점", f"학원을 지탱하는\n{KO.get(n,'핵심 힘')}",
          "학생의 성장을 이끄는 핵심 강점을 소개합니다")
    ac, ac2 = pal["accent"], pal["accent2"]
    cx0 = RIGHT_X + RIGHT_W/2
    cy0 = 3.55
    r = 0.72
    off = 1.55
    # 개수별 원 위치
    if n == 1:
        centers = [(cx0, cy0)]
    elif n == 2:
        centers = [(cx0-off, cy0),(cx0+off, cy0)]
    elif n == 3:
        centers = [(cx0-off, cy0-1.0),(cx0+off, cy0-1.0),(cx0, cy0+1.1)]
    else:
        centers = [(cx0-off, cy0-1.05),(cx0+off, cy0-1.05),(cx0-off, cy0+1.05),(cx0+off, cy0+1.05)]
    cols = [ac, ac2, ac2, ac]
    if n == 4:
        _circle_icon(s, cx0, cy0, 0.30, HEAD_INK, "+", "FFFFFF", 22)
    for i,(cx,cy) in enumerate(centers):
        _ring(s, cx, cy, r, cols[i], 2.4)
        _circle_icon(s, cx, cy-0.12, 0.13, cols[i])
        rect(s, cx-0.24, cy+0.04, 0.48, 0.26, cols[i], radius=True)
        txt(s, cx-r-0.15, cy+r+0.02, 2*r+0.3, 0.4, [(feats[i].get("title",""), 13, True, HEAD_INK)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


# 원형 아이콘 3단 (공용) — 정도와 정성 / 자기주도 / 교재학습도구
def _place_img(slide, url, x, y, w, h, cover=True):
    """place_image 확장판: base64 data URL(담당자가 PC에서 올린 이미지)도 지원."""
    if not url:
        return False
    u = str(url)
    if u.startswith("data:"):
        try:
            import base64 as _b64, io as _io2
            from pptx.util import Inches as _In
            head, _, b64 = u.partition(",")
            raw = _b64.b64decode(b64)
            try:
                from PIL import Image as _PIL
                im = _PIL.open(_io2.BytesIO(raw))
                iw, ih = im.size
                box_r = w/h; img_r = iw/ih
                pic = slide.shapes.add_picture(_io2.BytesIO(raw), _In(x), _In(y), _In(w), _In(h))
                if cover:
                    if img_r > box_r:
                        c = (1 - box_r/img_r)/2
                        pic.crop_left = c; pic.crop_right = c
                    elif img_r < box_r:
                        c = (1 - img_r/box_r)/2
                        pic.crop_top = c; pic.crop_bottom = c
                else:
                    # 비율 유지해 박스 안에 전체가 들어가게 맞춤(레터박스)
                    box_r = w/h
                    if img_r > box_r:      # 가로가 긴 이미지 → 폭 기준
                        nw = w; nh = w/img_r
                    else:                  # 세로가 긴 이미지 → 높이 기준
                        nh = h; nw = h*img_r
                    pic.left = _In(x + (w-nw)/2); pic.top = _In(y + (h-nh)/2)
                    pic.width = _In(nw); pic.height = _In(nh)
                return True
            except Exception:
                slide.shapes.add_picture(_io2.BytesIO(raw), _In(x), _In(y), _In(w), _In(h))
                return True
        except Exception:
            return False
    return place_image(slide, u, x, y, w, h, cover=cover)


def s_achieve_images(prs, pal, d, label_text, head, summary, urls, items=None):
    """실적 자료 이미지 슬라이드. 담당자가 올린 입결·성적·수상 이미지를 크게 배치.
    1장=크게 한 장, 2장=좌우, 3~4장=2x2. 이미지 옆/아래에 실적 텍스트 요약(있으면)."""
    urls = [u for u in (urls or []) if u]
    if not urls:
        return None
    s = _new(prs)
    _frame(s, pal, d)
    _left(s, pal, label_text, head, summary)
    n = min(len(urls), 4)
    x0, y0 = RIGHT_X, 1.45
    W, H = RIGHT_W, 4.7
    placed = 0
    # 실적 배너는 잘리면 정보가 사라지므로 cover=False(비율 유지, 전체 표시)
    if n == 1:
        placed += 1 if _place_img(s, urls[0], x0, y0, W, H, cover=False) else 0
    elif n == 2:
        ch2 = (H - 0.3)/2
        for i in range(2):
            if _place_img(s, urls[i], x0, y0 + i*(ch2+0.3), W, ch2, cover=False):
                placed += 1
    else:
        cw = (W - 0.3)/2
        ch = (H - 0.3)/2
        pos = [(x0, y0), (x0+cw+0.3, y0), (x0, y0+ch+0.3), (x0+cw+0.3, y0+ch+0.3)]
        for i in range(n):
            if _place_img(s, urls[i], pos[i][0], pos[i][1], cw, ch, cover=False):
                placed += 1
    if placed == 0:
        # 이미지 로드 실패 시 텍스트 슬라이드로 대체되도록 None 처리
        try:
            xml = prs.slides._sldIdLst[-1]
            prs.slides._sldIdLst.remove(xml)
        except Exception:
            pass
        return None
    return s


def _three_circles(prs, pal, d, label_text, head, summary, items, footer=None):
    # 실제 내용 있는 항목만 (빈 원 방지)
    items = [it for it in (items or []) if (it.get("title","").strip() or it.get("desc","").strip())]
    n = len(items)
    if n == 0:
        return None
    s = _new(prs)
    _frame(s, pal, d, footer_right=footer)
    _left(s, pal, label_text, head, summary)
    ac, ac2 = pal["accent"], pal["accent2"]
    n = min(n, 3)
    items = items[:n]
    gap = 0.6
    # 항목 수에 맞춰 카드폭: 1개면 넓게, 3개면 3등분
    if n == 1:
        cw = min(4.5, RIGHT_W)
        xs = [RIGHT_X + (RIGHT_W-cw)/2]
    else:
        cw = (RIGHT_W - (n-1)*gap)/n
        xs = [RIGHT_X + i*(cw+gap) for i in range(n)]
    y0 = 1.5
    r = 0.62 if n >= 2 else 0.7
    for i in range(n):
        x = xs[i]
        cxx = x + cw/2
        col = ac if i % 2 == 0 else ac2
        _circle_icon(s, cxx, y0+0.75, r, col, str(i+1), "FFFFFF", 24)
        raw_ttl = (items[i].get("title","") or "").strip()
        raw_dv  = _clean_desc(raw_ttl, items[i].get("desc",""))
        # 제목이 문장처럼 길면(라벨이 아니면) 본문으로 내리고 제목은 비움
        if len(raw_ttl) > 14 and not raw_dv:
            ttl3, dv = "", raw_ttl
        else:
            ttl3, dv = raw_ttl, raw_dv
        if ttl3:
            tsz = 15 if len(ttl3) <= 12 else 13
            txt(s, x-0.15, y0+1.6, cw+0.3, 0.7, [(ttl3, tsz, True, HEAD_INK)],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.1)
        if dv:
            cap = 120 if n == 1 else (90 if n == 2 else 70)
            dsz = 12.5 if n == 1 else 11.5
            dy = y0+2.4 if ttl3 else y0+1.7
            txt(s, x, dy, cw, 2.4, [(dv[:cap], dsz, False, INK)],
                align=PP_ALIGN.CENTER, line_spacing=1.26)
    return s


# 5) 번호원 (01~03) + 물결선 — 맞춤 수업 설계 / 커리큘럼 단계
def s_numbered(prs, pal, d, label_text, head, summary, items):
    items = _real_items(items)[:3]
    n = len(items)
    if n == 0:
        return None
    s = _new(prs)
    _frame(s, pal, d)
    _left(s, pal, label_text, head, summary)
    ac, ac2 = pal["accent"], pal["accent2"]
    y0 = 1.45
    rowh = min(1.7, (BODY_BOT-y0)/n)
    for i in range(n):
        y = y0 + i*rowh
        col = ac if i%2==0 else ac2
        _circle_icon(s, RIGHT_X+0.55, y+0.55, 0.5, col, f"{i+1:02d}", "FFFFFF", 18)
        if i < n-1:  # 물결/연결선
            vline(s, RIGHT_X+0.55, y+1.05, rowh-1.0, col, weight=2.0)
        txt(s, RIGHT_X+1.35, y+0.12, RIGHT_W-1.4, 0.5,
            [(items[i].get("title",""), 16, True, HEAD_INK)])
        dv = items[i].get("desc","")
        if dv:
            txt(s, RIGHT_X+1.35, y+0.62, RIGHT_W-1.4, 0.8, [(dv[:70], 12, False, INK)],
                line_spacing=1.24)
    return s


# 6) 회색바 리스트 (번호칩 + 제목 + 설명) — 네 가지 반 / 수업 대상
def s_barlist(prs, pal, d, label_text, head, summary, items):
    items = _real_items(items)[:4]
    if not items:
        return None
    n = len(items)
    s = _new(prs)
    _frame(s, pal, d)
    _left(s, pal, label_text, head, summary)
    ac = pal["accent"]
    y0 = 1.4
    avail = BODY_BOT - y0 - 0.1
    rowh = min(1.15, avail/n)
    for i,it in enumerate(items):
        y = y0 + i*rowh
        _card(s, RIGHT_X, y, RIGHT_W, rowh-0.18, CARD)
        # 번호칩
        chip = rect(s, RIGHT_X+0.18, y+0.18, 0.5, rowh-0.54, "333A44", radius=True)
        chip.shadow.inherit=False
        txt(s, RIGHT_X+0.18, y+0.18, 0.5, rowh-0.54, [(f"{i+1:02d}", 13, True, "FFFFFF")],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ttl = it.get("title","") or it.get("grade","")
        txt(s, RIGHT_X+0.9, y+0.14, RIGHT_W-1.0, 0.5, [(ttl, 15, True, HEAD_INK)],
            anchor=MSO_ANCHOR.MIDDLE)
        dv = it.get("desc","") or it.get("subj","")
        if dv:
            txt(s, RIGHT_X+0.9, y+0.56, RIGHT_W-1.1, rowh-0.7, [(dv[:60], 11.5, False, INK)],
                line_spacing=1.2)
    return s


# 7) 어두운 카드 (라운드 + 4단계) — 깊이 있는 학습
def s_darkcard(prs, pal, d, label_text, head, summary, items):
    s = _new(prs)
    _frame(s, pal, d)
    _left(s, pal, label_text, head, summary)
    ac, ac2 = pal["accent"], pal["accent2"]
    # 큰 어두운 카드
    x, y, w, h = RIGHT_X, 1.35, RIGHT_W, 4.9
    dc = rect(s, x, y, w, h, "3A3F47", radius=True); dc.shadow.inherit=False
    items = _real_items(items)[:4]
    n = max(1, len(items))
    ph = (h-0.6)/n
    for i,it in enumerate(items):
        yy = y+0.3 + i*ph
        col = ac if i%2==0 else ac2
        _circle_icon(s, x+0.65, yy+ph/2-0.1, 0.28, col, str(i+1), "FFFFFF", 15)
        txt(s, x+1.2, yy, w-1.4, 0.44, [(it.get("title",""), 14.5, True, "FFFFFF")],
            anchor=MSO_ANCHOR.MIDDLE)
        dv = it.get("desc","")
        if dv:
            txt(s, x+1.2, yy+0.44, w-1.5, ph-0.5, [(dv[:70], 11, False, "C7CDD6")],
                line_spacing=1.2)
    return s


# 9) 타임라인 STEP 점선 (STEP01~04) — 배우고 완성
def s_timeline(prs, pal, d, label_text, head, summary, items):
    items = _real_items(items)[:4]
    n = len(items)
    if n == 0:
        return None
    s = _new(prs)
    _frame(s, pal, d)
    _left(s, pal, label_text, head, summary)
    ac, ac2 = pal["accent"], pal["accent2"]
    y0 = 1.5
    rowh = (BODY_BOT - y0 - 0.1)/n
    lineX = RIGHT_X + 0.4
    for i,it in enumerate(items):
        y = y0 + i*rowh
        col = ac if i%2==0 else ac2
        _circle_icon(s, lineX, y+0.45, 0.20, col)
        if i < n-1:
            vline(s, lineX, y+0.68, rowh-0.4, DBL, weight=1.4)
        txt(s, lineX+0.55, y+0.06, 1.6, 0.4, [(f"STEP {i+1:02d}", 11, True, col)])
        txt(s, lineX+0.55, y+0.42, RIGHT_W-1.1, 0.45,
            [(it.get("title",""), 15, True, HEAD_INK)])
        dv = it.get("desc","")
        if dv:
            txt(s, lineX+0.55, y+0.86, RIGHT_W-1.1, rowh-0.9, [(dv[:65], 11.5, False, INK)],
                line_spacing=1.2)
    return s


# 11) 4박스 2x2 (제목+본문) — 데이터 기반 성장 / 특별 프로그램
def s_quad(prs, pal, d, label_text, head, summary, items):
    items = _real_items(items)[:4]
    n = len(items)
    if n == 0:
        return None
    s = _new(prs)
    _frame(s, pal, d)
    _left(s, pal, label_text, head, summary)
    ac, ac2 = pal["accent"], pal["accent2"]
    gx, gy = 0.44, 0.4
    base_y = 1.45
    # 개수별 배치: 1=한 칸 크게, 2=가로 2, 3=위2+아래1(중앙), 4=2x2
    if n == 1:
        cw = RIGHT_W; ch = 3.6
        pos = [(RIGHT_X, base_y)]
    elif n == 2:
        cw = RIGHT_W/2 - 0.22; ch = 3.6
        pos = [(RIGHT_X, base_y),(RIGHT_X+cw+gx, base_y)]
    else:
        cw = RIGHT_W/2 - 0.22; ch = 2.0
        pos = [(RIGHT_X, base_y),(RIGHT_X+cw+gx, base_y),
               (RIGHT_X, base_y+ch+gy),(RIGHT_X+cw+gx, base_y+ch+gy)]
        if n == 3:  # 세 번째는 아래 가운데
            pos[2] = (RIGHT_X + (RIGHT_W-cw)/2, base_y+ch+gy)
            pos = pos[:3]
    cols = [ac, ac2, ac2, ac]
    for i,(x,y) in enumerate(pos):
        _card(s, x, y, cw, ch)
        no = items[i].get("no","") or f"{i+1:02d}"
        txt(s, x+0.28, y+0.22, 1.4, 0.5, [(no, 20, True, cols[i])])
        ttl = items[i].get("title","")
        txt(s, x+0.28, y+0.78, cw-0.5, 0.6, [(ttl, 15, True, HEAD_INK)])
        dv = _clean_desc(ttl, items[i].get("desc",""))
        if dv:
            cap = 200 if ch > 3 else 75
            txt(s, x+0.28, y+1.3, cw-0.5, ch-1.5,
                [(_trim_sentence(dv, cap), 11.5, False, INK)], line_spacing=1.22)
    return s


# 12) 진단 회색 행 (회색 헤더바 + 화살표 + 설명) — 매트릭스 클리닉 / 과목별 관리
def s_diagrows(prs, pal, d, label_text, head, summary, cols, style="auto"):
    """과목별 관리. style: auto(과목 1개면 box, 2개면 row) / row(줄 타입) / box(박스 타입)"""
    s = _new(prs)
    _frame(s, pal, d)
    _left(s, pal, label_text, head, summary)
    ac, ac2 = pal["accent"], pal["accent2"]
    cols = (cols or [])[:2]
    if not cols:
        cols = [{"subject":"","items":[]}]
    st = (style or "auto").lower()
    if st == "auto":
        st = "box" if len(cols) == 1 else "row"

    if st == "box":
        # 박스 타입: 항목마다 카드를 2열로 배치 (과목이 하나일 때 보기 좋음)
        items = []
        for col in cols:
            for it in col.get("items", []):
                items.append((col.get("subject",""), it))
        items = items[:6]
        n = len(items)
        if n == 0:
            return s
        ncol = 2 if n > 2 else 1
        nrow = (n + ncol - 1)//ncol
        gx, gy = 0.35, 0.3
        cw = (RIGHT_W - (ncol-1)*gx)/ncol
        ch = min(1.55, (BODY_BOT - 1.4 - (nrow-1)*gy)/nrow)
        for i,(subj,it) in enumerate(items):
            r, c = divmod(i, ncol)
            x = RIGHT_X + c*(cw+gx)
            y = 1.4 + r*(ch+gy)
            _card(s, x, y, cw, ch)
            cc = ac if i % 2 == 0 else ac2
            rect(s, x, y, 0.09, ch, cc, radius=False)
            k = it.get("k","")
            txt(s, x+0.26, y+0.16, cw-0.4, 0.42, [(k, 13.5, True, HEAD_INK)])
            v = _trim_sentence(it.get("v",""), 70)
            if v:
                txt(s, x+0.26, y+0.62, cw-0.45, ch-0.75, [(v, 11.5, False, INK)],
                    line_spacing=1.24)
        return s

    # 줄 타입(기본): 과목별 좌우 2단 + 헤더바
    colw = (RIGHT_W - 0.4)/len(cols)
    for ci,col in enumerate(cols):
        x = RIGHT_X + ci*(colw+0.4)
        cc = ac if ci==0 else ac2
        hb = rect(s, x, 1.4, colw, 0.55, cc, radius=True); hb.shadow.inherit=False
        txt(s, x, 1.4, colw, 0.55, [(col.get("subject",""), 15, True, "FFFFFF")],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        its = col.get("items", [])[:5]
        y = 2.15
        for it in its:
            _card(s, x, y, colw, 0.78, CARD)
            k = it.get("k","")
            kb = rect(s, x+0.16, y+0.16, 0.9, 0.46, cc, radius=True); kb.shadow.inherit=False
            txt(s, x+0.16, y+0.16, 0.9, 0.46, [(k, 12, True, "FFFFFF")],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            txt(s, x+1.2, y+0.14, colw-1.3, 0.5,
                [("» "+_trim_sentence(it.get("v",""), 40), 11.5, False, INK)],
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
            y += 0.9
    return s


# 15) 불릿 (파란 반원 마커 + 제목 + 설명) — 학습 로드맵 / 정기 상담 / 규정
def s_bulletlist(prs, pal, d, label_text, head, summary, items):
    s = _new(prs)
    _frame(s, pal, d)
    _left(s, pal, label_text, head, summary)
    ac = pal["accent"]
    items = items[:5]
    y0 = 1.5
    n = max(1,len(items))
    rowh = min(1.15, (BODY_BOT-y0-0.1)/n)
    # 행 높이에 담길 본문 글자 상한(대략): 폭 RIGHT_W-0.7, 12pt, 담길 줄수
    cap = {4:66, 3:100, 2:150, 1:220}.get(n, 60)
    if n >= 5: cap = 52
    for i,it in enumerate(items):
        y = y0 + i*rowh
        _circle_icon(s, RIGHT_X+0.22, y+0.28, 0.16, ac)
        ttl = it.get("title","") or it.get("k","")
        txt(s, RIGHT_X+0.6, y+0.02, RIGHT_W-0.7, 0.5, [(ttl, 15, True, HEAD_INK)])
        dv = it.get("desc","") or it.get("v","")
        if dv:
            txt(s, RIGHT_X+0.6, y+0.5, RIGHT_W-0.7, rowh-0.55,
                [(_trim_sentence(dv, cap), 12, False, INK)], line_spacing=1.2)
    return s


# 16) SWOT 사분면 (S/W/O/T) — 마무리
def s16_closing(prs, pal, d):
    s = _new(prs)
    closing = d.get("closing", {}); ac = pal["accent"]; ac2 = pal["accent2"]
    academy = d.get("academy", {})
    _frame(s, pal, d)
    _left(s, pal, "마무리", closing.get("head","함께,\n끝까지").replace("\\n","\n"),
          closing.get("cta",""))
    # 약도 이미지가 있으면 우측을 연락처 카드 + 약도로 나눈다
    _cl_assets = d.get("assets", {}) or {}
    _map_url = _cl_assets.get("map", "") or ""
    x, y, w = RIGHT_X, 1.5, RIGHT_W
    map_h = 0
    if _map_url:
        # 위: 연락처 카드(2줄) + CTA / 아래: 약도
        card_h = 1.42
        map_y = y + card_h + 0.62
        map_h = 6.3 - map_y
        if map_h < 1.1:
            map_h = 0
        else:
            if not _place_img(s, _map_url, x, map_y, w, map_h, cover=False):
                map_h = 0
    card_h = 1.42 if map_h else 3.4
    _card(s, x, y, w, card_h, CARD)
    rows = [("학원", academy.get("name","")),
            ("과목", academy.get("subjects","")),
            ("전화", _fmt_phone(academy.get("phone",""))),
            ("위치", (academy.get("location","") or academy.get("address_short",""))),
            ("시간", academy.get("hours",""))]
    rows = [(k,v) for k,v in rows if v]
    if map_h:
        # 약도가 위치를 보여주므로 카드에는 짧은 핵심 정보만 (긴 주소 제외)
        rows = [(k,v) for k,v in rows if k in ("학원","전화")][:2]
    yy = y+0.28 if map_h else y+0.35
    gap = 0.56 if map_h else 0.58
    for k,v in rows:
        v = str(v).strip()
        vsz = 12.5 if len(v) <= 30 else (11 if len(v) <= 40 else 10)
        if len(v) > 52: v = v[:52].rstrip() + "…"
        kb = rect(s, x+0.3, yy, 0.85, 0.44, ac, radius=True); kb.shadow.inherit=False
        txt(s, x+0.3, yy, 0.85, 0.44, [(k, 12, True, "FFFFFF")],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x+1.3, yy, w-1.6, 0.44, [(v, vsz, False, INK)],
            anchor=MSO_ANCHOR.MIDDLE)
        yy += gap
    # CTA 배지 (약도가 있으면 카드 바로 아래 얇게, 없으면 하단)
    cta = closing.get("cta","지금 바로 상담을 예약하세요")
    cta_y = (y + card_h + 0.08) if map_h else 5.15
    cta_h = 0.45 if map_h else 0.66
    cb = rect(s, x, cta_y, w, cta_h, ac2, radius=True); cb.shadow.inherit=False
    txt(s, x, cta_y, w, cta_h, [(cta, 13.5 if map_h else 15, True, "FFFFFF")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


# ---- 시간표 (book 프레임 안에 v3 표 재사용) ----
def s_timetable(prs, pal, d, group, rows, part=None, parts=None):
    s = _new(prs)
    _frame(s, pal, d)
    v3pal = V3.PALETTES.get(_v3name(pal), V3.PALETTES["teal_blue"])
    gc = V3.group_color(v3pal, group)
    if group == "기타":
        gname = "특별 과정"
    elif group == "특강":
        gname = "특강"
    else:
        gname = group if group.endswith("부") else group+"부"
    part_txt = f" ({part}/{parts})" if parts and parts>1 else ""
    _left(s, pal, "시간표"+part_txt, f"{gname}\n시간표",
          "요일과 시간에 맞춰 반을 선택하세요")
    # 표: 우측 영역에 배치
    V3._draw_table_at(s, v3pal, gc, rows, RIGHT_X-3.4, 1.5, RIGHT_W+3.4) \
        if hasattr(V3, "_draw_table_at") else _draw_table_book(s, gc, rows)
    return s

def _draw_table_book(s, header_hex, rows):
    """book 프레임용 표: 우측 넓은 영역(RIGHT_X부터 오른쪽 여백까지)에 그림."""
    from pptx.oxml.ns import qn as _qn
    headers = ["반","요일 · 시간","강사","강의실"]
    cols_ratio = [0.30,0.34,0.18,0.18]
    x = MX; top = 1.5
    tbl_w = SLIDE_W - 2*MX
    n = len(rows); ncol = 4
    avail = 4.9
    rh = max(0.32, min(0.60, avail/(n+1))); hh = min(0.5, rh+0.04)
    body_font = 14 if n<=8 else (13 if n<=11 else 12)
    tbl_h = hh + rh*n
    # 시간표는 넓게: 중앙 접힘선 오른쪽 전체 폭 사용 (좌측 라벨/헤더는 유지)
    x = FOLD_X + 0.2
    tbl_w = SLIDE_W - x - MX
    gf = s.shapes.add_table(n+1, ncol, Inches(x), Inches(top), Inches(tbl_w), Inches(tbl_h))
    tbl = gf.table
    tbl.first_row=False; tbl.horz_banding=False; tbl.first_col=False
    try:
        tblPr = tbl._tbl.find(_qn('a:tblPr'))
        sid = tblPr.find(_qn('a:tableStyleId'))
        if sid is not None: tblPr.remove(sid)
    except Exception: pass
    for ci,ra in enumerate(cols_ratio):
        tbl.columns[ci].width = Inches(ra*tbl_w)
    tbl.rows[0].height = Inches(hh)
    for ri in range(1,n+1): tbl.rows[ri].height = Inches(rh)
    def _cell(cell, text, fs, bold, color, fill, left=True):
        cell.margin_left=Inches(0.12); cell.margin_right=Inches(0.06)
        cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.fill.solid(); cell.fill.fore_color.rgb=_rgb(fill)
        tf=cell.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if left else PP_ALIGN.CENTER
        r=p.add_run(); r.text=text; r.font.size=Pt(fs); r.font.bold=bold
        r.font.color.rgb=_rgb(color); r.font.name=FONT
    for ci,h in enumerate(headers):
        _cell(tbl.cell(0,ci), h, 13, True, "FFFFFF", header_hex)
    for ri,row in enumerate(rows):
        z = CARD if ri%2==1 else "FFFFFF"
        for ci in range(ncol):
            t = row[ci] if ci<len(row) else ""
            col = MUTED if t.strip()=="-" else INK
            _cell(tbl.cell(ri+1,ci), t, body_font, False, col, z)
    # 하단 라인
    for row in tbl.rows:
        for cell in row.cells:
            tcPr=cell._tc.get_or_add_tcPr()
            ex=tcPr.find(_qn("a:lnB"))
            if ex is not None: tcPr.remove(ex)
            ln=tcPr.makeelement(_qn("a:lnB"),{"w":"6350","cap":"flat"})
            fill=ln.makeelement(_qn("a:solidFill"),{})
            clr=fill.makeelement(_qn("a:srgbClr"),{"val":"E3E8EF"})
            fill.append(clr); ln.append(fill); tcPr.append(ln)
    return gf


def _v3name(pal):
    # book pal dict → v3 팔레트명 역매핑
    for name,p in BOOK_PALETTES.items():
        if p is pal: return name
    return "teal_blue"


MAX_ROWS_PER_TT = 12

# ======================================================================
#  build
# ======================================================================
def build(data, palette="teal_blue", out="brandbook_book.pptx"):
    pal = BOOK_PALETTES.get(palette, BOOK_PALETTES["teal_blue"])
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    ac = data.get("academy", {})
    intro = data.get("intro", {})
    feats = data.get("features", [])
    ach = data.get("achievements", {})
    targets = data.get("targets", {})
    curr = data.get("curriculum", {})
    specials = data.get("specials", {})
    mgmt = data.get("management", {})
    admission = data.get("admission", {})
    rules = data.get("rules", {})
    faq = data.get("faq", {})

    # 1 표지
    s01_cover(prs, pal, data)
    # 2 소개 (사분면)
    s02_intro(prs, pal, data)
    # 3 네 가지 힘 (원형연결 4)
    if feats:
        s03_four(prs, pal, data)
    # 4 정도와 정성 (원형 3단) — 실적 or 강점 상위 3
    ach_items = ach.get("items", []) if ach else []
    three = ([{"title":a.get("name",""),"desc":a.get("desc","")} for a in ach_items]
             or [{"title":f.get("title",""),"desc":f.get("desc","")} for f in feats[:3]])
    if three:
        _three_circles(prs, pal, data, "정도와 정성",
                       ach.get("head","우리가 만드는\n결과"),
                       "정성을 다한 수업이 결과로 이어집니다", three[:3])
    # 4-b 실적 자료 이미지 (담당자가 올린 입결·성적·수상 이미지가 있을 때만)
    _assets = data.get("assets", {}) or {}
    _ach_imgs = _assets.get("achievements", []) or []
    if _ach_imgs:
        s_achieve_images(prs, pal, data, "실적 자료",
                         "숫자로 확인하는\n성장의 결과",
                         "학원이 만들어낸 실제 결과입니다", _ach_imgs)
    # 5 맞춤 수업 설계 (번호원) — 커리큘럼 단계
    stages = curr.get("stages", []) if curr else []
    if stages:
        s_numbered(prs, pal, data, "맞춤 수업 설계",
                   curr.get("head","단계별로\n설계합니다"),
                   "수준에 맞춰 단계적으로 성장합니다",
                   [{"title":st.get("title",""),"desc":st.get("desc","")} for st in stages])
    # 6 네 가지 반 (회색바 리스트) — 수업 대상
    t_items = targets.get("items", []) if targets else []
    if t_items:
        s_barlist(prs, pal, data, "수업 대상",
                  targets.get("head","학년별\n맞춤 과정"),
                  "학년에 맞는 과정을 운영합니다",
                  [{"title":it.get("grade",""),"desc":it.get("subj","")} for it in t_items])
    # 7 깊이 있는 학습 (어두운 카드) — 특별 프로그램 상위 4
    sp_items = specials.get("items", []) if specials else []
    if sp_items:
        s_darkcard(prs, pal, data, "깊이 있는 학습",
                   specials.get("head","한 단계\n더 깊게"),
                   "심화 프로그램으로 실력을 완성합니다", sp_items[:4])
    # 8 시간표 (그룹 수만큼, 자동 분할)
    for tt in data.get("timetables", []):
        group = tt.get("group",""); rows = tt.get("rows", [])
        if not rows: continue
        if len(rows) <= MAX_ROWS_PER_TT:
            s_timetable(prs, pal, data, group, rows)
        else:
            chunks=[rows[i:i+MAX_ROWS_PER_TT] for i in range(0,len(rows),MAX_ROWS_PER_TT)]
            for pi,ch in enumerate(chunks,1):
                s_timetable(prs, pal, data, group, ch, part=pi, parts=len(chunks))
    # 9 배우고 완성 (타임라인) — 입학 절차
    adm_steps = admission.get("steps", []) if admission else []
    if adm_steps:
        s_timeline(prs, pal, data, "입학 절차",
                   admission.get("head","이렇게\n시작합니다"),
                   "간단한 절차로 시작할 수 있습니다",
                   [{"title":st.get("title",""),"desc":st.get("desc","")} for st in adm_steps])
    # 11 데이터 기반 성장 (4박스) — 특별 프로그램(있으면 이미 씀) → 강점 하위 4
    #    특별 프로그램이 없을 때만 강점으로 사분면 채움(중복 방지)
    if not sp_items and len(feats) >= 4:
        s_quad(prs, pal, data, "핵심 강점",
               "우리를 만드는\n네 가지 축",
               "학원의 핵심 강점을 정리했습니다",
               [{"title":f.get("title",""),"desc":f.get("desc","")} for f in feats[:4]])
    # 12 매트릭스 클리닉 (진단 회색행) — 과목별 관리
    mg_cols = (mgmt.get("cols") or mgmt.get("columns") or []) if mgmt else []
    # 앱은 {name, rows:[{k,v}]} 형태로 보내므로 {subject, items} 로 정규화
    _norm = []
    for c in mg_cols:
        if not isinstance(c, dict):
            continue
        _norm.append({
            "subject": c.get("subject") or c.get("name") or "",
            "items":   c.get("items")   or c.get("rows") or [],
        })
    mg_cols = _norm
    if mg_cols:
        s_diagrows(prs, pal, data, "과목별 관리",
                   mgmt.get("head","과목마다\n다르게 관리"),
                   "과목 특성에 맞춰 학생을 관리합니다", mg_cols,
                   style=(mgmt.get("style") or "auto"))
    # 15 학습 로드맵 / 정기 상담 (불릿) — 규정
    r_items = rules.get("items", []) if rules else []
    if r_items:
        s_bulletlist(prs, pal, data, "학원 지침",
                     rules.get("head","약속과\n지침"),
                     "안심하고 맡길 수 있도록 지침을 운영합니다",
                     [{"title":it.get("k",""),"desc":it.get("v","")} for it in r_items])
    # FAQ (불릿 리스트) — 있으면
    f_items = faq.get("items", []) if faq else []
    if f_items:
        s_bulletlist(prs, pal, data, "자주 묻는 질문",
                     faq.get("head","궁금한 점을\n확인하세요"),
                     "자주 묻는 질문을 정리했습니다",
                     [{"title":q.get("q",""),"desc":q.get("a","")} for q in f_items])
    # 16 마무리 (사분면/연락)
    s16_closing(prs, pal, data)

    prs.save(out)
    return out


if __name__ == "__main__":
    import json, sys
    src = sys.argv[1] if len(sys.argv)>1 else "sewon.json"
    pal = sys.argv[2] if len(sys.argv)>2 else "teal_blue"
    out = sys.argv[3] if len(sys.argv)>3 else "brandbook_book.pptx"
    data = json.load(open(src, encoding="utf-8"))
    build(data, pal, out)
    print("saved", out)
